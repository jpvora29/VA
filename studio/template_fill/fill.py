"""Fill engine — clone the active ``.pptx`` and write the resolved values in.

Generic over any template. Opens the template (so think-cell OLE objects, freeform
shapes, theme, masters and layouts are preserved), then writes each *filled* slot.

Correctness rules (a corrupt deck is worse than an unfilled one):
  * text is written **without deleting run elements** — set the first run's text and
    blank the rest, so think-cell field references and run formatting survive;
  * a chart is rewritten only after it is cut loose from its source workbook
    (:func:`_replace_chart_data`); a think-cell / externally-linked one never is, and keeps
    the template's authored data;
  * a final literal-substitution pass fills label tokens like ``Country (1)`` and
    ``xyz`` from the selection;
  * hidden / reordered slides are applied last by editing the slide-id list — so an
    unselected second country/product block is dropped cleanly.
"""
from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

from logger import get_logger
from studio.template_fill.model import materialize_fields

logger = get_logger(__name__)


# ── shape lookup (recursive, by stable cNvPr id) ─────────────────────────────


def _iter_leaves(shapes) -> Iterable[Any]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_leaves(sh.shapes)
        else:
            yield sh


def _index_by_id(slide) -> Dict[int, Any]:
    return {int(sh.shape_id): sh for sh in _iter_leaves(slide.shapes)}


# ── run-preserving text writes (never delete <a:r> elements) ─────────────────


_SOFT_BREAK = "\v"


def _set_paragraph_text(paragraph, text: str) -> None:
    """Write ``text`` into ``paragraph``, keeping its formatting.

    Fast path: reuse the existing runs (set run 0, blank the rest) so think-cell field
    references and run formatting survive. Two cases need the paragraph rebuilt instead:

      * the paragraph has NO ``<a:r>`` runs — its visible text (and all its formatting)
        lives in ``<a:fld>`` elements;
      * the new text carries a soft line break, which ``run.text`` would escape into a
        literal ``_x000B_`` instead of a ``<a:br>``.

    Both rebuild via ``paragraph.text`` and graft the original ``rPr`` onto the new runs,
    so nothing falls back to the theme default (product names keep their white bold, rank
    changes their green).
    """
    runs = list(paragraph.runs)
    if runs and _SOFT_BREAK not in text:
        runs[0].text = text
        for r in runs[1:]:
            r.text = ""
        return

    from copy import deepcopy

    from pptx.oxml.ns import qn

    source = runs[0]._r if runs else paragraph._p.find(qn("a:fld"))
    rPr = source.find(qn("a:rPr")) if source is not None else None
    rPr = deepcopy(rPr) if rPr is not None else None
    paragraph.text = text
    if rPr is not None:
        for run in paragraph.runs:
            old = run._r.find(qn("a:rPr"))
            if old is not None:
                run._r.remove(old)
            run._r.insert(0, deepcopy(rPr))


# The template's own trend palette (sampled from the authored delta runs):
# green ▲/↑ growth, red ▼/↓ decline, Marsh amber ►/flat.
_TREND_COLOR = {"▲": "6ABF30", "↑": "6ABF30", "▼": "C53532", "↓": "C53532", "►": "FFBF00"}
# A value line ending in a parenthesised trend, e.g. "$411M (+9.8%▲)" / "PY (+4▲)".
_DELTA_TAIL = re.compile(r"^(?P<head>.*?)\s*(?P<delta>\([^()]*[▲▼►]\))\s*$")


def _trend_color(text: str):
    from pptx.dml.color import RGBColor

    for arrow, hexv in _TREND_COLOR.items():
        if arrow in text:
            return RGBColor.from_string(hexv)
    return None


def _set_value_line(paragraph, line: str) -> None:
    """Write one value line, keeping the template's value-vs-delta run structure.

    The authored cells split "value (delta▲)" across runs — big white value runs, then
    coloured delta runs from the first "(". A whole-line write into run 0 destroyed
    that (everything came out in the value's white). Re-split on the same boundary and
    colour the delta by its OWN sign (the template's colour is a static example — a
    red-authored cell must still show green when the real number grew)."""
    m = _DELTA_TAIL.match(line)
    runs = list(paragraph.runs)
    if m and runs:
        split = next((i for i, r in enumerate(runs) if r.text.lstrip().startswith("(")), None)
        if split:                                   # a real delta run group exists
            runs[0].text = m.group("head") + " "
            for r in runs[1:split]:
                r.text = ""
            runs[split].text = m.group("delta")
            for r in runs[split + 1:]:
                r.text = ""
            color = _trend_color(m.group("delta"))
            if color is not None:
                runs[split].font.color.rgb = color
            return
    _set_paragraph_text(paragraph, line)


def _set_cell_text(cell, text: str) -> None:
    """Write ``text`` into a cell, mapping ``\\n``-separated lines onto the cell's own
    paragraphs — so a composite value keeps its caption paragraph's formatting (e.g.
    ``"$57M (+4.1%▲)\\nMarsh GWP"`` fills the value line and rewrites the 9pt caption
    in place). Surplus template paragraphs are blanked; surplus lines are appended."""
    tf = cell.text_frame
    paras = list(tf.paragraphs)
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_value_line(paras[i], line)
        else:
            tf.add_paragraph().text = line
    for p in paras[len(lines):]:
        for r in p.runs:
            r.text = ""


def _write_text_shape(shape, where: List[Any], text: str) -> None:
    if not shape.has_text_frame or not (where and where[0] == "para"):
        return
    idx = int(where[1])
    paras = shape.text_frame.paragraphs
    if 0 <= idx < len(paras):
        _set_paragraph_text(paras[idx], text)


def _write_table(shape, where: List[Any], text: str) -> None:
    if not shape.has_table or not (where and where[0] == "cell"):
        return
    r, c = int(where[1]), int(where[2])
    rows = shape.table.rows
    if 0 <= r < len(rows):
        cells = rows[r].cells
        if 0 <= c < len(cells):
            _set_cell_text(cells[c], text)


# ── commentary typography ─────────────────────────────────────────────────────

# Written commentary uses ONE face/size everywhere. The template's ellipsis
# placeholders carry whatever ad-hoc run formatting the author left (10pt here,
# 14pt there, theme default on appended paragraphs), so inheriting it makes the
# prose inconsistent — too big in one box (overlapping the lines below), too
# small in the next. Roles prefixed ``note:`` (prose slots) and ``fbnote:``
# (feedback/quadrant/highlights table cells) are restyled after the write.
_COMMENTARY_FONT_NAME = "Arial"
_COMMENTARY_FONT_PT = 11
_COMMENTARY_ROLE_PREFIXES = ("note:", "fbnote:")


def _is_commentary_role(role: Optional[str]) -> bool:
    return bool(role) and str(role).startswith(_COMMENTARY_ROLE_PREFIXES)


def _style_commentary_paragraphs(paragraphs, ink=None) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    for p in paragraphs:
        # A bullet list reads left-aligned. The author's EMPTY panels carry whatever
        # alignment happened to be on them (the quadrant's four columns are authored
        # centred, left, centred, centred), so inherited alignment makes one page's four
        # commentary columns look like four different slides.
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.name = _COMMENTARY_FONT_NAME
            r.font.size = Pt(_COMMENTARY_FONT_PT)
            if ink is not None:
                r.font.color.rgb = ink
            # Some template placeholders carry a yellow <a:highlight> — real
            # commentary must not inherit the author's "fill me" marker.
            rPr = r._r.find(qn("a:rPr"))
            hl = rPr.find(qn("a:highlight")) if rPr is not None else None
            if hl is not None:
                rPr.remove(hl)


def _authored_ink(frames) -> Optional[Any]:
    """The colour the author used for the commentary they actually WROTE on this page.

    An emptied prose box keeps the run colour of the text that was deleted from it — the
    quadrant's blank columns carry a leftover white, which is invisible the moment real
    commentary lands on their pale panel. The one column the author did write in is the
    honest source of the page's commentary ink, and using it everywhere keeps the four
    columns looking like one slide.
    """
    for frame in frames:
        for p in frame.paragraphs:
            for r in p.runs:
                if r.text.strip() and r.font.color.type is not None:
                    try:
                        return r.font.color.rgb
                    except AttributeError:      # a theme colour — already renders correctly
                        return None
    return None


def _anchor_commentary_top(frame) -> None:
    """Start commentary at the top of its box, never floating in the middle of it.

    A quadrant panel is as tall as the slide; middle-anchored bullets end up marooned in
    the centre of the column while the neighbouring column starts at the top.
    """
    from pptx.enum.text import MSO_ANCHOR

    try:
        frame.vertical_anchor = MSO_ANCHOR.TOP
    except (AttributeError, ValueError):        # a frame that cannot be anchored
        pass


# Grid metrics whose colour must follow the SIGN of the real value — the template
# authors them with static example colours (a red "-xx.x%" row), which a positive
# actual would wrongly inherit.
_SIGNED_ROLE_SUFFIXES = (":var", ":rank_change")


def _recolor_by_sign(shape, where: List[Any], text: str) -> None:
    """Colour a signed metric by the written value's own trend, in a text box OR a cell.

    The breakdown grid is authored as a TABLE, so a cell-only write has to be handled here
    too — otherwise a real ``-0.4%`` keeps the green the author's positive example was
    styled in, and the page reads a decline as growth.
    """
    t = str(text).strip()
    color = _trend_color(t)
    if color is None:
        if t.startswith("+"):
            color = _trend_color("▲")
        elif t.startswith("-"):
            color = _trend_color("▼")
    if color is None:
        return
    try:
        for paragraph in _signed_paragraphs(shape, where):
            for r in paragraph.runs:
                r.font.color.rgb = color
    except Exception as exc:  # noqa: BLE001 — colouring must never break the fill
        logger.warning("template_fill: sign recolor skipped: %s", exc)


def _signed_paragraphs(shape, where: List[Any]) -> List[Any]:
    """The paragraphs a signed value was just written into."""
    if not where:
        return []
    if where[0] == "para" and shape.has_text_frame:
        paras = shape.text_frame.paragraphs
        idx = int(where[1])
        return [paras[idx]] if 0 <= idx < len(paras) else []
    if where[0] == "cell" and shape.has_table:
        r, c = int(where[1]), int(where[2])
        rows = shape.table.rows
        if 0 <= r < len(rows):
            cells = rows[r].cells
            if 0 <= c < len(cells):
                return list(cells[c].text_frame.paragraphs)
    return []


# ── bullet-point commentary ──────────────────────────────────────────────────

# Commentary is written as a BULLET LIST — one point per ``\n``-separated line, each becoming
# its own bulleted paragraph. Where the author already bulleted the box we reuse THEIR bullet
# (the Trading Summary columns use a Wingdings "§"); where they turned bullets off we apply
# this standard one, because a point-wise list needs a visible marker.
_DEFAULT_BULLET_CHAR = "•"
_DEFAULT_BULLET_FONT = "Arial"
_BULLET_INDENT_EMU = 171450          # ~0.19", a conventional hanging indent
# A heading line inside a commentary block ("Key Highlights:") introduces the bullets rather
# than being one, so it keeps its own paragraph unbulleted.
_HEADING_LINE = re.compile(r":\s*$")


# `<a:pPr>`'s children are a SEQUENCE: the bullet-font slot precedes the bullet-character
# slot, which precedes tabLst/defRPr/extLst. Appending at the end instead produces invalid
# OOXML that PowerPoint silently discards on save — which is exactly what happened before
# these successor lists were used to insert each element in its proper place.
_AFTER_BU_FONT = ("a:buNone", "a:buAutoNum", "a:buChar", "a:tabLst", "a:defRPr", "a:extLst")
_AFTER_BU_CHAR = ("a:tabLst", "a:defRPr", "a:extLst")


def _drop(pPr, *tags: str) -> None:
    from pptx.oxml.ns import qn

    for tag in tags:
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)


def _set_bullet(paragraph, *, bulleted: bool) -> None:
    """Give ``paragraph`` a visible bullet, or explicitly none.

    Any bullet the author chose is left exactly as it is — only a paragraph with no bullet
    (or an explicit ``<a:buNone>``) gets the default one, so the deck keeps its own look.
    """
    from pptx.oxml.ns import qn

    pPr = paragraph._p.get_or_add_pPr()
    if not bulleted:
        _drop(pPr, "a:buChar", "a:buAutoNum")
        if pPr.find(qn("a:buNone")) is None:
            pPr.insert_element_before(pPr.makeelement(qn("a:buNone"), {}), *_AFTER_BU_CHAR)
        return

    if pPr.find(qn("a:buChar")) is not None or pPr.find(qn("a:buAutoNum")) is not None:
        return                                      # the author already bulleted this box
    _drop(pPr, "a:buNone", "a:buFontTx")
    if pPr.find(qn("a:buFont")) is None:
        font = pPr.makeelement(qn("a:buFont"), {"typeface": _DEFAULT_BULLET_FONT})
        pPr.insert_element_before(font, *_AFTER_BU_FONT)
    char = pPr.makeelement(qn("a:buChar"), {"char": _DEFAULT_BULLET_CHAR})
    pPr.insert_element_before(char, *_AFTER_BU_CHAR)
    pPr.set("marL", str(_BULLET_INDENT_EMU))
    pPr.set("indent", str(-_BULLET_INDENT_EMU))


def _commentary_frame(shape, where: List[Any]):
    """The text frame a commentary role writes into — a shape's own, or a table cell's."""
    if where and where[0] == "para" and getattr(shape, "has_text_frame", False):
        return shape.text_frame
    if where and where[0] == "cell" and getattr(shape, "has_table", False):
        r, c = int(where[1]), int(where[2])
        rows = shape.table.rows
        if 0 <= r < len(rows) and 0 <= c < len(rows[r].cells):
            return rows[r].cells[c].text_frame
    return None


def _write_bullets(frame, text: str, *, template_paragraph=None, ink=None) -> None:
    """Lay ``text``'s lines out over ``frame``'s paragraphs, one bulleted point per line.

    The frame's own paragraphs are reused first (so the author's formatting is inherited);
    extra points are appended as clones of ``template_paragraph`` — the first authored
    paragraph — so an appended bullet looks like the ones above it rather than falling back
    to the theme default. Surplus authored paragraphs are removed, not just blanked, so no
    empty bullet is left hanging.
    """
    from copy import deepcopy

    lines = [ln for ln in str(text).split("\n")]
    paras = list(frame.paragraphs)
    source = template_paragraph if template_paragraph is not None else (paras[0] if paras else None)

    for i, line in enumerate(lines):
        if i < len(paras):
            target = paras[i]
        elif source is not None:
            clone = deepcopy(source._p)
            source._p.getparent().append(clone)
            target = frame.paragraphs[-1]
        else:
            target = frame.add_paragraph()
        _set_paragraph_text(target, line)
        _set_bullet(target, bulleted=not _HEADING_LINE.search(line))

    for extra in paras[len(lines):]:
        extra._p.getparent().remove(extra._p)
    _style_commentary_paragraphs(frame.paragraphs, ink)
    _anchor_commentary_top(frame)


def _write_commentary(shape, where: List[Any], text: str, ink=None) -> None:
    """Write bullet-point commentary into a text shape or a table cell."""
    try:
        frame = _commentary_frame(shape, where)
        if frame is None:
            return
        _write_bullets(frame, text, ink=ink)
    except Exception as exc:  # noqa: BLE001 — commentary layout must never break the fill
        logger.warning("template_fill: commentary write fell back to a plain paragraph: %s", exc)
        if where and where[0] == "para":
            _write_text_shape(shape, where, str(text).replace("\n", " "))
        elif where and where[0] == "cell":
            _write_table(shape, where, str(text))


# ── literal label substitution (Country (n) / Region (n) / xyz) ──────────────


def _label_subs(values: Dict[str, Any]) -> List[Any]:
    """Literal placeholder substitutions applied to every run (label text), in order.

    Makes the static template wording follow the selection: ``Country (1)`` /
    ``Country xyz`` → the actual country, bare ``Carrier`` → the carrier name, and a
    year shift so ``2025`` / ``FY 2025`` track the selected reporting year.
    """
    subs: List[Any] = []
    countries = []
    i = 0
    while f"country_name[{i}]" in values:
        countries.append(str(values[f"country_name[{i}]"]))
        i += 1
    # "Country (n)" / "Country / Region(n)" / "Region (n)" → the nth country name.
    for idx, name in enumerate(countries, start=1):
        subs.append((re.compile(rf"\b(?:Country\s*/\s*)?(?:Country|Region)\s*\(\s*{idx}\s*\)", re.I), name))
    # Countries the selection doesn't include → blank, so no "Country (2)" placeholder
    # is left on the quadrant/portfolio/feedback slides.
    for idx in range(len(countries) + 1, 10):
        subs.append((re.compile(rf"\b(?:Country\s*/\s*)?(?:Country|Region)\s*\(\s*{idx}\s*\)", re.I), ""))
    # "Country xyz" / "Region xyz" / bare "xyz" mark the spotlight-YoY callout — name it after
    # the spotlight entity (a significant country, or a product when one country is in scope) so
    # the label agrees with the spotlight value. Falls back to the top country.
    spotlight = str(values.get("spotlight_name", "")).strip() or (countries[0] if countries else "")
    if spotlight:
        subs.append((re.compile(r"\b(?:Country|Region)\s+xyz\b", re.I), spotlight))
        subs.append((re.compile(r"\bxyz\b", re.I), spotlight))

    # Product deck: rewrite the template's authored example products (e.g. "Marine
    # Feedback", "Energy") to the single product this deck is for. Longest names first so a
    # multi-word product isn't half-matched by a shorter one.
    product = str(values.get("product_name", "")).strip()
    if product:
        vocab = [str(p).strip() for p in (values.get("product_vocab") or []) if str(p).strip()]
        for authored in sorted(set(vocab), key=len, reverse=True):
            if authored.casefold() != product.casefold():
                subs.append((re.compile(rf"\b{re.escape(authored)}\b"), product))

    carrier = str(values.get("subject_name", "")).strip()
    if carrier:
        # The cover's "Carrier X" title placeholder → the carrier. MUST precede the bare
        # "Carrier" rule below, which would otherwise leave a dangling " X".
        subs.append((re.compile(r"\bCarrier\s+X\b"), carrier))
        # The template's authored example carrier (e.g. "QBE GWP rank") → this deck's
        # subject. Longest names first so a multi-word carrier isn't half-matched. Also
        # upholds the no-named-peer rule: no other carrier's name survives the fill.
        authored = {str(c).strip() for c in (values.get("carrier_vocab") or []) if str(c).strip()}
        for name in sorted(authored, key=len, reverse=True):
            if name.casefold() != carrier.casefold():
                subs.append((re.compile(rf"\b{re.escape(name)}\b"), carrier))

        # Carrier / Carrier's / Carriers's (template typo) → the carrier (possessive kept).
        def _carrier(m, _c=carrier):
            return f"{_c}’s" if ("'" in m.group(0) or "’" in m.group(0)) else _c

        subs.append((re.compile(r"\bCarriers?(?:['’]s)?\b"), _carrier))

    ty, py = values.get("template_year"), values.get("period_year")
    if ty and py and int(ty) != int(py):
        delta = int(py) - int(ty)
        subs.append((re.compile(r"\b(20\d{2})\b"), lambda m: str(int(m.group(1)) + delta)))
    return subs


def _apply_subs(slide, subs: List[Any]) -> None:
    """Apply label substitutions at PARAGRAPH granularity — a phrase like
    ``Country xyz`` is usually split across runs, so per-run replacement would leave
    the leading word behind. We rewrite a paragraph only when a sub actually changes
    it, writing into run 0 and blanking the rest (run formatting of run 0 survives).

    The text is read via ``paragraph.text`` rather than by joining runs, so a paragraph
    whose text lives in ``<a:fld>`` elements — think-cell writes whole table rows that way —
    is substituted too instead of being silently skipped."""
    if not subs:
        return
    for sh in _iter_leaves(slide.shapes):
        frames = []
        if sh.has_text_frame:
            frames.append(sh.text_frame)
        elif sh.has_table:
            frames.extend(cell.text_frame for row in sh.table.rows for cell in row.cells)
        for tf in frames:
            for para in tf.paragraphs:
                original = para.text
                if not original:
                    continue
                new = original
                for pattern, repl in subs:
                    new = pattern.sub(repl, new)
                if new != original:
                    _set_paragraph_text(para, new)


# ── added free elements ──────────────────────────────────────────────────────


def _add_elements(slide, elements: List[Dict[str, Any]], width_emu: int, height_emu: int) -> None:
    from pptx.util import Emu, Pt

    for el in elements or []:
        tb = slide.shapes.add_textbox(Emu(int(el.get("x", 0))), Emu(int(el.get("y", 0))),
                                      Emu(int(el.get("w", width_emu // 4))),
                                      Emu(int(el.get("h", height_emu // 10))))
        tb.text_frame.text = str(el.get("text", ""))
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(int(el.get("size", 12)))


# ── charts (guarded, best-effort) ────────────────────────────────────────────


def _chart_is_external(chart) -> bool:
    """True for think-cell / externally-linked charts — never safe to rewrite.

    The link is either an external relationship or a ``<c:externalData>`` element in
    the chart XML (think-cell); treat either as external.
    """
    try:
        return (b"externalData" in chart.part.blob
                or any(getattr(r, "is_external", False) for r in chart.part.rels.values()))
    except Exception:  # noqa: BLE001
        return False


def _detach_external_data(chart) -> bool:
    """Cut a chart loose from its external (think-cell / linked-workbook) data source.

    Removes the ``<c:externalData>`` element and drops its relationship, turning the
    chart into a plain native chart whose cached plot XML remains — after which
    ``replace_data`` is safe (python-pptx creates a fresh embedded workbook).
    Returns True when the chart is (now) safe to rewrite."""
    from pptx.oxml.ns import qn

    try:
        cs = chart._chartSpace
        ext = cs.find(qn("c:externalData"))
        if ext is None:
            return not _chart_is_external(chart)
        rid = ext.get(qn("r:id"))
        cs.remove(ext)
        if rid:
            try:
                chart.part.drop_rel(rid)
            except (KeyError, AttributeError):
                pass
        return True
    except Exception as exc:  # noqa: BLE001
        logger.warning("template_fill: could not detach external chart data: %s", exc)
        return False


def _replace_chart_data(chart, chart_data) -> None:
    """``chart.replace_data``, with the chart first cut loose from its source workbook.

    Mandatory, not defensive. ``replace_data`` writes an .xlsx blob into whatever part
    ``c:externalData`` points at, and these templates were authored from BINARY workbooks
    (``.xlsb``) — so the xlsx bytes land in a part still declared as xlsb and PowerPoint
    refuses to open the deck at all. Detaching first makes python-pptx mint a fresh .xlsx.
    """
    _detach_external_data(chart)
    chart.replace_data(chart_data)


def _bubble_points(points: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """The plottable growth points — both YoY axes known. Axis values are FRACTIONS
    (the template charts carry ``0.0%``-formatted axes), sizes are raw premium."""
    return [p for p in points
            if p.get("carrier_yoy") is not None and p.get("marsh_yoy") is not None]


# ── chart number formats ─────────────────────────────────────────────────────

# ``replace_data`` rewrites a chart's cached values and, with them, their FORMAT CODE —
# which python-pptx sets to "General". Every axis and data label in the shipped templates
# is ``sourceLinked``, i.e. renders whatever the source data says, so a refilled chart
# prints ``207.8838507`` where the author had ``$207.9M`` and plots ``0.02`` where the axis
# reads ``2.0%``. Each writer therefore states the format its own numbers are in, and
# unlinks the axes that must not follow the source.
_MILLIONS_FORMAT = "#,##0.0"        # bars — already scaled to the unit the caption declares
_PERCENT_FORMAT = "0.0%"            # bubble axes — growth rates are written as fractions


def _pin_value_axis_format(chart, fmt: str) -> None:
    """Pin every value axis of ``chart`` to ``fmt``, unlinked from the source data."""
    from pptx.oxml.ns import qn

    for ax in chart._chartSpace.iter(qn("c:valAx")):
        numFmt = ax.find(qn("c:numFmt"))
        if numFmt is None:                      # an axis the author left unformatted
            continue
        numFmt.set("formatCode", fmt)
        numFmt.set("sourceLinked", "0")


def _show_series_names(chart) -> None:
    """Label every point with its own SERIES name, and nothing else.

    One series per line of business means the data label names the point — which is what a
    legend would do, except a legend of identically-coloured keys names nothing and steals
    a band of the plot area. Only flags the author already wrote are set, so the label's
    schema order (and the rest of its formatting) survives untouched.
    """
    from pptx.oxml.ns import qn

    flags = {"c:showSerName": "1", "c:showVal": "0", "c:showCatName": "0",
             "c:showPercent": "0", "c:showBubbleSize": "0", "c:showLegendKey": "0"}
    for dLbls in chart._chartSpace.iter(qn("c:dLbls")):
        for tag, val in flags.items():
            el = dLbls.find(qn(tag))
            if el is not None:
                el.set("val", val)


def _write_bubble_chart(chart, points: List[Dict[str, Any]]) -> None:
    """One series per line of business (x = Marsh YoY, y = carrier YoY, size = carrier GWP),
    each bubble labelled with its own line and the axes pinned to percentages."""
    from pptx.chart.data import BubbleChartData

    data = BubbleChartData(number_format=_PERCENT_FORMAT)
    for p in points:
        ser = data.add_series(str(p.get("lob") or ""))
        ser.add_data_point(p["marsh_yoy"] / 100.0, p["carrier_yoy"] / 100.0,
                           abs(p.get("size") or 0.0) or 1.0)
    _replace_chart_data(chart, data)
    _pin_value_axis_format(chart, _PERCENT_FORMAT)
    # Every bubble is cloned from the author's single authored series, so they all share one
    # colour: the legend comes off and each bubble names itself.
    chart.has_legend = False
    _show_series_names(chart)


def _blank_stale_point_labels(slide, chart_shape, points: List[Dict[str, Any]]) -> None:
    """Blank the hand-placed point-label textboxes over a refilled bubble chart —
    they were positioned for the template's authored dummy bubbles, and each bubble
    now carries its own line of business as a data label.

    A stale label is a text *placeholder* whose text matches a plotted name, or a
    short placeholder sitting wholly inside the chart frame (the authored labels for
    lines the data doesn't have). Quadrant captions/speech bubbles are rectangles,
    and the axis captions sit outside the frame — both survive. Labels often hold
    their text in ``<a:fld>`` elements (no runs), so the whole frame is replaced."""
    names = {str(p.get("lob") or "").strip() for p in points if p.get("lob")}

    def _inside(sh) -> bool:
        try:
            return (sh.left >= chart_shape.left and sh.top >= chart_shape.top
                    and sh.left + sh.width <= chart_shape.left + chart_shape.width
                    and sh.top + sh.height <= chart_shape.top + chart_shape.height)
        except TypeError:                       # a shape without geometry → leave it
            return False

    for sh in _iter_leaves(slide.shapes):
        if not getattr(sh, "has_text_frame", False) or sh is chart_shape:
            continue
        txt = sh.text_frame.text.strip()
        if not txt or "placeholder" not in sh.name.lower():
            continue
        if txt in names or (len(txt) <= 30 and _inside(sh)):
            sh.text_frame.text = ""


# ── the LC-ranking quadrant panels ───────────────────────────────────────────

_LC_LABEL_PT = 8            # matches the size the author set on the panel's rank axis
# "$120M" — the x axis states the pool size in millions, in place of the hand-drawn ticks.
_LC_SIZE_FORMAT = '"$"#,##0,,"M"'
# Headroom past the largest pool so the right-most point (and its label) sits inside the
# plot rather than on its edge. The author's own scaling: their axis max is data max × 1.04.
_LC_SIZE_HEADROOM = 1.04
# The author scales the rank axis 0–11 so their painted band boundary lands on rank 5 — the
# top-5 line the rest of the deck benchmarks against. Kept unless the data ranks deeper,
# in which case the axis grows to fit (and the boundary drifts with it).
_LC_RANK_AXIS_MAX = 11.0

_SHAPE_TYPE_FREEFORM, _SHAPE_TYPE_LINE, _SHAPE_TYPE_GROUP = 5, 9, 6
# The furniture the author hand-drew AROUND their example points, all of which is positioned
# for data that is about to be replaced: the text placeholders faking a broken money axis and
# labelling each dummy point, the tiny freeforms marking the axis breaks, and the connectors
# ruling the quadrant split (the painted bands meet along it, so the split still reads).
# The painted band GROUP itself stays — it is the page's design, not its data.
_LC_STALE_SHAPES = {_SHAPE_TYPE_LINE}
# A loose freeform smaller than this share of the panel is an axis-break mark, not artwork.
_LC_MARK_SHARE = 0.02


# How far outside its own plot frame a panel's furniture reaches, as a fraction of the
# chart's own size. Generous vertically (the axis captions and the panel's title sit above
# and below the frame), tight horizontally (the panels sit side by side). Anything two
# panels both reach is settled by whichever chart is nearer.
_PANEL_REACH_X, _PANEL_REACH_Y = 0.10, 0.40


def _centre(shape) -> Optional[Tuple[int, int]]:
    try:
        return shape.left + shape.width // 2, shape.top + shape.height // 2
    except TypeError:                           # a shape without geometry
        return None


def _panel_owner(shape, charts):
    """The quadrant panel a loose shape belongs to — the nearest chart it is in reach of.

    Reach is measured off each chart's own frame, so the page's title and KPI band (which
    sit well clear of every panel) belong to none of them and are never touched.
    """
    at = _centre(shape)
    if at is None:
        return None
    reachable = []
    for chart in charts:
        mid = _centre(chart)
        if mid is None:
            continue
        dx, dy = abs(at[0] - mid[0]), abs(at[1] - mid[1])
        if (dx <= chart.width * (0.5 + _PANEL_REACH_X)
                and dy <= chart.height * (0.5 + _PANEL_REACH_Y)):
            reachable.append((dx * dx + dy * dy, chart))
    return min(reachable, key=lambda r: r[0])[1] if reachable else None


def _is_stale_furniture(shape, panel) -> bool:
    """True for a shape drawn to serve the template's authored dummy points.

    Stale: the *placeholders* (the point labels and the text faking a broken money axis),
    the connectors ruling the quadrant split, and the tiny freeforms marking the axis
    breaks. Kept: the painted band group (the page's own design) and plain text boxes,
    which caption the panel rather than its data.
    """
    kind = getattr(shape, "shape_type", None)
    if kind == _SHAPE_TYPE_GROUP:               # the painted bands — the panel's backdrop
        return False
    if kind in _LC_STALE_SHAPES:
        return True
    if kind == _SHAPE_TYPE_FREEFORM:
        return _area(shape) <= _LC_MARK_SHARE * _area(panel)
    return (getattr(shape, "has_text_frame", False)
            and "placeholder" in (shape.name or "").lower())


def _area(shape) -> float:
    return float(shape.width or 0) * float(shape.height or 0)


def _strip_quadrant_furniture(slide, charts, *, blank: set) -> None:
    """Clear the stale hand-drawn furniture from every quadrant panel on ``slide``.

    A filled panel keeps what still tells the truth — its painted bands, axis captions and
    title; only what was drawn around the author's example points goes, because the refilled
    chart states its own axis and labels its own points. A panel in ``blank`` has no country
    left to report, so the whole panel goes with it.

    Top-level shapes only: a painted band is a GROUP, and it moves and stays whole.
    """
    for sh in list(slide.shapes):               # materialised: the loop deletes as it goes
        if getattr(sh, "has_chart", False):
            continue
        owner = _panel_owner(sh, charts)
        if owner is None:
            continue
        if int(owner.shape_id) in blank or _is_stale_furniture(sh, owner):
            sh._element.getparent().remove(sh._element)


def _reanchor(box, old_frame, new_frame) -> Tuple[int, int, int, int]:
    """Re-place a shape laid out around ``old_frame`` to suit ``new_frame``.

    Each axis keeps the relationship the author gave it: a caption centred on the panel
    stays centred, a rule spanning the panel spans the new one, and a caption hugging an
    edge keeps its distance from that edge. That is what lets a panel grow to half or all
    of the page and take its title and axis captions with it.

    ``box`` and both frames are ``(x, y, w, h)`` in EMU; so is the result.
    """
    x, w = _reanchor_axis((box[0], box[2]), (old_frame[0], old_frame[2]),
                          (new_frame[0], new_frame[2]))
    y, h = _reanchor_axis((box[1], box[3]), (old_frame[1], old_frame[3]),
                          (new_frame[1], new_frame[3]))
    return x, y, w, h


# A shape covering this much of the panel is a rule or a full-width caption: it stretches
# with the panel. One whose centre sits this close to the panel's stays centred on it.
_SPAN_SHARE, _CENTRED_SHARE = 0.60, 0.15


def _reanchor_axis(box, old, new) -> Tuple[int, int]:
    """One axis of :func:`_reanchor`: ``(lo, size)`` for the shape in the new frame."""
    lo, size = box
    o_lo, o_size = old
    n_lo, n_size = new
    d_lo, d_hi = lo - o_lo, (lo + size) - (o_lo + o_size)
    if size >= _SPAN_SHARE * o_size:            # spans the panel → stretches with it
        return n_lo + d_lo, max(1, (n_lo + n_size + d_hi) - (n_lo + d_lo))
    off = (lo + size / 2) - (o_lo + o_size / 2)
    if abs(off) <= _CENTRED_SHARE * o_size:     # centred on the panel → stays centred
        return int(n_lo + n_size / 2 + off - size / 2), size
    if abs(d_lo) <= abs(d_hi):                  # hugs the near edge → keeps its distance
        return n_lo + d_lo, size
    return n_lo + n_size + d_hi - size, size


def _fill_rank_scatter(chart, points: List[Dict[str, Any]], country: str) -> None:
    """Refill one country's quadrant panel IN PLACE with its lines of business.

    The author's chart is kept, not rebuilt: ``replace_data`` leaves every bit of series
    formatting alone, so the markers, the axis rules and the painted bands behind them stay
    exactly as drawn. Each point is a line of business at (size of the Marsh pool, the
    carrier's rank in it), labelled with the line's name — the labels the author hand-placed
    for their example points come off in :func:`_strip_quadrant_furniture`.

    Both axes are re-scaled to the country's own numbers: the money axis to the largest pool
    (its ticks now shown, since the hand-drawn broken axis goes), the rank axis to the
    author's own 0–11 — which is what puts their band boundary on the top-5 line — widened
    only if the carrier ranks deeper than that.
    """
    from pptx.chart.data import XyChartData
    from pptx.util import Pt

    data = XyChartData()
    series = data.add_series(str(country or ""))
    for p in points:
        series.add_data_point(float(p["size"]), float(p["rank"]))
    _replace_chart_data(chart, data)

    size_axis = chart.category_axis                 # the x (money) axis: both are value axes
    size_axis.minimum_scale = 0.0
    size_axis.maximum_scale = max(float(p["size"]) for p in points) * _LC_SIZE_HEADROOM
    size_axis.tick_labels.number_format = _LC_SIZE_FORMAT
    size_axis.tick_labels.number_format_is_linked = False
    size_axis.tick_labels.font.size = Pt(_LC_LABEL_PT)
    _show_tick_labels(size_axis)

    rank_axis = chart.value_axis                    # the y axis: rank 1 at the top (maxMin)
    rank_axis.minimum_scale = 0.0
    rank_axis.maximum_scale = max(_LC_RANK_AXIS_MAX,
                                  max(float(p["rank"]) for p in points) + 1.0)

    for point, row, side in zip(chart.plots[0].series[0].points, points,
                                _label_sides(points, size_axis.maximum_scale)):
        _name_point(point, str(row["name"]), side)
    chart.has_legend = False
    chart.has_title = False                         # the panel's own title names the country


# A point past this share of the money axis gets its name to the LEFT — a right-hand label
# would run off the plot. Lines of business sharing a rank stack their names instead.
_LC_INSET_SHARE = 0.6


def _label_sides(points: List[Dict[str, Any]], x_max: float) -> List[Any]:
    """Where each point's name sits, so no two names land on each other or off the plot.

    Rank is an integer, so several lines of business routinely share one row: the biggest
    pool of a row keeps the side its position calls for, and the rest stack above and below
    it. Deterministic — ``points`` arrive biggest pool first.
    """
    from pptx.enum.chart import XL_LABEL_POSITION as POS

    at_rank: Dict[int, int] = {}
    sides = []
    for p in points:
        rank = int(p["rank"])
        nth = at_rank[rank] = at_rank.get(rank, -1) + 1
        if nth:                                     # a line is already labelled on this row
            sides.append(POS.ABOVE if nth % 2 else POS.BELOW)
        elif float(p["size"]) > _LC_INSET_SHARE * x_max:
            sides.append(POS.LEFT)
        else:
            sides.append(POS.RIGHT)
    return sides


def _show_tick_labels(axis) -> None:
    """Turn an axis's tick labels back on (``c:tickLblPos``).

    The author hid this axis and drew its ticks by hand as text boxes, faking a break in the
    scale. Those go with the rest of the example furniture, so the axis must state its own.
    """
    from pptx.oxml.ns import qn

    pos = axis._element.find(qn("c:tickLblPos"))
    if pos is not None:
        pos.set("val", "nextTo")


def _name_point(point, text: str, side) -> None:
    """Label one scatter point with its line of business, on the given ``side`` of it.

    Written per point rather than as a plot-wide setting: a scatter plot carries no
    ``c:dLbls`` of its own in python-pptx, and the label text is the category name, which an
    XY series does not have — so each point states its own.
    """
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    label = point.data_label
    label.position = side
    label.text_frame.text = text
    label.text_frame.paragraphs[0].font.size = Pt(_LC_LABEL_PT)
    # The label IS the name: without this PowerPoint would append the plotted rank to it.
    show_val = label._dLbl.find(qn("c:showVal"))
    if show_val is not None:
        show_val.set("val", "0")


def _box(shape) -> Tuple[int, int, int, int]:
    return int(shape.left or 0), int(shape.top or 0), int(shape.width or 0), int(shape.height or 0)


def _panel_furniture(slide, charts) -> Dict[int, List[Any]]:
    """``{chart shape id: the loose shapes that belong to its panel}``.

    Decided for the whole slide up front, off the frames the AUTHOR drew: the panels are
    about to be re-laid-out, and once one has moved, "which panel is this caption nearest"
    no longer has the answer the page was authored with.
    """
    owned: Dict[int, List[Any]] = {int(sh.shape_id): [] for sh in charts}
    for sh in slide.shapes:
        if getattr(sh, "has_chart", False):
            continue
        owner = _panel_owner(sh, charts)
        if owner is not None:
            owned[int(owner.shape_id)].append(sh)
    return owned


def _plot_frame(chart, box) -> Optional[Tuple[int, int, int, int]]:
    """Where the plot area sits inside a chart occupying ``box``, or None if not pinned.

    The author lays the plot out manually, as FRACTIONS of the chart frame — so the plot's
    own rectangle follows from the frame, whatever size the panel is given.
    """
    from pptx.oxml.ns import qn

    layout = chart._chartSpace.plotArea.find(qn("c:layout"))
    manual = layout.find(qn("c:manualLayout")) if layout is not None else None
    if manual is None:
        return None
    try:
        f = [float(manual.find(qn(f"c:{k}")).get("val")) for k in ("x", "y", "w", "h")]
    except (AttributeError, TypeError, ValueError):
        return None
    x, y, w, h = box
    return int(x + f[0] * w), int(y + f[1] * h), int(f[2] * w), int(f[3] * h)


# How close a shape must sit to the plot area to count as drawn ON it, as a share of the
# panel. The painted bands are authored a rounding error off the plot's own edges.
_LC_SNAP_SHARE = 0.03


def _snaps_to_plot(box, plot) -> bool:
    """True when ``box`` is the author's plot-area backdrop rather than a loose caption."""
    tolerance = _LC_SNAP_SHARE * max(plot[2], plot[3])
    return all(abs(a - b) <= tolerance for a, b in zip(box, plot))


def _move_panel_furniture(furniture: List[Any], old, new, *, old_plot=None, new_plot=None) -> None:
    """Carry a panel's bands, title and axis captions across to its new frame.

    Anything the author drew ON the plot area — the painted priority bands — is re-placed on
    the NEW plot area rather than scaled with the frame: the plot's insets are a fraction of
    the chart, so a panel grown to the whole page would otherwise leave its bands sitting a
    third of an inch adrift of the points they sit behind.
    """
    if old == new:
        return
    for sh in furniture:
        try:
            box = _box(sh)
            if old_plot and new_plot and _snaps_to_plot(box, old_plot):
                sh.left, sh.top, sh.width, sh.height = new_plot
            else:
                sh.left, sh.top, sh.width, sh.height = _reanchor(box, old, new)
        except (TypeError, ValueError):         # a shape that cannot be placed → leave it
            continue


def _fill_ranking_panels(slide, sidx: int, panels: Dict[str, Any]) -> set:
    """Fill every LC-ranking panel on one slide; return the chart shape ids handled.

    Done a slide at a time rather than a chart at a time because the panels share the page:
    which loose shape belongs to which panel — and how much of the page each panel gets
    once the empty ones are dropped — can only be decided with all of them in hand.

    Each live panel is moved to its planned frame (with its bands, title and axis captions
    following it), then refilled in place; a panel with no country left comes off whole.
    """
    charts = [sh for sh in _iter_leaves(slide.shapes)
              if getattr(sh, "has_chart", False) and f"{sidx}:{int(sh.shape_id)}" in panels]
    if not charts:
        return set()
    handled = {int(sh.shape_id) for sh in charts}
    blank = {int(sh.shape_id) for sh in charts
             if not panels[f"{sidx}:{int(sh.shape_id)}"]["points"]}
    _strip_quadrant_furniture(slide, charts, blank=blank)
    furniture = _panel_furniture(slide, charts)
    for sh in charts:
        panel = panels[f"{sidx}:{int(sh.shape_id)}"]
        try:
            if int(sh.shape_id) in blank:
                # No country left to report: the panel comes off whole, rather than
                # shipping the template's example book under a title that has been erased.
                sh._element.getparent().remove(sh._element)
                continue
            rect = panel["rect"]
            old, new = _box(sh), (rect["x"], rect["y"], rect["w"], rect["h"])
            _move_panel_furniture(furniture[int(sh.shape_id)], old, new,
                                  old_plot=_plot_frame(sh.chart, old),
                                  new_plot=_plot_frame(sh.chart, new))
            sh.left, sh.top, sh.width, sh.height = new
            _fill_rank_scatter(sh.chart, panel["points"], panel.get("country") or "")
        except Exception as exc:  # noqa: BLE001 — a panel must never break the export
            logger.warning("template_fill: ranking panel skipped: %s", exc)
    logger.info("template_fill: %d ranking panel(s) on slide %d (%d blanked)",
                len(charts), sidx, len(blank))
    return handled


def _write_bar_chart(chart, series: Dict[str, Any]) -> None:
    """Write a CY-vs-PY category chart (``{categories, cy, py}``) into a clustered bar.

    The series NAMES come from the template's own authored series (``CY``/``PY``) so the
    legend, colours and ordering the author chose are preserved. The values arrive already
    scaled to the unit the page's caption declares ("GWP LoB (€M)"), so the labels state
    that unit's own plain format rather than the raw-currency scaling code the author's
    example data carried.
    """
    from pptx.chart.data import CategoryChartData

    authored = [str(s.name or "") for s in chart.series][:2]
    names = (authored + ["CY", "PY"])[:2]
    data = CategoryChartData(number_format=_MILLIONS_FORMAT)
    data.categories = [str(c) for c in series["categories"]]
    for name, key in zip(names, ("cy", "py")):
        data.add_series(name, [float(v) for v in series[key]])
    _replace_chart_data(chart, data)
    _pin_value_axis_format(chart, _MILLIONS_FORMAT)


def _fill_charts(prs, values: Dict[str, Any]) -> None:
    """Best-effort: write the computed data into the deck's charts.

      * **bubble** — the carrier-vs-Marsh growth scatter. Think-cell/externally linked
        bubbles are first detached from the link (native PowerPoint supports bubble
        charts, so the refilled chart stands alone);
      * **quadrant panels** — the LC-ranking page's per-country panels, addressed by
        ``slide:shape`` from the ``lc_ranking`` payload
        (:mod:`studio.template_fill.lc_page` decides which panel is which country, and how
        much of the page each gets). The authored scatter is refilled in place and moved to
        its planned frame; only the hand-drawn axis and dummy point labels come off;
      * **scatter** — the growth data, but only when NOT externally linked and not an
        LC-ranking panel;
      * **clustered bars** — the GWP-performance page's CY-vs-PY breakdowns, addressed by
        ``slide:shape`` from the ``gwp_bars`` payload
        (:mod:`studio.template_fill.gwp_page` decides which chart is which dimension).

    Off-switch: ``STUDIO_FILL_CHARTS=off``. Any per-chart failure is swallowed so a
    chart never breaks the export.
    """
    if os.getenv("STUDIO_FILL_CHARTS", "auto").strip().lower() in {"off", "0", "false", "no"}:
        return
    points = _bubble_points((values.get("growth_bubble") or {}).get("points") or [])
    bars = values.get("gwp_bars") or {}
    ranking = values.get("lc_ranking") or {}
    if not points and not bars and not ranking:
        return
    from pptx.chart.data import XyChartData

    for sidx, slide in enumerate(prs.slides):
        # The ranking panels rewrite the slide (their own furniture, and the chart shape
        # itself), so they run first and the plain per-chart pass skips what they took.
        done = _fill_ranking_panels(slide, sidx, ranking) if ranking else set()
        for sh in list(_iter_leaves(slide.shapes)):
            if not getattr(sh, "has_chart", False) or int(sh.shape_id) in done:
                continue
            chart = sh.chart
            ctype = str(chart.chart_type)
            series = bars.get(f"{sidx}:{int(sh.shape_id)}")
            try:
                if series:
                    if not _detach_external_data(chart):
                        continue
                    _write_bar_chart(chart, series)
                elif not points:
                    continue
                elif "BUBBLE" in ctype:
                    if not _detach_external_data(chart):
                        continue
                    _write_bubble_chart(chart, points)
                    _blank_stale_point_labels(slide, sh, points)
                elif "SCATTER" in ctype:
                    if _chart_is_external(chart):
                        continue
                    data = XyChartData()
                    ser = data.add_series("Carrier vs Marsh growth")
                    for p in points:
                        ser.add_data_point(p["marsh_yoy"] / 100.0, p["carrier_yoy"] / 100.0)
                    _replace_chart_data(chart, data)
            except Exception as exc:  # noqa: BLE001 — a chart must never break the export
                logger.warning("template_fill: chart fill skipped (%s): %s", ctype, exc)


# ── hidden / reordered slides ────────────────────────────────────────────────


def _apply_order(prs, order: List[int], hidden: List[int]) -> None:
    """Drop hidden slides and reorder the rest by editing the slide-id list."""
    keep = [i for i in order if i not in set(hidden)]
    if keep == list(range(len(prs.slides))):
        return
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for sld in ids:
        lst.remove(sld)
    for i in keep:
        if 0 <= i < len(ids):
            lst.append(ids[i])


# ── entry point ──────────────────────────────────────────────────────────────


def fill_template(doc: Dict[str, Any], *, out_path: Optional[str] = None) -> str:
    """Write the filled template to ``out_path`` (default: cwd) and return the path."""
    from pptx import Presentation

    prs = Presentation(doc["template_path"])
    fields = materialize_fields(doc)
    values = doc.get("values", {})
    width_emu = int(getattr(prs, "slide_width", 0) or 0)
    height_emu = int(getattr(prs, "slide_height", 0) or 0)
    subs = _label_subs(values)

    by_slide: Dict[int, List[Dict[str, Any]]] = {}
    for fld in fields.values():
        if fld["filled"]:
            by_slide.setdefault(fld["slide_idx"], []).append(fld)

    for sidx, slide in enumerate(prs.slides):
        # Labels first, values second: a slot's own write is already fully resolved, so the
        # literal substitutions must not run over it afterwards (a written "TTM April 2025"
        # would otherwise be shifted again by the template-year rule).
        _apply_subs(slide, subs)
        index = _index_by_id(slide)
        fields = by_slide.get(sidx, [])
        ink = _authored_ink(
            frame for fld in fields if _is_commentary_role(str(fld.get("role") or ""))
            for frame in [_commentary_frame(index.get(int(fld["shape_id"])), fld["where"])]
            if frame is not None)
        for fld in fields:
            shape = index.get(int(fld["shape_id"]))
            if shape is None:
                continue
            where, text = fld["where"], str(fld["text"])
            role = str(fld.get("role") or "")
            # Commentary owns its whole box/cell: it lays its points out as bulleted
            # paragraphs, so it replaces the single-paragraph write rather than following it.
            if _is_commentary_role(role):
                _write_commentary(shape, where, text, ink)
                continue
            if where and where[0] == "para":
                _write_text_shape(shape, where, text)
            elif where and where[0] == "cell":
                _write_table(shape, where, text)
            if role.endswith(_SIGNED_ROLE_SUFFIXES) and text.strip():
                _recolor_by_sign(shape, where, text)
        _add_elements(slide, doc.get("added", {}).get(str(sidx), []), width_emu, height_emu)

    _fill_charts(prs, values)

    n = len(prs.slides)
    _apply_order(prs, doc.get("order", list(range(n))), doc.get("hidden", []))

    if not out_path:
        out_path = str(Path.cwd() / "qbr_filled.pptx")
    prs.save(out_path)
    logger.info("template_fill: exported -> %s", out_path)
    return out_path
