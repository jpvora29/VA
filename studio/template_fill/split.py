"""Split a single deck into the fixed per-axis sub-templates.

A bootstrap/dev tool: the *final* sub-templates are author-made, but until then this carves
the existing ``qbr_template.pptx`` into runnable ``overall`` / ``country`` (/ ``product``)
sub-decks so the split→fill→merge pipeline works end-to-end today.

It is the inverse of :mod:`studio.template_fill.merge` — instead of appending whole decks it
extracts a chosen subset of slides into a fresh deck, reusing the very same part-graph clone
(so think-cell OLE + native charts are preserved exactly).

Section boundaries come from :func:`studio.template_fill.sections.classify_sections` (no
hard-coded indices): everything before the first ``country_divider`` is ``overall``; the
first country block (divider → its following slides up to the next divider) is ``country``.
"""
from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Sequence, Set, Tuple

from pptx import Presentation

from logger import get_logger
from studio.template_fill.analyze import analyze
from studio.template_fill.merge import _append_slide
from studio.template_fill.sections import Section, classify_sections

logger = get_logger(__name__)

# Section types that belong to a per-country block (filled once per selected country).
_COUNTRY_BLOCK = {
    Section.COUNTRY_DIVIDER,
    Section.BREAKDOWN,
    Section.CARRIER_TITLE,
    Section.GROWTH,
}


def extract_slides(src_path: str, indices: Sequence[int]) -> Presentation:
    """A fresh |Presentation| holding only ``indices`` of ``src_path`` (in that order)."""
    src = Presentation(src_path)
    dst = Presentation()
    dst.slide_width = src.slide_width
    dst.slide_height = src.slide_height
    # A blank Presentation() ships zero slides; if that ever changes, start clean.
    for _ in range(len(dst.slides._sldIdLst)):
        del dst.slides._sldIdLst[0]

    cloned: Dict[int, object] = {}
    reserved: Set[str] = set()
    shared: Dict[Tuple[str, str], object] = {}
    for i in indices:
        _append_slide(dst, src.slides[i], cloned, reserved, shared)
    return dst


def extract_to_file(src_path: str, indices: Sequence[int], out_path: str) -> str:
    prs = extract_slides(src_path, indices)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    logger.info("split: %d slide(s) %s -> %s", len(indices), list(indices), out_path)
    return out_path


def plan_axes(src_path: str) -> Dict[str, List[int]]:
    """``{axis_name: [slide indices]}`` for ``src_path`` from its section structure.

    ``overall`` = every slide before the first country divider.
    ``country`` = the first country block (its divider through the slide before the next
    divider), keeping only the per-country section types (drops trailing blanks).
    """
    template = analyze(src_path)
    secs = classify_sections(template)
    n = len(template.slides)

    dividers = [i for i in range(n) if secs.get(i) == Section.COUNTRY_DIVIDER]
    axes: Dict[str, List[int]] = {}

    first_divider = dividers[0] if dividers else n
    axes["overall"] = list(range(first_divider))

    if dividers:
        start = dividers[0]
        end = dividers[1] if len(dividers) > 1 else n
        axes["country"] = [i for i in range(start, end) if secs.get(i) in _COUNTRY_BLOCK]
    return axes


def split_template(src_path: str, out_dir: str = "template") -> Dict[str, str]:
    """Write each planned axis to ``out_dir/<axis>_template.pptx``; return ``{axis: path}``.

    The dict is keyed by axis name (``overall``/``country``); the file it points at carries
    the ``_template`` suffix to match the author's naming convention.
    """
    axes = plan_axes(src_path)
    out: Dict[str, str] = {}
    for axis, indices in axes.items():
        if not indices:
            continue
        out[axis] = extract_to_file(src_path, indices, str(Path(out_dir) / f"{axis}_template.pptx"))
    return out


def main(argv: Sequence[str]) -> int:
    src = argv[0] if argv else "template/qbr_template.pptx"
    out_dir = argv[1] if len(argv) > 1 else "template"
    written = split_template(src, out_dir)
    for axis, path in written.items():
        print(f"{axis:<8} -> {path}")
    if not written:
        print("no axes found (no country divider?) — nothing written")
        return 1
    return 0


if __name__ == "__main__":
    import sys

    raise SystemExit(main(sys.argv[1:]))
