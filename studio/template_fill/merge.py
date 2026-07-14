"""Merge engine — concatenate several filled ``.pptx`` files into one deck.

The split-template model fills small author-made sub-decks separately (``overall``,
one ``product`` deck per product, one ``country`` deck per country) and then stitches
them, in order, into a single presentation. python-pptx has no native "append the
slides of deck B onto deck A", so this does it at the OPC layer.

How it works (and why it survives think-cell / native charts):
  * a slide is copied by cloning its *part graph* — the slide part plus every part it
    reaches (layout → master → theme, images, charts + their embedded workbooks, and
    think-cell OLE objects) — into the destination package;
  * each clone keeps the **original bytes** of its source part, and its relationships are
    rebuilt with the **same rId keys**, so the in-XML ``r:embed`` / ``r:id`` references
    inside the copied blob stay valid without ever editing the XML — nothing is
    re-serialized, so OLE objects and externally-linked charts are not disturbed;
  * fresh, collision-free partnames are allocated in the destination so two decks that
    both ship ``/ppt/slides/slide1.xml`` don't clash.

Entry points:
    ``merge_pptx(paths) -> Presentation``     stitch in order, return the live object
    ``merge_to_file(paths, out_path) -> str``  …and save it
"""
from __future__ import annotations

import hashlib
import os
import re
import shutil
import sys
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part, _Relationship
from pptx.opc.packuri import PackURI

from logger import get_logger

logger = get_logger(__name__)

# A slide part's own forward relationship — every appended slide gets exactly one of
# these from the destination presentation part.
_SLIDE_RELTYPE = RT.SLIDE

_TRAILING_INT = re.compile(r"^(.*?)(\d+)(\.[^.]+)$")


def _name_template(partname: PackURI) -> str:
    """A printf template for new partnames in the same folder/extension as ``partname``.

    ``/ppt/media/image3.png`` → ``/ppt/media/image%d.png``;
    ``/ppt/slides/slide1.xml`` → ``/ppt/slides/slide%d.xml``.
    """
    s = str(partname)
    m = _TRAILING_INT.match(s)
    if m:
        return f"{m.group(1)}%d{m.group(3)}"
    base, dot, ext = s.rpartition(".")
    return f"{base}%d.{ext}" if dot else f"{s}%d"


def _alloc_partname(package, template: str, reserved: Set[str]) -> PackURI:
    """Next free partname for ``template`` considering BOTH the package and names this
    merge has already handed out but not yet linked into the part graph.

    ``package.next_partname`` only sees parts reachable through relationships, so during
    the clone recursion freshly-created (not-yet-linked) siblings are invisible to it and
    would be assigned the same name — hence the explicit ``reserved`` set.
    """
    existing = {str(p.partname) for p in package.iter_parts()} | reserved
    n = 1
    while (template % n) in existing:
        n += 1
    name = template % n
    reserved.add(name)
    return PackURI(name)


def _clone_part(
    package,
    src_part: Part,
    cloned: Dict[int, Part],
    reserved: Set[str],
    shared: Dict[Tuple[str, str], Part],
) -> Part:
    """Deep-clone ``src_part`` (and everything it reaches) into ``package``.

    Returns the destination clone. Two levels of de-duplication:
      * ``cloned`` (per source deck, by object identity) — a part reused by several slides
        of the same deck (a master, a theme) is copied once;
      * ``shared`` (whole merge, by content hash) — a **leaf** part with no relationships
        (an image, a font, an embedded blob) that is byte-identical across sub-decks is
        stored ONCE. This is what stops a 3 MB template image from being copied per product
        block. Only relationship-free parts are shared, so nothing with divergent rels (a
        master/layout) is ever wrongly merged.
    """
    memo = cloned.get(id(src_part))
    if memo is not None:
        return memo

    leaf = len(src_part.rels) == 0
    if leaf:
        digest = (src_part.content_type, hashlib.sha1(src_part.blob).hexdigest())
        existing = shared.get(digest)
        if existing is not None:
            cloned[id(src_part)] = existing
            return existing

    dst_part = Part(
        _alloc_partname(package, _name_template(src_part.partname), reserved),
        src_part.content_type,
        package,
        src_part.blob,            # original bytes — never re-serialized
    )
    cloned[id(src_part)] = dst_part
    if leaf:
        shared[digest] = dst_part

    base_uri = dst_part.partname.baseURI
    dst_rels = dst_part.rels._rels  # the {rId: _Relationship} backing dict
    for rId, rel in src_part.rels.items():
        if rel.is_external:
            dst_rels[rId] = _Relationship(base_uri, rId, rel.reltype, RTM.EXTERNAL, rel.target_ref)
        else:
            target = _clone_part(package, rel.target_part, cloned, reserved, shared)
            dst_rels[rId] = _Relationship(base_uri, rId, rel.reltype, RTM.INTERNAL, target)
    return dst_part


def _append_slide(
    base: Presentation,
    src_slide,
    cloned: Dict[int, Part],
    reserved: Set[str],
    shared: Dict[Tuple[str, str], Part],
) -> None:
    """Clone one source slide's part graph into ``base`` and register it as a new slide."""
    dst_slide_part = _clone_part(base.part.package, src_slide.part, cloned, reserved, shared)
    rId = base.part.relate_to(dst_slide_part, _SLIDE_RELTYPE)
    base.slides._sldIdLst.add_sldId(rId)


_SLD_LAYOUT_ID = re.compile(rb'(<p:sldLayoutId[^>]*\bid=")(\d+)(")')


def _renumber_layout_ids(base: Presentation) -> None:
    """Give every ``p:sldLayoutId`` in the merged package a document-unique id.

    Each source deck numbers its layout ids from 2147483649, so the cloned masters
    collide with the base's — one of the invalidities that made PowerPoint refuse
    the merged file. The ids live only inside each master's ``sldLayoutIdLst``
    (references are by r:id), so a byte-level renumber of the master blobs is safe.
    """
    masters = [p for p in base.part.package.iter_parts()
               if p.content_type.endswith("slideMaster+xml")]
    used = [int(m.group(2)) for p in masters for m in _SLD_LAYOUT_ID.finditer(p.blob)]
    counter = max(used, default=2147483648) + 1
    for part in masters:
        # Only CLONED masters are raw ``Part``s holding literal blob bytes; the base
        # deck's own masters are typed XmlParts whose ids are already consistent.
        if type(part) is not Part:
            continue

        def fresh(m):
            nonlocal counter
            out = m.group(1) + str(counter).encode() + m.group(3)
            counter += 1
            return out

        part._blob = _SLD_LAYOUT_ID.sub(fresh, part.blob)


def _register_cloned_masters(base: Presentation) -> None:
    """Register every reachable slide/notes master in the presentation's master lists.

    Cloned slides bring their layout → master part-graphs across, but a master that is
    only *reachable* (via slide rels) and not LISTED in ``p:sldMasterIdLst`` /
    ``p:notesMasterIdLst`` makes the package invalid to PowerPoint ("could not open
    the file" / repair prompt) even though python-pptx reads it happily. Walk the
    merged package and add a presentation-level relationship + id-list entry for any
    master PowerPoint doesn't know about yet.
    """
    from pptx.oxml.ns import qn

    pres_part = base.part
    pres_elm = pres_part._element

    def listed(list_tag: str, id_tag: str, reltype: str, content_type_suffix: str) -> None:
        lst = pres_elm.find(qn(list_tag))
        known = {id(rel.target_part) for rel in pres_part.rels.values() if rel.reltype == reltype}
        masters = [p for p in pres_part.package.iter_parts()
                   if p.content_type.endswith(content_type_suffix) and id(p) not in known]
        if not masters:
            return
        if lst is None:
            # Schema order: sldMasterIdLst, notesMasterIdLst, … precede sldIdLst.
            lst = pres_elm.makeelement(qn(list_tag), {})
            anchor = pres_elm.find(qn("p:sldMasterIdLst"))
            pres_elm.insert(list(pres_elm).index(anchor) + 1 if anchor is not None else 0, lst)
        next_id = max([int(e.get("id")) for e in lst if e.get("id")] + [2147483647]) + 1
        for part in masters:
            rId = pres_part.relate_to(part, reltype)
            entry = lst.makeelement(qn(id_tag), {})
            if id_tag == "p:sldMasterId":                 # notesMasterId carries no id attr
                entry.set("id", str(next_id))
                next_id += 1
            entry.set(qn("r:id"), rId)
            lst.append(entry)
        logger.info("merge_pptx: registered %d %s master(s)", len(masters), id_tag)

    listed("p:sldMasterIdLst", "p:sldMasterId", RT.SLIDE_MASTER, "slideMaster+xml")
    listed("p:notesMasterIdLst", "p:notesMasterId", RT.NOTES_MASTER, "notesMaster+xml")


def merge_pptx(paths: List[str]) -> Presentation:
    """Concatenate the slides of ``paths`` (in order) into one |Presentation|.

    The first path is opened as the base; every slide of each later deck is appended.
    Raises ``ValueError`` if ``paths`` is empty.
    """
    if not paths:
        raise ValueError("merge_pptx: no input paths")
    base = Presentation(paths[0])
    reserved: Set[str] = set()

    # Seed the content-hash pool with the base deck's own leaf parts, so a media blob the
    # appended decks share with the base is stored once, not re-copied per sub-deck.
    shared: Dict[Tuple[str, str], Part] = {}
    for part in base.part.package.iter_parts():
        if len(part.rels) == 0:
            shared[(part.content_type, hashlib.sha1(part.blob).hexdigest())] = part

    for path in paths[1:]:
        src = Presentation(path)
        cloned: Dict[int, Part] = {}   # per-source-deck identity memo
        for slide in src.slides:
            _append_slide(base, slide, cloned, reserved, shared)
    _renumber_layout_ids(base)
    _register_cloned_masters(base)
    logger.info("merge_pptx: stitched %d deck(s) -> %d slide(s)", len(paths), len(base.slides._sldIdLst))
    return base


def merge_to_file(paths: List[str], out_path: str) -> str:
    """Merge ``paths`` and save the result to ``out_path`` (returned).

    Windows desktop exports prefer PowerPoint itself for the final append. That
    keeps opaque PowerPoint/vendor package parts intact, which is safer than
    asking python-pptx to re-save a complex merged package. Set
    ``STUDIO_PPT_MERGE_ENGINE=opc`` to force the old pure-python path for CI.
    """
    engine = os.getenv("STUDIO_PPT_MERGE_ENGINE", "auto").strip().lower()
    if not paths:
        raise ValueError("merge_pptx: no input paths")

    if len(paths) == 1 and engine != "opc":
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(paths[0], out_path)
        logger.info("merge_pptx: copied single deck -> %s", out_path)
        return out_path

    if engine in {"auto", "", "powerpoint", "win32", "com"} and sys.platform == "win32":
        # NO OPC fallback on Windows: for the real (think-cell) templates the OPC
        # merge currently produces a package PowerPoint refuses to open — serving it
        # corrupts the preview render and the export. Retry once (the usual failure
        # is a transiently busy COM instance), then fail loudly so the caller can
        # fall back to a valid single-template deck instead.
        try:
            return _merge_to_file_powerpoint(paths, out_path)
        except Exception as exc:  # noqa: BLE001
            logger.warning("merge_pptx: PowerPoint merge failed (%s) — retrying once", exc)
        return _merge_to_file_powerpoint(paths, out_path)
    if engine not in {"auto", "", "opc"}:
        raise ValueError(f"unknown STUDIO_PPT_MERGE_ENGINE={engine!r}")

    prs = merge_pptx(paths)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    logger.info("merge_pptx: exported -> %s", out_path)
    return out_path


def _merge_to_file_powerpoint(paths: List[str], out_path: str) -> str:
    """Append decks using local Microsoft PowerPoint COM automation."""
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("pywin32 is not available") from exc

    out = Path(out_path).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    abs_paths = [str(Path(p).resolve()) for p in paths]
    if out.exists():
        out.unlink()

    app: Optional[object] = None
    prs: Optional[object] = None
    pythoncom.CoInitialize()
    try:
        try:
            app = win32com.client.DispatchEx("PowerPoint.Application")
        except Exception:
            app = win32com.client.Dispatch("PowerPoint.Application")
        try:
            app.DisplayAlerts = 0
        except Exception:
            pass
        prs = app.Presentations.Open(abs_paths[0], WithWindow=False)
        for src in abs_paths[1:]:
            prs.Slides.InsertFromFile(src, prs.Slides.Count)
        prs.SaveAs(str(out))
        logger.info("merge_pptx: PowerPoint merged %d deck(s) -> %s", len(paths), out)
        return str(out)
    finally:
        try:
            if prs is not None:
                prs.Close()
        except Exception:
            pass
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()
