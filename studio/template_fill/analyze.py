"""Dynamic template understanding — introspect any ``.pptx`` into a ``Template``.

Pure ``python-pptx`` reading: every slide becomes a list of ``Shape`` records with
absolute geometry (EMU), kind, the text of each paragraph, table cells, and chart
type/series. Grouped shapes are flattened to their leaves with absolute geometry
(group transforms applied), so the downstream slot detector and the geometry
preview see every fillable unit regardless of nesting.

No knowledge of any particular template lives here — this is the same code for the
starter QBR deck or a template a user uploads. Results are cached by file hash in
``registry`` so analysis runs once per template.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Tuple

from logger import get_logger
from studio.template_fill import preview_assets

logger = get_logger(__name__)


# ── data model ───────────────────────────────────────────────────────────────


@dataclass
class Shape:
    """One leaf shape on a slide, with absolute geometry and its content."""

    shape_id: int                       # cNvPr id — unique within the slide
    name: str
    kind: str                           # text | table | chart | picture | ole | other
    x: int = 0                          # absolute EMU
    y: int = 0
    w: int = 0
    h: int = 0
    paragraphs: List[str] = field(default_factory=list)          # text shapes
    table: Optional[List[List[str]]] = None                       # table cells [r][c]
    chart_type: Optional[str] = None                              # chart shapes
    chart_categories: List[str] = field(default_factory=list)
    chart_series: List[Tuple[str, List[float]]] = field(default_factory=list)
    chart_external: bool = False                                   # think-cell/linked data → can't auto-fill
    fill_color: Optional[str] = None                              # solid fill hex (no #)
    font_color: Optional[str] = None                              # first-run font hex
    font_face: Optional[str] = None
    font_size_pt: Optional[float] = None
    bold: bool = False
    italic: bool = False
    align: Optional[str] = None                                   # left|center|right
    image_url: Optional[str] = None                               # Dash-served asset URL

    @property
    def text(self) -> str:
        return "\n".join(self.paragraphs)


@dataclass
class Slide:
    index: int
    layout: str
    background_url: Optional[str] = None
    shapes: List[Shape] = field(default_factory=list)

    def title(self) -> str:
        """Best-effort slide title — the first non-empty 'Title' shape, else first text."""
        for sh in self.shapes:
            if "title" in sh.name.lower() and sh.text.strip():
                return sh.text.strip()
        for sh in self.shapes:
            if sh.kind == "text" and sh.text.strip():
                return sh.text.strip()
        return ""


@dataclass
class Template:
    path: str
    width_emu: int
    height_emu: int
    slides: List[Slide] = field(default_factory=list)

    def shape(self, slide_idx: int, shape_id: int) -> Optional[Shape]:
        if 0 <= slide_idx < len(self.slides):
            for sh in self.slides[slide_idx].shapes:
                if sh.shape_id == shape_id:
                    return sh
        return None


# ── geometry: a transform maps local EMU coords → absolute slide coords ───────


def _identity(left, top, width, height):
    return int(left or 0), int(top or 0), int(width or 0), int(height or 0)


def _group_transform(parent_tf, grp):
    """Build a child→absolute transform for a group, composed with ``parent_tf``."""
    ax, ay, aw, ah = parent_tf(grp.left, grp.top, grp.width, grp.height)
    try:
        xfrm = grp._element.grpSpPr.xfrm
        ch_off = xfrm.chOff
        ch_ext = xfrm.chExt
        cox, coy = int(ch_off.x), int(ch_off.y)
        cw, chh = int(ch_ext.cx), int(ch_ext.cy)
        sx = aw / cw if cw else 1.0
        sy = ah / chh if chh else 1.0
    except Exception:  # noqa: BLE001 — a malformed group must not break analysis
        cox = coy = 0
        sx = sy = 1.0

    def tf(left, top, width, height):
        return (
            int(ax + (int(left or 0) - cox) * sx),
            int(ay + (int(top or 0) - coy) * sy),
            int(int(width or 0) * sx),
            int(int(height or 0) * sy),
        )

    return tf


# ── shape extraction ─────────────────────────────────────────────────────────


def _kind(shape) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    st = shape.shape_type
    if shape.has_table:
        return "table"
    if shape.has_chart:
        return "chart"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if st in (MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT):
        return "ole"
    if shape.has_text_frame:
        return "text"
    return "other"


def _read_table(shape) -> List[List[str]]:
    t = shape.table
    return [[cell.text for cell in row.cells] for row in t.rows]


def _read_chart(shape):
    ch = shape.chart
    ctype = str(ch.chart_type)
    cats: List[str] = []
    try:
        cats = [str(c) for c in ch.plots[0].categories]
    except Exception:  # noqa: BLE001
        cats = []
    series: List[Tuple[str, List[float]]] = []
    try:
        for s in ch.series:
            vals = [v for v in s.values if v is not None]
            series.append((str(s.name or ""), [float(v) for v in vals]))
    except Exception:  # noqa: BLE001
        series = []
    # A think-cell / linked-workbook chart cannot be safely repopulated (replace_data
    # corrupts the link), so flag it for a manual-fill cue instead of touching it. The
    # link shows up EITHER as an external relationship OR as a ``<c:externalData>``
    # element in the chart XML (think-cell uses the latter) — treat either as external.
    external = False
    try:
        external = (b"externalData" in ch.part.blob
                    or any(getattr(r, "is_external", False) for r in ch.part.rels.values()))
    except Exception:  # noqa: BLE001
        external = False
    return ctype, cats, series, external


def _walk(shapes, tf):
    """Yield ``(leaf_shape, absolute_box)`` for every leaf, recursing into groups."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk(shape.shapes, _group_transform(tf, shape))
            continue
        yield shape, tf(shape.left, shape.top, shape.width, shape.height)


# MSO_THEME_COLOR name → the clrMap attribute that points at a theme-scheme slot.
_THEME_TO_MAP = {
    "DARK_1": "tx1", "TEXT_1": "tx1", "LIGHT_1": "bg1", "BACKGROUND_1": "bg1",
    "DARK_2": "tx2", "TEXT_2": "tx2", "LIGHT_2": "bg2", "BACKGROUND_2": "bg2",
    "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
    "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
    "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink",
}


def _theme_palette(prs) -> dict:
    """Resolve the template's theme colours → ``{MSO_THEME_COLOR name: hex}``.

    Reads the theme ``clrScheme`` and the master ``clrMap`` so e.g. ``TEXT_1`` → the
    deck's navy and ``ACCENT_5`` → its yellow — the colours the preview must reproduce."""
    from lxml import etree
    from pptx.oxml.ns import qn

    out: dict = {}
    try:
        master = prs.slide_masters[0]
        theme = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
        scheme = etree.fromstring(theme.blob).find(".//" + qn("a:clrScheme"))
        raw = {}
        for child in scheme:
            tag = etree.QName(child).localname
            srgb, sysc = child.find(qn("a:srgbClr")), child.find(qn("a:sysClr"))
            if srgb is not None:
                raw[tag] = srgb.get("val")
            elif sysc is not None:
                raw[tag] = sysc.get("lastClr")
        cm = master.element.find(qn("p:clrMap"))
        cmap = dict(cm.attrib) if cm is not None else {}
        for name, attr in _THEME_TO_MAP.items():
            hexv = raw.get(cmap.get(attr, attr))
            if hexv:
                out[name] = hexv.upper()
    except Exception as exc:  # noqa: BLE001
        logger.warning("analyze: theme palette unresolved: %s", exc)
    return out


def _solid_fill(shape, palette: Optional[dict] = None) -> Optional[str]:
    try:
        from pptx.enum.dml import MSO_FILL

        if shape.fill.type != MSO_FILL.SOLID:
            return None
        fc = shape.fill.fore_color
        try:
            return str(fc.rgb)
        except Exception:  # noqa: BLE001 — a theme colour: resolve via the palette
            tc = getattr(fc, "theme_color", None)
            if tc is not None and palette:
                name = str(tc).split("(")[0].strip().split(".")[-1]
                return palette.get(name)
    except Exception:  # noqa: BLE001 — gradient/picture/none fills aren't a hex
        pass
    return None


def _first_run_style(shape):
    color = None
    face = None
    size_pt = None
    bold = False
    italic = False
    align = None
    try:
        para = shape.text_frame.paragraphs[0]
        a = para.alignment
        align = {1: "left", 2: "center", 3: "right"}.get(int(a)) if a is not None else None
        run = para.runs[0] if para.runs else None
        if run is not None:
            bold = bool(run.font.bold)
            italic = bool(run.font.italic)
            face = run.font.name
            if run.font.size is not None:
                size_pt = float(run.font.size.pt)
            if run.font.color and run.font.color.type is not None:
                color = str(run.font.color.rgb)
    except Exception:  # noqa: BLE001
        pass
    return color, face, size_pt, bold, italic, align


def _extract(shape, box, template_path: str, slide_idx: int, palette: Optional[dict] = None) -> Shape:
    x, y, w, h = box
    kind = _kind(shape)
    rec = Shape(shape_id=int(shape.shape_id), name=str(shape.name), kind=kind,
                x=x, y=y, w=w, h=h)
    rec.fill_color = _solid_fill(shape, palette)
    if kind == "table":
        rec.table = _read_table(shape)
    elif kind == "chart":
        rec.chart_type, rec.chart_categories, rec.chart_series, rec.chart_external = _read_chart(shape)
    elif kind == "picture":
        try:
            rec.image_url = preview_assets.cache_picture(
                template_path,
                slide_idx,
                rec.shape_id,
                str(getattr(shape.image, "ext", "png")),
                bytes(shape.image.blob),
            )
        except Exception as exc:  # noqa: BLE001
            if "no embedded image" not in str(exc).lower():
                logger.warning("analyze: slide %d picture %r asset extraction failed: %s",
                               slide_idx, getattr(shape, "name", "?"), exc)
    elif shape.has_text_frame:
        rec.paragraphs = [p.text for p in shape.text_frame.paragraphs]
        rec.font_color, rec.font_face, rec.font_size_pt, rec.bold, rec.italic, rec.align = _first_run_style(shape)
    return rec


# ── entry point ──────────────────────────────────────────────────────────────


def analyze(template_path: str) -> Template:
    """Introspect ``template_path`` into a ``Template`` (geometry + content)."""
    from pptx import Presentation

    p = Path(template_path)
    prs = Presentation(str(p))
    palette = _theme_palette(prs)
    backgrounds = preview_assets.ensure_slide_backgrounds(str(p), len(prs.slides))
    slides: List[Slide] = []
    for i, s in enumerate(prs.slides):
        recs: List[Shape] = []
        for shape, box in _walk(s.shapes, _identity):
            try:
                recs.append(_extract(shape, box, str(p), i, palette))
            except Exception as exc:  # noqa: BLE001 — skip an unreadable shape, keep going
                logger.warning("analyze: slide %d shape %r unreadable: %s", i, getattr(shape, "name", "?"), exc)
        bg = backgrounds[i] if i < len(backgrounds) else None
        slides.append(Slide(index=i, layout=str(getattr(s.slide_layout, "name", "")), background_url=bg, shapes=recs))
    return Template(
        path=str(p),
        width_emu=int(prs.slide_width or 12192000),
        height_emu=int(prs.slide_height or 6858000),
        slides=slides,
    )
