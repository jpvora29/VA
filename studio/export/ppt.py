"""Export a `DeckSpec` to a native, editable .pptx — optionally following a template.

Baseline strategy (robust across any template): open the template so the deck
inherits its theme (master background, colours, fonts), then lay each slide's
typed blocks out as real PowerPoint objects (text, native charts, tables) on a
blank layout. With no template, a clean 16:9 brand deck is produced.

Everything is editable in PowerPoint (no screenshots). Charts are native and built
from the same numbers shown on screen, so the export never diverges from the deck.

Incremental complex-template support lives in `template.py` (placeholder discovery
+ `LayoutMap`); this module already honours the template's theme and slide size.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from logger import get_logger
from studio.deck.model import DeckSpec, SlideSpec
from studio.export.template import TemplateProfile
from studio.page.format import money, pct
from ui.color_pallet import ColorPalette

logger = get_logger(__name__)

_EMU_PER_IN = 914_400
_MARGIN = 0.55
_HEADER_H = 1.15
_FOOTER_H = 0.4
FONT = "Arial"


def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr.lstrip("#"))


def _hex(value: Any, fallback: str) -> str:
    text = str(value or "").strip().lstrip("#")
    if len(text) == 6:
        try:
            int(text, 16)
            return text.upper()
        except ValueError:
            pass
    return fallback.lstrip("#").upper()


def _is_dark(hexstr: str) -> bool:
    value = _hex(hexstr, "FFFFFF")
    red, green, blue = (int(value[i : i + 2], 16) for i in (0, 2, 4))
    luminance = (0.2126 * red + 0.7152 * green + 0.0722 * blue) / 255
    return luminance < 0.48


def _slide_dims(prof: TemplateProfile) -> Tuple[float, float]:
    return prof.width_emu / _EMU_PER_IN, prof.height_emu / _EMU_PER_IN


def _text(
    slide, x, y, w, h, text, *, size=12, bold=False, color="1C2636", align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP, spacing=None, font=None,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font or FONT
    run.font.color.rgb = _rgb(color)
    if spacing:
        p.line_spacing = spacing
    return tb


def _rect(slide, x, y, w, h, fill_hex, *, line_hex=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _rgb(line_hex)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ── block renderers (region = x, y, w, h in inches) ──────────────────────────


def _kpis(slide, region, items, prof):
    x, y, w, h = region
    cols = 2
    gap = 0.18
    cw = (w - gap) / cols
    ch = 0.92
    for i, k in enumerate(items[:4]):
        cx = x + (i % cols) * (cw + gap)
        cy = y + (i // cols) * (ch + gap)
        _rect(slide, cx, cy, cw, ch, "FFFFFF", line_hex=prof.colors["line"])
        _text(slide, cx + 0.18, cy + 0.12, cw - 0.3, 0.3, str(k["label"]).upper(),
              size=8, bold=True, color=prof.colors["muted"])
        _text(slide, cx + 0.18, cy + 0.34, cw - 0.3, 0.4, str(k["value"]),
              size=17, bold=True, color=prof.colors["navy"])
        if k.get("delta"):
            tone = {"good": "good", "danger": "danger", "warn": "warn"}.get(k.get("tone"), "muted")
            _text(slide, cx + 0.18, cy + 0.7, cw - 0.3, 0.22, str(k["delta"]),
                  size=8.5, bold=True, color=prof.colors[tone])


def _commentary(slide, region, block, prof):
    x, y, w, h = region
    _rect(slide, x, y, w, h, "F7F9FD", line_hex=prof.colors["line"])
    _rect(slide, x, y, 0.07, h, prof.colors["blue"], radius=False)
    pad = 0.22
    _text(slide, x + pad, y + 0.16, w - 2 * pad, 0.3, "EXECUTIVE COMMENTARY",
          size=10, bold=True, color=prof.colors["navy"])
    _text(slide, x + pad, y + 0.5, w - 2 * pad, 0.9, block.headline,
          size=12, bold=True, color=prof.colors["ink"], spacing=1.1)
    cy = y + 1.5
    for pt in list(block.points)[:4]:
        tone = {"good": "good", "danger": "danger", "warn": "warn"}.get(pt.get("tone"), "blue")
        _rect(slide, x + pad, cy + 0.07, 0.1, 0.1, prof.colors[tone], radius=False)
        label = (pt.get("label") or "").strip()
        text = (label + " " if label else "") + pt.get("text", "")
        _text(slide, x + pad + 0.22, cy, w - 2 * pad - 0.22, 0.5, text,
              size=10, color=prof.colors["ink"], spacing=1.05)
        cy += 0.52
    if block.actions:
        cy += 0.05
        _text(slide, x + pad, cy, w - 2 * pad, 0.25, "RECOMMENDED ACTIONS",
              size=9, bold=True, color=prof.colors["navy"])
        cy += 0.3
        for a in list(block.actions)[:3]:
            _text(slide, x + pad, cy, w - 2 * pad, 0.4, "•  " + a, size=9.5, color=prof.colors["ink"])
            cy += 0.34


def _bar(slide, region, block, prof):
    x, y, w, h = region
    data = CategoryChartData()
    data.categories = list(block.labels)
    data.add_series("Premium", list(block.values))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data
    )
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    pts = list(series.points)
    ramp = ColorPalette.sequential(len(pts)) or ["#0B4BFF"]
    for i, point in enumerate(pts):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _rgb(ramp[i % len(ramp)])
    for ax in (chart.category_axis, chart.value_axis):
        ax.tick_labels.font.size = Pt(9)
        ax.tick_labels.font.name = FONT


def _line(slide, region, block, prof):
    x, y, w, h = region
    data = CategoryChartData()
    data.categories = list(block.labels)
    data.add_series("Performance", list(block.values))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), data
    )
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    series = chart.series[0]
    series.format.line.color.rgb = _rgb(prof.colors["blue"])
    series.format.line.width = Pt(2.25)
    for ax in (chart.category_axis, chart.value_axis):
        ax.tick_labels.font.size = Pt(8)
        ax.tick_labels.font.name = FONT
    chart.value_axis.has_major_gridlines = True


def _donut(slide, region, block, prof):
    x, y, w, h = region
    data = CategoryChartData()
    data.categories = list(block.labels)
    data.add_series("Share", list(block.values))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), data
    )
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    pts = list(chart.plots[0].series[0].points)
    ramp = ColorPalette.sequential(len(pts)) or ["#0B4BFF"]
    for i, point in enumerate(pts):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _rgb(ramp[i % len(ramp)])


def _fmt_cell(value, kind):
    if value is None or value == "":
        return "—"
    if kind == "money":
        return money(value)
    if kind == "pct":
        return pct(value)
    if kind in ("delta", "pct_signed"):
        return pct(value, signed=True)
    return str(value)


def _table(slide, region, block, prof, *, header_style=None, body_style=None):
    x, y, w, h = region
    cols = list(block.columns)
    rows = list(block.rows)
    n = len(rows) + 1
    gf = slide.shapes.add_table(n, len(cols), Inches(x), Inches(y), Inches(w), Inches(min(h, 0.34 * n)))
    table = gf.table
    header_style = header_style or {
        "font": FONT,
        "size": 9,
        "color": prof.colors["muted"],
    }
    body_style = body_style or {
        "font": FONT,
        "size": 9,
        "color": prof.colors["ink"],
    }
    for j, c in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = c["label"]
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(header_style["size"])
        para.font.bold = True
        para.font.name = header_style["font"]
        para.font.color.rgb = _rgb(header_style["color"])
        if c.get("align") == "right":
            para.alignment = PP_ALIGN.RIGHT
    for i, row in enumerate(rows, start=1):
        for j, c in enumerate(cols):
            cell = table.cell(i, j)
            cell.text = _fmt_cell(row.get(c["key"]), c.get("kind", "text"))
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(body_style["size"])
            para.font.name = body_style["font"]
            para.font.color.rgb = _rgb(body_style["color"])
            if c.get("align") == "right":
                para.alignment = PP_ALIGN.RIGHT


def _callout(slide, region, block, prof):
    x, y, w, h = region
    tone = "warn" if block.tone == "warn" else "blue"
    _rect(slide, x, y, w, h, "F7F9FD", line_hex=prof.colors["line"])
    _rect(slide, x, y, 0.07, h, prof.colors[tone], radius=False)
    _text(slide, x + 0.22, y + 0.1, w - 0.4, h - 0.2, block.text,
          size=12, bold=True, color=prof.colors["navy"], anchor=MSO_ANCHOR.MIDDLE)


def _rail(slide, region, takeaways, prof, *, heading="KEY TAKEAWAYS"):
    """Left commentary rail: heading + tone-dotted bullets."""
    x, y, w, h = region
    _rect(slide, x, y, w, h, "F6F9FF", line_hex=prof.colors["line"])
    pad = 0.2
    _text(slide, x + pad, y + 0.16, w - 2 * pad, 0.3, heading, size=10, bold=True, color=prof.colors["navy"])
    cy = y + 0.6
    for pt in list(takeaways)[:4]:
        tone = {"good": "good", "danger": "danger", "warn": "warn"}.get(pt.get("tone"), "blue")
        _rect(slide, x + pad, cy + 0.06, 0.09, 0.09, prof.colors[tone], radius=False)
        label = (pt.get("label") or "").strip()
        text = (label + " " if label else "") + pt.get("text", "")
        _text(slide, x + pad + 0.2, cy - 0.04, w - 2 * pad - 0.2, 0.7, text,
              size=10, color=prof.colors["ink"], spacing=1.05)
        cy += 0.72


def _stat_band(
    slide,
    region,
    items,
    prof,
    *,
    fill="FFFFFF",
    font=None,
    font_color=None,
    font_size=None,
    label_style=None,
    value_style=None,
    delta_style=None,
):
    x, y, w, h = region
    _rect(slide, x, y, w, h, fill, line_hex=prof.colors["line"])
    cw = w / max(len(items), 1)
    body_color = font_color or prof.colors["navy"]
    label_size = max(6.5, (font_size or 11) * 0.75)
    value_size = max(12, (font_size or 11) * 1.75)
    label_style = label_style or {
        "font": font,
        "size": label_size,
        "color": font_color or prof.colors["muted"],
    }
    value_style = value_style or {
        "font": font,
        "size": value_size,
        "color": body_color,
    }
    for i, k in enumerate(items):
        cx = x + i * cw
        if i:
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(y + 0.15), Inches(0.012), Inches(h - 0.3))
            ln.fill.solid(); ln.fill.fore_color.rgb = _rgb(prof.colors["line"]); ln.line.fill.background()
        _text(
            slide,
            cx + 0.18,
            y + 0.16,
            cw - 0.3,
            0.25,
            str(k["label"]).upper(),
            size=label_style["size"],
            bold=True,
            color=label_style["color"],
            font=label_style["font"],
        )
        _text(
            slide,
            cx + 0.18,
            y + 0.42,
            cw - 0.3,
            0.5,
            str(k["value"]),
            size=value_style["size"],
            bold=True,
            color=value_style["color"],
            font=value_style["font"],
        )
        if k.get("delta"):
            tone = {"good": "good", "danger": "danger", "warn": "warn"}.get(k.get("tone"), "muted")
            active_delta_style = delta_style or {
                "font": font,
                "size": max(7, (font_size or 11) * 0.82),
                "color": font_color or prof.colors[tone],
            }
            _text(
                slide,
                cx + 0.18,
                y + h - 0.32,
                cw - 0.3,
                0.25,
                str(k["delta"]),
                size=active_delta_style["size"],
                bold=True,
                color=active_delta_style["color"],
                font=active_delta_style["font"],
            )
        if k.get("trend_values") and h >= 0.9:
            data = CategoryChartData()
            data.categories = list(k.get("trend_labels", []))
            data.add_series("Trend", list(k["trend_values"]))
            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE,
                Inches(cx + 0.18),
                Inches(y + h - 0.28),
                Inches(max(0.4, cw - 0.34)),
                Inches(0.22),
                data,
            )
            chart = chart_frame.chart
            chart.has_legend = False
            chart.has_title = False
            chart.category_axis.visible = False
            chart.value_axis.visible = False
            chart.series[0].format.line.color.rgb = _rgb(prof.colors["blue"])
            chart.series[0].format.line.width = Pt(1.25)


def _swot(slide, region, block, prof):
    x, y, w, h = region
    gap = 0.25
    cw = (w - gap) / 2
    ch = (h - gap) / 2
    quad = [
        ("strengths", "STRENGTHS", "good"), ("weaknesses", "WEAKNESSES", "danger"),
        ("opportunities", "OPPORTUNITIES", "blue"), ("threats", "THREATS", "warn"),
    ]
    for i, (key, label, tone) in enumerate(quad):
        cx = x + (i % 2) * (cw + gap)
        cy = y + (i // 2) * (ch + gap)
        _rect(slide, cx, cy, cw, ch, "FFFFFF", line_hex=prof.colors["line"])
        _rect(slide, cx, cy, cw, 0.06, prof.colors[tone], radius=False)
        _text(slide, cx + 0.18, cy + 0.14, cw - 0.3, 0.28, label, size=11, bold=True, color=prof.colors[tone])
        ty = cy + 0.5
        for b in (getattr(block, key, []) or [])[:4]:
            _text(slide, cx + 0.22, ty, cw - 0.4, 0.4, "•  " + b, size=9.5, color=prof.colors["ink"], spacing=1.0)
            ty += 0.34


def _cards(slide, region, block, prof):
    x, y, w, h = region
    cards = list(block.cards)[:3]
    gap = 0.3
    cw = (w - gap * (len(cards) - 1)) / max(len(cards), 1)
    ch = min(h, 2.6)
    cy = y + (h - ch) / 2
    for i, c in enumerate(cards):
        cx = x + i * (cw + gap)
        tone = {"good": "good", "danger": "danger", "warn": "warn"}.get(c.get("tone"), "blue")
        _rect(slide, cx, cy, cw, ch, "FFFFFF", line_hex=prof.colors["line"])
        _rect(slide, cx, cy, cw, 0.07, prof.colors[tone], radius=False)
        _text(slide, cx + 0.2, cy + 0.22, cw - 0.4, 0.25, c.get("tag", ""), size=9, bold=True, color=prof.colors[tone])
        _text(slide, cx + 0.2, cy + 0.55, cw - 0.4, 0.5, c.get("title", ""), size=15, bold=True, color=prof.colors["navy"])
        _text(slide, cx + 0.2, cy + 1.1, cw - 0.4, ch - 1.2, c.get("body", ""), size=10.5, color=prof.colors["muted"], spacing=1.1)


def _waterfall(slide, region, block, prof):
    """Waterfall via the invisible-base stacked-column trick (python-pptx has no
    native waterfall). Bars accumulate to a final Total."""
    x, y, w, h = region
    labels = list(block.labels) + ["Total"]
    vals = [float(v) for v in block.values]
    base: List[float] = []
    delta: List[float] = []
    running = 0.0
    for v in vals:
        base.append(running)
        delta.append(v)
        running += v
    base.append(0.0)
    delta.append(running)  # total column starts at 0

    data = CategoryChartData()
    data.categories = labels
    data.add_series("base", base)
    data.add_series("delta", delta)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 50
    s_base, s_delta = plot.series[0], plot.series[1]
    s_base.format.fill.background()       # invisible stepping base
    s_base.format.line.fill.background()
    pts = list(s_delta.points)
    for i, pt in enumerate(pts):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _rgb(prof.colors["navy"] if i == len(pts) - 1 else prof.colors["blue"])
    for ax in (chart.category_axis, chart.value_axis):
        ax.tick_labels.font.size = Pt(9)
        ax.tick_labels.font.name = FONT


def _visual(slide, region, block, prof):
    if block.kind == "chart":
        if block.chart == "donut":
            _donut(slide, region, block, prof)
        elif block.chart == "waterfall":
            _waterfall(slide, region, block, prof)
        elif block.chart == "line":
            _line(slide, region, block, prof)
        else:
            _bar(slide, region, block, prof)
    elif block.kind == "table":
        _table(slide, region, block, prof)
    elif block.kind == "matrix":
        _matrix_to_ppt(slide, region, {"title": getattr(block, "title", ""), "points": block.points}, prof)
    elif block.kind == "heatmap":
        _heatmap_to_ppt(slide, region, {"title": getattr(block, "title", ""), "rows": block.rows,
                                        "columns": block.columns, "values": block.values}, prof)
    elif block.kind == "radar":
        _radar_to_ppt(slide, region, {"title": getattr(block, "title", ""), "labels": block.labels,
                                      "values": block.values}, prof)
    elif block.kind == "timeline":
        _timeline_to_ppt(slide, region, {"title": getattr(block, "title", ""), "tasks": block.tasks}, prof)


# ── slide assembly ────────────────────────────────────────────────────────────


def _add_slide(prs, prof, spec: SlideSpec, idx: int, total: int):
    blank = prof.layout_for("blank")
    layout = prs.slide_layouts[blank if blank is not None else len(prs.slide_layouts) - 1]
    slide = prs.slides.add_slide(layout)
    sw, sh = _slide_dims(prof)

    if spec.layout == "cover":
        _rect(slide, 0, 0, sw, sh, prof.colors["navy"], radius=False)
        _text(slide, 1.0, sh / 2 - 1.5, sw - 2, 0.4, spec.eyebrow, size=13, bold=True, color="5CC8FF")
        _text(slide, 1.0, sh / 2 - 1.05, sw - 2, 1.2, spec.title, size=40, bold=True, color="FFFFFF")
        _text(slide, 1.0, sh / 2 + 0.25, sw - 2, 0.5, spec.subtitle, size=16, color="9FC3E8")
        _footer(slide, prof, idx, dark=True)
        return

    if spec.layout == "divider":
        # Flat Marsh navy, no gradient; a single solid accent rule.
        _rect(slide, 0, 0, sw, sh, "000F47", radius=False)
        _text(slide, 1.0, sh / 2 - 0.85, sw - 2, 0.35, spec.eyebrow, size=13, bold=True, color="5CC8FF")
        _text(slide, 1.0, sh / 2 - 0.4, sw - 2, 1.0, spec.title, size=34, bold=True, color="FFFFFF")
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(sh / 2 + 0.62), Inches(1.4), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb("009DE0"); bar.line.fill.background()
        _footer(slide, prof, idx, dark=True)
        return

    # Header: eyebrow (+ question) + action title + accent rule.
    eyebrow = spec.eyebrow + (f"   ·   Q: {spec.question}" if spec.question else "")
    _text(slide, _MARGIN, 0.4, sw - 2 * _MARGIN, 0.28, eyebrow, size=10.5, bold=True, color=prof.colors["blue"])
    _text(slide, _MARGIN, 0.66, sw - 2 * _MARGIN, 0.7, spec.title, size=22, bold=True, color=prof.colors["navy"])
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(_MARGIN), Inches(1.42), Inches(1.0), Inches(0.04))
    rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(prof.colors["blue"]); rule.line.fill.background()

    cx, cy = _MARGIN, 1.7
    cw, ch = sw - 2 * _MARGIN, sh - 1.7 - _FOOTER_H - 0.25
    has_reco = bool(spec.recommendation)
    if has_reco:
        ch -= 0.75  # leave room for the recommendation strip
    blocks = list(spec.blocks)

    if spec.layout == "swot" and blocks:
        _swot(slide, (cx, cy, cw, ch), blocks[0], prof)
    elif spec.layout == "initiatives" and blocks:
        _cards(slide, (cx, cy, cw, ch), blocks[0], prof)
    elif spec.layout == "agenda":
        ty = cy + 0.1
        for p in list(spec.takeaways):
            _rect(slide, cx, ty, 0.5, 0.5, prof.colors["navy"])
            _text(slide, cx, ty + 0.06, 0.5, 0.4, p.get("label", ""), size=14, bold=True,
                  color="FFFFFF", align=PP_ALIGN.CENTER)
            _text(slide, cx + 0.72, ty, cw - 0.9, 0.5, p.get("text", ""), size=15, bold=True,
                  color=prof.colors["navy"], anchor=MSO_ANCHOR.MIDDLE)
            ty += 0.72
    elif spec.layout == "methodology":
        _methodology(slide, (cx, cy, cw, ch), spec, prof)
    elif spec.layout == "exec":
        kpis = next((b for b in blocks if b.kind == "kpis"), None)
        cards = next((b for b in blocks if b.kind == "cards"), None)
        if kpis:
            _stat_band(slide, (cx, cy, cw, 1.3), kpis.items, prof)
        low_y = cy + 1.55
        low_h = ch - 1.55
        gap = 0.3
        railw = (cw - gap) * 0.55
        _rail(slide, (cx, low_y, railw, low_h), spec.takeaways, prof, heading="WHAT IT MEANS")
        if cards:
            ax = cx + railw + gap
            _rect(slide, ax, low_y, cw - railw - gap, low_h, "F7F9FD", line_hex=prof.colors["line"])
            _text(slide, ax + 0.2, low_y + 0.16, cw - railw - gap - 0.4, 0.3, "PRIORITY ACTIONS", size=10, bold=True, color=prof.colors["navy"])
            ty = low_y + 0.55
            for c in cards.cards[:3]:
                _text(slide, ax + 0.2, ty, cw - railw - gap - 0.4, 0.5, "•  " + c.get("title", ""), size=10.5, color=prof.colors["ink"], spacing=1.05)
                ty += 0.5
    else:
        # content (insight / decision): the dense LayoutPlan — stat band ▸
        # (rail | primary visual) ▸ full-width secondary visual — so the slide fills.
        from studio.deck.compose import compose

        plan = compose(spec)
        ay, ah = cy, ch
        if plan.stat_band:
            _stat_band(slide, (cx, ay, cw, 1.0), list(plan.stat_band)[:4], prof)
            ay, ah = ay + 1.15, ah - 1.15
        sec_h = 1.7 if plan.secondary is not None else 0.0
        main_h = ah - (sec_h + 0.18 if sec_h else 0)
        gap = 0.35
        railw = (cw - gap) * 0.42
        _rail(slide, (cx, ay, railw, main_h), spec.takeaways, prof)
        if plan.primary is not None:
            _visual(slide, (cx + railw + gap, ay, cw - railw - gap, main_h), plan.primary, prof)
        if plan.secondary is not None:
            _visual(slide, (cx, ay + main_h + 0.18, cw, sec_h), plan.secondary, prof)

    if has_reco:
        _reco_strip(slide, (cx, cy + ch + 0.18, cw, 0.55), spec, prof)
    _footer(slide, prof, idx)


def _reco_strip(slide, region, spec, prof):
    x, y, w, h = region
    _rect(slide, x, y, w, h, "EEF3FF", line_hex=prof.colors["line"])
    _rect(slide, x, y, 0.06, h, prof.colors["blue"], radius=False)
    meta = []
    if spec.owner:
        meta.append(spec.owner)
    if spec.due_date:
        meta.append(spec.due_date)
    if spec.confidence:
        meta.append(f"{spec.confidence} confidence")
    _text(slide, x + 0.2, y + 0.06, w * 0.66, h - 0.12, f"RECOMMENDATION   {spec.recommendation}",
          size=10.5, bold=True, color=prof.colors["navy"], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, x + w * 0.67, y + 0.06, w * 0.32 - 0.2, h - 0.12, "   ·   ".join(meta),
          size=9.5, color=prof.colors["muted"], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _role_style(props, role, *, font, size, color):
    stored = ((props.get("text_styles") or {}).get(role) or {})
    return {
        "font": stored.get("font_family") or font,
        "size": float(stored.get("font_size") or size),
        "color": _hex(stored.get("font_color"), color),
    }


def _styled_commentary(
    slide,
    region,
    points,
    heading,
    prof,
    *,
    fill,
    font,
    color,
    size,
    heading_style=None,
    body_style=None,
):
    """Render every edited commentary line without the polished rail's four-item cap."""
    x, y, w, h = region
    _rect(slide, x, y, w, h, fill, line_hex=prof.colors["line"])
    pad = 0.2
    heading_style = heading_style or {"font": font, "size": max(7, size * 0.86), "color": color}
    body_style = body_style or {"font": font, "size": size, "color": color}
    _text(
        slide,
        x + pad,
        y + 0.14,
        w - 2 * pad,
        min(0.36, max(0.24, h * 0.12)),
        str(heading or "NOTES").upper(),
        size=heading_style["size"],
        bold=True,
        color=heading_style["color"],
        font=heading_style["font"],
    )
    lines = []
    for point in points or []:
        if isinstance(point, str):
            lines.append(point)
            continue
        label = str(point.get("label", "") or "").strip()
        text = str(point.get("text", "") or "").strip()
        if label or text:
            lines.append((label + " " if label else "") + text)
    _text(
        slide,
        x + pad,
        y + 0.52,
        w - 2 * pad,
        max(0.2, h - 0.64),
        "\n".join(f"•  {line}" for line in lines),
        size=body_style["size"],
        color=body_style["color"],
        font=body_style["font"],
        spacing=1.05,
    )


def _styled_recommendation(
    slide,
    region,
    props,
    prof,
    *,
    fill,
    font,
    color,
    size,
    label_style=None,
    body_style=None,
    meta_style=None,
):
    x, y, w, h = region
    _rect(slide, x, y, w, h, fill, line_hex=prof.colors["line"])
    _rect(slide, x, y, 0.06, h, prof.colors["blue"], radius=False)
    meta = [
        str(props.get(key, "") or "").strip()
        for key in ("owner", "due", "confidence")
        if str(props.get(key, "") or "").strip()
    ]
    label_style = label_style or {"font": font, "size": size, "color": color}
    body_style = body_style or {"font": font, "size": size, "color": color}
    meta_style = meta_style or {
        "font": font,
        "size": max(7, size * 0.9),
        "color": color,
    }
    label_w = min(1.35, w * 0.22)
    _text(
        slide,
        x + 0.2,
        y + 0.06,
        label_w,
        h - 0.12,
        "RECOMMENDATION",
        size=label_style["size"],
        bold=True,
        color=label_style["color"],
        font=label_style["font"],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        x + 0.2 + label_w,
        y + 0.06,
        max(0.4, w * 0.66 - label_w),
        h - 0.12,
        props.get("text", ""),
        size=body_style["size"],
        bold=True,
        color=body_style["color"],
        font=body_style["font"],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        x + w * 0.67,
        y + 0.06,
        w * 0.32 - 0.2,
        h - 0.12,
        "   ·   ".join(meta),
        size=meta_style["size"],
        color=meta_style["color"],
        font=meta_style["font"],
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _methodology(slide, region, spec, prof):
    x, y, w, h = region
    gaps = list(spec.takeaways)
    half = (len(gaps) + 1) // 2
    colw = (w - 0.4) / 2
    for ci, col in enumerate((gaps[:half], gaps[half:])):
        gx = x + ci * (colw + 0.4)
        gy = y
        for g in col:
            _rect(slide, gx, gy, 0.05, 0.5, prof.colors["warn"], radius=False)
            _text(slide, gx + 0.15, gy, colw - 0.2, 0.25, g.get("label", ""), size=11, bold=True, color=prof.colors["navy"])
            _text(slide, gx + 0.15, gy + 0.24, colw - 0.2, 0.6, g.get("text", ""), size=9.5, color=prof.colors["muted"], spacing=1.05)
            gy += 0.78


def _footer(slide, prof, idx, *, dark=False):
    sw, sh = _slide_dims(prof)
    color = "FFFFFF" if dark else prof.colors["muted"]
    _text(slide, _MARGIN, sh - 0.42, 4, 0.25, "Marsh ICG", size=9, bold=True, color=color)
    _text(slide, sw / 2 - 2.5, sh - 0.42, 5, 0.25, "Strictly Private & Confidential  ·  Source: GPR, Carrier Survey",
          size=8, color=color, align=PP_ALIGN.CENTER)
    _text(slide, sw - _MARGIN - 1, sh - 0.42, 1, 0.25, str(idx),
          size=10, bold=True, color=color, align=PP_ALIGN.RIGHT)


def export_deck(
    deck: DeckSpec, *, template_path: Optional[str] = None, out_path: Optional[str] = None
) -> str:
    """Write `deck` to a .pptx (following `template_path`'s theme if given). Returns the path."""
    prof = TemplateProfile.load(template_path)
    prs = Presentation(template_path) if template_path and Path(template_path).exists() else Presentation()
    prs.slide_width = Emu(prof.width_emu)
    prs.slide_height = Emu(prof.height_emu)

    total = len(deck.slides)
    for i, spec in enumerate(deck.slides, start=1):
        _add_slide(prs, prof, spec, i, total)

    if not out_path:
        title = (deck.meta.get("title") or "qbr").replace(" ", "_").replace("—", "-")
        out_path = str(Path.cwd() / f"{title}.pptx")
    prs.save(out_path)
    logger.info("studio: exported deck -> %s", out_path)
    return out_path


# ── canvas-layout export (widget geometry → native PowerPoint objects) ────────


def _grid_region(prof, x, y, w, h):
    """Map a widget's 12×8 grid box to an inches region on the slide."""
    sw, sh = _slide_dims(prof)
    m, top, bot = 0.4, 0.32, 0.5
    cw = (sw - 2 * m) / 12.0
    ch = (sh - top - bot) / 8.0
    return (m + x * cw, top + y * ch, w * cw, h * ch)


def _advanced_title(slide, region, props, prof):
    x, y, w, h = region
    title = str(props.get("title", "") or "")
    if not title:
        return region
    style = _role_style(
        props, "title", font=FONT, size=10, color=prof.colors["navy"]
    )
    _text(
        slide,
        x + 0.08,
        y + 0.03,
        w - 0.16,
        0.28,
        title,
        size=style["size"],
        bold=True,
        color=style["color"],
        font=style["font"],
    )
    return (x, y + 0.32, w, max(0.2, h - 0.32))


def _matrix_to_ppt(slide, region, props, prof):
    x, y, w, h = _advanced_title(slide, region, props, prof)
    _rect(slide, x, y, w, h, "FFFFFF", line_hex=prof.colors["line"])
    ax_x, ax_y = x + 0.38, y + h - 0.32
    ax_w, ax_h = w - 0.55, h - 0.55
    for ratio in (0.5,):
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(ax_x + ax_w * ratio),
            Inches(y + 0.12),
            Inches(0.008),
            Inches(ax_h),
        )
        ln.fill.solid(); ln.fill.fore_color.rgb = _rgb("AAB4C4"); ln.line.fill.background()
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(ax_x),
            Inches(y + 0.12 + ax_h * ratio),
            Inches(ax_w),
            Inches(0.008),
        )
        ln.fill.solid(); ln.fill.fore_color.rgb = _rgb("AAB4C4"); ln.line.fill.background()
    colors = ["007A78", "0B4BFF", "F2A900", "7A61D1"]
    for index, point in enumerate(props.get("points", [])):
        diameter = min(0.85, max(0.32, float(point.get("size", 30)) / 75))
        px = ax_x + float(point.get("x", 0)) / 100 * ax_w - diameter / 2
        py = y + 0.12 + (1 - float(point.get("y", 0)) / 100) * ax_h - diameter / 2
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(px),
            Inches(py),
            Inches(diameter),
            Inches(diameter),
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = _rgb(
            point.get("color", colors[index % len(colors)])
        )
        bubble.line.color.rgb = _rgb("FFFFFF")
        bubble.line.width = Pt(1)
        _text(
            slide,
            px,
            py + diameter * 0.2,
            diameter,
            diameter * 0.6,
            str(point.get("label", "")),
            size=7,
            bold=True,
            color="FFFFFF",
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    _text(slide, ax_x, ax_y, ax_w, 0.2, "Ease to win", size=7.5, color=prof.colors["muted"], align=PP_ALIGN.CENTER)


def _heatmap_to_ppt(slide, region, props, prof):
    x, y, w, h = _advanced_title(slide, region, props, prof)
    rows, cols, values = list(props.get("rows", [])), list(props.get("columns", [])), list(props.get("values", []))
    left = 0.9
    cw = (w - left) / max(len(cols), 1)
    ch = (h - 0.28) / max(len(rows), 1)
    for j, label in enumerate(cols):
        _text(slide, x + left + j * cw, y, cw, 0.24, str(label), size=7, bold=True, color=prof.colors["muted"], align=PP_ALIGN.CENTER)
    for i, label in enumerate(rows):
        _text(slide, x, y + 0.28 + i * ch, left - 0.08, ch, str(label), size=7.5, color=prof.colors["ink"], anchor=MSO_ANCHOR.MIDDLE)
        for j in range(len(cols)):
            value = float(values[i][j]) if i < len(values) and j < len(values[i]) else 0
            fill = "007A78" if value >= 80 else "56C5C1" if value >= 65 else "DDE7FF" if value >= 50 else "F1F3F7"
            _rect(slide, x + left + j * cw, y + 0.28 + i * ch, cw - 0.02, ch - 0.02, fill, radius=False)
            _text(slide, x + left + j * cw, y + 0.28 + i * ch, cw - 0.02, ch - 0.02, f"{value:.0f}", size=8, bold=True, color="FFFFFF" if value >= 65 else prof.colors["navy"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _radar_to_ppt(slide, region, props, prof):
    region = _advanced_title(slide, region, props, prof)
    x, y, w, h = region
    data = CategoryChartData()
    data.categories = list(props.get("labels", []))
    data.add_series("Risk", list(props.get("values", [])))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_FILLED, Inches(x), Inches(y), Inches(w), Inches(h), data
    )
    chart = frame.chart
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _rgb("DDE7FF")
    chart.series[0].format.line.color.rgb = _rgb(prof.colors["blue"])


def _timeline_to_ppt(slide, region, props, prof):
    x, y, w, h = _advanced_title(slide, region, props, prof)
    tasks = list(props.get("tasks", []))
    max_end = max([float(t.get("start", 0)) + float(t.get("duration", 1)) for t in tasks] or [1])
    label_w = min(1.65, w * 0.34)
    row_h = h / max(len(tasks), 1)
    tone = {"on_track": prof.colors["good"], "at_risk": prof.colors["warn"], "planned": prof.colors["blue"]}
    for i, task in enumerate(tasks):
        cy = y + i * row_h
        _text(slide, x, cy, label_w - 0.08, row_h, str(task.get("task", "")), size=7.5, color=prof.colors["ink"], anchor=MSO_ANCHOR.MIDDLE)
        start = float(task.get("start", 0))
        duration = float(task.get("duration", 1))
        bar_x = x + label_w + start / max_end * (w - label_w)
        bar_w = max(0.18, duration / max_end * (w - label_w))
        _rect(slide, bar_x, cy + row_h * 0.22, bar_w, row_h * 0.56, tone.get(str(task.get("status", "")).lower(), prof.colors["muted"]), radius=False)
        _text(slide, bar_x, cy + row_h * 0.22, bar_w, row_h * 0.56, str(task.get("owner", "")), size=6.8, color="FFFFFF", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _callout_to_ppt(slide, region, props, prof):
    x, y, w, h = region
    tone = {"good": prof.colors["good"], "warn": prof.colors["warn"]}.get(props.get("tone"), prof.colors["blue"])
    _rect(slide, x, y, w, h, "F3F6FF", line_hex=prof.colors["line"])
    _rect(slide, x, y, 0.07, h, tone, radius=False)
    _text(slide, x + 0.2, y + 0.14, w - 0.35, 0.22, str(props.get("label", "EXECUTIVE TAKEAWAY")), size=8, bold=True, color=tone)
    _text(slide, x + 0.2, y + 0.42, w - 0.35, min(0.55, h * 0.42), str(props.get("title", "")), size=15, bold=True, color=prof.colors["navy"])
    _text(slide, x + 0.2, y + 0.95, w - 0.35, max(0.2, h - 1.08), str(props.get("body", "")), size=9.5, color=prof.colors["muted"])


def _actions_to_ppt(slide, region, props, prof):
    x, y, w, h = region
    _rect(slide, x, y, w, h, "FFFFFF", line_hex=prof.colors["line"])
    _text(slide, x + 0.18, y + 0.12, w - 0.36, 0.25, str(props.get("title", "Priority decisions")), size=10, bold=True, color=prof.colors["navy"])
    items = list(props.get("items", []))
    row_h = max(0.32, (h - 0.48) / max(len(items), 1))
    tone = {"on_track": prof.colors["good"], "at_risk": prof.colors["warn"], "planned": prof.colors["blue"]}
    for i, item in enumerate(items):
        cy = y + 0.44 + i * row_h
        _rect(slide, x + 0.18, cy + 0.08, 0.08, 0.08, tone.get(str(item.get("status", "")).lower(), prof.colors["muted"]), radius=False)
        _text(slide, x + 0.34, cy, w * 0.58, row_h, str(item.get("action", "")), size=8, color=prof.colors["ink"], anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, x + w * 0.63, cy, w * 0.2, row_h, str(item.get("owner", "")), size=7.5, color=prof.colors["muted"], anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, x + w * 0.84, cy, w * 0.13, row_h, str(item.get("due", "")), size=7.5, color=prof.colors["muted"], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _widget_to_ppt(slide, prof, w, region, *, dark=False):
    kind = w["kind"]
    p = w.get("props", {})
    x, y, ww, hh = region
    font = p.get("font_family") or FONT
    size = float(
        p.get("font_size")
        or (26 if p.get("hero") else (19 if kind == "headline" else 10))
    )
    color = _hex(
        p.get("font_color"),
        "FFFFFF" if dark else prof.colors["ink"],
    )
    fill = _hex(
        p.get("background_color"),
        prof.colors["navy"]
        if dark
        else ("EEF3FF" if kind == "reco" else "FFFFFF"),
    )
    if p.get("background_color"):
        _rect(slide, x, y, ww, hh, fill, line_hex=prof.colors["line"])
    if kind == "headline":
        eyebrow_style = _role_style(
            p,
            "eyebrow",
            font=font,
            size=max(7, size * 0.42),
            color=color,
        )
        title_style = _role_style(
            p, "title", font=font, size=size, color=color
        )
        subtitle_style = _role_style(
            p,
            "subtitle",
            font=font,
            size=max(8, size * 0.45),
            color=color,
        )
        if p.get("eyebrow"):
            _text(
                slide,
                x,
                y,
                ww,
                0.25,
                p["eyebrow"],
                size=eyebrow_style["size"],
                bold=True,
                color=eyebrow_style["color"],
                font=eyebrow_style["font"],
            )
            y, hh = y + 0.28, hh - 0.28
        title_h = hh * 0.58 if p.get("subtitle") else hh
        _text(
            slide,
            x,
            y,
            ww,
            title_h,
            p.get("text", ""),
            size=title_style["size"],
            bold=True,
            color=title_style["color"],
            font=title_style["font"],
            anchor=MSO_ANCHOR.TOP,
            spacing=1.05,
        )
        if p.get("subtitle"):
            _text(
                slide,
                x,
                y + title_h,
                ww,
                max(0.2, hh - title_h),
                p["subtitle"],
                size=subtitle_style["size"],
                color=subtitle_style["color"],
                font=subtitle_style["font"],
            )
    elif kind == "text":
        if p.get("swot"):
            _swot(slide, region, type("B", (), {k: p["swot"].get(k, []) for k in ("strengths", "weaknesses", "opportunities", "threats")})(), prof)
        else:
            heading_style = _role_style(
                p,
                "heading",
                font=font,
                size=max(7, size * 0.86),
                color=color,
            )
            body_style = _role_style(
                p, "body", font=font, size=size, color=color
            )
            _styled_commentary(
                slide,
                region,
                p.get("points", []),
                p.get("heading", ""),
                prof,
                fill=fill,
                font=font,
                color=color,
                size=size,
                heading_style=heading_style,
                body_style=body_style,
            )
    elif kind in ("kpiband", "kpi"):
        label_style = _role_style(
            p, "label", font=font, size=max(6.5, size * 0.75), color=color
        )
        value_style = _role_style(
            p, "value", font=font, size=max(12, size * 1.75), color=color
        )
        delta_style = _role_style(
            p, "delta", font=font, size=max(7, size * 0.82), color=color
        )
        _stat_band(
            slide,
            region,
            list(p.get("items", []))[:6],
            prof,
            fill=fill,
            font=font,
            font_color=color,
            font_size=size,
            label_style=label_style,
            value_style=value_style,
            delta_style=delta_style,
        )
    elif kind == "chart":
        blk = type("B", (), {"kind": "chart", "chart": p.get("chart", "bar"), "labels": p.get("labels", []), "values": p.get("values", [])})()
        chart_region = region
        if p.get("title"):
            title_style = _role_style(
                p, "title", font=font, size=max(8, size), color=color
            )
            _text(
                slide,
                x,
                y,
                ww,
                0.3,
                p["title"],
                size=title_style["size"],
                bold=True,
                color=title_style["color"],
                font=title_style["font"],
            )
            chart_region = (x, y + 0.32, ww, max(0.2, hh - 0.32))
        if blk.values:
            _visual(slide, chart_region, blk, prof)
        else:
            _placeholder_box(slide, chart_region, prof, "Chart — bind data")
    elif kind == "table":
        blk = type("B", (), {"kind": "table", "columns": p.get("columns", []), "rows": p.get("rows", [])})()
        if blk.rows and blk.columns:
            header_style = _role_style(
                p, "header", font=font, size=max(8, size * 0.82), color=color
            )
            body_style = _role_style(
                p, "body", font=font, size=max(8, size * 0.82), color=color
            )
            _table(
                slide,
                region,
                blk,
                prof,
                header_style=header_style,
                body_style=body_style,
            )
        else:
            _placeholder_box(slide, region, prof, "Table — bind data")
    elif kind == "matrix":
        _matrix_to_ppt(slide, region, p, prof)
    elif kind == "heatmap":
        _heatmap_to_ppt(slide, region, p, prof)
    elif kind == "radar":
        _radar_to_ppt(slide, region, p, prof)
    elif kind == "radial":
        blk = type("B", (), {"labels": p.get("labels", []), "values": p.get("values", [])})()
        _donut(slide, _advanced_title(slide, region, p, prof), blk, prof)
    elif kind == "bridge":
        blk = type("B", (), {"labels": p.get("labels", []), "values": p.get("values", [])})()
        _waterfall(slide, _advanced_title(slide, region, p, prof), blk, prof)
    elif kind == "timeline":
        _timeline_to_ppt(slide, region, p, prof)
    elif kind == "callout":
        _callout_to_ppt(slide, region, p, prof)
    elif kind == "actions":
        _actions_to_ppt(slide, region, p, prof)
    elif kind == "reco":
        label_style = _role_style(
            p, "label", font=font, size=max(8, size * 0.82), color=color
        )
        body_style = _role_style(
            p, "body", font=font, size=size, color=color
        )
        meta_style = _role_style(
            p, "meta", font=font, size=max(7, size * 0.9), color=color
        )
        _styled_recommendation(
            slide,
            region,
            p,
            prof,
            fill=fill,
            font=font,
            color=color,
            size=size,
            label_style=label_style,
            body_style=body_style,
            meta_style=meta_style,
        )
    elif kind == "divider":
        label_style = _role_style(
            p, "label", font=font, size=size, color=color
        )
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + hh / 2), Inches(ww), Inches(0.03))
        ln.fill.solid(); ln.fill.fore_color.rgb = _rgb(prof.colors["line"]); ln.line.fill.background()
        if p.get("text"):
            _text(
                slide,
                x,
                y,
                ww,
                hh,
                p["text"],
                size=label_style["size"],
                bold=True,
                color=label_style["color"],
                font=label_style["font"],
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
    else:
        _placeholder_box(slide, region, prof, p.get("label", kind))


def _placeholder_box(slide, region, prof, label):
    x, y, w, h = region
    _rect(slide, x, y, w, h, "F4F7FE", line_hex=prof.colors["line"])
    _text(slide, x, y, w, h, label, size=12, bold=True, color=prof.colors["blue"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _add_canvas_slide(prs, prof, widgets, idx, spec=None, page_style=None):
    blank = prof.layout_for("blank")
    layout = prs.slide_layouts[blank if blank is not None else len(prs.slide_layouts) - 1]
    slide = prs.slides.add_slide(layout)
    background = _hex(
        (page_style or {}).get("background_color"),
        prof.colors["navy"]
        if spec and spec.layout in ("cover", "divider")
        else "FFFFFF",
    )
    dark = _is_dark(background)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(background)
    for w in sorted(widgets, key=lambda w: (w["y"], w["x"])):
        _widget_to_ppt(
            slide,
            prof,
            w,
            _grid_region(prof, w["x"], w["y"], w["w"], w["h"]),
            dark=dark,
        )
    _footer(slide, prof, idx, dark=dark)


def export_document(doc, *, template_path: Optional[str] = None, out_path: Optional[str] = None) -> str:
    """Export the shared document. Pages composed on the canvas (custom widget
    layout) export by widget geometry; untouched pages use the polished renderer —
    so what you arranged on the canvas is what lands in PowerPoint."""
    from studio.page import document as D

    prof = TemplateProfile.load(template_path)
    prs = Presentation(template_path) if template_path and Path(template_path).exists() else Presentation()
    prs.slide_width = Emu(prof.width_emu)
    prs.slide_height = Emu(prof.height_emu)

    hidden = set(doc.get("hidden", []))
    order = [sid for sid in doc.get("order", []) if sid not in hidden]
    meta = dict(doc.get("meta", {}))
    for i, sid in enumerate(order, start=1):
        spec = D.materialize_slide(doc, sid)
        has_page_style = sid in doc.get("page_style", {})
        if D.has_custom_layout(doc, sid) or has_page_style:
            _add_canvas_slide(
                prs,
                prof,
                D.page_widgets(doc, sid),
                i,
                spec,
                D.effective_page_style(doc, sid),
            )
        else:
            _add_slide(prs, prof, spec, i, len(order))

    if not out_path:
        title = (meta.get("title") or "qbr").replace(" ", "_").replace("—", "-")
        out_path = str(Path.cwd() / f"{title}.pptx")
    prs.save(out_path)
    logger.info("studio: exported document -> %s", out_path)
    return out_path
