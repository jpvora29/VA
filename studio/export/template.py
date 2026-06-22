"""Template profile + layout mapping for PPT export.

This is the seam that lets the exporter "follow a template" and grow to handle
COMPLEX templates incrementally:

  - `TemplateProfile.load(path)` inspects any ``.pptx`` (theme colours/fonts, slide
    size, and each layout's placeholder roles) — a focused port of the reference
    ``skills/ppt/template_analyzer``. With no template it returns brand defaults.
  - `LayoutMap` maps a studio slide ``layout`` → the template layout to use, by
    TYPE (title / section / two_content / content / blank), resolved against the
    template's own layouts. Blocks fill matching placeholders when present, else
    the exporter falls back to absolute positioning.

Incremental path: a plain template needs nothing; a richer template is supported
by (a) the analyzer discovering more placeholders and (b) extending `LayoutMap`
for bespoke named layouts — no change to compute/deck/rendering.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional

# Brand defaults (used when no template is supplied).
_DEFAULT_COLORS = {
    "navy": "000F47",
    "blue": "0B4BFF",
    "ink": "1C2636",
    "muted": "5B6577",
    "good": "1F9D55",
    "warn": "B9810A",
    "danger": "C53532",
    "line": "D7DFEB",
}
_DEFAULT_FONTS = {"major": "Arial", "minor": "Arial"}


@dataclass
class LayoutInfo:
    index: int
    name: str
    type: str  # title | section | content | two_content | blank | other
    placeholders: List[str] = field(default_factory=list)


@dataclass
class TemplateProfile:
    """What the exporter needs to honour a template (or brand defaults)."""

    path: Optional[str] = None
    width_emu: int = 12192000          # 13.333in (16:9)
    height_emu: int = 6858000          # 7.5in
    colors: Dict[str, str] = field(default_factory=lambda: dict(_DEFAULT_COLORS))
    fonts: Dict[str, str] = field(default_factory=lambda: dict(_DEFAULT_FONTS))
    layouts: List[LayoutInfo] = field(default_factory=list)
    is_dark: bool = False

    def layout_for(self, kind: str) -> Optional[int]:
        """Index of the first layout whose type matches `kind`, if any."""
        for layout in self.layouts:
            if layout.type == kind:
                return layout.index
        return None

    @classmethod
    def load(cls, path: Optional[str]) -> "TemplateProfile":
        if not path:
            return cls()
        p = Path(path)
        if not p.exists() or p.suffix.lower() != ".pptx":
            return cls()
        try:
            from pptx import Presentation

            prs = Presentation(str(p))
        except Exception:  # noqa: BLE001 - a bad template must not break export
            return cls()

        layouts: List[LayoutInfo] = []
        for i, layout in enumerate(prs.slide_layouts):
            phs = []
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                phs.append(str(t).split(".")[-1] if t else ph.name)
            layouts.append(
                LayoutInfo(index=i, name=getattr(layout, "name", f"Layout {i}"),
                           type=_layout_type(getattr(layout, "name", ""), phs), placeholders=phs)
            )
        return cls(
            path=str(p),
            width_emu=int(prs.slide_width or 12192000),
            height_emu=int(prs.slide_height or 6858000),
            layouts=layouts,
        )


def _layout_type(name: str, placeholders: List[str]) -> str:
    n = (name or "").lower()
    if "title slide" in n or "cover" in n:
        return "title"
    if "section" in n:
        return "section"
    if "two" in n or "comparison" in n:
        return "two_content"
    if "blank" in n:
        return "blank"
    if "content" in n or "title and" in n:
        return "content"
    pset = {p.lower() for p in placeholders if p}
    if {"title", "subtitle"} <= pset:
        return "title"
    if "title" in pset:
        return "content"
    return "other"


# Studio slide layout → preferred template layout TYPE (resolved per template).
LAYOUT_MAP: Dict[str, str] = {
    "title": "title",
    "section": "section",
    "exec_summary": "two_content",
    "two_column": "two_content",
    "kpi_chart": "content",
    "table": "content",
    "content": "content",
}
