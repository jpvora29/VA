"""Phase 2 — the QBR deck becomes structured JSON. No LLM calls.

Rule-based throughout: shapes, tables and charts are read straight out of the file,
sections are detected from the deck's own divider slides, and ownership (Carrier vs
Marsh) is inferred from the slide. Deterministic parsing is the point — the model
never decides what the deck says.

Granularity is the one knob: ``"slide"`` writes one JSON per slide, ``"section"``
writes one per section. :mod:`mom.modes` picks it.
"""

import json
import os
import re
import threading

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from logger import get_logger

log = get_logger(__name__)

# A paragraph shorter than this is chrome (a bullet glyph, a stray character), not content.
MIN_TEXT_LEN = 2


# ── Entry counter ─────────────────────────────────────────────────────────────
# Entry ids are sequential within one extraction. The counter is thread-local so two
# MoM runs extracting at the same time each get their own sequence rather than
# interleaving into one another's slide files.
_ids = threading.local()


def _new_entry_id() -> str:
    _ids.n = getattr(_ids, "n", 0) + 1
    return f"entry_{_ids.n:04d}"


def _reset_counter() -> None:
    _ids.n = 0


# ── Shape helpers ─────────────────────────────────────────────────────────────

def _get_position(shape) -> dict | None:
    try:
        return {
            "left_in":   round(shape.left   / 914400, 2) if shape.left   is not None else None,
            "top_in":    round(shape.top    / 914400, 2) if shape.top    is not None else None,
            "width_in":  round(shape.width  / 914400, 2) if shape.width  is not None else None,
            "height_in": round(shape.height / 914400, 2) if shape.height is not None else None,
        }
    except Exception:
        return None


def _is_title_shape(shape) -> bool:
    try:
        if shape.is_placeholder:
            pt = shape.placeholder_format.type
            return pt in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        pass
    return False


def _flatten(text: str) -> str:
    """A shape's text as one line.

    A soft line break (Shift+Enter) inside a text frame comes back from python-pptx as
    a vertical tab, so a two-line title reads as one word joined by an invisible
    control character. That character is illegal in both a worksheet and a Word
    document, and it travels a long way from here — into the slide JSON, the tagging
    prompt and the run log — before anything complains about it.
    """
    return re.sub(r"\s+", " ", (text or "").replace("\x0b", " ")).strip()


def _get_slide_title(slide) -> str | None:
    for shape in slide.shapes:
        if _is_title_shape(shape) and shape.has_text_frame:
            text = _flatten(shape.text_frame.text)
            if text:
                return text
    return None


def _normalize_numeric(text: str) -> dict | None:
    if not text:
        return None
    t = text.strip()
    m = re.match(r'^([+-]?)(\$)?([\d,]+(?:\.\d+)?)\s*(%|k|K|m|M|b|B)?$', t)
    if not m:
        return None
    sign, currency, number, suffix = m.groups()
    try:
        value = float(number.replace(",", ""))
    except ValueError:
        return None
    if sign == "-":
        value = -value
    multiplier = {"k": 1e3, "K": 1e3, "m": 1e6, "M": 1e6, "b": 1e9, "B": 1e9}.get(suffix, 1)
    if suffix not in ("%", None):
        value *= multiplier
    unit = "percent" if suffix == "%" else ("currency" if currency else ("count" if not suffix else "scaled_count"))
    return {"numeric_value": value, "unit": unit, "raw": text}


def _table_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    md = ["| " + " | ".join(rows[0]) + " |"]
    md.append("|" + "|".join(["---"] * len(rows[0])) + "|")
    for row in rows[1:]:
        md.append("| " + " | ".join(row) + " |")
    return "\n".join(md)


# ── Slide owner determination ─────────────────────────────────────────────────

def _determine_slide_owner(slide, section_title: str, carrier_name: str) -> str:
    """Apply the 5-rule cascade. Returns 'Marsh', 'Carrier', or 'Unknown'."""
    carrier_lower = (carrier_name or "").lower().strip()
    sect_lower    = (section_title or "").lower()

    if "marsh" in sect_lower or "icg" in sect_lower:
        return "Marsh"
    if carrier_lower and carrier_lower in sect_lower:
        return "Carrier"
    try:
        if "marsh" in slide.slide_layout.name.lower():
            return "Marsh"
    except Exception:
        pass
    for shape in slide.shapes:
        if _is_title_shape(shape):
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                if text.startswith("Marsh"):
                    return "Marsh"
                break
    return "Unknown"


# ── Table extraction ──────────────────────────────────────────────────────────

def _is_subheader_row(row: list[str]) -> bool:
    non_empty = [c.strip() for c in row if c.strip()]
    if not non_empty:
        return False
    for c in non_empty:
        if len(c) > 10:
            return False
        num = _normalize_numeric(c)
        if num is not None:
            try:
                iv = int(float(c.replace(",", "")))
                if not (1900 <= iv <= 2099):
                    return False
            except ValueError:
                return False
    return True


def _merge_two_level_headers(row0: list[str], row1: list[str]) -> list[str]:
    filled = []
    last = ""
    for c in row0:
        if c.strip():
            last = c.strip()
        filled.append(last)
    merged = []
    for f, s in zip(filled, row1):
        f, s = f.strip(), s.strip()
        merged.append(f"{f} {s}" if f and s else (f or s))
    return merged


def _process_table(shape, slide_number: int, slide_title: str,
                   slide_layout: str, slide_owner: str,
                   position: dict, out_list: list,
                   table_label: str = "") -> None:
    table  = shape.table
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    grid   = [[table.cell(r, c).text for c in range(n_cols)] for r in range(n_rows)]
    markdown = _table_to_markdown(grid)

    if n_rows == 1 and n_cols == 1:
        out_list.append({
            "entry_id":      _new_entry_id(),
            "slide_number":  slide_number,
            "slide_title":   slide_title,
            "slide_layout":  slide_layout,
            "slide_owner":   slide_owner,
            "content_type":  "table_label",
            "text":          grid[0][0].strip(),
            "data":          None,
            "position":      position,
            "tag": None, "umbrella_tag": None,
            "table_context": None,
            "parent_table_id": None,
        })
        return

    table_entry_id = _new_entry_id()
    out_list.append({
        "entry_id":    table_entry_id,
        "slide_number": slide_number,
        "slide_title":  slide_title,
        "slide_layout": slide_layout,
        "slide_owner":  slide_owner,
        "content_type": "table",
        "text":         markdown,
        "data":         {"rows": grid, "n_rows": n_rows, "n_cols": n_cols},
        "position":     position,
        "tag": None, "umbrella_tag": None,
        "table_context": {"all_column_headers": grid[0] if grid else [], "full_table_markdown": markdown},
        "parent_table_id": None,
    })

    if n_rows == 0 or n_cols == 0:
        return

    header_row = grid[0]
    data_start = 1
    if n_rows >= 3 and _is_subheader_row(grid[1]):
        header_row = _merge_two_level_headers(grid[0], grid[1])
        data_start = 2

    start_col   = 1 if n_cols > 1 else 0
    col_headers = [h.strip() for h in header_row[start_col:]]
    context_prefix = (table_label.strip() or slide_title or "").strip()

    for r in range(data_start, n_rows):
        row_label  = grid[r][0].strip() if n_cols > 1 else None
        row_values = grid[r][start_col:]
        if not any(v.strip() for v in row_values):
            continue

        row_cells = []
        for c_off, value in enumerate(row_values):
            v = value.strip()
            if not v:
                continue
            col_name = col_headers[c_off] if c_off < len(col_headers) else ""
            row_cells.append({
                "column_name": col_name or None,
                "value":       v,
                "numeric":     _normalize_numeric(v),
            })

        parts    = [f"{c['column_name']}: {c['value']}" if c["column_name"] else c["value"]
                    for c in row_cells]
        row_core = "; ".join(parts)
        if row_label:
            row_core = f"{row_label} — {row_core}"
        row_text = (
            f"{context_prefix}: {row_core}"
            if context_prefix and not row_label
            else row_core
        )

        out_list.append({
            "entry_id":     _new_entry_id(),
            "slide_number": slide_number,
            "slide_title":  slide_title,
            "slide_layout": slide_layout,
            "slide_owner":  slide_owner,
            "content_type": "table_row",
            "text":         row_text,
            "row_label":    row_label,
            "cells":        row_cells,
            "position":     position,
            "tag": None, "umbrella_tag": None,
            "table_context": {"all_column_headers": col_headers, "full_table_markdown": markdown},
            "parent_table_id": table_entry_id,
        })


# ── Text extraction ───────────────────────────────────────────────────────────

def _process_text_shape(shape, slide_number: int, slide_title: str,
                        slide_layout: str, slide_owner: str,
                        position: dict, out_list: list, is_title: bool) -> None:
    tf = shape.text_frame
    for para in tf.paragraphs:
        text = para.text.strip()
        if not text or len(text) < MIN_TEXT_LEN:
            continue
        out_list.append({
            "entry_id":     _new_entry_id(),
            "slide_number": slide_number,
            "slide_title":  slide_title,
            "slide_layout": slide_layout,
            "slide_owner":  slide_owner,
            "content_type": "title" if is_title else "text",
            "text":         text,
            "data":         None,
            "position":     position,
            "tag": None, "umbrella_tag": None,
        })


# ── Per-slide extractor ───────────────────────────────────────────────────────

def _extract_slide_entries(slide, slide_number: int, section_title: str,
                           carrier_name: str) -> dict:
    layout_name = ""
    try:
        layout_name = slide.slide_layout.name
    except Exception:
        pass

    title      = _get_slide_title(slide)
    slide_owner = _determine_slide_owner(slide, section_title, carrier_name)
    entries    = []

    pending_table_label = ""
    for shape in slide.shapes:
        pos      = _get_position(shape)
        is_title = _is_title_shape(shape)

        if shape.has_table:
            tbl = shape.table
            if len(tbl.rows) == 1 and len(tbl.columns) == 1:
                cell_text = tbl.cell(0, 0).text.strip()
                if cell_text:
                    pending_table_label = cell_text
                _process_table(shape, slide_number, title, layout_name, slide_owner, pos, entries,
                               table_label=pending_table_label)
            else:
                _process_table(shape, slide_number, title, layout_name, slide_owner, pos, entries,
                               table_label=pending_table_label)
                pending_table_label = ""
            continue

        try:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart = shape.chart
                chart_title = ""
                try:
                    chart_title = chart.chart_title.text_frame.text.strip()
                except Exception:
                    pass
                series_texts = []
                try:
                    for series in chart.series:
                        s_name = getattr(series, "name", "") or ""
                        for pt in series.values:
                            if pt is not None:
                                series_texts.append(f"{s_name}: {pt}")
                except Exception:
                    pass
                if chart_title or series_texts:
                    text = (chart_title + " — " if chart_title else "") + "; ".join(series_texts[:20])
                    entries.append({
                        "entry_id":     _new_entry_id(),
                        "slide_number": slide_number,
                        "slide_title":  title,
                        "slide_layout": layout_name,
                        "slide_owner":  slide_owner,
                        "content_type": "chart_point",
                        "text":         text,
                        "data":         None,
                        "position":     pos,
                        "tag": None, "umbrella_tag": None,
                    })
                continue
        except Exception:
            pass

        if shape.has_text_frame:
            _process_text_shape(shape, slide_number, title, layout_name,
                                slide_owner, pos, entries, is_title)

    return {"title": title, "layout": layout_name, "slide_owner": slide_owner, "entries": entries}


# ── Section detection ─────────────────────────────────────────────────────────

def detect_sections(per_slide: dict, n_slides: int,
                    meeting_topics: list[str]) -> list[dict]:
    """Returns [{section_title, start_slide, end_slide}]. Slides 1-2 = Introduction."""
    sections = [{"section_title": "Introduction", "start_slide": 1, "end_slide": 2}]
    if n_slides <= 2:
        return sections

    topic_words_per_topic = {}
    for t in meeting_topics:
        words = set(re.sub(r"[^a-z0-9]", " ", t.lower()).split())
        words -= {"and", "the", "of", "in", "to", "a", "an", "for", "with", "by", "from"}
        if words:
            topic_words_per_topic[t] = words

    divider_slides = []
    for slide_num in range(3, n_slides + 1):
        slide_data   = per_slide.get(slide_num, {})
        layout_name  = slide_data.get("layout", "")
        title        = slide_data.get("title") or ""
        entry_count  = len(slide_data.get("entries", []))
        text_entries = [e for e in slide_data.get("entries", [])
                        if e.get("content_type") in ("title", "text")]

        is_div = any(kw in layout_name.lower() for kw in
                     {"section", "divider", "break", "header only", "chapter"})
        if not is_div:
            has_table_entry = any(e.get("content_type") == "table"       for e in slide_data.get("entries", []))
            has_chart_entry = any(e.get("content_type") == "chart_point" for e in slide_data.get("entries", []))
            short_text_only = (
                not has_table_entry and not has_chart_entry and
                len(text_entries) <= 2 and
                all(len(e.get("text", "")) <= 30 for e in text_entries)
            )
            if short_text_only and entry_count <= 2:
                is_div = True

        if is_div:
            slide_title_words = set(re.sub(r"[^a-z0-9]", " ", title.lower()).split())
            best_topic, best_overlap = None, 0
            for t, twords in topic_words_per_topic.items():
                overlap = len(slide_title_words & twords)
                if overlap > best_overlap:
                    best_overlap = overlap
                    best_topic = t

            if not best_topic and not title:
                text_entries_on_slide = [
                    e.get("text", "").strip()
                    for e in slide_data.get("entries", [])
                    if e.get("content_type") in ("title", "text")
                    and e.get("text", "").strip()
                    and not re.fullmatch(r"0?\d{1,2}", e.get("text", "").strip())
                ]
                title = text_entries_on_slide[0] if text_entries_on_slide else ""

            raw_label  = best_topic if best_topic else title
            safe_label = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", " ", raw_label)
            safe_label = re.sub(r"\s+", " ", safe_label).strip()
            label      = safe_label if safe_label else f"Section {len(divider_slides) + 1}"
            divider_slides.append((slide_num, label))

    if not divider_slides:
        sections.append({"section_title": "Main Content", "start_slide": 3, "end_slide": n_slides})
        return sections

    divider_slides.append((n_slides + 1, "__END__"))
    for i, (div_slide, label) in enumerate(divider_slides[:-1]):
        next_div_slide = divider_slides[i + 1][0]
        content_start  = div_slide + 1
        content_end    = next_div_slide - 1
        if content_start <= content_end:
            sections.append({"section_title": label, "start_slide": content_start, "end_slide": content_end})

    return sections


# ── Output: save by SECTION (Carrier & Marsh format mode) ─────────────────────

def save_sections(per_slide: dict, sections: list[dict], output_dir: str) -> list[dict]:
    """
    Writes one JSON per section to output_dir.
    Returns manifest: [{section_number, section_title, slide_range, file, entry_count}]
    Used by the Carrier & Marsh format pipeline.
    """
    os.makedirs(output_dir, exist_ok=True)
    manifest = []

    for idx, sec in enumerate(sections, start=1):
        title = sec["section_title"]
        start = sec["start_slide"]
        end   = sec["end_slide"]

        entries = []
        for sn in range(start, end + 1):
            slide_data = per_slide.get(sn)
            if slide_data:
                entries.extend(slide_data["entries"])

        safe_title = re.sub(r"[^a-zA-Z0-9]+", "_", title).strip("_").lower()[:40]
        filename   = f"section_{idx:02d}_{safe_title}.json"
        filepath   = os.path.join(output_dir, filename)

        section_doc = {
            "section_number": idx,
            "section_title":  title,
            "slide_range":    [start, end],
            "entries":        entries,
        }
        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(section_doc, f, indent=2, ensure_ascii=False)

        manifest.append({
            "section_number": idx,
            "section_title":  title,
            "slide_range":    [start, end],
            "file":           filepath,
            "entry_count":    len(entries),
        })

    return manifest


# ── Output: save by SLIDE (Skill approach mode) ───────────────────────────────

def save_slides(per_slide: dict, sections: list[dict], output_dir: str) -> list[dict]:
    """
    Writes one JSON per slide to output_dir.
    Each file embeds its section context.
    Returns manifest: [{slide_number, slide_title, section_number, section_title, file, entry_count}]
    Used by the Skill approach pipeline.
    """
    os.makedirs(output_dir, exist_ok=True)

    slide_to_section: dict[int, dict] = {}
    for idx, sec in enumerate(sections, start=1):
        for sn in range(sec["start_slide"], sec["end_slide"] + 1):
            slide_to_section[sn] = {
                "section_number": idx,
                "section_title":  sec["section_title"],
                "slide_range":    [sec["start_slide"], sec["end_slide"]],
            }

    manifest = []
    for slide_num in sorted(per_slide.keys()):
        slide_data = per_slide[slide_num]
        entries    = slide_data.get("entries", [])
        sec_info   = slide_to_section.get(slide_num, {
            "section_number": 0,
            "section_title":  "Unknown",
            "slide_range":    [slide_num, slide_num],
        })

        filename = f"slide_{slide_num:03d}.json"
        filepath = os.path.join(output_dir, filename)

        slide_doc = {
            "slide_number":   slide_num,
            "slide_title":    slide_data.get("title"),
            "slide_layout":   slide_data.get("layout"),
            "slide_owner":    slide_data.get("slide_owner"),
            "section_number": sec_info["section_number"],
            "section_title":  sec_info["section_title"],
            "slide_range":    sec_info["slide_range"],
            "entries":        entries,
        }
        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(slide_doc, f, indent=2, ensure_ascii=False)

        manifest.append({
            "slide_number":   slide_num,
            "slide_title":    slide_data.get("title"),
            "section_number": sec_info["section_number"],
            "section_title":  sec_info["section_title"],
            "file":           filepath,
            "entry_count":    len(entries),
        })

    return manifest


# ── PPT metadata extractor (slides 1–2) — shared by both modes ───────────────

def extract_ppt_metadata(pptx_path: str) -> dict:
    """
    Extract client name, meeting date, attendees, and agenda items
    from slides 1 and 2 of the QBR deck. Identical in both original pipelines.
    """
    import re as _re

    prs      = Presentation(pptx_path)
    n_slides = len(prs.slides)

    result = {"client": None, "subject": None, "meeting_date": None,
               "attendees": [], "agenda_items": []}

    _DATE_PATTERN = _re.compile(
        r'\b(?:Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|'
        r'Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|'
        r'Dec(?:ember)?)\b.{0,6}\d{4}'
        r'|\b\d{1,2}[\s/\-]\w+[\s/\-]\d{2,4}\b'
        r'|\b\d{1,2}/\d{1,2}/\d{2,4}\b',
        _re.IGNORECASE
    )

    def _cell_text(cell) -> str:
        return cell.text.strip().replace("\x0b", " ").replace("\n", " ").strip()

    def _shape_texts(slide) -> list[tuple[bool, str]]:
        out = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = False
            try:
                if shape.is_placeholder:
                    is_title = shape.placeholder_format.type in (
                        PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
            except Exception:
                pass
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    out.append((is_title, t))
        return out

    if n_slides >= 1:
        slide1 = prs.slides[0]
        texts  = _shape_texts(slide1)
        for is_title, t in texts:
            if is_title and result["subject"] is None:
                result["subject"] = _re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", " ", t).strip()
                result["subject"] = _re.sub(r"\s+", " ", result["subject"]).strip()
                m = _re.match(r'^([^\-–/]+?)\s*[\-–]', t)
                if m:
                    candidate = m.group(1).strip()
                    if candidate.lower() not in ("marsh", "icg", "marsh icg"):
                        result["client"] = candidate
            elif not is_title and result["meeting_date"] is None:
                dm = _DATE_PATTERN.search(t)
                if dm:
                    result["meeting_date"] = dm.group(0).strip()

        if result["meeting_date"] is None:
            for shape in slide1.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            ct = _cell_text(cell)
                            dm = _DATE_PATTERN.search(ct)
                            if dm:
                                result["meeting_date"] = dm.group(0).strip()
                                break

    if n_slides >= 2:
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            if not shape.has_table:
                continue
            tbl    = shape.table
            n_rows = len(tbl.rows)
            n_cols = len(tbl.columns)
            if n_rows < 1 or n_cols < 1:
                continue

            grid = [[_cell_text(tbl.cell(r, c)) for c in range(n_cols)] for r in range(n_rows)]
            header_text = " ".join(grid[0]).lower()

            if "agenda" in header_text:
                for row in grid[1:]:
                    parts = [c for c in row if c]
                    if not parts:
                        continue
                    if len(parts) == 1 and _re.fullmatch(r"\d+", parts[0]):
                        continue
                    item = parts[-1] if len(parts) > 1 else parts[0]
                    item = _re.sub(r"\s+", " ", item).strip()
                    if item:
                        result["agenda_items"].append(item)
                continue

            company = ""
            for row in grid:
                if not any(row):
                    continue
                if row[0]:
                    company = row[0]
                if n_cols >= 3:
                    name = _re.sub(r"\s+", " ", row[1] if len(row) > 1 else "").strip()
                    role = _re.sub(r"\s+", " ", row[2] if len(row) > 2 else "").strip()
                    if name and name.lower() not in ("name", "attendees", ""):
                        entry = f"{name} ({company}" + (f" – {role}" if role else "") + ")"
                        result["attendees"].append(entry)
                elif n_cols == 2:
                    name = _re.sub(r"\s+", " ", row[0] or (row[1] if len(row) > 1 else "")).strip()
                    if name and name.lower() not in ("name", "attendees", ""):
                        result["attendees"].append(f"{name} ({company})" if company else name)
                elif n_cols == 1:
                    name = _re.sub(r"\s+", " ", row[0]).strip()
                    if name and name.lower() not in ("name", "attendees", ""):
                        result["attendees"].append(name)

    return result


# ── Public run function ───────────────────────────────────────────────────────

def run(pptx_path: str, carrier_name: str, meeting_topics: list[str],
        output_dir: str, *, granularity: str = "slide") -> list[dict]:
    """
    Full Phase 2 pipeline.
    Extracts all slides, detects sections, saves output files.

    Args:
        granularity: "slide"   → one JSON per slide   (skill approach mode)
                     "section" → one JSON per section  (carrier/marsh format mode)

    Returns manifest list.
    """
    _reset_counter()
    prs      = Presentation(pptx_path)
    n_slides = len(prs.slides)
    log.info("MoM: loaded deck with %d slide(s)", n_slides)

    per_slide: dict[int, dict] = {}
    for i, slide in enumerate(prs.slides, start=1):
        per_slide[i] = _extract_slide_entries(
            slide, i, section_title="", carrier_name=carrier_name
        )

    sections = detect_sections(per_slide, n_slides, meeting_topics)

    for sec in sections:
        title = sec["section_title"]
        for sn in range(sec["start_slide"], sec["end_slide"] + 1):
            if sn in per_slide:
                slide = prs.slides[sn - 1]
                owner = _determine_slide_owner(slide, title, carrier_name)
                per_slide[sn]["slide_owner"] = owner
                for entry in per_slide[sn]["entries"]:
                    entry["slide_owner"] = owner

    if granularity == "section":
        return save_sections(per_slide, sections, output_dir)
    else:
        return save_slides(per_slide, sections, output_dir)
