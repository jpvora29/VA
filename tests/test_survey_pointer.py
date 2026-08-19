"""The survey line on the summary page — when it is written, and what it may claim.

Unit tests run on a stub grid (no warehouse); the end-to-end test generates both a
premium-basis and a survey-basis deck from the real templates and checks the line is on
the page that carries the score tile, and only on the survey basis.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from typing import Dict, Optional, Tuple

import pytest

from studio.compute import DATA_BASIS_PREMIUM, DATA_BASIS_WITH_SURVEY
from studio.template_fill.survey import pointer as P


@pytest.fixture(autouse=True)
def _no_llm(monkeypatch):
    monkeypatch.setenv("STUDIO_AI", "off")


@dataclass
class _Result:
    subject: str = "Zurich"
    data_basis: str = DATA_BASIS_WITH_SURVEY
    resolved_filters: Dict = field(default_factory=lambda: {"Country": ("Singapore",)})


@dataclass
class _Grid:
    """The fields :mod:`studio.template_fill.survey.pointer` reads off a ScoreGrid."""

    year: int = 2025
    prior_year: Optional[int] = 2024
    overall: Optional[float] = 6.95
    prior_overall: Optional[float] = 6.93
    sections: Tuple[str, ...] = ()
    totals: Dict[str, float] = field(default_factory=dict)

    def overall_delta(self):
        if self.overall is None or self.prior_overall is None:
            return None
        return self.overall - self.prior_overall

    def section_total(self, section):
        return self.totals.get(section)


def _with_grid(monkeypatch, grid):
    from studio.template_fill.survey import facts as survey_facts

    monkeypatch.setattr(survey_facts, "load_grid", lambda r, c: grid)


# ── the data-basis gate ──────────────────────────────────────────────────────


def test_a_premium_basis_run_gets_no_survey_line(monkeypatch):
    _with_grid(monkeypatch, _Grid())
    assert P.overall_point(_Result(data_basis=DATA_BASIS_PREMIUM)) is None


def test_a_survey_basis_run_reports_the_score_and_its_year(monkeypatch):
    _with_grid(monkeypatch, _Grid())
    said = P.overall_point(_Result())
    assert said.startswith("Brokers scored Zurich 7.0 in the 2025 carrier survey")


def test_a_scope_the_survey_book_cannot_answer_gets_no_line(monkeypatch):
    from studio.template_fill.survey import facts as survey_facts

    _with_grid(monkeypatch, None)
    monkeypatch.setattr(survey_facts, "load_overall_score", lambda r, c: None)
    assert P.overall_point(_Result()) is None


def test_a_failing_survey_book_never_breaks_the_deck(monkeypatch):
    from studio.template_fill.survey import facts as survey_facts

    def boom(*a, **k):
        raise RuntimeError("survey warehouse down")

    monkeypatch.setattr(survey_facts, "load_grid", boom)
    monkeypatch.setattr(survey_facts, "load_overall_score", boom)
    assert P.overall_point(_Result()) is None


# ── movement, in the table's own terms ───────────────────────────────────────


def test_a_move_inside_the_neutral_band_is_not_called_a_rise(monkeypatch):
    """The Carrier Survey table paints |delta| <= 0.2 as no material change; the prose
    must not contradict the colour on the next page."""
    _with_grid(monkeypatch, _Grid(overall=6.95, prior_overall=6.93))
    assert "broadly unchanged on 2024" in P.overall_point(_Result())


def test_a_material_rise_is_named_with_its_size(monkeypatch):
    _with_grid(monkeypatch, _Grid(overall=7.4, prior_overall=6.9))
    assert "up 0.5 on 2024" in P.overall_point(_Result())


def test_a_material_fall_is_named(monkeypatch):
    _with_grid(monkeypatch, _Grid(overall=6.4, prior_overall=6.9))
    assert "down 0.5 on 2024" in P.overall_point(_Result())


def test_no_prior_surveyed_year_means_no_movement_claim(monkeypatch):
    _with_grid(monkeypatch, _Grid(prior_year=None, prior_overall=None))
    said = P.overall_point(_Result())
    assert said == "Brokers scored Zurich 7.0 in the 2025 carrier survey."


# ── the practice spread ──────────────────────────────────────────────────────


_SECTIONS = ("Underwriting", "Claims", "Policy Administration")


def test_a_trivial_spread_between_practices_is_not_a_finding(monkeypatch):
    """0.11 across every practice is noise — naming a 'strongest' one is the survey-side
    version of naming a $145K move on a $208M book."""
    _with_grid(monkeypatch, _Grid(sections=_SECTIONS, totals={
        "Underwriting": 6.92, "Claims": 6.95, "Policy Administration": 7.01}))
    assert "strongest practice" not in P.overall_point(_Result())


def test_a_real_spread_names_both_ends_with_their_scores(monkeypatch):
    _with_grid(monkeypatch, _Grid(sections=_SECTIONS, totals={
        "Underwriting": 6.1, "Claims": 6.9, "Policy Administration": 7.6}))
    said = P.overall_point(_Result())
    assert "Policy Administration its strongest practice at 7.6" in said
    assert "Underwriting its weakest at 6.1" in said


# ── scope ────────────────────────────────────────────────────────────────────


def test_a_multi_country_run_reports_the_score_without_a_practice_breakdown(monkeypatch):
    """The grid is loaded per country, so a multi-market run must not report one market's
    practices as if they were the whole book's."""
    from studio.template_fill.survey import facts as survey_facts

    monkeypatch.setattr(survey_facts, "load_grid",
                        lambda r, c: pytest.fail("a multi-country run loaded one country's grid"))
    monkeypatch.setattr(survey_facts, "load_overall_score", lambda r, c: 7.2)
    said = P.overall_point(_Result(resolved_filters={"Country": ("Singapore", "Japan")}))
    assert said == "Brokers scored Zurich 7.2 in the carrier survey."


# ── end to end ───────────────────────────────────────────────────────────────


def _slide_texts(path):
    from pptx import Presentation

    return [" ".join(sh.text_frame.text for sh in s.shapes
                     if sh.has_text_frame and sh.text_frame.text.strip())
            for s in Presentation(path).slides]


def test_the_survey_line_lands_on_the_page_that_carries_the_score_tile(tmp_path):
    from studio.compute import compute_overall
    from studio.template_fill import assemble as A
    from studio.template_fill.binding_map import available

    if "overall" not in set(available()):
        pytest.skip("split templates not present")

    result = compute_overall(
        filters={"Carrier_Group": "Zurich", "Country": "Singapore", "Year": 2025})

    with_survey = _slide_texts(A.assemble_deck(
        result, out_path=str(tmp_path / "survey.pptx"), scope="overall",
        data_basis=DATA_BASIS_WITH_SURVEY))
    premium = _slide_texts(A.assemble_deck(
        result, out_path=str(tmp_path / "premium.pptx"), scope="overall",
        data_basis=DATA_BASIS_PREMIUM))

    carrying = [t for t in with_survey if "Brokers scored" in t]
    assert len(carrying) == 1, "the survey line should appear on exactly one page"
    assert "Survey" in carrying[0], "the line is not on the page that reports the score"
    assert not any("Brokers scored" in t for t in premium), \
        "a premium-basis deck was not asked about the survey book"
