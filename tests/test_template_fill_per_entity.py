"""Per-entity role resolution — the split-template re-scoping seam.

Hermetic: exercises the filter re-scoping (which product/country a sub-deck is built for)
without touching the database, by stubbing ``resolve_roles`` to capture the scoped filters.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from studio import compute as C
from studio.template_fill import bindings as B
from studio.template_fill.fill import fill_template


def _result(**filters):
    return C.OverallResult(subject="ACME", flow="gpr", resolved_filters=dict(filters))


def test_selected_products_and_countries_read_filters():
    r = _result(Product_Line=("Aviation", "Marine"), Country="UK")
    assert B.selected_products(r) == ("Aviation", "Marine")
    assert B.selected_countries(r) == ("UK",)          # scalar → single-item tuple
    assert B.selected_products(_result()) == ()        # unfiltered → empty


def test_rescope_pins_one_value_without_mutating_original():
    r = _result(Product_Line=("Aviation", "Marine"), Country="UK", Year=2025)
    scoped = B._rescope(r, "Product_Line", "Aviation")
    assert scoped.resolved_filters["Product_Line"] == "Aviation"
    assert scoped.resolved_filters["Country"] == "UK"   # other filters preserved
    assert scoped.resolved_filters["Year"] == 2025
    # original untouched
    assert r.resolved_filters["Product_Line"] == ("Aviation", "Marine")


def test_resolve_for_product_scopes_filters(monkeypatch):
    seen = {}
    monkeypatch.setattr(B, "resolve_roles", lambda res: seen.update(res.resolved_filters) or {})
    r = _result(Product_Line=("Aviation", "Marine"), Country="UK")
    B.resolve_roles_for_product(r, "Marine")
    assert seen["Product_Line"] == "Marine"
    assert seen["Country"] == "UK"


def test_resolve_for_country_scopes_filters(monkeypatch):
    seen = {}
    monkeypatch.setattr(B, "resolve_roles", lambda res: seen.update(res.resolved_filters) or {})
    r = _result(Product_Line="Aviation", Country=("UK", "France"))
    B.resolve_roles_for_country(r, "France")
    assert seen["Country"] == "France"
    assert seen["Product_Line"] == "Aviation"


def test_resolve_for_product_injects_product_name(monkeypatch):
    monkeypatch.setattr(B, "resolve_roles", lambda res: {})
    values = B.resolve_roles_for_product(_result(Product_Line=("Marine", "Property")), "Property")
    assert values["product_name"] == "Property"      # so fill can rewrite the authored word


def test_fill_rewrites_authored_products_to_deck_product(tmp_path):
    # A product deck authored with "Marine"/"Energy" example words, filled for "Property".
    src = str(tmp_path / "product.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text_frame.text = (
        "Marine Feedback — Energy portfolio"
    )
    prs.save(src)

    doc = {
        "template_path": src,
        "values": {"product_name": "Property", "product_vocab": ["Marine", "Energy", "Property"]},
        "manifest": [],
        "overrides": {}, "map_overrides": {}, "added": {},
    }
    out = fill_template(doc, out_path=str(tmp_path / "filled.pptx"))

    text = "\n".join(
        sh.text_frame.text for s in Presentation(out).slides for sh in s.shapes if sh.has_text_frame
    )
    assert "Property Feedback — Property portfolio" in text
    assert "Marine" not in text and "Energy" not in text
