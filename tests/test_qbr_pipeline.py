"""End-to-end QBR pipeline (plan §Pipeline Style + regression tests).

Runs the real business flow against the seed DB with STUDIO_AI=off:

    selection -> structured data load -> evidence -> report plan
      -> template descriptor/intent/binding -> render plan
      -> commentary -> QA -> exported PPTX

and asserts the plan's regression guarantees: deterministic output across
repeated runs, slide order preserved under parallel drafting, sub-functions
testable independently, and export blocked only on critical QA failures.
"""
from __future__ import annotations

import asyncio
import os
from pathlib import Path

import pytest

os.environ["STUDIO_AI"] = "off"          # deterministic: no LLM anywhere

from studio.pipeline import (  # noqa: E402
    StudioSelection,
    build_qbr_deck,
    build_report_plan,
    bounded_gather,
    gather_ordered,
)
from studio.pipeline.qbr_pipeline import build_evidence, commentary_targets  # noqa: E402

FILTERS = {"carrier": "Zurich", "year": 2025}
TEMPLATE = "template/qbr_template.pptx"


@pytest.fixture(scope="module")
def result(tmp_path_factory):
    out = tmp_path_factory.mktemp("deck") / "qbr_e2e.pptx"
    sel = StudioSelection(filters=FILTERS, template_path=TEMPLATE, out_path=str(out))
    return build_qbr_deck(sel)


# ── the end-to-end path ───────────────────────────────────────────────────────


def test_pipeline_exports_a_real_deck(result):
    assert result.deck_path and Path(result.deck_path).exists()
    assert Path(result.deck_path).stat().st_size > 10_000

    from pptx import Presentation

    prs = Presentation(result.deck_path)          # exported file opens cleanly
    assert len(prs.slides) > 0


def test_qa_report_present_and_not_blocking(result):
    counts = result.qa_report.counts()
    assert counts["total"] > 0                     # QA actually ran and recorded notes
    assert counts["critical"] == 0
    assert not result.qa_report.blocking


def test_every_commentary_number_traces_to_a_fact(result):
    pack = result.evidence_pack
    assert any(c.sentences for c in result.commentary)
    for slide in result.commentary:
        for sent in slide.sentences:
            assert sent.fact_ids, f"uncited sentence on slide {slide.slide_idx}"
            for fid in sent.fact_ids:
                assert fid in pack.items, f"unknown fact {fid}"


def test_report_plan_claims_cite_existing_facts(result):
    from studio.content.qa import check_report_plan, has_errors

    violations = check_report_plan(result.report_plan, result.evidence_pack)
    assert not has_errors(violations)
    assert result.report_plan.findings
    assert result.report_plan.data_gaps            # honest gaps, not filler


def test_binding_map_and_layout_intent_cover_the_template(result):
    assert result.binding_map is not None and result.binding_map.bindings
    assert len(result.layout_intent.slides) == len(result.render_plan.doc["order"])


# ── determinism regressions ───────────────────────────────────────────────────


def test_repeated_runs_are_deterministic(result, tmp_path):
    sel = StudioSelection(filters=FILTERS, template_path=TEMPLATE,
                          out_path=str(tmp_path / "again.pptx"))
    again = build_qbr_deck(sel)

    first = [(c.slide_idx, c.purpose, [(s.text, s.fact_ids) for s in c.sentences])
             for c in result.commentary]
    second = [(c.slide_idx, c.purpose, [(s.text, s.fact_ids) for s in c.sentences])
              for c in again.commentary]
    assert first == second
    assert result.qa_report.counts() == again.qa_report.counts()
    assert sorted(result.evidence_pack.items) == sorted(again.evidence_pack.items)


def test_parallel_commentary_preserves_slide_order(result):
    indices = [c.slide_idx for c in result.commentary]
    assert indices == sorted(indices)              # layout order, not completion order


# ── sub-functions are independently testable ──────────────────────────────────


def test_evidence_and_report_plan_steps_run_standalone():
    from studio.compute import compute_overall
    from studio.content.evidence_graph import build_evidence_graph

    computed = compute_overall(filters=FILTERS)
    pack = build_evidence(computed)
    assert pack.items and pack.subject == "Zurich"

    graph = build_evidence_graph(pack)
    assert graph.nodes and graph.edges

    plan = build_report_plan(pack, graph, StudioSelection(filters=FILTERS))
    assert plan.findings


def test_commentary_targets_skip_hidden_slides(result):
    targets = commentary_targets(result.layout_intent, result.render_plan)
    hidden = set(result.render_plan.hidden_slides)
    assert targets
    assert not [idx for idx, _ in targets if idx in hidden]


# ── async helpers ─────────────────────────────────────────────────────────────


def test_gather_ordered_never_reorders():
    async def slow(i):
        await asyncio.sleep(0.03 - i * 0.01)       # later items finish first
        return i

    async def main():
        return await gather_ordered([slow(i) for i in range(3)])

    assert asyncio.run(main()) == [0, 1, 2]


def test_bounded_gather_limits_concurrency_and_keeps_order():
    running = 0
    peak = 0

    async def work(i):
        nonlocal running, peak
        running += 1
        peak = max(peak, running)
        await asyncio.sleep(0.01)
        running -= 1
        return i

    async def main():
        return await bounded_gather([work(i) for i in range(8)], limit=2)

    assert asyncio.run(main()) == list(range(8))
    assert peak <= 2
