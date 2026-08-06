"""Assembly orchestration — plan → fill each sub-deck → merge.

Hermetic: registers tiny in-memory split templates and stubs the DB-backed role resolvers,
so it proves the *composition* (overall + N products + M countries, in order) without a
warehouse. Real role values are covered by the per-entity and compute tests.
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from studio import compute as C
from studio.template_fill import assemble as A
from studio.template_fill import binding_map as BM


@pytest.fixture(autouse=True)
def _force_opc_merge(monkeypatch):
    monkeypatch.setenv("STUDIO_PPT_MERGE_ENGINE", "opc")


def _tiny_template(path: str, n_slides: int) -> str:
    prs = Presentation()
    for i in range(n_slides):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = f"s{i}"
    prs.save(path)
    return path


@pytest.fixture
def split_templates(tmp_path, monkeypatch):
    overall = _tiny_template(str(tmp_path / "overall.pptx"), 2)
    product = _tiny_template(str(tmp_path / "product.pptx"), 1)
    country = _tiny_template(str(tmp_path / "country.pptx"), 1)

    # Register the three fixed templates (empty bindings — composition test only).
    BM._REGISTRY.clear()
    BM.get_binding_map.cache_clear()
    BM._REGISTRY["overall"] = lambda: BM.BindingMap("overall", overall, ())
    BM._REGISTRY["product"] = lambda: BM.BindingMap("product", product, ())
    BM._REGISTRY["country"] = lambda: BM.BindingMap("country", country, ())
    monkeypatch.setattr(BM, "_discover_json_maps", lambda: None)  # ignore on-disk maps

    # Stub the DB-backed resolvers (assemble imported them by name into its namespace).
    monkeypatch.setattr(A, "resolve_roles", lambda r: {})
    monkeypatch.setattr(A, "resolve_roles_for_product", lambda r, p: {})
    monkeypatch.setattr(A, "resolve_roles_for_country", lambda r, c: {})
    yield
    BM._REGISTRY.clear()
    BM.get_binding_map.cache_clear()


def test_assemble_composes_overall_products_countries(split_templates, tmp_path):
    result = C.OverallResult(
        subject="ACME", flow="gpr",
        resolved_filters={"Product_Line": ("Aviation", "Marine"), "Country": ("UK", "France")},
    )
    out = A.assemble_deck(result, out_path=str(tmp_path / "deck.pptx"), work_dir=str(tmp_path / "work"))

    merged = Presentation(out)
    # overall(2) + 2 products(1 each) + 2 countries(1 each) = 6
    assert len(merged.slides._sldIdLst) == 6


def test_assemble_overall_only_when_no_selection_and_no_data(split_templates, tmp_path):
    # No selection AND no engine → the carrier-book fallback resolves to nothing → overall only.
    result = C.OverallResult(subject="ACME", flow="gpr", resolved_filters={})
    out = A.assemble_deck(result, out_path=str(tmp_path / "deck.pptx"), work_dir=str(tmp_path / "work"))
    assert len(Presentation(out).slides._sldIdLst) == 2


def test_assemble_falls_back_to_full_carrier_book_when_unselected(split_templates, tmp_path, monkeypatch):
    # Nothing pinned → build a block for every product/country the carrier writes in.
    monkeypatch.setattr(A, "product_vocab", lambda r: ("Marine", "Property", "Energy"))
    monkeypatch.setattr(A, "carrier_countries", lambda r: ("Singapore", "Japan"))
    result = C.OverallResult(subject="ACME", flow="gpr", resolved_filters={})

    decks = A.plan_subdecks(result)
    kinds = [d.template for d in decks]
    assert kinds == ["overall", "product", "product", "product", "country", "country"]
    assert [d.label for d in decks[1:]] == ["Marine", "Property", "Energy", "Singapore", "Japan"]

    out = A.assemble_deck(result, out_path=str(tmp_path / "deck.pptx"), work_dir=str(tmp_path / "work"))
    # overall(2) + 3 products(1) + 2 countries(1)
    assert len(Presentation(out).slides._sldIdLst) == 7


# ── the overall block reports on the Setup selection ─────────────────────────


def _scoped_results(monkeypatch, result):
    """``{sub-deck label: the result its PAGES are computed from}``.

    Captured at ``_build_subdeck``, which is exactly what the page providers
    (grids, gwp_page, lc_page…) are handed.
    """
    seen = {}
    build = A._build_subdeck

    def capture(template_name, scoped_result, values, label):
        seen[label] = scoped_result
        return build(template_name, scoped_result, values, label)

    monkeypatch.setattr(A, "_build_subdeck", capture)
    A.plan_subdecks(result)
    return seen


def test_the_overall_block_keeps_the_pinned_products(split_templates, monkeypatch):
    """Regression: the overall summary dropped the product filter outright, so a run
    pinned to two lines still reported the carrier's whole book on page 2."""
    result = C.OverallResult(
        subject="ACME", flow="gpr",
        resolved_filters={"Product_Line": ("Aviation", "Marine"), "Country": ("UK",)},
    )
    seen = _scoped_results(monkeypatch, result)
    assert seen["overall"].resolved_filters["Product_Line"] == ("Aviation", "Marine")
    assert seen["overall"].resolved_filters["Country"] == ("UK",)


def test_every_subdeck_carries_the_runs_whole_product_selection(split_templates, monkeypatch):
    """A per-product sub-deck narrows to one line, so the SELECTION has to travel
    separately — a portfolio page inside it ranks against the run, not the book."""
    result = C.OverallResult(
        subject="ACME", flow="gpr",
        resolved_filters={"Product_Line": ("Aviation", "Marine")},
    )
    seen = _scoped_results(monkeypatch, result)
    for scoped in seen.values():
        assert scoped.scope_products == ("Aviation", "Marine")
    assert seen["Marine"].resolved_filters["Product_Line"] == "Marine"


def test_an_unfiltered_run_carries_no_product_selection(split_templates, monkeypatch):
    """The fallback book is not a selection — nothing to narrow a portfolio page to."""
    monkeypatch.setattr(A, "product_vocab", lambda r: ("Marine", "Property"))
    seen = _scoped_results(monkeypatch, C.OverallResult(subject="ACME", flow="gpr"))
    assert seen["overall"].scope_products == ()


def test_plan_skips_unregistered_axes(tmp_path, monkeypatch):
    # Only 'overall' registered → product/country selections are ignored.
    overall = _tiny_template(str(tmp_path / "overall.pptx"), 3)
    BM._REGISTRY.clear()
    BM.get_binding_map.cache_clear()
    BM._REGISTRY["overall"] = lambda: BM.BindingMap("overall", overall, ())
    monkeypatch.setattr(BM, "_discover_json_maps", lambda: None)
    monkeypatch.setattr(A, "resolve_roles", lambda r: {})

    result = C.OverallResult(resolved_filters={"Product_Line": ("Aviation",), "Country": ("UK",)})
    decks = A.plan_subdecks(result)
    assert [d.template for d in decks] == ["overall"]
    BM._REGISTRY.clear()
    BM.get_binding_map.cache_clear()
