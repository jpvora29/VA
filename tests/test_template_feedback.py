"""Feedback/quadrant/highlights table fill — detection, roles and value composition.

Hermetic: the detection/composition tests run on synthetic ``Template`` objects with the
compute layer stubbed; the bubble-chart test runs against the real country template when
present (skip otherwise). ``STUDIO_AI=off`` keeps every path deterministic.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from studio import compute as C
from studio.template_fill import feedback as F
from studio.template_fill import roles as R
from studio.template_fill.analyze import Shape, Slide, Template


@pytest.fixture(autouse=True)
def _no_llm(monkeypatch):
    monkeypatch.setenv("STUDIO_AI", "off")


# ── synthetic template (the product feedback page + the country quadrant page) ─


def _feedback_table() -> Shape:
    return Shape(shape_id=2, name="Table 1", kind="table", table=[
        ["", "What’s working well", "What’s not", "Growth Opportunities", "", ""],
        ["Country / Region(1)", "Relationship: …………", "", "",
         "$xxx.xM (+x.x%▲)\nMarsh GWP", "$xx.xM (-x.x▼)\nCarrier GWP"],
        ["", "", "", "", "x (=0►)\nCarrier Rank", "x.x% (+x.x%▲)\nCarrier SoW"],
        ["Country / Region (2)", "", "", ".",
         "$xxx.xM (+x.x%▲) Marsh GWP", "$xx.xM (-x.x▼)\nCarrier GWP"],
        ["", "", "", "", "x (+x▲)\nCarrier Rank", "x.x% (-x.x▼)\nCarrier SoW"],
    ])


def _quadrant_table() -> Shape:
    return Shape(shape_id=4, name="Table 3", kind="table", table=[
        ["Successes", "Challenges", "Opportunities for 2026", "Key Messages"],
        ["Performance", ". ", "", ""],
    ])


def _highlights_table() -> Shape:
    return Shape(shape_id=3, name="Table 2", kind="table", table=[["Key Highlights:"]])


def _template() -> Template:
    return Template(path="synthetic", width_emu=12192000, height_emu=6858000, slides=[
        Slide(index=0, layout="", shapes=[_feedback_table()]),
        Slide(index=1, layout="", shapes=[_quadrant_table(), _highlights_table()]),
    ])


_FACTS = {
    "carrier": {"current": 48e6, "pct": 52.5},
    "marsh": {"current": 411e6, "pct": 9.8},
    "rank": {"current": 2, "delta": 4},
    "sow": {"current": 11.6, "delta": 3.2},
}


def _result():
    return C.OverallResult(subject="ACME", flow="gpr",
                           resolved_filters={"Carrier_Group": "ACME", "Year": 2025})


@pytest.fixture
def _stub_compute(monkeypatch):
    monkeypatch.setattr(F, "_reporting_filters", lambda r: dict(r.resolved_filters))
    monkeypatch.setattr(F, "_countries_in_scope", lambda r: ["Japan"])
    monkeypatch.setattr(F, "_facts", lambda r, f: _FACTS)


# ── detection + role binding ──────────────────────────────────────────────────


def test_augment_binds_feedback_quadrant_and_highlight_cells():
    bindings = F.augment(_template(), [])
    roles = {b.role for b in bindings}
    # Country (1) block: 3 commentary columns (fbnote:) + 4 KPI callouts (fb:).
    assert {"fbnote:0:2:1:1", "fbnote:0:2:1:2", "fbnote:0:2:1:3",
            "fb:0:2:1:4", "fb:0:2:1:5", "fb:0:2:2:4", "fb:0:2:2:5"} <= roles
    # Country (2) block binds too, plus the quadrant row and the highlights cell.
    assert {"fb:0:2:3:4", "fbnote:1:4:1:0", "fbnote:1:4:1:3", "fbnote:1:3:0:0"} <= roles
    assert all(not b.placeholder for b in bindings)


def test_augment_rebinds_existing_manifest_slots_in_place():
    from studio.template_fill.slots import Slot

    existing = R.Binding(Slot(0, 2, ["cell", 1, 4], "$xxx.xM", "money", ""), None, True)
    out = F.augment(_template(), [existing])
    assert existing.role == "fb:0:2:1:4" and existing.placeholder is False
    assert existing in out


def test_augment_ignores_unrelated_tables():
    plain = Template(path="p", width_emu=1, height_emu=1, slides=[
        Slide(index=0, layout="", shapes=[Shape(shape_id=9, name="T", kind="table",
                                                table=[["A", "B"], ["1", "2"]])]),
    ])
    assert F.augment(plain, []) == []


# ── value composition (numbers on the page) ───────────────────────────────────


def test_values_fill_kpi_cells_in_template_style(_stub_compute):
    vals = F.values(_template(), _result())
    assert vals["fb:0:2:1:4"] == "$411M (+9.8%▲)\nMarsh GWP"
    assert vals["fb:0:2:1:5"] == "$48M (+52.5%▲)\nCarrier GWP"
    assert vals["fb:0:2:2:4"] == "2 (+4▲)\nCarrier Rank"
    assert vals["fb:0:2:2:5"] == "11.6% (+3.2%▲)\nCarrier SoW"


def test_values_blank_country_rows_beyond_scope(_stub_compute):
    # Only one country in scope → the Country (2) block is blanked entirely.
    vals = F.values(_template(), _result())
    assert vals["fb:0:2:3:4"] == "" and vals["fb:0:2:4:5"] == ""


def test_values_commentary_carries_figures(_stub_compute):
    vals = F.values(_template(), _result())
    assert "+52.5%" in vals["fbnote:0:2:1:1"] and "$48M" in vals["fbnote:0:2:1:1"]  # working well
    assert "$363M" in vals["fbnote:0:2:1:3"]                                        # headroom
    assert vals["fbnote:1:3:0:0"].startswith("Key Highlights:")
    assert "rank #2" in vals["fbnote:1:4:1:3"]                                      # key messages


def test_declines_flow_to_challenges():
    facts = {"carrier": {"current": 30e6, "pct": -12.0}, "marsh": {"current": 100e6, "pct": 2.0},
             "rank": {"current": 6, "delta": -2}, "sow": {"current": 4.0, "delta": -1.1}}
    text = F._compose("challenges", facts, F._PANEL_BULLETS)
    assert "-12.0%" in text and "slipped" in text
    assert F._kpi_cell("rank", facts) == "6 (-2▼)\nCarrier Rank"


def test_commentary_written_in_consistent_arial_11(tmp_path):
    # Commentary inherits ad-hoc template run formatting (18pt here, 10pt there) —
    # every note:/fbnote: write must come out Arial 11, KPI cells keep their own style.
    from pptx import Presentation
    from pptx.util import Inches, Pt

    from studio.template_fill.fill import fill_template
    from studio.template_fill.slots import Slot

    src = str(tmp_path / "t.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "…………"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    frame = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(2))
    frame.table.rows[1].cells[1].text = "…………"
    tb_id, tbl_id = tb.shape_id, frame.shape_id
    prs.save(src)

    note_role, cell_role = f"note:0:{tb_id}:0", f"fbnote:0:{tbl_id}:1:1"
    manifest = [
        R.Binding(Slot(0, tb_id, ["para", 0], "…………", "text", ""), note_role, False).to_dict(),
        R.Binding(Slot(0, tbl_id, ["cell", 1, 1], "…………", "text", ""), cell_role, False).to_dict(),
    ]
    doc = {"template_path": src, "manifest": manifest,
           "values": {note_role: "Growth of +5.0% YoY.", cell_role: "Line one.\nLine two."},
           "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "styled.pptx"))

    prs2 = Presentation(out)
    shapes = {sh.shape_id: sh for sh in prs2.slides[0].shapes}
    note_runs = [r for p in shapes[tb_id].text_frame.paragraphs for r in p.runs if r.text]
    cell_runs = [r for p in shapes[tbl_id].table.rows[1].cells[1].text_frame.paragraphs
                 for r in p.runs if r.text]
    assert note_runs and cell_runs
    for run in note_runs + cell_runs:
        assert run.font.name == "Arial" and run.font.size.pt == 11


def test_commentary_columns_are_laid_out_alike_across_a_page(tmp_path):
    # The quadrant's four panels are authored inconsistently — one left-aligned and top
    # anchored with real navy text, the rest centred, middle anchored and carrying the
    # white run colour of the example text the author deleted. Filled, they must read as
    # four columns of one slide.
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches

    from studio.template_fill.fill import fill_template
    from studio.template_fill.slots import Slot

    navy = RGBColor(0x00, 0x0F, 0x47)
    src = str(tmp_path / "quad.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    authored = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(4))
    authored.text_frame.text = "Performance"
    authored.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
    blank = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(3), Inches(4))
    blank.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    blank.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    a_id, b_id = authored.shape_id, blank.shape_id
    prs.save(src)

    roles = {a_id: f"note:0:{a_id}:0", b_id: f"note:0:{b_id}:0"}
    doc = {"template_path": src,
           "manifest": [R.Binding(Slot(0, sid, ["para", 0], "", "text", ""), role, False).to_dict()
                        for sid, role in roles.items()],
           "values": {roles[a_id]: "Successes point.", roles[b_id]: "Opportunity point."},
           "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "quad_out.pptx"))

    shapes = {sh.shape_id: sh for sh in Presentation(out).slides[0].shapes}
    for sid in (a_id, b_id):
        frame = shapes[sid].text_frame
        assert frame.vertical_anchor == MSO_ANCHOR.TOP
        for p in frame.paragraphs:
            assert p.alignment == PP_ALIGN.LEFT
            # The emptied panel takes the ink of the column the author actually wrote,
            # not the leftover colour of the text deleted from it.
            assert all(r.font.color.rgb == navy for r in p.runs)


# ── bullet-point commentary ──────────────────────────────────────────────────


def _bullet_char(paragraph):
    """The bullet a paragraph will actually show: its char, ``None``, or inherited."""
    from pptx.oxml.ns import qn

    pPr = paragraph._p.find(qn("a:pPr"))
    if pPr is None:
        return "inherit"
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        el = pPr.find(qn(tag))
        if el is not None:
            return el.get("char") or tag.split(":")[1]
    return "inherit"


def _fill_commentary(tmp_path, text, *, authored=("…………",), bullet_char=None):
    """Fill one text box with ``text`` as commentary; return its exported paragraphs."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    from studio.template_fill.fill import fill_template
    from studio.template_fill.slots import Slot

    src = str(tmp_path / "src.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    tb.text_frame.text = authored[0]
    for line in authored[1:]:
        tb.text_frame.add_paragraph().text = line
    if bullet_char:                                   # the author's own bullet
        for p in tb.text_frame.paragraphs:
            pPr = p._p.get_or_add_pPr()
            pPr.append(pPr.makeelement(qn("a:buChar"), {"char": bullet_char}))
    shape_id = tb.shape_id
    prs.save(src)

    role = f"note:0:{shape_id}:0"
    doc = {"template_path": src,
           "manifest": [R.Binding(Slot(0, shape_id, ["para", 0], authored[0], "text", ""),
                                  role, False).to_dict()],
           "values": {role: text}, "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "out.pptx"))
    shape = next(sh for sh in Presentation(out).slides[0].shapes
                 if sh.shape_id == shape_id)
    return list(shape.text_frame.paragraphs)


def test_commentary_is_written_one_bullet_per_line(tmp_path):
    paras = _fill_commentary(tmp_path, "First point.\nSecond point.\nThird point.")
    assert [p.text for p in paras] == ["First point.", "Second point.", "Third point."]
    assert [_bullet_char(p) for p in paras] == ["•", "•", "•"]


def test_surplus_authored_paragraphs_are_removed_not_left_empty(tmp_path):
    # Four authored lines, two real points → no empty bulleted paragraph may survive.
    paras = _fill_commentary(tmp_path, "One point.\nTwo points.",
                             authored=("…", "…", "…", "…"))
    assert [p.text for p in paras] == ["One point.", "Two points."]


def test_the_authors_own_bullet_is_kept(tmp_path):
    # The Trading Summary columns are bulleted with a Wingdings "§" — the deck's look wins.
    paras = _fill_commentary(tmp_path, "First point.\nSecond point.", bullet_char="§")
    assert [_bullet_char(p) for p in paras] == ["§", "§"]


def test_a_heading_line_is_not_bulleted(tmp_path):
    # "Key Highlights:" introduces the bullets rather than being one of them.
    paras = _fill_commentary(tmp_path, "Key Highlights:\nFirst point.\nSecond point.")
    assert [_bullet_char(p) for p in paras] == ["buNone", "•", "•"]


_GROWING = {"carrier": {"current": 48e6, "pct": 52.5}, "marsh": {"current": 411e6, "pct": 9.8},
            "rank": {"current": 2, "delta": 4}, "sow": {"current": 11.6, "delta": 3.2},
            "peer": {"current": 45e6, "pct": 7.6, "sow": 10.8, "sow_delta": -0.2}}
_SHRINKING = {"carrier": {"current": 30e6, "pct": -12.0}, "marsh": {"current": 100e6, "pct": 2.0},
              "rank": {"current": 6, "delta": -2}, "sow": {"current": 4.0, "delta": -1.1},
              "peer": {"current": 20e6, "pct": 1.0, "sow": 5.0, "sow_delta": 0.1}}


# Marsh outgrowing the carrier gives the Growth column its second point ("capture the flow").
_LAGGING = {**_GROWING, "carrier": {"current": 48e6, "pct": 2.0}}


@pytest.mark.parametrize("composer,facts", [
    ("working", _GROWING), ("growth", _LAGGING), ("key_messages", _GROWING),
    ("challenges", _SHRINKING),
])
def test_composers_return_one_point_per_line(composer, facts):
    lines = F._compose(composer, facts, F._PANEL_BULLETS).split("\n")
    assert len(lines) > 1, "a commentary cell should carry several points, not one paragraph"
    assert all(line.strip() for line in lines)


def test_a_composer_with_nothing_to_say_returns_no_bullets():
    # A book that grew, climbed the rank and already writes above the peer share average has
    # no evidenced challenge — the column must stay empty rather than invent one (the
    # template's own fill-me cue then shows through).
    assert F._compose("challenges", _GROWING, F._PANEL_BULLETS) == ""


def test_a_growing_book_below_peer_share_still_reports_a_challenge():
    # Growth alone is not success: writing 7.1% of the wallet against a 10.8% peer average
    # is a real, evidenced gap, and the quadrant must say so.
    behind = {**_GROWING, "sow": {"current": 7.1, "delta": 3.2}}
    text = F._compose("challenges", behind, F._PANEL_BULLETS)
    assert "3.7pp below the top-5 peer average of 10.8%" in text
    assert "$15M" in text                       # 3.7pp of the $411M Marsh book in scope


def test_a_panel_carries_more_of_the_argument_than_a_table_cell():
    facts = {**_GROWING, "carrier": {"current": 48e6, "pct": 52.5, "delta": 16e6},
             "movers": [{"name": "Cyber", "delta": 12e6, "pct": 40.0},
                        {"name": "Marine", "delta": -1e6, "pct": -4.0}]}
    panel = F._compose("working", facts, F._PANEL_BULLETS).split("\n")
    cell = F._compose("working", facts, F._CELL_BULLETS).split("\n")
    assert len(panel) > len(cell) and panel[:len(cell)] == cell


def test_a_single_value_dimension_does_not_restate_the_headline():
    # A product page scoped to one country decomposes into that country and nothing else —
    # "the increase was led by Singapore" on a Singapore row says nothing.
    facts = {**_GROWING, "carrier": {"current": 48e6, "pct": 52.5, "delta": 16e6},
             "movers": [{"name": "Singapore", "delta": 16e6, "pct": 52.5}],
             "pool": [{"name": "Singapore", "delta": 40e6, "pct": 9.8}]}
    assert "led by" not in F._compose("working", facts, F._PANEL_BULLETS)
    assert "capture gap" not in F._compose("growth", facts, F._PANEL_BULLETS)


def test_highlights_keeps_its_heading_then_bullets():
    lines = F._compose("highlights", _GROWING, F._PANEL_BULLETS).split("\n")
    assert lines[0] == "Key Highlights:"
    assert len(lines) > 1 and all(lines[1:])


# ── the real bubble chart (native fill after detaching the think-cell link) ───


COUNTRY_TEMPLATE = "template/country_template.pptx"


def test_bubble_chart_filled_from_growth_points(tmp_path):
    if not Path(COUNTRY_TEMPLATE).exists():
        pytest.skip("country template not present")
    from pptx import Presentation

    from studio.template_fill.fill import fill_template

    points = [{"lob": "Property", "carrier_yoy": 12.0, "marsh_yoy": 5.0, "size": 40e6},
              {"lob": "Cyber", "carrier_yoy": -3.0, "marsh_yoy": 8.0, "size": 15e6}]
    doc = {"template_path": COUNTRY_TEMPLATE, "manifest": [],
           "values": {"growth_bubble": {"points": points}},
           "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "bubble.pptx"))

    prs = Presentation(out)
    charts = [sh.chart for s in prs.slides for sh in s.shapes if getattr(sh, "has_chart", False)]
    bubble = next(c for c in charts if "BUBBLE" in str(c.chart_type))
    assert [s.name for s in bubble.plots[0].series] == ["Property", "Cyber"]
    assert list(bubble.plots[0].series[0].values) == [12.0 / 100.0]
    # Every bubble is cloned from one authored series, so they share a colour: a legend of
    # identical keys names nothing. Each bubble carries its own line of business instead,
    # and both axes read as percentages rather than the source data's raw fractions.
    assert not bubble.has_legend
    xml = bubble._chartSpace.xml
    assert '<c:showSerName val="1"/>' in xml
    assert xml.count('<c:numFmt formatCode="0.0%" sourceLinked="0"/>') == 2
    # The hand-placed labels for the authored dummy bubbles are blanked.
    slide = next(s for s in prs.slides
                 for sh in s.shapes if getattr(sh, "has_chart", False) and sh.chart is bubble)
    texts = [sh.text_frame.text.strip() for sh in slide.shapes
             if getattr(sh, "has_text_frame", False)]
    assert "Property" not in texts and "FINPRO" not in texts


# ── the growth quadrant reports on the SELECTION, and never on examples ──────


def test_an_empty_growth_payload_clears_the_authored_bubbles(tmp_path):
    """A page titled "<Carrier> vs Marsh growth rates" must never plot the author's
    example lines of business: with nothing to say, the chart is emptied.

    Reachable in practice — a single year of data supports no YoY at all, so every
    point comes back with no growth to plot.
    """
    if not Path(COUNTRY_TEMPLATE).exists():
        pytest.skip("country template not present")
    from pptx import Presentation

    from studio.template_fill.fill import fill_template

    doc = {"template_path": COUNTRY_TEMPLATE, "manifest": [],
           "values": {"growth_bubble": {"points": [
               {"lob": "Property", "carrier_yoy": None, "marsh_yoy": None, "size": 40e6}]}},
           "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "empty_bubble.pptx"))

    prs = Presentation(out)
    charts = [sh.chart for s in prs.slides for sh in s.shapes if getattr(sh, "has_chart", False)]
    bubble = next(c for c in charts if "BUBBLE" in str(c.chart_type))
    plotted = [v for s in bubble.plots[0].series for v in s.values]
    assert not [v for v in plotted if v is not None], "authored example bubbles survived"


def test_no_growth_role_at_all_leaves_the_chart_untouched(tmp_path):
    """Absent is not the same as empty: a doc that never computed the quadrant (a
    preview, a partial fill) must not have its chart wiped."""
    if not Path(COUNTRY_TEMPLATE).exists():
        pytest.skip("country template not present")
    from pptx import Presentation

    from studio.template_fill.fill import fill_template

    doc = {"template_path": COUNTRY_TEMPLATE, "manifest": [], "values": {},
           "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "untouched.pptx"))

    prs = Presentation(out)
    charts = [sh.chart for s in prs.slides for sh in s.shapes if getattr(sh, "has_chart", False)]
    bubble = next(c for c in charts if "BUBBLE" in str(c.chart_type))
    assert list(bubble.plots[0].series), "the authored chart was cleared without being asked"


def test_the_growth_quadrant_resolves_a_year_when_none_is_pinned():
    """Regression: both axes are period comparisons, and ``movement_by_dim`` returns
    nothing without a year — so an unpinned-year run produced an EMPTY payload and the
    chart kept the template's own example bubbles under the carrier's name."""
    from studio.compute import compute_overall
    from studio.template_fill.bindings import resolve_roles

    unpinned = resolve_roles(compute_overall(
        filters={"carrier": "Zurich", "country": ["Singapore"]}))
    points = (unpinned["growth_bubble"] or {}).get("points") or []
    assert points, "no growth points without a pinned year"
    assert all(p["carrier_yoy"] is not None for p in points)


def test_the_growth_quadrant_plots_only_the_selected_lines_of_business():
    """The chart answers the Setup selection: pinning two lines plots those two."""
    from studio.compute import compute_overall
    from studio.template_fill.bindings import resolve_roles

    picked = ["Cyber", "Marine"]
    roles = resolve_roles(compute_overall(
        filters={"carrier": "Zurich", "country": ["Singapore"], "product_line": picked}))
    plotted = {p["lob"] for p in roles["growth_bubble"]["points"]}
    assert plotted == set(picked)

    # …and an unfiltered run still covers the carrier's whole book.
    everything = resolve_roles(compute_overall(
        filters={"carrier": "Zurich", "country": ["Singapore"]}))
    assert len({p["lob"] for p in everything["growth_bubble"]["points"]}) > len(picked)
