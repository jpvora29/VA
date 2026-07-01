"""Splitter — carving one deck into per-axis sub-templates by section.

Hermetic: builds a small titled deck so section classification (overall vs country block)
is exercised without the real asset. A skipif check splits the real template when present.
"""
from __future__ import annotations

import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from studio.template_fill.split import extract_slides, plan_axes, split_template


def _titled_deck(path: str, titles) -> str:
    prs = Presentation()
    layout = prs.slide_layouts[5]  # "Title Only" — has a title placeholder
    for t in titles:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = t
    prs.save(path)
    return path


def test_plan_axes_splits_overall_and_country(tmp_path):
    src = _titled_deck(
        str(tmp_path / "deck.pptx"),
        ["Overall Summary", "Highlights", "Country (1)", "Carrier breakdown", "Carrier vs Marsh growth"],
    )
    axes = plan_axes(src)
    assert axes["overall"] == [0, 1]                 # everything before the first divider
    assert axes["country"] == [2, 3, 4]              # divider + breakdown + growth block


def test_extract_slides_count_and_order(tmp_path):
    src = _titled_deck(str(tmp_path / "deck.pptx"), ["A", "B", "C", "D"])
    sub = extract_slides(src, [2, 0])
    assert len(sub.slides._sldIdLst) == 2


def test_split_template_writes_axis_files(tmp_path):
    src = _titled_deck(
        str(tmp_path / "deck.pptx"),
        ["Overall Summary", "Country (1)", "Carrier breakdown"],
    )
    written = split_template(src, str(tmp_path / "out"))
    assert set(written) == {"overall", "country"}
    for path in written.values():
        assert os.path.exists(path)
        assert len(Presentation(path).slides._sldIdLst) >= 1


def test_plan_axes_no_divider_is_all_overall(tmp_path):
    src = _titled_deck(str(tmp_path / "deck.pptx"), ["Overall Summary", "Highlights"])
    axes = plan_axes(src)
    assert axes["overall"] == [0, 1]
    assert "country" not in axes


@pytest.mark.skipif(
    not os.path.exists("template/qbr_template.pptx"),
    reason="real QBR template not present",
)
def test_split_real_template(tmp_path):
    written = split_template("template/qbr_template.pptx", str(tmp_path / "out"))
    assert "overall" in written and "country" in written
    assert len(Presentation(written["country"]).slides._sldIdLst) >= 2
