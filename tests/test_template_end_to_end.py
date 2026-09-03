"""End-to-end: a Setup selection → the assembled, populated ``.pptx``.

Runs the REAL pipeline over the REAL templates under ``template/`` against the seed DB:

    selection → compute_overall → per-axis sub-decks → fill (values · charts · commentary)
              → merge → one deck

and then reads the exported file back to assert the deck a user would open is populated:
no ``x``-placeholder survives, no authored example entity (Carrier X · QBE · Europe ·
Germany · Country (1)) survives, every page's charts carry real categories, and the
closing back cover is present. Deterministic: seed DB, no LLM.

Slow by nature (it fills and merges ten-odd decks), so it is marked ``e2e``:

    pytest -m "not e2e"      # skip
"""
from __future__ import annotations

import re

import pytest
from pptx import Presentation

from studio.compute import compute_overall
from studio.template_fill.analyze import analyze
from studio.template_fill.assemble import assemble_deck, plan_subdecks
from studio.template_fill.binding_map import available

pytestmark = pytest.mark.e2e

# Placeholder tokens ("$xx.xm", "x.x%", "#x") must all be gone from a populated deck.
_PLACEHOLDER = re.compile(r"(?<![A-Za-z$])\$?[xX]+(?:[.,][xX]+)*(?:[MBKmbk])?%?(?![A-Za-z])")
# Example entities the author baked into the templates; every one must be renamed.
_AUTHORED = ("Carrier X", "QBE", "Europe", "Germany", "Country (1)", "Region (1)", "xyz")
# Only "…" the author left as an explicit fill-me cue may remain, and only where premium
# data honestly cannot fill it (relationship feedback columns).
_SELECTION = {"carrier": "Zurich", "country": ["Singapore", "Japan"], "year": 2025}


def _deck_text(path: str) -> str:
    template = analyze(path)
    parts = []
    for slide in template.slides:
        for shape in slide.shapes:
            parts.append(shape.text)
            for row in (shape.table or []):
                parts.extend(row)
    return "\n".join(p for p in parts if p)


def _unfilled(path: str):
    """``[(slide_idx, shape text)]`` for every shape still showing a placeholder token."""
    out = []
    for slide in analyze(path).slides:
        for shape in slide.shapes:
            blocks = [shape.text] + [c for row in (shape.table or []) for c in row]
            for block in blocks:
                if any(_PLACEHOLDER.search(line) for line in block.splitlines()):
                    out.append((slide.index, " ".join(block.split())))
                    break
    return out


@pytest.fixture(scope="module")
def assembled(tmp_path_factory):
    tmp = tmp_path_factory.mktemp("e2e")
    result = compute_overall(filters=_SELECTION)
    assert result.subject == "Zurich"
    path = assemble_deck(result, out_path=str(tmp / "deck.pptx"), work_dir=str(tmp / "work"))
    return path, _deck_text(path)


def test_every_registered_axis_is_assembled():
    result = compute_overall(filters=_SELECTION)
    decks = plan_subdecks(result)
    axes = [d.template for d in decks]
    assert axes[0] == "overall"                       # the deck opens on the overall block
    assert "country" in available() and axes.count("country") == 2   # one per selected country
    assert axes[-1] == "end"                          # …and closes on the back cover


def test_no_placeholder_token_survives(assembled):
    """Every tile the deck's own database can answer must be filled.

    The seed database carries the premium book only, so the broker-survey tile has no
    source and correctly keeps its placeholder — anything ELSE left unfilled is a bug.
    """
    path, _ = assembled
    leftovers = [(idx, text) for idx, text in _unfilled(path) if "survey" not in text.lower()]
    assert not leftovers, f"unfilled placeholders in the export: {leftovers}"


def test_no_authored_example_entity_survives(assembled):
    _, text = assembled
    assert [name for name in _AUTHORED if name in text] == []


def test_the_carrier_and_its_countries_are_named(assembled):
    _, text = assembled
    assert "Zurich" in text
    for country in _SELECTION["country"]:
        assert country in text


def test_commentary_is_written_and_cites_figures(assembled):
    _, text = assembled
    # The authored example narrative (about Casualty at 2.3%) must be replaced by prose
    # built from this carrier's own facts.
    assert "Casualty at 2.3%" not in text
    assert "share of wallet" in text.lower()
    assert re.search(r"\$\d", text), "commentary should carry real currency figures"


def _commentary_boxes(path):
    """Every multi-point commentary box in the deck, as ``[(bullet, text)]`` per paragraph."""
    from pptx import Presentation
    from pptx.oxml.ns import qn

    def bullet(paragraph):
        pPr = paragraph._p.find(qn("a:pPr"))
        if pPr is None:
            return "inherit"
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
            el = pPr.find(qn(tag))
            if el is not None:
                return el.get("char") or tag.split(":")[1]
        return "inherit"

    boxes = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(paras) > 1 and any(len(p.text.split()) > 6 for p in paras):
                boxes.append([(bullet(p), p.text) for p in paras])
    return boxes


def test_commentary_is_written_as_bullet_points(assembled):
    path, _ = assembled
    boxes = _commentary_boxes(path)
    assert boxes, "expected multi-point commentary boxes in the deck"
    for box in boxes:
        bullets = [b for b, _ in box]
        # Every point carries a visible bullet — the author's own, or the standard one.
        assert all(b not in ("buNone", "inherit") for b in bullets), box
        # A box uses ONE marker throughout rather than mixing styles.
        assert len(set(bullets)) == 1, box
        # And no point is a run-on paragraph of several sentences joined together.
        assert all(txt.count(". ") <= 1 for _, txt in box), box


def test_no_empty_bullet_is_left_behind(assembled):
    """A box authored with four example lines but filled with two points must end up with
    two paragraphs — a blanked-but-kept paragraph would render as an empty bullet."""
    from pptx import Presentation

    path, _ = assembled
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paras = list(shape.text_frame.paragraphs)
            if len(paras) < 2 or not any(len(p.text.split()) > 6 for p in paras):
                continue
            assert all(p.text.strip() for p in paras), \
                f"blank paragraph left in {shape.name!r}: {[p.text for p in paras]}"


def test_every_chart_carries_real_categories(assembled):
    path, _ = assembled
    charts = [sh for slide in analyze(path).slides for sh in slide.shapes if sh.kind == "chart"]
    assert charts, "the deck should still contain its charts"
    for chart in charts:
        assert "Germany" not in chart.chart_categories
        assert not any(_PLACEHOLDER.fullmatch(str(c)) for c in chart.chart_categories)


def test_the_growth_bubble_is_series_per_line_of_business(assembled):
    path, _ = assembled
    bubbles = [sh for slide in analyze(path).slides for sh in slide.shapes
               if sh.kind == "chart" and "BUBBLE" in (sh.chart_type or "")]
    assert bubbles
    names = [name for name, _ in bubbles[0].chart_series]
    assert len(names) > 1 and all(names), "each bubble should be its own named LoB series"


def test_the_deck_opens_on_the_carrier_and_closes_on_the_back_cover(assembled):
    path, _ = assembled
    slides = analyze(path).slides
    assert slides[0].title().strip() == "Zurich"
    assert "back cover" in slides[-1].layout.lower()


def test_the_export_is_a_readable_presentation(assembled):
    path, _ = assembled
    assert len(Presentation(path).slides) == len(analyze(path).slides) > 1


def test_the_survey_tile_is_bound_even_though_this_database_cannot_fill_it():
    """The tile is MAPPED (so a warehouse with the survey flow fills it) — it is only the
    seed database's missing Carriers table that leaves the placeholder standing."""
    from studio.template_fill.binding_map import get_binding_map

    roles = {b.role for b in get_binding_map("overall").bindings}
    assert "survey_score" in roles


# ── the Setup selection governs page 2, not only the pages after it ──────────


def test_pinned_products_narrow_the_overall_summary_and_the_ranking_page():
    """The overall block used to drop the product filter outright, so page 2 reported the
    carrier's WHOLE book however narrow the Setup selection was — and the portfolio page
    behind it ranked every line of business rather than the ones that were picked.

    Runs the real plan over the real overall template against the seed DB.
    """
    picked = ["Cyber", "Marine"]
    whole = plan_subdecks(compute_overall(filters=_SELECTION))[0]
    narrow = plan_subdecks(compute_overall(filters={**_SELECTION, "product_line": picked}))[0]
    assert whole.template == narrow.template == "overall"

    # The headline figures follow the selection, so they are strictly smaller.
    assert 0 < narrow.values["carrier_gwp"] < whole.values["carrier_gwp"]
    assert 0 < narrow.values["marsh_gwp"] < whole.values["marsh_gwp"]

    # …and the portfolio/ranking page plots exactly the lines that were picked.
    ranked = {point["name"]
              for panel in narrow.values["lc_ranking"].values()
              for point in panel["points"]}
    assert ranked and ranked <= set(picked)


# ── the decomposition reaches the page a reader opens ───────────────────────


def test_the_deck_argues_from_named_industries_or_segments(assembled):
    """The change this suite exists to prove.

    Before, every commentary column was written from the same six scope-level figures, so
    a product page could only restate its own headline. A deck that never names an industry
    or a client segment has fallen back to that.
    """
    _, text = assembled
    named = [v for v in ("Renewable Energy", "Pharmaceuticals", "Technology & Telecom",
                         "Healthcare & Life Sciences", "Manufacturing", "Financial Services",
                         "Commercial", "Corporate", "Risk Management")
             if v in text]
    assert named, "no industry or client segment reached the deck's commentary"


def test_a_portfolio_wide_finding_is_not_printed_on_every_page(assembled):
    """The same industries are unwritten in every market, so the finding belongs to the
    page above. Said on each one it is the same sentence with a different figure, which the
    claim ledger cannot dedup because the string differs every time."""
    _, text = assembled
    absences = [line for line in text.split("\n") if "wrote none of it" in line]
    assert len(absences) <= 2, f"the same absence is argued {len(absences)} times:\n" + \
                               "\n".join(absences)


def test_no_column_closes_on_a_bare_product_instruction(assembled):
    """"Across the book the call is to defend Cyber, scale Financial Lines, fix Casualty"
    named six products, carried no figure, and was true of any carrier with six lines."""
    _, text = assembled
    low = text.lower()
    for slogan in ("across the book the call is to", "defend cyber", "scale financial lines",
                   "fix casualty", "selectively pursue property"):
        assert slogan not in low, slogan


def test_every_opportunity_bullet_carries_a_figure(assembled):
    """An opportunity with nothing at stake behind it is a heading, not a finding."""
    _, text = assembled
    for line in text.split("\n"):
        stripped = line.strip()
        if "wrote none of it" in stripped or "of premium is on the table" in stripped:
            assert "$" in stripped, stripped
