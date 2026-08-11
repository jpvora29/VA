"""Where the Carrier Survey slide lands, and when it is generated at all.

Hermetic, like ``test_template_assemble.py``: tiny in-memory templates and stubbed
resolvers, so it proves the COMPOSITION without a warehouse.
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from studio import compute as C
from studio.template_fill import assemble as A
from studio.template_fill import binding_map as BM


@pytest.fixture(autouse=True)
def _force_opc_merge(monkeypatch):
    monkeypatch.setenv("STUDIO_PPT_MERGE_ENGINE", "opc")


def _tiny(path: str, n: int) -> str:
    prs = Presentation()
    for i in range(n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = f"s{i}"
    prs.save(path)
    return path


@pytest.fixture
def axes(tmp_path, monkeypatch):
    BM._REGISTRY.clear()
    BM.get_binding_map.cache_clear()
    for name, n in (("overall", 2), ("country", 5), ("survey", 1), ("end", 1)):
        path = _tiny(str(tmp_path / f"{name}.pptx"), n)
        BM._REGISTRY[name] = (lambda name=name, path=path: BM.BindingMap(name, path, ()))
    monkeypatch.setattr(BM, "_discover_json_maps", lambda: None)
    monkeypatch.setattr(A, "resolve_roles", lambda r: {})
    monkeypatch.setattr(A, "resolve_roles_for_product", lambda r, p: {})
    monkeypatch.setattr(A, "resolve_roles_for_country", lambda r, c: {})
    monkeypatch.setattr(A.survey_facts, "has_survey_data", lambda r, c: c != "Atlantis")
    yield
    BM._REGISTRY.clear()
    BM.get_binding_map.cache_clear()


def _result(countries=("Singapore", "Japan")):
    return C.OverallResult(subject="Zurich", flow="gpr",
                           resolved_filters={"Country": tuple(countries)})


def test_premium_only_generates_no_survey_slide(axes):
    decks = A.plan_subdecks(_result(), data_basis="premium")
    assert "survey" not in [d.template for d in decks]


def test_omitting_data_basis_keeps_todays_deck(axes):
    assert "survey" not in [d.template for d in A.plan_subdecks(_result())]


def test_survey_follows_each_country_block(axes):
    decks = A.plan_subdecks(_result(), data_basis="premium_survey")
    assert [d.template for d in decks] == [
        "overall", "country", "survey", "country", "survey", "end"]


def test_the_survey_deck_is_labelled_and_scoped_to_its_country(axes):
    decks = A.plan_subdecks(_result(), data_basis="premium_survey")
    surveys = [d for d in decks if d.template == "survey"]
    assert [d.label for d in surveys] == ["Singapore survey", "Japan survey"]
    assert surveys[0].values["country_name[0]"] == "Singapore"


def test_a_country_with_no_survey_book_is_skipped(axes):
    decks = A.plan_subdecks(_result(("Singapore", "Atlantis")), data_basis="premium_survey")
    assert [d.template for d in decks] == ["overall", "country", "survey", "country", "end"]


def test_an_overall_only_scope_generates_no_survey_slide(axes):
    decks = A.plan_subdecks(_result(), scope="overall", data_basis="premium_survey")
    assert [d.template for d in decks] == ["overall", "end"]


def test_the_merged_deck_has_six_slides_per_country(axes, tmp_path):
    out = A.assemble_deck(_result(), out_path=str(tmp_path / "deck.pptx"),
                          work_dir=str(tmp_path / "work"), data_basis="premium_survey")
    # overall(2) + 2 x (country 5 + survey 1) + end(1)
    assert len(Presentation(out).slides._sldIdLst) == 15


def test_generate_passes_the_selection_s_data_basis_through():
    from pathlib import Path

    src = Path("studio/authoring/generate.py").read_text(encoding="utf-8")
    assert 'data_basis=selection.get("data_basis")' in src


# ── the summary page's overall survey-score tile ─────────────────────────────
# The tile is sourced from the survey book, so it is subject to the same DATA BASIS gate as
# the Carrier Survey page: filled when the run asked for that book, off the page otherwise.


def _tile_template(tmp_path, caption: str = "Overall Carrier Survey"):
    """A one-slide template carrying a score tile with ``caption`` under its placeholder."""
    from studio.template_fill.analyze import analyze

    path = str(tmp_path / "tile.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.text = "x.x"
    box.text_frame.add_paragraph().text = caption
    prs.save(path)
    return analyze(path)


def _tile_result(basis: str):
    return C.OverallResult(subject="Zurich", flow="gpr", data_basis=basis,
                           resolved_filters={"Country": ("Singapore",)})


def test_the_tile_is_filled_on_the_survey_basis(tmp_path, monkeypatch):
    from studio.template_fill.survey import kpi

    monkeypatch.setattr(kpi.survey_facts, "load_overall_score", lambda r, c=(): 7.2)
    values = kpi.values(_tile_template(tmp_path), _tile_result("premium_survey"))
    assert values == {"survey_score": 7.2}


def test_the_tile_comes_off_the_page_on_the_premium_basis(tmp_path, monkeypatch):
    """A premium-basis run was not asked about the survey book, so the tile must not be
    filled from it — and must not ship the template's own ``x.x`` either."""
    from studio.template_fill.survey import kpi

    monkeypatch.setattr(kpi.survey_facts, "load_overall_score", lambda r, c=(): 7.2)
    template = _tile_template(tmp_path)
    values = kpi.values(template, _tile_result("premium"))
    assert values == {"drop_shapes": kpi.tiles(template)} and values["drop_shapes"]


def test_the_tile_comes_off_when_the_survey_book_has_no_score(tmp_path, monkeypatch):
    from studio.template_fill.survey import kpi

    monkeypatch.setattr(kpi.survey_facts, "load_overall_score", lambda r, c=(): None)
    template = _tile_template(tmp_path)
    assert kpi.values(template, _tile_result("premium_survey")) == {
        "drop_shapes": kpi.tiles(template)}


def test_the_tile_is_scoped_to_the_countries_the_run_pinned(tmp_path, monkeypatch):
    from studio.template_fill.survey import kpi

    seen = []

    def record(result, countries=()):
        seen.append(tuple(countries))
        return 7.2

    monkeypatch.setattr(kpi.survey_facts, "load_overall_score", record)
    kpi.values(_tile_template(tmp_path), _tile_result("premium_survey"))
    assert seen == [("Singapore",)]


def test_a_template_with_no_such_tile_is_left_alone(tmp_path):
    """Detection is by caption, so the product/country decks — and the Carrier Survey page,
    whose TITLE is "Carrier Survey" — are untouched."""
    from studio.template_fill.survey import kpi

    assert kpi.tiles(_tile_template(tmp_path, caption="Carrier Survey")) == []
    assert kpi.values(_tile_template(tmp_path, caption="Carrier Rank"),
                      _tile_result("premium")) == {}


def test_a_dropped_shape_is_gone_from_the_exported_file(tmp_path):
    """The fill engine honours ``drop_shapes`` — the tile leaves the .pptx, it is not blanked."""
    from studio.template_fill.fill import fill_template

    src = str(tmp_path / "src.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    keep = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    keep.text_frame.text = "keep me"
    drop = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
    drop.text_frame.text = "x.x Overall Carrier Survey"
    prs.save(src)

    out = fill_template({"template_path": src, "manifest": [], "overrides": {}, "added": {},
                         "values": {"drop_shapes": [f"0:{drop.shape_id}"]}},
                        out_path=str(tmp_path / "out.pptx"))
    texts = [sh.text_frame.text for sh in Presentation(out).slides[0].shapes
             if sh.has_text_frame]
    assert texts == ["keep me"]


# ── payloads several providers contribute to ─────────────────────────────────
#
# A sub-deck's values are built by folding each provider's output in turn. Most keys are one
# role's own text and belong to one provider, but the fill-engine payloads are addressed by
# SHAPE, so any page can contribute — and a plain dict update would let whichever provider
# ran last silently erase the others' entries. Two providers already write `drop_shapes`
# (the survey tile, the GWP page's country chart); today's templates keep them on separate
# decks, so these pin the behaviour before a re-authored template makes the clash live.


def test_two_providers_writing_the_same_payload_are_combined():
    merged = A._merge_values({"drop_shapes": ["0:1"], "cell_fills": {"0:9": [{"r": 0}]}},
                             {"drop_shapes": ["0:2"], "cell_fills": {"0:8": [{"r": 1}]}})
    assert merged["drop_shapes"] == ["0:1", "0:2"]
    assert set(merged["cell_fills"]) == {"0:9", "0:8"}


def test_the_same_shape_named_twice_is_dropped_once():
    merged = A._merge_values({"drop_shapes": ["0:1", "0:2"]}, {"drop_shapes": ["0:2", "0:3"]})
    assert merged["drop_shapes"] == ["0:1", "0:2", "0:3"]


def test_an_ordinary_role_is_still_overwritten_by_the_later_provider():
    """Only the shared payloads combine — a role's text has exactly one author."""
    assert A._merge_values({"title": "old"}, {"title": "new"})["title"] == "new"


def test_a_provider_that_writes_nothing_leaves_the_payload_alone():
    base = {"drop_shapes": ["0:1"], "resize_shapes": {"0:5": {"x": 1}}}
    assert A._merge_values(base, {"title": "t"})["drop_shapes"] == ["0:1"]
    assert A._merge_values({}, base)["resize_shapes"] == {"0:5": {"x": 1}}


def test_the_real_provider_chain_keeps_both_providers_drops(axes, monkeypatch):
    """The fix where it actually bites: two providers on one sub-deck, both dropping a
    shape the page cannot honestly fill. Before, the second erased the first."""
    monkeypatch.setattr(A, "prune", type("P", (), {"hidden_country_pages": staticmethod(
        lambda template, n: ())})())
    providers = (lambda template, result: {"drop_shapes": ["0:11"], "gwp_bars": {"0:1": {}}},
                 lambda template, result: {"drop_shapes": ["0:22"]})

    sub = A._build_subdeck("overall", C.OverallResult(subject="Zurich"), {}, "overall",
                           providers=providers)

    assert sub.values["drop_shapes"] == ["0:11", "0:22"]
    assert sub.values["gwp_bars"] == {"0:1": {}}
