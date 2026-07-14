"""Regression tests for the generated QBR: temporal views, advanced widgets,
country/SWOT pages, and screen + PPTX rendering of the new deck blocks.

All deterministic — they run against the seed DB (no DB_PATH, no LLM).
"""
from __future__ import annotations

from core.analytics.library import compute_period_series
from core.analytics.types import PrimitiveArgs
from pptx import Presentation

from studio.compute import compute_overall, mom, opportunity_matrix, qoq, ttm
from studio.data import get_engine
from studio.deck import build_deck
from studio.deck.model import ChartBlock, HeatmapBlock, MatrixBlock, RadarBlock
from studio.export import export_deck
from studio.page import slide as SL

ENG = get_engine()
GPR = "gpr"
_SCOPE = {"Carrier_Group": "Zurich", "Country": "Singapore"}


# ── temporal engine (TTM / QoQ / MoM) ────────────────────────────────────────


def test_seed_has_monthly_granularity():
    series = compute_period_series(
        PrimitiveArgs(flow=GPR, metric="premium", filters={"Carrier_Group": "Zurich"}),
        grain="month", engine=ENG,
    )
    assert len(series) == 36  # 3 years × 12 months, derived from Billing_Date
    assert series[0].dims["period"] == "2023-01"
    assert series[-1].dims["period"] == "2025-12"


def test_ttm_qoq_mom_views_compute():
    t = ttm(GPR, _SCOPE, ENG)
    assert t and len(t["values"]) == 36 - 11  # rolling 12-month windows
    assert t["ttm_pct"] is not None
    assert t["current"] > 0

    m = mom(GPR, _SCOPE, ENG)
    assert len(m["values"]) == 35 and m["latest"] is not None

    q = qoq(GPR, _SCOPE, ENG)
    assert len(q["values"]) >= 8  # ~11 quarter-over-quarter changes across 3 years


def test_monthly_split_preserves_annual_totals():
    """The monthly seed must sum back to the annual figure YoY/breakdowns rely on."""
    from core.analytics.library import compute_breakdown

    facts = compute_breakdown(
        PrimitiveArgs(flow=GPR, metric="premium", filters={"Carrier_Group": "Zurich", "Year": 2025}),
        engine=ENG,
    )
    monthly = sum(
        f.value
        for f in compute_period_series(
            PrimitiveArgs(flow=GPR, metric="premium", filters={"Carrier_Group": "Zurich"}),
            grain="month", engine=ENG,
        )
        if str(f.dims["period"]).startswith("2025")
    )
    assert abs(facts[0].value - monthly) < 1.0  # rounding tolerance


# ── deterministic, premium-derived opportunity matrix (never random) ─────────


def test_movers_table_shows_both_years_and_yoy_not_absolute_change():
    from studio.content.evidence import _MOVE_COLS, _mover_label

    labels = [c["label"] for c in _MOVE_COLS]
    assert "Prior yr" in labels and "Current yr" in labels and "YoY" in labels
    assert "Change" not in labels  # the absolute-change column is gone
    # >100% growth shows BOTH years; ≤100% does not.
    hot = _mover_label({"name": "Cyber", "pct": 150.0, "prior": 10_000_000, "current": 25_000_000})
    assert "150" in hot and "→" in hot
    mild = _mover_label({"name": "Property", "pct": 12.0, "prior": 1, "current": 1})
    assert "→" not in mild and "12" in mild


def test_no_hhi_anywhere_in_the_deck():
    res = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"], "year": 2025}, engine=ENG)
    deck = build_deck(res, carrier="Zurich", country="Singapore", year=2025)
    blob = " ".join(
        s.title + " ".join(t.get("text", "") for t in s.takeaways) + " " + " ".join(str(e) for e in s.evidence)
        for s in deck.slides
    )
    assert "HHI" not in blob and "hhi" not in blob


def test_year_multiselect_uses_latest_for_comparison():
    from studio.compute import _current_year, _resolve_filters

    assert _current_year(_resolve_filters({"year": [2023, 2025, 2024]})) == 2025
    assert _current_year(_resolve_filters({"year": 2024})) == 2024
    assert _current_year({}) is None
    deck = build_deck(
        compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"], "year": [2024, 2025]}, engine=ENG),
        carrier="Zurich", country="Singapore", year=2025,
    )
    assert len(deck.slides) > 5


def test_swot_strengths_include_premium_scale_and_share_of_portfolio():
    from studio.narrate.commentary import build_swot

    res = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"], "year": 2025}, engine=ENG)
    strengths = " ".join(build_swot(res).strengths)
    assert "of the book" in strengths  # high premium + share-of-portfolio strength


def test_peer_average_is_mean_of_peer_totals_not_per_row():
    """Regression for the ~899x bug: peer benchmark = average of each peer's TOTAL."""
    from core.analytics.library import compute_breakdown
    from core.analytics.types import PrimitiveArgs
    from studio.compute import _resolve_filters, peer_gap

    f = _resolve_filters({"carrier": "Zurich", "country": ["Singapore"], "year": 2025})
    pg = peer_gap("gpr", f, ENG)
    assert pg and pg["n_peers"] == 4  # Zurich's seed peer set
    assert 0.1 < pg["ratio"] < 10  # like-for-like totals, not the old per-row blow-up

    # peer_avg must equal the mean of each peer carrier's independently-computed total.
    totals = []
    for peer in ("AXA XL", "Allianz", "Chubb", "AIG"):
        facts = compute_breakdown(
            PrimitiveArgs(flow="gpr", metric="premium", group_by=(), filters={**f, "Carrier_Group": peer}),
            engine=ENG,
        )
        if facts:
            totals.append(facts[0].value)
    expected = sum(totals) / len(totals)
    assert abs(pg["peer_avg"] - round(expected, 2)) < 1.0


def test_opportunity_matrix_is_deterministic_and_premium_derived():
    res = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"], "year": 2025}, engine=ENG)
    p1 = opportunity_matrix(res.flow, res.resolved_filters, res.engine, res.whitespace)
    p2 = opportunity_matrix(res.flow, res.resolved_filters, res.engine, res.whitespace)
    assert p1 and p1 == p2  # identical on repeat → no randomness
    for pt in p1:
        assert 0 <= pt["x"] <= 100 and 0 <= pt["y"] <= 100 and pt["size"] > 0
    # potential axis is market-premium-scaled: the biggest candidate sits at 100.
    assert max(pt["y"] for pt in p1) == 100


# ── country pages + SWOT tick are gated by the cuts ──────────────────────────


def _deck(cuts):
    res = compute_overall(
        filters={"carrier": "Zurich", "country": ["Singapore", "Hong Kong"], "year": 2025}, engine=ENG
    )
    return build_deck(res, carrier="Zurich", country="Singapore", year=2025, report="qbr", cuts=cuts)


def test_country_and_swot_slides_only_appear_when_ticked():
    base = _deck(())
    assert not any((s.eyebrow or "").startswith("GEOGRAPHIC") for s in base.slides)
    assert not any(s.layout == "swot" for s in base.slides)

    full = _deck(("country_breakdown", "country_swot"))
    geo = [s for s in full.slides if (s.eyebrow or "").startswith("GEOGRAPHIC")]
    swot = [s for s in full.slides if s.layout == "swot"]
    assert len(geo) == 2  # one breakdown slide per selected country
    assert len(swot) == 2  # one SWOT slide per selected country
    assert {s.eyebrow for s in swot} == {"SWOT · SINGAPORE", "SWOT · HONG KONG"}


def test_performance_slide_carries_ttm_trend_and_kpi_band():
    deck = _deck(())
    perf = next(s for s in deck.slides if (s.eyebrow or "").startswith("PERFORMANCE"))
    kinds = [b.kind for b in perf.blocks]
    assert "kpis" in kinds
    line = next(b for b in perf.blocks if b.kind == "chart" and b.chart == "line")
    assert line.labels and line.values  # the trailing-12-month trend


# ── new deck blocks render on screen and export as native PPTX ───────────────


def test_advanced_blocks_render_on_screen():
    blocks = [
        MatrixBlock([{"label": "Cyber", "x": 50, "y": 90, "size": 40}], "matrix"),
        HeatmapBlock(["Property"], ["SG", "HK"], [[80, 40]], "heatmap"),
        RadarBlock(["Rank", "Growth", "Wallet"], [70, 60, 50], "radar"),
        ChartBlock("line", ["Jan", "Feb", "Mar"], [1, 2, 3], "trend"),
    ]
    for block in blocks:
        assert SL._visual(block) is not None


def test_generated_advanced_widgets_export_as_native_powerpoint(tmp_path):
    deck = _deck(("country_breakdown", "country_swot", "rank_similar", "peer_average"))
    out = tmp_path / "qbr.pptx"
    export_deck(deck, out_path=str(out))
    prs = Presentation(str(out))
    assert len(prs.slides) == len(deck.slides)
    # Every content slide (matrix/heatmap/radar included) lays down real shapes,
    # not a single blank placeholder.
    contentful = [sl for sl in prs.slides if len(sl.shapes) > 3]
    assert len(contentful) >= 10


def test_stale_persisted_template_doc_is_dropped(tmp_path):
    # qs-tdoc persists in browser local storage; an assembled doc points at a temp
    # .pptx that Windows cleans between sessions. The master render must drop it
    # (crash regression: FileNotFoundError "<carrier>_<hash>_QBR.pptx" on app open).
    from studio.authoring.generate import usable_tdoc

    stale = {"template_path": str(tmp_path / "AIG_GROUP_deadbeef_QBR.pptx"), "assembled": True}
    assert usable_tdoc(stale) is None
    assert usable_tdoc(None) is None

    live_path = tmp_path / "live.pptx"
    live_path.write_bytes(b"x")
    live = {"template_path": str(live_path)}
    assert usable_tdoc(live) is live
