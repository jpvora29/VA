"""The Minutes-of-Meeting workspace, from the pure helpers up to the whole workflow.

The integration tests at the bottom run the REAL pipeline over a real .pptx and a real
note file, with the model replaced by a stub. That is the point of injecting
``call_json``: the extraction, tagging, scoring, verification, DOCX writing and the
files on disk are all exercised without a credential.
"""
from __future__ import annotations

import json
import re
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from mom import config, jobs, tagger, verifier
from mom.minutes_docx import agenda_items, write_minutes, xml_safe
from mom.modes import AI_SUMMARY, DEFAULT_MODE, MODES, SELF_NOTES, resolve_mode
from mom.notes import UnsupportedNoteFormat, read_note
from mom.pipeline import (
    MoMRequest,
    minutes_filename,
    run_mom_pipeline,
    sections_from_manifest,
)
from mom.progress import PHASES, label_for, percent_done
from mom.summariser import SHAPES, build_evidence_block, clean_bullets
from mom.uploads import UploadRejected, safe_name, stage_upload


# ── the stub model ───────────────────────────────────────────────────────────


class StubModel:
    """Answers each pipeline prompt by its label, and records what it was asked.

    It reads the ids out of the prompt and tags them round-robin from the real tag
    list, so the answers are always pairs the pipeline will accept — which is what
    makes the whole run reach the DOCX.
    """

    def __init__(self, tags, verdict: str = "pass"):
        self.tags = list(tags)
        self.verdict = verdict
        self.labels: list[str] = []

    def __call__(self, prompt: str, *, label: str = "", phase: str = "unknown") -> dict:
        self.labels.append(label)
        if label == "parse_meeting_note":
            return self._meeting_note()
        if label == "tag_meeting_note":
            return self._assign(prompt, r"\[(mn_\d+)\]", "id")
        if label.startswith("tag_"):
            return self._assign(prompt, r"\[(entry_\d+)\]", "entry_id")
        if label == "verify_tagging":
            return {"verdict": self.verdict, "issues": [], "notes": "", "corrections": []}
        if label == "summary_generation":
            return self._summary(prompt)
        raise AssertionError(f"the pipeline asked for an unexpected label: {label!r}")

    def _meeting_note(self) -> dict:
        return {
            "client": None,
            "subject": None,
            "meeting_date": None,
            "location": "London",
            "author": "Marsh ICG",
            "attendees": [],
            "discussion_points": [
                {"topic": "Performance", "bullets": [
                    "Global GWP reached $2.0bn, up 4.2% year on year.",
                    "Share of wallet improved to 2.7%.",
                    "Rank moved up one place to 8th globally.",
                ]},
                {"topic": "Marsh Update", "bullets": [
                    "Broker Workbench rollout continues across renewals.",
                    "A dedicated London team now services wholesale business.",
                    "Cargo LATAM facility participation was secured.",
                ]},
            ],
            "action_items": [
                {"action": "Share updated pricing", "owner": "Marsh", "due_date": "Q3 2026"},
                {"action": "Arrange a cyber roundtable", "owner": None, "due_date": None},
            ],
        }

    def _assign(self, prompt: str, pattern: str, key: str) -> dict:
        ids = re.findall(pattern, prompt)
        # Only ids from the "Now tag the following" half — the few-shot block above it
        # uses ex_NN, so nothing there matches either pattern.
        results = []
        for index, item_id in enumerate(dict.fromkeys(ids)):
            tag = self.tags[index % len(self.tags)]
            results.append({key: item_id, "umbrella_tag": tag.umbrella_tag,
                            "sub_tag": tag.sub_tag})
        return {"results": results}

    def _summary(self, prompt: str) -> dict:
        if "carrier_update" in prompt:
            return {
                "carrier_update": ["Carrier: GWP grew 4.2% to $2.0bn.", "Rank improved to 8th."],
                "marsh_update": ["Broker Workbench rollout continues.", ""],
            }
        return {
            "strategy_and_initiatives": ["Broker Workbench rollout continues."],
            "country_product_region": ["Cargo LATAM facility participation secured."],
            "key_takeaways": ["GWP grew 4.2% to $2.0bn; rank improved to 8th."],
        }


@pytest.fixture
def tags():
    return tagger.load_tag_list(config.tag_list_path())


@pytest.fixture
def model(tags):
    return StubModel(tags)


# ── fixtures on disk ─────────────────────────────────────────────────────────


# Layout 5 is "Title Only" and layout 2 is "Section Header" in the default template.
# The extractor reads a real title PLACEHOLDER (not a text box) for the deck's client
# and slide titles, and treats a "section"-named layout as a divider — so the fixture
# has to use the placeholders a real deck uses.
_TITLE_ONLY, _SECTION_HEADER = 5, 2


def _add_text_slide(presentation, title: str, bullets=(), layout: int = _TITLE_ONLY):
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout])
    slide.shapes.title.text = title
    if bullets:
        body = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(4))
        frame = body.text_frame
        frame.text = bullets[0]
        for line in bullets[1:]:
            frame.add_paragraph().text = line
    return slide


def _add_divider(presentation, title: str):
    return _add_text_slide(presentation, title, layout=_SECTION_HEADER)


def _add_table(slide, rows, left_in=0.5, top_in=1.6):
    shape = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(left_in), Inches(top_in), Inches(8), Inches(2)
    )
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            shape.table.cell(r, c).text = value
    return shape


@pytest.fixture
def deck_path(tmp_path) -> Path:
    """A small but realistic QBR deck: cover, agenda + attendees, three sections."""
    presentation = Presentation()

    cover = _add_text_slide(presentation, "Zurich - Marsh ICG QBR")
    note = cover.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.6))
    note.text_frame.text = "March 2026, London"

    front = _add_text_slide(presentation, "Agenda and attendees")
    _add_table(front, [["Agenda", ""], ["1", "Performance"], ["2", "Marsh Update"],
                       ["3", "Key Takeaways"]])
    _add_table(front, [["Company", "Name", "Role"],
                       ["Zurich", "A. Underwriter", "Head of Specialty"],
                       ["Marsh", "B. Broker", "Client Executive"]], top_in=4.0)

    _add_divider(presentation, "Performance")
    _add_text_slide(presentation, "Global performance", [
        "Global GWP with Marsh reached $2.0bn, up 4.2% year on year.",
        "Share of wallet 2.7%, up 1pp; ranked 8th globally, up one place.",
        "North America GWP $1.0bn, up 7.7% year on year.",
    ])
    _add_divider(presentation, "Marsh Update")
    _add_text_slide(presentation, "Placement priorities", [
        "Broker Workbench rollout continues across new business and renewals.",
        "A dedicated London team now services wholesale business.",
    ])
    _add_divider(presentation, "Key Takeaways")
    _add_text_slide(presentation, "Headlines", [
        "Commercial insurance rates decreased 5% in the quarter.",
        "Cyber and Crisis Management are outperforming the Marsh portfolio.",
    ])

    path = tmp_path / "Zurich UK QBR.pptx"
    presentation.save(str(path))
    return path


@pytest.fixture
def note_docx(tmp_path) -> Path:
    document = Document()
    document.add_heading("QBR meeting note", level=1)
    for line in ("Global GWP reached $2.0bn, up 4.2%.",
                 "Broker Workbench rollout continues.",
                 "Rates decreased 5% in the quarter."):
        document.add_paragraph(line, style="List Bullet")
    table = document.add_table(rows=2, cols=3)
    for index, value in enumerate(("Action", "Owner", "Timeline")):
        table.rows[0].cells[index].text = value
    for index, value in enumerate(("Share updated pricing", "Marsh", "Q3 2026")):
        table.rows[1].cells[index].text = value
    path = tmp_path / "self notes.docx"
    document.save(str(path))
    return path


@pytest.fixture
def note_pdf(tmp_path) -> Path:
    pymupdf = pytest.importorskip("pymupdf")
    document = pymupdf.open()
    page = document.new_page()
    page.insert_text((72, 96), "QBR meeting note\nGlobal GWP reached $2.0bn, up 4.2%.")
    path = tmp_path / "ai summary.pdf"
    document.save(str(path))
    document.close()
    return path


@pytest.fixture
def run_paths(tmp_path) -> config.RunPaths:
    return config.RunPaths(tmp_path / "run").create()


# ── modes and progress ───────────────────────────────────────────────────────


def test_the_two_modes_differ_only_in_their_three_knobs():
    """Everything before the last phase is one code path; the mode is data."""
    assert {m.id for m in MODES} == {"ai_summary", "self_notes"}
    assert AI_SUMMARY.granularity == "slide" and AI_SUMMARY.summary == "skill"
    assert SELF_NOTES.granularity == "section" and SELF_NOTES.summary == "carrier_marsh"
    assert set(SHAPES) == {m.summary for m in MODES}


@pytest.mark.parametrize("value", [None, "", "bogus", "AI_SUMMARY", 7])
def test_an_unknown_mode_falls_back_rather_than_raising(value):
    """A stale store value must not take the whole workspace down."""
    assert resolve_mode(value).id == DEFAULT_MODE


def test_self_notes_accepts_the_formats_the_reader_can_actually_read():
    """The upload zone and the reader must agree, or a .docx crashes the run."""
    from mom.notes import READERS

    for mode in MODES:
        accepted = {ext.strip() for ext in mode.accept.split(",")}
        assert accepted <= set(READERS), f"{mode.id} accepts what nothing can read"


def test_progress_reads_the_phase_not_the_prose():
    """The old bar grepped stdout for keywords, so a reworded print broke it."""
    percents = [percent_done(phase.id) for phase in PHASES]
    assert percents == sorted(percents)
    assert percents[0] == 0
    assert percent_done("summary", finished=True) == 100
    assert percent_done("not-a-phase") == 0
    assert label_for("tagging") and label_for(None) == "Working"


# ── the tag vocabulary ───────────────────────────────────────────────────────


def test_the_tag_list_ships_with_the_package(tags):
    """A fresh checkout must be able to run; the CSV is not an external input."""
    assert config.tag_list_path().is_file()
    assert len(tags) == 10
    assert len({tag.pair for tag in tags}) == len(tags), "duplicate sub-tags"
    assert tagger.scored_umbrellas(tags) == {
        "Strategy & Initiatives", "Country / Product / Region", "Key Takeaways"
    }


def test_the_tag_list_and_the_few_shot_examples_agree(tags):
    """Calibrating the model against a sub-tag it can never assign teaches it nothing."""
    from mom.examples import FEW_SHOT_EXAMPLES

    for tag in tags:
        assert tag.sub_tag in FEW_SHOT_EXAMPLES, tag.sub_tag


def test_a_tag_list_without_the_required_columns_is_refused(tmp_path):
    path = tmp_path / "bad.csv"
    path.write_text("theme,detail\na,b\n", encoding="utf-8")
    with pytest.raises(ValueError, match="umbrella_tag"):
        tagger.load_tag_list(path)


def test_only_pairs_from_the_tag_list_survive(tags):
    """The model may not invent a category; an unknown pair becomes Unclassified."""
    answer = {"results": [
        {"id": "mn_0001", "umbrella_tag": "Key Takeaways",
         "sub_tag": "KPIs & Performance Headlines"},
        {"id": "mn_0002", "umbrella_tag": "Vibes", "sub_tag": "Made Up"},
    ]}
    assigned = tagger._read_assignments(answer, "id", tags)
    assert assigned["mn_0001"] == ("Key Takeaways", "KPIs & Performance Headlines")
    assert assigned["mn_0002"] == (config.UNCLASSIFIED, config.UNCLASSIFIED)


def test_every_bullet_gets_an_id_and_an_owner():
    meeting = {"discussion_points": [
        {"topic": "Zurich update", "bullets": ["One.", "  ", "Two."]},
        {"topic": "Market", "bullets": ["Three."]},
    ]}
    bullets = tagger.collect_bullets(meeting, "Zurich")
    assert [b["id"] for b in bullets] == ["mn_0001", "mn_0002", "mn_0003"]
    assert bullets[0]["bullet_owner"] == "Zurich"
    assert bullets[2]["bullet_owner"] == "Marsh"


def test_an_entry_with_nothing_in_it_is_never_sent_to_the_model():
    assert tagger.compress_entry({"content_type": "table", "text": "x"}) is None
    assert tagger.compress_entry({"content_type": "text", "text": " a "}) is None
    assert tagger.compress_entry({"content_type": "text", "text": "GWP up 4.2%"}) == "GWP up 4.2%"
    row = {"content_type": "table_row", "text": "", "row_label": "Property",
           "cells": [{"column_name": "GWP", "value": "$2.0bn"}]}
    assert tagger.compress_entry(row) == "Property — GWP: $2.0bn"


# ── scoring ──────────────────────────────────────────────────────────────────


def _bullet(sub_tag, umbrella="Key Takeaways", index=1):
    return {"id": f"mn_{index:04d}", "topic": "T", "bullet": f"b{index}",
            "bullet_owner": "Marsh", "umbrella_tag": umbrella, "sub_tag": sub_tag}


def test_a_meeting_note_bullet_outweighs_a_deck_entry(tmp_path, tags):
    """The note records what was discussed; the deck only shows what was prepared."""
    tagged = tmp_path / "tagged"
    tagged.mkdir()
    (tagged / "slide_001.json").write_text(json.dumps({"entries": [
        {"umbrella_tag": "Key Takeaways", "sub_tag": "Outperformers & Underperformers",
         "text": "Cyber up 39.5%"},
    ]}), encoding="utf-8")

    pairs = tagger.score_priorities(
        [_bullet("KPIs & Performance Headlines")], tagged, tagger.scored_umbrellas(tags)
    )
    by_sub = {p["sub_tag"]: p["score"] for p in pairs}
    assert by_sub["KPIs & Performance Headlines"] == config.MEETING_NOTE_WEIGHT
    assert by_sub["Outperformers & Underperformers"] == config.PPT_WEIGHT
    assert config.MEETING_NOTE_WEIGHT > config.PPT_WEIGHT


def test_one_umbrella_cannot_fill_the_priority_list(tmp_path, tags):
    """Without the cap, ten Key Takeaways sub-tags would crowd out every other theme."""
    tagged = tmp_path / "tagged"
    tagged.mkdir()
    key_subs = [t.sub_tag for t in tags if t.umbrella_tag == "Key Takeaways"]
    bullets = [_bullet(sub, index=i) for i, sub in enumerate(key_subs * 2, 1)]

    pairs = tagger.score_priorities(bullets, tagged, tagger.scored_umbrellas(tags))
    assert len(pairs) <= config.UMBRELLA_CAP


def test_unclassified_content_never_reaches_the_minutes(tmp_path, tags):
    tagged = tmp_path / "tagged"
    tagged.mkdir()
    bullets = [_bullet(config.UNCLASSIFIED, umbrella=config.UNCLASSIFIED)]
    assert tagger.score_priorities(bullets, tagged, tagger.scored_umbrellas(tags)) == []


def test_a_corrupt_tagged_file_does_not_lose_the_run(tmp_path, tags):
    tagged = tmp_path / "tagged"
    tagged.mkdir()
    (tagged / "slide_001.json").write_text("{ not json", encoding="utf-8")
    pairs = tagger.score_priorities(
        [_bullet("KPIs & Performance Headlines")], tagged, tagger.scored_umbrellas(tags)
    )
    assert [p["sub_tag"] for p in pairs] == ["KPIs & Performance Headlines"]


# ── verification ─────────────────────────────────────────────────────────────


def test_too_few_priority_pairs_stops_the_run():
    result = verifier.verify_rules([{"umbrella_tag": "a", "sub_tag": "b", "score": 2}], [])
    assert not result.passed and result.errors


def test_mostly_unclassified_bullets_stop_the_run():
    """Thin minutes are worse than none: if nothing matched, say so instead."""
    pairs = [{"umbrella_tag": "u", "sub_tag": f"s{i}", "score": 2,
              "meeting_note_bullets": [1], "ppt_entries": []} for i in range(5)]
    bullets = [_bullet(config.UNCLASSIFIED, index=i) for i in range(10)]
    assert not verifier.verify_rules(pairs, bullets).passed


def test_a_healthy_run_passes_the_rules():
    pairs = [{"umbrella_tag": "u", "sub_tag": f"s{i}", "score": 4,
              "meeting_note_bullets": [1], "ppt_entries": [1]} for i in range(5)]
    bullets = [_bullet("KPIs & Performance Headlines", index=i) for i in range(10)]
    result = verifier.verify_rules(pairs, bullets)
    assert result.passed and not result.warnings


def test_a_correction_is_applied_only_when_it_names_a_real_pair(tags):
    bullets = [_bullet("KPIs & Performance Headlines", index=1)]
    good = {"id": "mn_0001", "umbrella_tag": "Strategy & Initiatives",
            "sub_tag": "Digital & Innovation"}
    assert verifier.apply_corrections([good], bullets, tags) == 1
    assert bullets[0]["sub_tag"] == "Digital & Innovation"

    invented = {"id": "mn_0001", "umbrella_tag": "Vibes", "sub_tag": "Made Up"}
    unknown = {"id": "mn_9999", "umbrella_tag": "Key Takeaways",
               "sub_tag": "KPIs & Performance Headlines"}
    assert verifier.apply_corrections([invented, unknown], bullets, tags) == 0
    assert bullets[0]["sub_tag"] == "Digital & Innovation"


def test_an_audit_call_that_fails_does_not_fail_the_run():
    """The audit is advisory. Losing it must not lose an otherwise good run."""
    def broken(prompt, *, label="", phase=""):
        raise RuntimeError("no credentials")

    result = verifier.verify_with_llm([], {}, [], broken)
    assert result.passed and result.warnings


def test_the_checkpoint_raises_instead_of_exiting_the_process(tags):
    """The standalone version raised SystemExit here, which inside a web app would
    have killed the worker thread with nothing to show the user."""
    thin = tagger.TaggingResult(top_pairs=[], tagged_bullets=[], tags=list(tags))
    with pytest.raises(verifier.VerificationFailed):
        verifier.run_checkpoint(
            thin, {}, tagged_dir=Path("."), priority_path=Path("x.json"), llm_check=False
        )


# ── reading the note ─────────────────────────────────────────────────────────


def test_a_word_note_reads_its_paragraphs_and_tables(note_docx):
    """The standalone app offered .docx on the Self Notes zone and then crashed on it."""
    pages = read_note(note_docx)
    assert len(pages) == 1
    assert "Broker Workbench" in pages[0].text
    assert any("Owner" in table for table in pages[0].tables)
    assert "Owner" in pages[0].as_prompt_block()


def test_a_pdf_note_reads_its_pages(note_pdf):
    pages = read_note(note_pdf)
    assert pages and "2.0bn" in pages[0].text


def test_a_format_we_cannot_read_is_refused_by_name(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_text("hello", encoding="utf-8")
    with pytest.raises(UnsupportedNoteFormat, match="notes.txt"):
        read_note(path)


# ── uploads ──────────────────────────────────────────────────────────────────


def test_an_upload_cannot_write_outside_its_directory():
    """The filename comes from the client, so it is never joined onto a path as-is."""
    assert safe_name("../../etc/passwd") == "passwd"
    assert safe_name("deck/../../x.pptx") == "x.pptx"
    assert safe_name("") == "upload"
    assert safe_name("Zurich UK QBR.pptx") == "Zurich UK QBR.pptx"


def test_an_upload_of_the_wrong_type_is_rejected(monkeypatch, tmp_path):
    monkeypatch.setenv("MOM_RUNS_DIR", str(tmp_path))
    with pytest.raises(UploadRejected):
        stage_upload("data:text/plain;base64,aGk=", "notes.txt", ".pdf,.docx")


def test_a_staged_upload_lands_on_disk_not_in_the_browser(monkeypatch, tmp_path):
    """The store must carry a path — a whole deck as base64 would round-trip on
    every callback that reads it."""
    monkeypatch.setenv("MOM_RUNS_DIR", str(tmp_path))
    staged = stage_upload("data:application/pdf;base64,aGk=", "note.pdf", ".pdf")
    assert staged.exists() and staged.path.read_bytes() == b"hi"
    assert set(staged.as_store()) == {"name", "path"}
    assert "base64" not in json.dumps(staged.as_store())


# ── the deck, and the shape it hands on ──────────────────────────────────────


def test_the_deck_reads_its_metadata_and_agenda(deck_path):
    from mom import deck as deck_module

    metadata = deck_module.extract_ppt_metadata(str(deck_path))
    assert metadata["client"] == "Zurich"
    assert metadata["meeting_date"] == "March 2026"
    assert metadata["agenda_items"] == ["Performance", "Marsh Update", "Key Takeaways"]
    assert any("A. Underwriter" in person for person in metadata["attendees"])


def test_a_slide_manifest_collapses_to_the_same_sections_as_a_section_one():
    """The agenda is written from sections either way, so both granularities must
    produce the same section list from the same deck."""
    slides = [
        {"slide_number": 3, "section_number": 1, "section_title": "Performance"},
        {"slide_number": 4, "section_number": 1, "section_title": "Performance"},
        {"slide_number": 6, "section_number": 2, "section_title": "Marsh Update"},
    ]
    assert sections_from_manifest(slides, "slide") == [
        {"section_number": 1, "section_title": "Performance", "slide_range": [3, 4]},
        {"section_number": 2, "section_title": "Marsh Update", "slide_range": [6, 6]},
    ]

    sections = [
        {"section_number": 1, "section_title": "Performance", "slide_range": [3, 4]},
        {"section_number": 2, "section_title": "Marsh Update", "slide_range": [6, 6]},
    ]
    assert sections_from_manifest(sections, "section") == sections


def test_the_document_is_named_after_the_deck():
    assert minutes_filename(Path("/x/Zurich UK QBR.pptx")) == "Zurich UK QBR_Meeting_Notes.docx"


# ── writing the document ─────────────────────────────────────────────────────


def _headings(path: Path) -> list[str]:
    return [p.text for p in Document(str(path)).paragraphs if p.style.name.startswith("Heading")]


def _all_text(path: Path) -> str:
    document = Document(str(path))
    rows = [c.text for t in document.tables for r in t.rows for c in r.cells]
    return "\n".join([p.text for p in document.paragraphs] + rows)


def test_the_skill_shape_writes_its_three_sections(tmp_path):
    path = write_minutes(
        meeting_data={"client": "Zurich", "subject": "Zurich QBR",
                      "meeting_date": "March 2026", "attendees": ["A (Zurich)"],
                      "action_items": [{"action": "Share pricing", "owner": "Marsh",
                                        "due_date": "Q3 2026"}]},
        summary={"strategy_and_initiatives": ["BWB rollout."],
                 "country_product_region": ["LATAM facility."],
                 "key_takeaways": ["GWP up 4.2%."]},
        top_pairs=[{"umbrella_tag": "Key Takeaways", "sub_tag": "KPIs & Performance Headlines"}],
        ppt_sections=[{"section_number": 2, "section_title": "Performance"}],
        mode="skill",
        path=tmp_path / "out.docx",
    )
    assert _headings(path) == [
        "Agenda", "Strategy & Initiatives", "Country / Product / Region",
        "Key Takeaways", "Action Items",
    ]
    text = _all_text(path)
    assert "GWP up 4.2%." in text and "Share pricing" in text and "Q3 2026" in text


def test_the_carrier_shape_names_its_first_section_after_the_carrier(tmp_path):
    path = write_minutes(
        meeting_data={"client": "Zurich", "subject": "QBR", "action_items": []},
        summary={"carrier_update": ["Rank improved."], "marsh_update": ["BWB rollout."]},
        top_pairs=[],
        ppt_sections=[],
        mode="carrier_marsh",
        path=tmp_path / "out.docx",
    )
    assert _headings(path) == ["Agenda", "Zurich", "Marsh Update"]


def test_an_unknown_output_shape_is_refused(tmp_path):
    with pytest.raises(ValueError, match="mode must be"):
        write_minutes(meeting_data={}, summary={}, top_pairs=[], ppt_sections=[],
                      mode="freestyle", path=tmp_path / "x.docx")


def test_the_agenda_prefers_the_decks_own_agenda_slide():
    meeting = {"ppt_agenda_items": ["Performance", "Marsh Update"],
               "discussion_points": [{"topic": "Something else"}]}
    sections = [{"section_number": 2, "section_title": "Performance"}]
    assert agenda_items(meeting, sections) == ["Performance", "Marsh Update"]
    assert agenda_items({"discussion_points": [{"topic": "Fallback"}]}, []) == ["Fallback"]
    assert agenda_items({}, [{"section_number": 1, "section_title": "Introduction"},
                             {"section_number": 2, "section_title": "Performance"}]) == ["Performance"]


def test_control_characters_never_reach_the_document():
    """python-docx will happily write XML that Word then refuses to open."""
    assert xml_safe("a\x0bb\n c   d") == "a b c d"


def test_a_party_label_the_model_prepended_anyway_is_stripped():
    assert clean_bullets(["Carrier: GWP grew.", "", "  "]) == ["GWP grew."]
    assert clean_bullets([]) == ["No material content identified."]


def test_the_evidence_block_carries_both_sources():
    block = build_evidence_block([{
        "umbrella_tag": "Key Takeaways", "sub_tag": "KPIs & Performance Headlines",
        "score": 6,
        "meeting_note_bullets": [{"bullet": "GWP up 4.2%.", "bullet_owner": "Marsh"}],
        "ppt_entries": [{"text": "SoW 2.7%", "slide_owner": "Carrier", "slide_title": "KPIs"},
                        {"text": "dropped", "slide_owner": "Carrier", "slide_title": "KPIs"}],
    }], max_ppt_lines=1)
    assert "Priority 1: Key Takeaways / KPIs & Performance Headlines" in block
    assert "GWP up 4.2%." in block and "SoW 2.7%" in block
    assert "dropped" not in block, "the per-tag line cap is not applied"


# ── run directories ──────────────────────────────────────────────────────────


def test_a_run_owns_its_directory(monkeypatch, tmp_path):
    """The standalone pipeline wrote into one shared data/ folder, so a second run
    overwrote the first. Two runs must now sit side by side."""
    monkeypatch.setenv("MOM_RUNS_DIR", str(tmp_path))
    first, second = config.new_run_paths(), config.new_run_paths()
    assert first.root != second.root
    for paths in (first, second):
        assert paths.inputs.is_dir() and paths.raw_data.is_dir()
        assert paths.tagged_data.is_dir() and paths.output_dir.is_dir()
        assert paths.root in paths.summary_json.parents


# ── end to end ───────────────────────────────────────────────────────────────


@pytest.mark.parametrize("mode", MODES, ids=[m.id for m in MODES])
def test_the_whole_workflow_produces_a_document(mode, deck_path, note_docx, run_paths, model):
    """upload -> note -> deck -> tag -> score -> verify -> .docx, over real files."""
    result = run_mom_pipeline(
        MoMRequest(note_path=note_docx, deck_path=deck_path, mode=mode, paths=run_paths),
        call_json=model,
    )

    assert result.docx_path.is_file()
    assert result.docx_path.name == "Zurich UK QBR_Meeting_Notes.docx"
    assert result.client == "Zurich", "the deck is the source of truth for the client"

    # Every phase actually ran.
    assert "parse_meeting_note" in model.labels
    assert "tag_meeting_note" in model.labels
    assert any(label.startswith(f"tag_{mode.granularity}") for label in model.labels)
    assert "verify_tagging" in model.labels and "summary_generation" in model.labels

    # The intermediate artefacts are on disk, in this run's directory.
    assert json.loads(run_paths.meeting_note_json.read_text(encoding="utf-8"))["client"] == "Zurich"
    assert json.loads(run_paths.priority_data.read_text(encoding="utf-8"))
    assert list(run_paths.raw_data.glob(f"{mode.granularity}_*.json"))
    assert list(run_paths.tagged_data.glob("*.json"))
    assert result.run_log_path.is_file()

    # The document reads like minutes.
    headings = _headings(result.docx_path)
    assert headings[0] == "Agenda" and headings[-1] == "Action Items"
    text = _all_text(result.docx_path)
    assert "Zurich" in text and "Share updated pricing" in text
    assert "Performance" in text, "the deck's agenda did not reach the document"


def test_a_failed_audit_stops_before_the_document_is_written(
    deck_path, note_docx, run_paths, tags
):
    """A run that the audit judges misaligned must not leave a half-good .docx behind."""
    with pytest.raises(verifier.VerificationFailed):
        run_mom_pipeline(
            MoMRequest(note_path=note_docx, deck_path=deck_path,
                       mode=AI_SUMMARY, paths=run_paths),
            call_json=StubModel(tags, verdict="fail"),
        )
    assert not list(run_paths.output_dir.glob("*.docx"))


def test_two_runs_do_not_write_into_each_others_directories(
    deck_path, note_docx, run_paths, model
):
    """The standalone pipeline wrote into the folder it was launched from, so a second
    run overwrote the first one's slide files, priorities and document."""
    request = MoMRequest(note_path=note_docx, deck_path=deck_path,
                         mode=AI_SUMMARY, paths=run_paths)
    first = run_mom_pipeline(request, call_json=model)

    second_paths = config.RunPaths(run_paths.root.parent / "run2").create()
    second = run_mom_pipeline(
        MoMRequest(note_path=note_docx, deck_path=deck_path, mode=AI_SUMMARY,
                   paths=second_paths),
        call_json=model,
    )

    assert first.docx_path != second.docx_path
    assert first.docx_path.is_file() and second.docx_path.is_file()
    assert first.run_log_path != second.run_log_path
    assert run_paths.root not in second_paths.root.parents


class _FakeResponse:
    def __init__(self, payload: dict, tokens: int):
        self.content = json.dumps(payload)
        self.usage_metadata = {
            "input_tokens": tokens,
            "output_tokens": tokens,
            "output_token_details": {"reasoning_tokens": 1},
        }


class _FakeChat:
    def __init__(self, tokens: int):
        self.tokens = tokens

    def invoke(self, prompt: str):
        return _FakeResponse({"results": []}, self.tokens)


def test_a_run_log_counts_its_own_calls_and_no_others():
    """The standalone logger was a module singleton, so two runs shared its totals."""
    from mom.llm import LlmJsonCaller
    from mom.run_log import RunLog

    first, second = RunLog("first"), RunLog("second")
    for _ in range(3):
        LlmJsonCaller(_FakeChat(10), first)("p", label="tag", phase="tagging")
    LlmJsonCaller(_FakeChat(5), second)("p", label="tag", phase="tagging")

    assert first.summary().total_llm_calls == 3
    assert second.summary().total_llm_calls == 1
    assert second.summary().total_tokens == 10
    assert first.summary().phase_breakdown[0].reasoning_tokens == 3


def test_a_response_that_is_not_json_names_the_call_that_produced_it():
    from mom.llm import LlmJsonCaller
    from mom.run_log import RunLog

    class _Broken:
        def invoke(self, prompt):
            response = _FakeResponse({}, 1)
            response.content = "I'm afraid I can't do that."
            return response

    with pytest.raises(ValueError, match="summary_generation"):
        LlmJsonCaller(_Broken(), RunLog())("p", label="summary_generation", phase="summary")


def test_the_run_log_workbook_is_written(tmp_path):
    from mom.run_log import RunLog
    from mom.run_log_excel import write_run_log

    log = RunLog("run-1")
    log.start_phase("tagging")
    log.record_call(call_index=1, label="tag", phase="tagging", input_tokens=10,
                    output_tokens=5, reasoning_tokens=2, duration_s=0.5)
    log.end_phase("tagging")

    path = write_run_log(log.summary(), tmp_path / "run_log.xlsx")
    assert path.is_file()

    import openpyxl

    workbook = openpyxl.load_workbook(path)
    assert workbook.sheetnames == ["Run Summary", "Call Log"]
    assert workbook["Call Log"].cell(2, 3).value == "tag"


# ── the job the workspace polls ──────────────────────────────────────────────


def test_a_job_reports_progress_then_the_finished_document(deck_path, note_docx, run_paths, model):
    request = MoMRequest(note_path=note_docx, deck_path=deck_path,
                         mode=AI_SUMMARY, paths=run_paths)
    job = jobs.start_run(request, runner=lambda req, report: run_mom_pipeline(
        req, call_json=model, report=report
    ))

    for _ in range(600):  # the stub run is fast; this is a ceiling, not a wait
        if job.done:
            break
        import time

        time.sleep(0.05)
    assert job.done, "the run never finished"

    state = job.snapshot()
    assert state["error"] is None
    assert state["percent"] == 100 and state["done"]
    assert state["filename"] == "Zurich UK QBR_Meeting_Notes.docx"
    assert state["client"] == "Zurich" and state["topics"]
    assert Path(job.docx_path()).is_file()
    assert jobs.get_job(job.job_id) is job
    jobs.clear_job(job.job_id)
    assert jobs.get_job(job.job_id) is None


def test_a_failing_run_reaches_the_user_as_a_message(run_paths, note_docx, deck_path):
    """A crashed thread with no message is the one outcome the user cannot act on."""
    def explode(request, report):
        report("tagging", "half way")
        raise RuntimeError("the tag list is missing")

    job = jobs.start_run(
        MoMRequest(note_path=note_docx, deck_path=deck_path, mode=AI_SUMMARY, paths=run_paths),
        runner=explode,
    )
    for _ in range(200):
        if job.done:
            break
        import time

        time.sleep(0.05)

    state = job.snapshot()
    assert state["done"] and state["error"] == "the tag list is missing"
    assert state["percent"] != 100, "a failed run must not read as complete"
    assert job.docx_path() is None
