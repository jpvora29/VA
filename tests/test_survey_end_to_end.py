"""End-to-end: a Premium + survey selection → the assembled deck's Carrier Survey pages.

Runs the REAL pipeline over the REAL templates against the seed DB:

    selection → compute_overall → per-country sub-decks (+ survey) → fill (cells · colours
              · ribbon) → merge → one deck

then reads the exported file back. Deterministic: seed DB, no LLM.

    pytest -m "not e2e"      # skip
"""
from __future__ import annotations

import re

import pytest
from pptx import Presentation

from studio import seed as S
from studio.compute import compute_overall
from studio.template_fill.assemble import assemble_deck, plan_subdecks

pytestmark = pytest.mark.e2e

_SELECTION = {"carrier": "Zurich", "country": ["Singapore", "Japan"], "year": 2025}


@pytest.fixture(scope="module")
def result():
    S.ensure_seed_db()
    out = compute_overall(filters=_SELECTION)
    assert out.subject == "Zurich"
    return out


@pytest.fixture(scope="module")
def survey_deck(result, tmp_path_factory):
    tmp = tmp_path_factory.mktemp("survey_e2e")
    path = assemble_deck(result, out_path=str(tmp / "deck.pptx"),
                         work_dir=str(tmp / "work"), data_basis="premium_survey")
    return path


def _survey_slides(path):
    """Every slide in the exported deck whose title is the Carrier Survey page."""
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        # The page's TITLE, not any mention: the summary page's survey pointer sentence
        # ("Brokers scored ... in the 2025 carrier survey") is prose about the score, not
        # the Carrier Survey page, and a substring match counted it as one.
        if any(t.strip().lower().startswith("carrier survey") for t in texts):
            out.append(slide)
    return out


def test_one_survey_page_per_selected_country(survey_deck):
    assert len(_survey_slides(survey_deck)) == len(_SELECTION["country"])


def test_each_survey_page_follows_its_own_country_block(result, tmp_path):
    # scope="country": _SELECTION pins no product, and against the REAL template registry
    # (unlike the hermetic test_survey_assemble.py fixture, which never registers a product
    # axis) an unpinned product falls back to the carrier's full product vocabulary — so the
    # default scope="all" would legitimately interleave 6 product sub-decks here. That is
    # correct, unrelated pre-existing behaviour (see plan_subdecks' docstring), not a defect
    # in the survey feature. Scoping to "country" isolates exactly what this test names:
    # survey placement relative to its own country block.
    axes = [d.template for d in plan_subdecks(result, scope="country", data_basis="premium_survey")]
    assert axes == ["overall", "country", "survey", "country", "survey", "end"]


def test_the_survey_page_is_titled_for_its_country(survey_deck):
    titles = []
    for slide in _survey_slides(survey_deck):
        titles.extend(sh.text_frame.text.strip() for sh in slide.shapes if sh.has_text_frame)
    for country in _SELECTION["country"]:
        assert country in titles
    assert "Country (1)" not in titles


def test_no_x_placeholder_survives_in_the_score_table(survey_deck):
    for slide in _survey_slides(survey_deck):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            cells = [shape.table.cell(r, c).text
                     for r in range(len(shape.table.rows))
                     for c in range(len(shape.table.columns))]
            assert "x.x" not in cells


def test_the_scores_are_real_numbers_on_the_survey_scale(survey_deck):
    for slide in _survey_slides(survey_deck):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for r in range(1, len(shape.table.rows)):
                for c in range(1, len(shape.table.columns)):
                    value = float(shape.table.cell(r, c).text)
                    assert 1.0 <= value <= 10.0


def test_cells_are_coloured_by_their_move_against_last_year(survey_deck):
    from studio.template_fill.survey import bands

    legend = set(bands.LEGEND)
    seen = set()
    for slide in _survey_slides(survey_deck):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for r in range(1, len(shape.table.rows)):
                for c in range(1, len(shape.table.columns)):
                    cell = shape.table.cell(r, c)
                    try:
                        seen.add(str(cell.fill.fore_color.rgb))
                    except (AttributeError, TypeError):
                        pass          # a cell with no comparable prior year
    assert seen & legend, "no cell took a band colour"
    assert seen <= legend, f"unexpected colours on the page: {seen - legend}"
    assert len(seen & legend) >= 3, "the seeded drifts should span several bands"


def test_the_ribbon_picture_is_ours_not_the_authored_one(survey_deck):
    from studio.template_fill.survey import ribbon

    if not ribbon.available():
        pytest.skip("kaleido/Chrome not available on this host")
    authored = Presentation("template/survey_template.pptx").slides[0]
    original = max((sh for sh in authored.shapes if sh.shape_type == 13),
                   key=lambda sh: sh.height).image.blob
    for slide in _survey_slides(survey_deck):
        pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
        biggest = max(pictures, key=lambda sh: sh.height)
        assert biggest.image.blob != original


def test_the_think_cell_object_is_gone_from_the_filled_page(survey_deck):
    from studio.template_fill.survey import ribbon

    if not ribbon.available():
        pytest.skip("kaleido/Chrome not available on this host")
    for slide in _survey_slides(survey_deck):
        assert not [sh for sh in slide.shapes if "think-cell" in (sh.name or "").lower()]


def test_no_peer_carrier_is_named_anywhere_on_the_page(survey_deck):
    peers = [c for c in S.CARRIERS if c != "Zurich"]
    for slide in _survey_slides(survey_deck):
        blocks = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        for shape in slide.shapes:
            if shape.has_table:
                blocks.extend(shape.table.cell(r, c).text
                              for r in range(len(shape.table.rows))
                              for c in range(len(shape.table.columns)))
        text = " ".join(blocks)
        for peer in peers:
            assert peer not in text


# ── regression: the premium-only deck must be untouched ──────────────────────


def test_premium_only_deck_is_unchanged_by_this_feature(result, survey_deck, tmp_path):
    """The premium deck carries no survey page, and the survey deck adds EXACTLY one
    slide per country to it — derived from the two decks, so neither count is hardcoded
    against a template that may gain or lose slides for unrelated reasons."""
    premium = assemble_deck(result, out_path=str(tmp_path / "premium.pptx"),
                            work_dir=str(tmp_path / "w1"), data_basis="premium")
    assert not _survey_slides(premium)
    premium_n = len(Presentation(premium).slides._sldIdLst)
    survey_n = len(Presentation(survey_deck).slides._sldIdLst)
    assert survey_n - premium_n == len(_SELECTION["country"])


def test_default_data_basis_still_produces_the_premium_deck(result, tmp_path):
    default = assemble_deck(result, out_path=str(tmp_path / "default.pptx"),
                            work_dir=str(tmp_path / "w2"))
    assert not _survey_slides(default)


# ── the summary page's overall survey-score tile ─────────────────────────────
# Same gate as the page itself, on the page the deck OPENS with: filled from the survey
# book when the run asked for it, off the slide when it did not.


# The caption reaches the exported deck as "Overall <carrier> Survey" — the fill engine
# rewrites the authored "Carrier" to the subject's own name, like every other label.
_TILE = re.compile(r"overall\s+\S+\s+survey", re.I)


def _survey_tiles(path):
    """The text of every overall survey-score tile in the exported deck."""
    return [sh.text_frame.text for slide in Presentation(path).slides
            for sh in slide.shapes
            if sh.has_text_frame and _TILE.search(sh.text_frame.text)]


def test_the_overall_survey_tile_carries_a_real_score_on_the_survey_basis(survey_deck):
    tiles = _survey_tiles(survey_deck)
    assert len(tiles) == 1
    score = float(re.search(r"\d+\.\d", tiles[0]).group(0))
    assert 1.0 <= score <= 10.0


def test_the_overall_survey_tile_comes_off_the_page_on_the_premium_basis(result, tmp_path):
    """A premium-basis deck must not carry a survey number — nor the template's own "x.x"
    where one would have gone. Scoped to the overall block: that is the page under test."""
    path = assemble_deck(result, out_path=str(tmp_path / "premium.pptx"),
                         work_dir=str(tmp_path / "work"), scope="overall")
    assert _survey_tiles(path) == []
