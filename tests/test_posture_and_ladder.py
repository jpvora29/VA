"""Strategic postures and the opportunity evidence ladder.

The postures are pinned against the book the plan was written from (Cyber defends,
Financial Lines scales, Casualty is a fix, Marine and Energy are selective), so a
threshold change that quietly re-labels the portfolio fails here rather than in a deck.

The ladder's tests are mostly about what the deck may NOT say: premium data evidences an
observation, never an addressable market, and the recommended verb has to stop there.
"""
from __future__ import annotations

import pytest

from studio import opportunity as O
from studio.posture import (
    LEAD_RANK, MATERIAL_SHARE_MOVE, SCALE_GROWTH_MARGIN, Posture, PostureInput, posture_for,
)


@pytest.fixture(autouse=True)
def _no_llm(monkeypatch):
    monkeypatch.setenv("STUDIO_AI", "off")


def _p(**kw) -> PostureInput:
    base = dict(name="Line", premium=40e6, pool=400e6, growth_pct=8.0,
                pool_growth_pct=8.0, rank=6, rank_change=0, share=10.0, share_change=0.0)
    base.update(kw)
    return PostureInput(**base)


# ── the portfolio the plan describes ─────────────────────────────────────────


@pytest.mark.parametrize("name,kw,expected", [
    # Leading on rank, growing hard — a leader's first job is holding the lead.
    ("Cyber", dict(rank=1, rank_change=5, growth_pct=97.3, pool_growth_pct=13.5,
                   share=14.7, peer_share=11.5), Posture.DEFEND),
    # Not leading, but outgrowing the pool with the rank moving behind it.
    ("Financial Lines", dict(rank=3, rank_change=3, growth_pct=57.1, pool_growth_pct=13.2,
                             share=10.7, peer_share=10.8), Posture.SCALE),
    # Shrinking while the pool grows.
    ("Casualty", dict(rank=6, rank_change=0, growth_pct=-0.4, pool_growth_pct=10.2,
                      share=7.6, share_change=-0.8), Posture.FIX),
    # Ahead of the pool, but not by enough to call it a trend, and the rank did not move.
    ("Marine", dict(rank=7, rank_change=0, growth_pct=12.0, pool_growth_pct=5.4,
                    share=7.9, share_change=0.5), Posture.SELECTIVE),
    ("Energy", dict(rank=7, rank_change=0, growth_pct=8.9, pool_growth_pct=8.5,
                    share=7.6, share_change=0.1), Posture.SELECTIVE),
    # A pool the carrier is absent from is not a position to fix or scale.
    ("Renewable Energy", dict(premium=0.0, pool=184.9e6, growth_pct=None,
                              rank=None, share=None), Posture.VALIDATE),
])
def test_the_book_gets_the_posture_the_plan_expects(name, kw, expected):
    call = posture_for(_p(name=name, **kw))
    assert call is not None and call.posture is expected


def test_every_posture_carries_a_reason():
    call = posture_for(_p(rank=1))
    assert call.because and not call.because.endswith(".")


def test_a_book_with_no_figures_gets_no_posture():
    assert posture_for(PostureInput(name="Unknown")) is None


# ── precedence, stated as rules rather than as examples ──────────────────────


def test_a_leader_that_is_also_growing_defends_rather_than_scales():
    """Both tests hold; the lead is the one that decides."""
    call = posture_for(_p(rank=LEAD_RANK, growth_pct=60.0, pool_growth_pct=10.0,
                          rank_change=2, share_change=5.0))
    assert call.posture is Posture.DEFEND


def test_a_barely_written_pool_validates_rather_than_fixes():
    """It is trailing the pool, but there is no position there to fix yet."""
    call = posture_for(_p(premium=1e6, pool=400e6, growth_pct=-5.0, pool_growth_pct=10.0,
                          share_change=-1.0))
    assert call.posture is Posture.VALIDATE


def test_a_shrinking_book_never_defends():
    call = posture_for(_p(rank=1, growth_pct=-8.0, pool_growth_pct=4.0, share_change=-2.0))
    assert call.posture is Posture.FIX


def test_outgrowing_the_pool_without_the_rank_moving_is_not_scale():
    call = posture_for(_p(growth_pct=8.0 + SCALE_GROWTH_MARGIN + 1, pool_growth_pct=8.0,
                          rank_change=0, share_change=0.0))
    assert call.posture is Posture.SELECTIVE


def test_a_material_share_gain_is_scale_on_its_own():
    call = posture_for(_p(share_change=MATERIAL_SHARE_MOVE, rank=6, rank_change=0))
    assert call.posture is Posture.SCALE


# ── the evidence ladder ──────────────────────────────────────────────────────


def test_premium_data_alone_only_ever_evidences_an_observation():
    assert O.level_from_premium_only() is O.Evidence.OBSERVED


def test_an_observed_gap_earns_validate_and_nothing_stronger():
    assert O.action_for(O.Evidence.OBSERVED) == "Validate"
    assert O.action_for(O.Evidence.VALIDATED) == "Enter"


def test_an_observed_gap_says_what_is_still_unconfirmed():
    assert "appetite and capacity" in O.qualifier_for(O.Evidence.OBSERVED)
    assert O.qualifier_for(O.Evidence.VALIDATED) == ""


def _money(v):
    return f"${v/1e6:.1f}M"


def test_a_gap_is_described_as_placed_elsewhere_not_as_a_market_to_be_had():
    said = O.describe(O.Opportunity("Renewable Energy", 184.9e6), money=_money)
    assert "placed with other carriers" in said
    for banned in ("addressable", "headroom", "whitespace", "market the account"):
        assert banned not in said.lower()


def test_the_recommendation_for_an_observed_gap_never_says_enter():
    said = O.recommend(O.Opportunity("Renewable Energy", 184.9e6), money=_money)
    assert said.startswith("Validate")
    assert "appetite and capacity are unconfirmed" in said
    assert "enter" not in said.lower()


def test_a_validated_opportunity_may_finally_say_enter():
    opp = O.Opportunity("Renewable Energy", 184.9e6, level=O.Evidence.VALIDATED)
    assert O.recommend(opp, money=_money).startswith("Enter")


def test_what_the_carrier_already_writes_is_not_counted_as_unwritten():
    opp = O.Opportunity("Marine", 300e6, written=24e6)
    assert opp.unwritten == pytest.approx(276e6)


# ── on the page ──────────────────────────────────────────────────────────────


def test_the_deck_states_a_portfolio_call_and_a_book_call(tmp_path):
    """End to end: the overall page opens its priorities on the portfolio's stance, and a
    key-messages column opens on the book's own."""
    from pptx import Presentation

    from studio.compute import compute_overall
    from studio.template_fill import assemble as A
    from studio.template_fill.binding_map import available

    if "overall" not in set(available()):
        pytest.skip("split templates not present")

    result = compute_overall(
        filters={"Carrier_Group": "Zurich", "Country": "Singapore", "Year": 2025})
    out = A.assemble_deck(result, out_path=str(tmp_path / "deck.pptx"), scope="overall")
    text = "\n".join(sh.text_frame.text for s in Presentation(out).slides
                     for sh in s.shapes if sh.has_text_frame)

    assert "The book's first calls are to" in text, "no portfolio stance on the deck"
    assert "The call here is to" in text, "no book-level stance on the deck"
    # The stance vocabulary is closed: every call uses one of the five verbs.
    portfolio = next(l for l in text.split("\n") if "The book's first calls are to" in l)
    assert any(v in portfolio for v in
               ("defend", "scale", "fix", "selectively pursue", "validate"))
    # ...and an instruction now names what is at stake. "Across the book the call is to
    # defend Cyber, scale Financial Lines, fix Casualty and selectively pursue Property,
    # Marine and Energy" listed six products, carried no figure, and would have been true
    # of any carrier with six lines.
    assert "Across the book the call is to" not in text
    assert "$" in portfolio, f"a call with nothing at stake behind it: {portfolio}"


def test_the_deck_never_calls_placed_premium_an_addressable_market(tmp_path):
    from pptx import Presentation

    from studio.compute import compute_overall
    from studio.template_fill import assemble as A
    from studio.template_fill.binding_map import available

    if "overall" not in set(available()):
        pytest.skip("split templates not present")

    result = compute_overall(
        filters={"Carrier_Group": "Zurich", "Country": "Singapore", "Year": 2025})
    out = A.assemble_deck(result, out_path=str(tmp_path / "deck.pptx"), scope="all")
    text = " ".join(sh.text_frame.text for s in Presentation(out).slides
                    for sh in s.shapes if sh.has_text_frame).lower()

    assert "addressable market" not in text
    assert "of whitespace" not in text
