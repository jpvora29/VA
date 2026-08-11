"""The two generic capabilities the Carrier Survey page needs from the fill engine:
painting table-cell backgrounds, and swapping a picture's image for a rendered one.

Hermetic — builds its own one-slide deck, so it proves the mechanics without a template.
"""
from __future__ import annotations

import base64
import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from studio.template_fill import fill as F

# Two valid, distinct 1x1 PNGs. Inline rather than drawn with Pillow: Pillow is only a
# transitive dependency here, and the test needs nothing more than "two different PNGs".
_RED = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg==")
_BLUE = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg==")


@pytest.fixture
def deck():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    picture = slide.shapes.add_picture(io.BytesIO(_RED), Inches(1), Inches(4),
                                       Inches(2), Inches(1))
    return prs, table.shape_id, picture.shape_id


def test_cell_background_is_painted_from_the_payload(deck):
    prs, table_id, _ = deck
    F._fill_cell_backgrounds(prs, {"cell_fills": {f"0:{table_id}": [
        {"r": 1, "c": 1, "hex": "CF3638"}]}})
    cell = prs.slides[0].shapes[0].table.cell(1, 1)
    assert str(cell.fill.fore_color.rgb) == "CF3638"


def test_a_none_hex_clears_the_cell_to_no_fill(deck):
    from pptx.enum.dml import MSO_FILL

    prs, table_id, _ = deck
    key = f"0:{table_id}"
    F._fill_cell_backgrounds(prs, {"cell_fills": {key: [{"r": 0, "c": 0, "hex": "008542"}]}})
    F._fill_cell_backgrounds(prs, {"cell_fills": {key: [{"r": 0, "c": 0, "hex": None}]}})
    assert prs.slides[0].shapes[0].table.cell(0, 0).fill.type == MSO_FILL.BACKGROUND


def test_an_out_of_range_cell_is_skipped_not_raised(deck):
    prs, table_id, _ = deck
    F._fill_cell_backgrounds(prs, {"cell_fills": {f"0:{table_id}": [
        {"r": 9, "c": 9, "hex": "CF3638"}]}})   # must not raise


def test_no_payload_is_a_no_op(deck):
    prs, _, _ = deck
    F._fill_cell_backgrounds(prs, {})           # must not raise


def test_picture_blob_is_replaced_in_place(deck):
    prs, _, picture_id = deck
    before = prs.slides[0].shapes[1].image.blob
    replacement = _BLUE
    F._replace_pictures(prs, {"pictures": {f"0:{picture_id}": replacement}})
    after = prs.slides[0].shapes[1]
    assert after.image.blob == replacement
    assert after.image.blob != before


def test_picture_replacement_keeps_the_authored_frame(deck):
    prs, _, picture_id = deck
    shape = prs.slides[0].shapes[1]
    frame = (shape.left, shape.top, shape.width, shape.height)
    F._replace_pictures(prs, {"pictures": {f"0:{picture_id}": _BLUE}})
    after = prs.slides[0].shapes[1]
    assert (after.left, after.top, after.width, after.height) == frame


def test_swapping_one_picture_does_not_corrupt_a_sibling_sharing_its_image_part(deck):
    # python-pptx dedupes image parts by content hash: two ``add_picture`` calls with
    # identical bytes share ONE part and rId. A naive blob-in-place swap would silently
    # rewrite both shapes; the fix mints a new part and repoints only the target's blip.
    prs, _, picture_id = deck
    slide = prs.slides[0]
    sibling = slide.shapes.add_picture(io.BytesIO(_RED), Inches(4), Inches(4),
                                       Inches(2), Inches(1))
    F._replace_pictures(prs, {"pictures": {f"0:{picture_id}": _BLUE}})
    swapped = next(s for s in slide.shapes if s.shape_id == picture_id)
    untouched = next(s for s in slide.shapes if s.shape_id == sibling.shape_id)
    assert swapped.image.blob == _BLUE
    assert untouched.image.blob == _RED


def test_think_cell_object_is_removed_from_a_refilled_slide(deck):
    prs, _, picture_id = deck
    slide = prs.slides[0]
    box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
    box.name = "think-cell data - do not delete"
    F._replace_pictures(prs, {"pictures": {f"0:{picture_id}": _BLUE}})
    assert not [s for s in slide.shapes if "think-cell" in s.name.lower()]


def test_think_cell_survives_a_slide_we_did_not_refill(deck):
    # An empty ``pictures`` dict would trip the top-level no-op return before the slide
    # loop ever runs, which would prove nothing about the per-slide ``touched`` scoping.
    # Targeting a shape id that is not on this slide keeps the loop live but the flag False.
    prs, _, picture_id = deck
    slide = prs.slides[0]
    box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
    box.name = "think-cell data - do not delete"
    F._replace_pictures(prs, {"pictures": {f"0:{picture_id + 999}": _BLUE}})
    assert [s for s in slide.shapes if "think-cell" in s.name.lower()]


# ── trimming a table to the size of what there is to say ─────────────────────
#
# A page whose axes come from the DATA (the Carrier Survey table's practices are the ones
# the carrier is surveyed on) has to be able to come out shorter than the template. Left in,
# the surplus lines are blank strips between the numbers and the Total.


@pytest.fixture
def grid_deck():
    """A 4x4 table whose cells name their own coordinates, so a trim is checkable."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(4, 4, Inches(1), Inches(1), Inches(8), Inches(4))
    for r in range(4):
        for c in range(4):
            shape.table.cell(r, c).text = f"{r}{c}"
    return prs, shape.shape_id


def _grid(prs):
    table = prs.slides[0].shapes[0].table
    return [[c.text for c in row.cells] for row in table.rows]


def test_dropped_columns_and_rows_leave_the_table(grid_deck):
    prs, table_id = grid_deck
    F._drop_table_lines(prs, {"drop_table_lines": {f"0:{table_id}": {"rows": [1], "cols": [1, 2]}}})
    assert _grid(prs) == [["00", "03"], ["20", "23"], ["30", "33"]]


def test_the_trimmed_table_still_spans_its_frame(grid_deck):
    """The freed width and height are shared over the survivors — a table that shrank to
    a third of the page would read as a mistake, not as a shorter table."""
    prs, table_id = grid_deck
    table = prs.slides[0].shapes[0].table
    before = (sum(c.width for c in table.columns), sum(r.height for r in table.rows))
    F._drop_table_lines(prs, {"drop_table_lines": {f"0:{table_id}": {"rows": [2], "cols": [0, 3]}}})
    table = prs.slides[0].shapes[0].table
    assert (sum(c.width for c in table.columns), sum(r.height for r in table.rows)) == before


def test_indices_are_the_original_ones_however_many_are_dropped(grid_deck):
    """Callers address cells by the template's own indices — dropping high-to-low is what
    keeps a two-column trim from removing the wrong second column."""
    prs, table_id = grid_deck
    F._drop_table_lines(prs, {"drop_table_lines": {f"0:{table_id}": {"cols": [0, 1, 2]}}})
    assert _grid(prs) == [["03"], ["13"], ["23"], ["33"]]


def test_a_table_with_nothing_to_drop_is_untouched(grid_deck):
    prs, table_id = grid_deck
    before = _grid(prs)
    F._drop_table_lines(prs, {"drop_table_lines": {f"0:{table_id}": {"rows": [], "cols": []}}})
    F._drop_table_lines(prs, {})
    assert _grid(prs) == before


# ── moving and resizing a shape ──────────────────────────────────────────────
#
# A page that DROPS a visual it cannot honestly fill has to hand its space to the one beside
# it, or the slide ships a hole where the author drew content. The GWP page does exactly
# this when a single-country run leaves the country-vs-country chart nothing to compare.


@pytest.fixture
def boxed():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(4))
    return prs, box


def test_a_shape_is_moved_and_resized_from_the_payload(boxed):
    prs, box = boxed
    F._resize_shapes(prs, {"resize_shapes": {f"0:{box.shape_id}": {
        "x": 100_000, "y": 200_000, "w": 300_000, "h": 400_000}}})
    assert (box.left, box.top, box.width, box.height) == (100_000, 200_000, 300_000, 400_000)


def test_only_the_named_edges_move(boxed):
    """A page widening a chart must not also stretch it downward through the table."""
    prs, box = boxed
    before = (box.top, box.height)
    F._resize_shapes(prs, {"resize_shapes": {f"0:{box.shape_id}": {"x": 42, "w": 4242}}})
    assert (box.left, box.width) == (42, 4242)
    assert (box.top, box.height) == before


def test_a_shape_the_payload_does_not_name_is_untouched(boxed):
    prs, box = boxed
    before = (box.left, box.top, box.width, box.height)
    F._resize_shapes(prs, {"resize_shapes": {f"0:{box.shape_id + 999}": {"x": 1}}})
    F._resize_shapes(prs, {})
    assert (box.left, box.top, box.width, box.height) == before


def test_a_malformed_box_costs_that_shape_only(boxed):
    """One bad entry must never break the export."""
    prs, box = boxed
    F._resize_shapes(prs, {"resize_shapes": {f"0:{box.shape_id}": {"x": "not an emu"}}})
    assert box.left == Inches(1)
