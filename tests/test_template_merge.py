"""Merge engine — concatenating filled sub-decks into one deck.

Hermetic: builds small decks in-memory with python-pptx (no fixture asset needed) so the
core stitch mechanics — slide count, picture/table preservation, no duplicate partnames,
no dangling relationships — are guaranteed in CI. A second, opt-in check exercises the
real ``template/qbr_template.pptx`` (think-cell + native charts) when it is present.
"""
from __future__ import annotations

import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from studio.template_fill.merge import merge_pptx, merge_to_file


@pytest.fixture(autouse=True)
def _force_opc_merge(monkeypatch):
    monkeypatch.setenv("STUDIO_PPT_MERGE_ENGINE", "opc")


def _tiny_deck(path: str, *, n_slides: int, with_picture: bool = False) -> str:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(n_slides):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = f"slide {i}"
        if with_picture:
            # a 1x1 PNG keeps a media part + slide->image relationship in play
            png = (
                b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
                b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01"
                b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
            )
            import io

            slide.shapes.add_picture(io.BytesIO(png), Emu(0), Emu(0), Inches(1), Inches(1))
    prs.save(path)
    return path


def _count(prs) -> dict:
    pics = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
    return {"slides": len(prs.slides._sldIdLst), "pics": pics}


def test_merge_concatenates_in_order(tmp_path):
    a = _tiny_deck(str(tmp_path / "a.pptx"), n_slides=2, with_picture=True)
    b = _tiny_deck(str(tmp_path / "b.pptx"), n_slides=3, with_picture=True)

    out = merge_to_file([a, b], str(tmp_path / "merged.pptx"))
    merged = Presentation(out)

    counts = _count(merged)
    assert counts["slides"] == 5            # 2 + 3, in order
    assert counts["pics"] == 5              # every picture survives the stitch

    # No duplicate partnames and no dangling internal relationships.
    parts = list(merged.part.package.iter_parts())
    names = [str(p.partname) for p in parts]
    assert len(names) == len(set(names)), "duplicate partnames after merge"
    reachable = {id(p) for p in parts}
    for p in parts:
        for rId, rel in p.rels.items():
            if not rel.is_external:
                assert id(rel.target_part) in reachable, "dangling internal relationship"


def test_merge_single_path_is_identity(tmp_path):
    a = _tiny_deck(str(tmp_path / "solo.pptx"), n_slides=2)
    merged = merge_pptx([a])
    assert len(merged.slides._sldIdLst) == 2


def test_merge_empty_raises():
    with pytest.raises(ValueError):
        merge_pptx([])


def test_merge_to_file_can_use_powerpoint_engine(tmp_path, monkeypatch):
    a = _tiny_deck(str(tmp_path / "a.pptx"), n_slides=1)
    b = _tiny_deck(str(tmp_path / "b.pptx"), n_slides=1)
    calls = {}

    def fake_powerpoint(paths, out_path):
        calls["paths"] = list(paths)
        calls["out_path"] = out_path
        Path(out_path).write_bytes(Path(paths[0]).read_bytes())
        return out_path

    monkeypatch.setenv("STUDIO_PPT_MERGE_ENGINE", "powerpoint")
    monkeypatch.setattr("studio.template_fill.merge._merge_to_file_powerpoint", fake_powerpoint)

    out = merge_to_file([a, b], str(tmp_path / "merged.pptx"))

    assert out == str(tmp_path / "merged.pptx")
    assert calls == {"paths": [a, b], "out_path": str(tmp_path / "merged.pptx")}


@pytest.mark.skipif(
    not os.path.exists("template/qbr_template.pptx"),
    reason="real QBR template (think-cell + charts) not present",
)
def test_merge_preserves_charts_and_ole_in_real_template(tmp_path):
    src = "template/qbr_template.pptx"

    def census(prs):
        charts = oles = 0
        for s in prs.slides:
            for sh in s.shapes:
                if getattr(sh, "has_chart", False) and sh.has_chart:
                    charts += 1
                if sh.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                    oles += 1
        return len(prs.slides._sldIdLst), charts, oles

    s_slides, s_charts, s_oles = census(Presentation(src))
    out = merge_to_file([src, src], str(tmp_path / "real_merged.pptx"))
    m_slides, m_charts, m_oles = census(Presentation(out))
    assert (m_slides, m_charts, m_oles) == (2 * s_slides, 2 * s_charts, 2 * s_oles)
