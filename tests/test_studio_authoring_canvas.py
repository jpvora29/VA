"""Focused regression tests for the Studio Boardroom Canvas page rail."""
from __future__ import annotations

import json

from dash.development.base_component import Component
from pptx import Presentation

from studio.deck.model import DeckSpec, SlideSpec
from studio.export import export_document
from studio.page import authoring as A
from studio.page import canvas as CV
from studio.page import document as D


def _deck() -> DeckSpec:
    return DeckSpec(
        slides=(
            SlideSpec(layout="cover", title="Zurich Singapore", eyebrow="QBR"),
            SlideSpec(layout="exec", title="Executive summary", eyebrow="SUMMARY"),
            SlideSpec(layout="divider", title="Performance", eyebrow="SECTION 01"),
            SlideSpec(layout="insight", title="Premium grew", eyebrow="PERFORMANCE"),
        ),
        meta={"carrier": "Zurich", "country": "Singapore", "year": 2025},
    )


def _walk(node):
    if node is None:
        return
    if isinstance(node, (list, tuple)):
        for child in node:
            yield from _walk(child)
        return
    yield node
    if isinstance(node, Component):
        yield from _walk(getattr(node, "children", None))


def _classes(node) -> set[str]:
    out: set[str] = set()
    for item in _walk(node):
        class_name = getattr(item, "className", None)
        if class_name:
            out.update(str(class_name).split())
    return out


def test_canvas_uses_pages_panel_and_live_preview_surface():
    deck = _deck()
    doc = D.new_document(deck)

    body = A.canvas_body(deck, 0, doc, None)
    classes = _classes(body)

    assert "qs-pages-panel" in classes
    assert "qs-pg-preview-surface" in classes
    assert "qs-cv-surface" in classes

    ids = [
        getattr(item, "id", None)
        for item in _walk(body)
        if getattr(item, "id", None) is not None
    ]
    stable_ids = [json.dumps(value, sort_keys=True) if isinstance(value, dict) else str(value) for value in ids]
    assert len(stable_ids) == len(set(stable_ids))


def test_page_sections_are_numbered_and_follow_divider_widget_edits():
    deck = _deck()
    doc = D.new_document(deck)
    divider_sid = D.sid_at(doc, 2)
    assert divider_sid is not None

    doc = D.set_widget_prop(doc, divider_sid, "w1", "text", "Growth opportunities")
    sections = A._page_sections(deck, doc)

    assert [(s["number"], s["label"]) for s in sections] == [
        (1, "Executive summary"),
        (2, "Growth opportunities"),
    ]
    assert sections[1]["idxs"] == [2, 3]


def test_insert_section_adds_real_divider_page_and_layout():
    doc = D.new_document(_deck())

    updated, index = D.add_divider_slide(doc, 1)
    sid = D.sid_at(updated, index)

    assert index == 2
    assert sid is not None
    assert updated["slides"][sid]["layout"] == "divider"
    assert updated["slides"][sid]["title"] == "New section"
    widgets = D.page_widgets(updated, sid)
    assert widgets == [
        {
            "id": "w1",
            "kind": "headline",
            "x": 0,
            "y": 2,
            "w": 12,
            "h": 3,
            "props": {
                "text": "New section",
                "eyebrow": "SECTION 02",
                "subtitle": "",
                "hero": True,
            },
        }
    ]


def test_inserted_section_exports_with_dark_divider_treatment(tmp_path):
    doc, index = D.add_divider_slide(D.new_document(_deck()), 1)
    out = tmp_path / "canvas-divider.pptx"

    export_document(doc, out_path=str(out))
    prs = Presentation(out)
    slide = prs.slides[index]
    text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))

    assert "New section" in text
    assert str(slide.background.fill.fore_color.rgb) == "000F47"


def test_thumbnail_surface_is_read_only_and_has_no_canvas_dom_ids():
    doc = D.new_document(_deck())
    sid = D.sid_at(doc, 0)
    assert sid is not None

    thumb = CV.thumbnail_surface(D.page_widgets(doc, sid), layout="cover", accent="navy")
    ids = [getattr(item, "id", None) for item in _walk(thumb)]

    assert "qs-cv-surface" not in ids
    assert "qs-cv-layer" not in ids
    assert "layout-cover" in _classes(thumb)


def test_commentary_editor_round_trips_every_line():
    doc = D.new_document(_deck())
    sid = D.sid_at(doc, 1)
    assert sid is not None
    text_widget = next(w for w in D.page_widgets(doc, sid) if w["kind"] == "text")
    edited = "\n".join(
        [
            "[good] Growth: Premium increased",
            "[warn] Retention: Renewal pressure remains",
            "Broker engagement improved",
            "[danger] Gap: Two products are underpenetrated",
            "Next review is scheduled",
        ]
    )

    updated = D.set_widget_prop(doc, sid, text_widget["id"], "points_text", edited)
    points = D.get_widget(updated, sid, text_widget["id"])["props"]["points"]

    assert len(points) == 5
    assert points[0] == {
        "label": "Growth.",
        "text": "Premium increased",
        "tone": "good",
    }
    assert "Next review is scheduled" in D.commentary_to_text(points)


def test_widget_and_slide_appearance_render_in_canvas_and_inspector():
    doc = D.new_document(_deck())
    sid = D.sid_at(doc, 1)
    assert sid is not None
    wid = D.page_widgets(doc, sid)[0]["id"]
    for prop, value in (
        ("font_family", "Georgia"),
        ("font_size", 28),
        ("font_color", "#112233"),
        ("background_color", "#DDEEFF"),
    ):
        doc = D.set_widget_prop(doc, sid, wid, prop, value)
    doc = D.set_widget_text_style(doc, sid, wid, "title", "font_family", "Georgia")
    doc = D.set_widget_text_style(doc, sid, wid, "title", "font_size", 30)
    doc = D.set_widget_text_style(doc, sid, wid, "title", "font_color", "#334455")
    doc = D.set_page_style(doc, sid, "background_color", "#ABCDEF")

    body = A.canvas_body(_deck(), 1, doc, wid, zoom=70, inspector_tab="style")
    components = list(_walk(body))
    surface = next(
        item
        for item in components
        if "qs-cv-surface" in str(getattr(item, "className", "")).split()
    )
    widget = next(
        item
        for item in components
        if "qs-cv-bg-custom" in str(getattr(item, "className", "")).split()
        and "selected" in str(getattr(item, "className", "")).split()
    )
    pattern_ids = [
        item.id for item in components if isinstance(getattr(item, "id", None), dict)
    ]

    assert surface.style["backgroundColor"] == "#ABCDEF"
    assert widget.style["backgroundColor"] == "#DDEEFF"
    assert widget.style["--qs-widget-font-family"] == "Georgia"
    assert widget.style["--qs-widget-font-size"] == "28px"
    assert widget.style["--qs-widget-font-color"] == "#112233"
    assert widget.style["--qs-title-font-family"] == "Georgia"
    assert widget.style["--qs-title-font-size"] == "30px"
    assert widget.style["--qs-title-font-color"] == "#334455"
    assert "qs-inspector-full" in _classes(body)
    assert {item["prop"] for item in pattern_ids if item["type"] == "qs-tstyle"} == {
        "font_family",
        "font_size",
    }
    assert {
        item["role"] for item in pattern_ids if item["type"] == "qs-tstyle"
    } == {"eyebrow", "title", "subtitle"}
    assert any(item["type"] == "qs-color-swatch" for item in pattern_ids)
    assert any(
        item["type"] == "qs-color-custom" and item["scope"] == "page"
        for item in pattern_ids
    )
    assert all(
        getattr(item, "type", None) != "color"
        for item in components
        if getattr(item, "_type", "") == "Input"
    )
    assert "qs-cv-gridlines" not in _classes(body)
    assert "qs-cv-grid-cell" not in _classes(body)
    frame = next(
        item
        for item in components
        if "qs-slide-frame" in str(getattr(item, "className", "")).split()
        and "canvas" in str(getattr(item, "className", "")).split()
    )
    assert frame.style["transform"] == f"scale({A.zoom_scale(70)})"
    assert {item["op"] for item in pattern_ids if item["type"] == "qs-zoom"} == {
        "out",
        "in",
        "fit",
    }
    assert A.adjusted_zoom(70, "in") == 75
    assert A.adjusted_zoom(70, "out") == 65
    assert A.adjusted_zoom(70, "fit") == A.ZOOM_FIT
    assert "qs-canvas-footer" in _classes(body)


def test_qbr_widget_library_renders_advanced_visuals_and_data_editing():
    doc = D.new_document(_deck())
    sid = D.sid_at(doc, 1)
    assert sid is not None
    advanced = (
        "heatmap",
        "bridge",
        "matrix",
        "radar",
        "radial",
        "timeline",
        "callout",
        "actions",
        "table",
        "kpi",
    )
    for kind in advanced:
        doc, wid = D.add_widget(doc, sid, kind)
        widget = D.get_widget(doc, sid, wid)
        assert widget is not None
        rendered = CV.thumbnail_surface([widget])
        assert "qs-cv-placeholder" not in _classes(rendered)

    matrix = next(w for w in D.page_widgets(doc, sid) if w["kind"] == "matrix")
    body = A.canvas_body(_deck(), 1, doc, matrix["id"], inspector_tab="data")
    ids = [
        item.id
        for item in _walk(body)
        if isinstance(getattr(item, "id", None), dict)
    ]
    assert any(
        item["type"] == "qs-wprop" and item["prop"] == "data_json"
        for item in ids
    )

    updated = D.set_widget_prop(
        doc,
        sid,
        matrix["id"],
        "data_json",
        json.dumps(
            {
                "title": "Updated opportunity matrix",
                "points": [{"label": "Cyber", "x": 70, "y": 80, "size": 40}],
            }
        ),
    )
    assert D.get_widget(updated, sid, matrix["id"])["props"]["title"] == (
        "Updated opportunity matrix"
    )


def test_component_library_uses_visual_browser_controls_and_real_add_actions():
    library = CV.palette(
        category="risk",
        tab="recommended",
        view="list",
        search="",
    )
    components = list(_walk(library))
    ids = [
        item.id
        for item in components
        if isinstance(getattr(item, "id", None), dict)
    ]
    add_kinds = {
        item["kind"] for item in ids if item.get("type") == "qs-addw"
    }

    assert add_kinds == {"heatmap", "radar"}
    assert "qs-lib-list" in _classes(library)
    assert "qs-lib-preview" in _classes(library)
    assert "qs-lib-sidebar" in _classes(library)
    assert "qs-lib-tabs" in _classes(library)
    assert "qs-lib-view" in _classes(library)
    assert {item["category"] for item in ids if item.get("type") == "qs-libcat"} == {
        "all",
        "financial",
        "customer",
        "growth",
        "risk",
        "operations",
        "market",
    }
    assert {item["tab"] for item in ids if item.get("type") == "qs-libtab"} == {
        "all",
        "recommended",
        "mine",
        "governed",
        "recent",
    }
    assert {item["view"] for item in ids if item.get("type") == "qs-libview"} == {
        "grid",
        "list",
    }
    assert "qs-lib-search" in {
        getattr(item, "id", None) for item in components
        if isinstance(getattr(item, "id", None), str)
    }


def test_component_library_modal_passes_browser_state_when_open():
    doc = D.new_document(_deck())
    body = A.canvas_body(
        _deck(),
        1,
        doc,
        None,
        library_category="operations",
        library_tab="mine",
        library_view="grid",
        library_search="tracker",
        library_open=True,
    )
    ids = [
        item.id
        for item in _walk(body)
        if isinstance(getattr(item, "id", None), dict)
    ]
    add_kinds = {
        item["kind"] for item in ids if item.get("type") == "qs-addw"
    }

    assert add_kinds == {"actions"}
    assert "qs-lib-grid" in _classes(body)
    assert "qs-lib-modal" in _classes(body)  # palette lives in the popup now
    active_categories = [
        item.id["category"]
        for item in _walk(body)
        if isinstance(getattr(item, "id", None), dict)
        and item.id.get("type") == "qs-libcat"
        and "active" in str(getattr(item, "className", "")).split()
    ]
    assert active_categories == ["operations"]


def test_inspector_collapses_and_library_opens_as_a_modal():
    doc = D.new_document(_deck())

    # Inspector still collapses on the canvas; the library is no longer a panel.
    collapsed = A.canvas_body(_deck(), 1, doc, None, inspector_collapsed=True)
    panel_toggles = {
        item.id["panel"]
        for item in _walk(collapsed)
        if isinstance(getattr(item, "id", None), dict)
        and item.id.get("type") == "qs-panel-toggle"
    }
    assert panel_toggles == {"inspector"}
    assert "inspector-collapsed" in _classes(collapsed)

    # Closed by default: the + button is present, the modal is not.
    closed = A.canvas_body(_deck(), 1, doc, None)
    lib_ops = {
        item.id["op"]
        for item in _walk(closed)
        if isinstance(getattr(item, "id", None), dict)
        and item.id.get("type") == "qs-lib-toggle"
    }
    assert "open" in lib_ops
    assert "qs-lib-modal" not in _classes(closed)
    assert "library-open" not in _classes(closed)

    # Open: the modal renders with close + backdrop controls.
    opened = A.canvas_body(_deck(), 1, doc, None, library_open=True)
    assert "library-open" in _classes(opened)
    assert "qs-lib-modal" in _classes(opened)
    open_ops = {
        item.id["op"]
        for item in _walk(opened)
        if isinstance(getattr(item, "id", None), dict)
        and item.id.get("type") == "qs-lib-toggle"
    }
    assert {"close", "close-bd"} <= open_ops


def test_widget_inspector_tabs_are_specific_to_widget_type():
    doc = D.new_document(_deck())
    sid = D.sid_at(doc, 1)
    assert sid is not None
    widgets = D.page_widgets(doc, sid)
    headline = next(w for w in widgets if w["kind"] == "headline")
    commentary = next(w for w in widgets if w["kind"] == "text")

    headline_body = A.canvas_body(
        _deck(), 1, doc, headline["id"], inspector_tab="data"
    )
    headline_tabs = {
        item.id["tab"]
        for item in _walk(headline_body)
        if isinstance(getattr(item, "id", None), dict)
        and item.id.get("type") == "qs-insp-tab"
    }
    assert headline_tabs == {"setup", "style", "rules"}
    assert "qs-inspector-field" in _classes(headline_body)

    text_body = A.canvas_body(
        _deck(), 1, doc, commentary["id"], inspector_tab="data"
    )
    text_tabs = {
        item.id["tab"]
        for item in _walk(text_body)
        if isinstance(getattr(item, "id", None), dict)
        and item.id.get("type") == "qs-insp-tab"
    }
    assert text_tabs == {"setup", "data", "style", "rules"}
    assert "Commentary points" in [
        getattr(item, "children", None) for item in _walk(text_body)
    ]


def test_canvas_mode_contains_scrolling_to_inner_panels():
    css = (
        __import__("pathlib").Path(__file__).parents[1]
        / "assets"
        / "studio_authoring.css"
    ).read_text(encoding="utf-8")

    assert ".qs-root.mode-canvas {" in css
    # The frame is a FIXED height so the panels inside it scroll instead of the page.
    # Since Studio became a tab under the app navbar it is the shell height, not the
    # viewport: `--va-shell-h` is `100vh` minus the navbar (assets/theme_tokens.css).
    assert "height: var(--va-shell-h);" in css
    assert "100vh" not in css, "Studio no longer owns the viewport"
    assert ".qs-root.mode-canvas .qs-canvas-viewport" in css
    viewport_rule = css.split(
        ".qs-root.mode-canvas .qs-canvas-viewport", 1
    )[1].split("}", 1)[0]
    component_rule = css.split(
        ".qs-comp-grid {", 1
    )[1].split("}", 1)[0]
    assert "overflow: hidden" in viewport_rule
    assert "overflow: auto hidden" in component_rule
    assert "padding: 42px 48px" in viewport_rule
    assert "grid-template-columns: 294px minmax(0,1fr) 42px" in css
    assert ".qs-canvas.inspector-collapsed" in css
    # Component library is now an on-demand modal opened by the + button (it no
    # longer occupies a permanent bottom row), so the stage gets the full height.
    assert ".qs-lib-modal" in css
    assert ".qs-lib-open-btn" in css
    assert ".qs-cv-gridlines" not in css


def test_custom_appearance_and_full_commentary_export_to_powerpoint(tmp_path):
    doc = D.new_document(_deck())
    sid = D.sid_at(doc, 1)
    assert sid is not None
    widgets = D.page_widgets(doc, sid)
    headline = next(w for w in widgets if w["kind"] == "headline")
    commentary = next(w for w in widgets if w["kind"] == "text")

    doc = D.set_page_style(doc, sid, "background_color", "#ABCDEF")
    for prop, value in (
        ("font_family", "Georgia"),
        ("font_size", 28),
        ("font_color", "#112233"),
        ("background_color", "#DDEEFF"),
    ):
        doc = D.set_widget_prop(doc, sid, headline["id"], prop, value)
    doc = D.set_widget_prop(doc, sid, headline["id"], "subtitle", "Role-specific subtitle")
    doc = D.set_widget_text_style(
        doc, sid, headline["id"], "title", "font_family", "Aptos"
    )
    doc = D.set_widget_text_style(
        doc, sid, headline["id"], "title", "font_size", 31
    )
    doc = D.set_widget_text_style(
        doc, sid, headline["id"], "title", "font_color", "#445566"
    )
    doc = D.set_widget_text_style(
        doc, sid, headline["id"], "subtitle", "font_color", "#778899"
    )
    doc = D.set_widget_prop(
        doc,
        sid,
        commentary["id"],
        "points_text",
        "\n".join(f"Point {index}" for index in range(1, 7)),
    )

    out = tmp_path / "styled-canvas.pptx"
    export_document(doc, out_path=str(out))
    slide = Presentation(out).slides[1]
    all_text = "\n".join(
        shape.text for shape in slide.shapes if hasattr(shape, "text")
    )
    title_shape = next(
        shape
        for shape in slide.shapes
        if hasattr(shape, "text") and "Executive summary" in shape.text
    )
    title_run = title_shape.text_frame.paragraphs[0].runs[0]

    assert str(slide.background.fill.fore_color.rgb) == "ABCDEF"
    assert title_run.font.name == "Aptos"
    assert round(title_run.font.size.pt) == 31
    assert str(title_run.font.color.rgb) == "445566"
    subtitle_shape = next(
        shape
        for shape in slide.shapes
        if hasattr(shape, "text") and "Role-specific subtitle" in shape.text
    )
    subtitle_run = subtitle_shape.text_frame.paragraphs[0].runs[0]
    assert str(subtitle_run.font.color.rgb) == "778899"
    assert "Point 6" in all_text
    assert any(
        shape.fill.type is not None
        and getattr(shape.fill.fore_color, "rgb", None) is not None
        and str(shape.fill.fore_color.rgb) == "DDEEFF"
        for shape in slide.shapes
    )


def test_advanced_qbr_widgets_export_as_native_powerpoint_content(tmp_path):
    deck = DeckSpec(
        slides=(SlideSpec(layout="insight", title="Advanced QBR"),),
        meta={},
    )
    for kind in (
        "heatmap",
        "bridge",
        "matrix",
        "radar",
        "radial",
        "timeline",
        "callout",
        "actions",
        "table",
        "kpi",
        "chart",
    ):
        doc = D.new_document(deck)
        sid = D.sid_at(doc, 0)
        assert sid is not None
        doc.setdefault("layouts", {})[sid] = []
        doc, _wid = D.add_widget(doc, sid, kind)
        out = tmp_path / f"{kind}.pptx"

        export_document(doc, out_path=str(out))
        slide = Presentation(out).slides[0]

        assert out.exists()
        assert len(slide.shapes) > 1
