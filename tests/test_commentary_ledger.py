"""The claim ledger — a deck makes each claim once.

Unit tests pin the ledger's own rules; the end-to-end test generates a real deck from the
split templates and counts repeated sentences across every slide, which is the regression
that mattered: one sentence about the whole book landed on four slides.
"""
from __future__ import annotations

import re
from collections import Counter

import pytest

from studio.template_fill.ledger import ClaimLedger, signature


@pytest.fixture(autouse=True)
def _no_llm(monkeypatch):
    monkeypatch.setenv("STUDIO_AI", "off")


# ── signatures ───────────────────────────────────────────────────────────────


def test_signature_ignores_casing_spacing_and_the_full_stop():
    assert signature("Premium grew 57.1%.") == signature("  premium   grew 57.1%  ")


def test_signature_keeps_the_figures():
    """Two products saying the same thing about their own books are DIFFERENT claims."""
    assert signature("Cyber grew 97.3%.") != signature("Cyber grew 57.1%.")


# ── taking claims ────────────────────────────────────────────────────────────


def test_a_claim_already_made_gives_way_to_the_next_point():
    led = ClaimLedger()
    assert led.take(["A grew 10%.", "B fell 2%."], limit=1) == ["A grew 10%."]
    assert led.take(["A grew 10%.", "B fell 2%."], limit=1) == ["B fell 2%."]


def test_dropping_a_used_claim_promotes_the_next_into_the_space():
    """The point of deduping BEFORE the trim: the page still fills its column."""
    led = ClaimLedger()
    led.take(["first.", "second."], limit=2)
    assert led.take(["first.", "second.", "third.", "fourth."], limit=2) == ["third.", "fourth."]


def test_a_page_with_nothing_fresh_repeats_rather_than_blanks():
    led = ClaimLedger()
    led.take(["only claim."], limit=1)
    assert led.take(["only claim."], limit=1) == ["only claim."]


def test_a_point_trimmed_off_the_end_is_not_recorded_as_said():
    """It never reached a slide, so the next page may still make it."""
    led = ClaimLedger()
    assert led.take(["shown.", "trimmed."], limit=1) == ["shown."]
    assert led.take(["trimmed."], limit=1) == ["trimmed."]


def test_blank_lines_are_never_recorded_or_returned():
    led = ClaimLedger()
    assert led.take(["", "   ", "real claim."], limit=3) == ["real claim."]


def test_two_decks_do_not_share_a_memory():
    """The ledger is per-deck state, not a module global."""
    first = ClaimLedger()
    first.take(["A grew 10%."], limit=1)
    assert ClaimLedger().take(["A grew 10%.", "B fell 2%."], limit=1) == ["A grew 10%."]


# ── end to end: the deck does not repeat itself ──────────────────────────────


def _repeated_sentences(path) -> dict:
    """``{sentence: [slide numbers]}`` for every sentence appearing on more than one slide."""
    from pptx import Presentation

    counts: Counter = Counter()
    where: dict = {}
    for idx, slide in enumerate(Presentation(path).slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                texts += [c.text for r in shape.table.rows for c in r.cells if c.text.strip()]
        for block in texts:
            for line in block.split("\n"):
                for sentence in re.split(r"(?<=[.])\s+", line.strip()):
                    sentence = sentence.strip()
                    if len(sentence.split()) >= 8:      # a claim, not a label or a KPI
                        counts[sentence] += 1
                        where.setdefault(sentence, []).append(idx)
    return {s: where[s] for s, n in counts.items() if n > 1}


def test_a_generated_deck_does_not_repeat_its_claims(tmp_path):
    """Regression: the same opening sentence about the whole book reached the highlights
    page, both the trading summary's columns and the ranking page — four slides, one claim.

    A single-country run describes ONE book on every page, so the ledger's keep-one floor
    still allows a handful of repeats where a page has nothing else true to say; the bar
    here is the order of magnitude, which was 17 extra copies before the ledger.
    """
    from studio.compute import compute_overall
    from studio.template_fill import assemble as A
    from studio.template_fill.binding_map import available

    if not {"overall", "product"} <= set(available()):
        pytest.skip("split templates not present")

    result = compute_overall(
        filters={"Carrier_Group": "Zurich", "Country": "Singapore", "Year": 2025})
    out = A.assemble_deck(result, out_path=str(tmp_path / "deck.pptx"), scope="all")

    repeated = _repeated_sentences(out)
    extra = sum(len(slides) - 1 for slides in repeated.values())
    assert extra <= 5, f"the deck repeats itself {extra} times: {repeated}"
    assert not any(len(slides) > 3 for slides in repeated.values()), \
        f"a claim reached four or more slides: {repeated}"


# ── shape repetition across a wide deck ──────────────────────────────────────
#
# Reported: "Lot of repetition of pointers, even when 8-10 products are
# included." Every claim WAS new — `signature` keeps the figures, so ten products
# saying "grew X%" about their own books are ten different claims and all ten
# shipped. What repeated was the sentence they were poured into.

from studio.template_fill.ledger import MAX_SHAPE_USES, shape


def _product_lines(product: str, i: int) -> list:
    return [
        f"The book grew {i + 4}% to ${i + 2}.1M, and that growth sits in {product}.",
        f"Share of wallet in {product} rose to {i + 9}.1% from {i + 7}.4%.",
        f"Marsh placed ${i * 7 + 20}M in {product} the carrier did not write.",
    ]


PRODUCTS = ["Property", "Casualty", "Cyber", "Marine",
            "Energy", "Aviation", "Construction", "Financial Lines"]


def test_two_products_pour_into_the_same_shape():
    a = "The book grew 12% to $4.1M, and that growth sits in Property."
    b = "The book grew 8% to $2.9M, and that growth sits in Casualty."
    assert shape(a) == shape(b)
    assert signature(a) != signature(b), "they are still two different claims"


def test_the_entity_is_masked_but_the_sentence_is_not():
    assert "@" in shape("Growth sits in Financial Lines.")
    assert "#" in shape("The book grew 12%.")


def test_a_wide_deck_does_not_repeat_one_shape_on_every_page():
    """The regression: eight products, eight identical sentences."""
    ledger = ClaimLedger()
    chosen = [ledger.take(_product_lines(p, i), limit=1)[0]
              for i, p in enumerate(PRODUCTS)]
    assert len({shape(line) for line in chosen}) > 1, "one mould served every page"


def test_a_shape_gets_its_allowance_then_gives_way():
    ledger = ClaimLedger()
    chosen = [ledger.take(_product_lines(p, i), limit=1)[0]
              for i, p in enumerate(PRODUCTS[:MAX_SHAPE_USES + 1])]
    shapes = [shape(line) for line in chosen]
    assert shapes[0] == shapes[1], "a shape may be used more than once"
    assert shapes[-1] != shapes[0], "but not indefinitely"


def test_a_page_still_fills_its_column_when_every_shape_is_worn_out():
    """Deduplication may never blank a cell."""
    ledger = ClaimLedger()
    for i, p in enumerate(PRODUCTS):
        assert ledger.take(_product_lines(p, i), limit=1), f"{p} got nothing"


def test_an_exact_repeat_is_still_dropped_outright():
    ledger = ClaimLedger()
    line = "The book grew 12% to $4.1M."
    ledger.take([line, "Rank held at 4th."], limit=2)
    assert not ledger.take([line, "Cyber remains thin."], limit=1)[0] == line
