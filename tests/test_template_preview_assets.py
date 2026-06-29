from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from studio.template_fill.analyze import analyze


def _make_deck(path: Path) -> None:
    img_path = path.with_suffix(".png")
    Image.new("RGB", (40, 20), "#cc3344").save(img_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(0.6))
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "Template headline"
    run.font.name = "Aptos"
    run.font.size = Pt(18)
    run.font.bold = True

    slide.shapes.add_picture(str(img_path), Inches(1), Inches(1.5), Inches(1.2), Inches(0.6))
    prs.save(path)


def test_analyze_extracts_picture_asset_and_text_style(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    deck = tmp_path / "template.pptx"
    _make_deck(deck)

    template = analyze(str(deck))

    text = next(sh for sh in template.slides[0].shapes if sh.kind == "text")
    assert text.font_face == "Aptos"
    assert text.font_size_pt == 18
    assert text.bold is True

    picture = next(sh for sh in template.slides[0].shapes if sh.kind == "picture")
    assert picture.image_url and picture.image_url.startswith("/assets/studio_template_previews/")
    assert (tmp_path / picture.image_url.lstrip("/")).exists()


def test_rendered_backgrounds_are_optional(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    monkeypatch.setenv("STUDIO_TEMPLATE_RENDERER", "none")
    deck = tmp_path / "template.pptx"
    _make_deck(deck)

    template = analyze(str(deck))

    assert template.slides[0].background_url is None
