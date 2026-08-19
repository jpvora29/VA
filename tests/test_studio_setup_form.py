"""Setup-form regressions: the controls, the peer set, and the deck-sections panel.

All deterministic — they run against the seed DB (no DB_PATH, no LLM). The last
test walks the real business flow the form drives:

    Setup selection → structured load → evidence → binding map → filled template
    → assembled .pptx

so a change to the form is proved end to end, not just at the widget level.
"""
from __future__ import annotations

import json

import pytest

from studio.authoring.setup import carriers_in_scope
from studio.page.authoring.setup import (
    peer_set_body,
    scope_axes,
    setup_body,
    template_sections_panel,
)


def _rendered(component) -> str:
    """A Dash component tree flattened to JSON so tests can assert on its content."""
    return json.dumps(component, default=lambda o: getattr(o, "__dict__", str(o)))


@pytest.fixture(scope="module")
def seeded_survey():
    """The seed DB, which carries the survey book the panel reads."""
    from studio import seed as S

    S.ensure_seed_db()
    return True


def _form() -> str:
    return _rendered(setup_body([], filter_options={}, filter_values={}))


# ── report type: Full QBR is the only deliverable, so there is no chooser ─────


def test_setup_form_has_no_report_type_control():
    form = _form()
    assert "studio-report-type" not in form
    assert "Full QBR" not in form


def test_generate_pins_the_report_to_qbr():
    """The control is gone, so the callback must supply the value itself."""
    import inspect

    from studio.authoring import setup as S

    src = inspect.getsource(S.register_setup)
    assert '"report": "qbr"' in src
    assert 'State("studio-report-type"' not in src


# ── AI assist is on by default ───────────────────────────────────────────────


# ── data basis: premium only, or premium + survey ────────────────────────────


def test_the_form_offers_gpr_or_gpr_plus_carrier_survey():
    """The labels name the BOOKS an author knows by name, not the internal measure."""
    form = _form()
    assert "studio-data-basis" in form
    assert "GPR" in form and "GPR + Carrier Survey" in form
    assert "Premium only" not in form and "Premium + survey" not in form
    # It sits with the source question — both answer "what is this deck built from?".
    assert "studio-data-source" in form
    assert "GPR / Survey" in form and "Custom data" in form
    assert "Existing database" not in form and "My uploaded data" not in form


def test_data_basis_defaults_to_premium_only():
    """Today's deck is premium-only, so the default must not imply otherwise."""
    from studio.page.authoring.setup import DATA_BASIS_DEFAULT, _data_basis_control

    radio = _data_basis_control().children[0].children[1]
    assert radio.id == "studio-data-basis"
    assert radio.value == DATA_BASIS_DEFAULT == "premium"
    assert [o["value"] for o in radio.options] == ["premium", "premium_survey"]


def test_the_choice_reaches_the_selection_generate_builds_from():
    """A deck must record what it was asked for, and every cached build has to re-key
    when the answer changes."""
    import inspect

    from studio.authoring import setup as S

    src = inspect.getsource(S.register_setup)
    assert 'State("studio-data-basis", "value")' in src
    assert '"data_basis": data_basis or A.DATA_BASIS_DEFAULT' in src
    # …and so do the survey identities chosen beside it.
    assert 'State("studio-survey-carrier", "value")' in src
    assert '"survey_carrier": survey_carrier or None' in src
    assert '"survey_peers":' in src


# ── the survey selections ────────────────────────────────────────────────────
#
# The survey book keeps its own carrier vocabulary. The deck resolves the premium subject
# into it, but the author is the one who can see both lists — so the match is SHOWN, and
# can be overridden, before a page of someone else's scores is built.


def test_the_survey_panel_is_in_the_form_but_hidden_by_default():
    form = _form()
    assert "studio-survey-carrier" in form and "studio-survey-peer" in form
    assert "studio-survey-section" in form


def test_the_panel_is_hidden_on_the_premium_basis():
    from studio.authoring.setup import survey_panel_state

    style, options, value, _note = survey_panel_state({"carrier": "Zurich"}, "premium", None, None)
    assert style == {"display": "none"}
    assert options == [] and value is None


def test_the_panel_offers_the_survey_book_s_own_carrier_names(seeded_survey):
    from studio.authoring.setup import survey_panel_state

    style, options, value, _note = survey_panel_state(
        {"carrier": "Zurich", "country": "Singapore"}, "premium_survey", None, None)
    assert style == {}
    names = [o["value"] for o in options]
    assert "Zurich" in names and len(names) > 1
    assert value == "Zurich"                       # the match it would make on its own


def test_the_panel_shows_how_the_carrier_was_matched(seeded_survey):
    from studio.authoring.setup import survey_panel_state

    _s, _o, _v, note = survey_panel_state(
        {"carrier": "Zurich", "country": "Singapore"}, "premium_survey", None, None)
    assert "Surveyed as Zurich" in _rendered(note)


def test_the_panel_warns_when_the_carrier_is_not_surveyed(seeded_survey):
    """The page will be skipped — better said on the form than discovered in the deck."""
    from studio.authoring.setup import survey_panel_state

    _s, _o, value, note = survey_panel_state(
        {"carrier": "Nobody At All", "country": "Singapore"}, "premium_survey", None, None)
    assert value is None
    assert "could not be matched" in _rendered(note)


def test_a_pinned_carrier_is_kept_over_the_match(seeded_survey):
    from studio.authoring.setup import survey_panel_state

    _s, _o, value, _n = survey_panel_state(
        {"carrier": "Zurich", "country": "Singapore"}, "premium_survey", "Chubb", None)
    assert value == "Chubb"


def test_an_uploaded_dataset_has_no_survey_book(seeded_survey):
    from studio.authoring.setup import survey_panel_state

    class _Record:
        dataset_id = "x"

    style, options, value, note = survey_panel_state(
        {"carrier": "Zurich"}, "premium_survey", None, _Record())
    assert style == {} and options == [] and value is None
    assert "governed data" in _rendered(note)


# ── the survey peer group ────────────────────────────────────────────────────
#
# The premium panel has always shown the peer group the deck WOULD use before it is built.
# The survey panel started empty, so the survey page's peer set was invisible until the
# deck existed — and it is keyed differently (``Carrier``, not ``Carrier_Group``), so it is
# not the group shown above it either.


def _peer_state(selected, basis="premium_survey", carrier="Zurich", chosen=None, record=None,
                keep_choice=True):
    """``([(country, options, peers)], read-out)`` — the per-market survey peer groups."""
    from studio.authoring.setup import survey_peer_state

    return survey_peer_state(selected, basis, carrier, chosen, record,
                             keep_choice=keep_choice)


def _one_market(selected, **kwargs):
    """``(options, peers, read-out)`` for a scope that covers a single market."""
    groups, note = _peer_state(selected, **kwargs)
    if not groups:
        return [], [], note
    (_country, options, peers), = groups
    return options, list(peers), note


def test_the_peer_panel_is_prefilled_from_the_peers_table(seeded_survey):
    from studio.data import peer_members

    options, value, note = _one_market({"carrier": "Zurich", "country": "Singapore"})
    assert value, "the survey peer group must be shown before the deck is built"
    # Keyed on Carrier and scoped to the country — the survey flow's own peer columns.
    assert set(value) == set(peer_members("survey", "Zurich", country=("Singapore",)))
    assert "Peers table" in _rendered(note)
    assert "Zurich" not in [o["value"] for o in options], "a carrier is not its own peer"


def test_the_peer_panel_follows_the_survey_carrier_not_the_premium_one(seeded_survey):
    """Two carriers keep two different peer groups; the panel must track the survey one."""
    zurich = _one_market({"country": "Singapore"}, carrier="Zurich")[1]
    aig = _one_market({"country": "Singapore"}, carrier="AIG")[1]
    assert zurich and aig and set(zurich) != set(aig)
    assert "AIG" not in aig and "Zurich" not in zurich


def test_an_edited_peer_selection_survives_a_filter_change(seeded_survey):
    options, value, note = _one_market({"carrier": "Zurich", "country": "Singapore"},
                                       chosen={"Singapore": ["Chubb"]})
    assert value == ["Chubb"]
    assert "your selection" in _rendered(note)


def test_a_stale_peer_selection_is_replaced_by_the_peers_table(seeded_survey):
    """A name this scope does not survey is not a selection worth keeping."""
    _o, value, _n = _one_market({"carrier": "Zurich", "country": "Singapore"},
                                chosen={"Singapore": ["Someone Not Surveyed"]})
    assert value and "Someone Not Surveyed" not in value


def test_the_peer_panel_is_silent_off_the_survey_basis(seeded_survey):
    options, value, note = _one_market({"carrier": "Zurich"}, basis="premium")
    assert options == [] and value == []
    assert "Peers table" not in _rendered(note)


def test_the_peer_panel_waits_for_a_survey_carrier(seeded_survey):
    options, value, _n = _one_market({"country": "Singapore"}, carrier=None)
    assert options == [] and value == []


def test_the_peer_panel_says_so_when_the_peers_table_has_no_row(seeded_survey):
    """Silence would read as "no peers needed"; the page would then rank against the whole
    surveyed field, which is a different statement and worth saying out loud."""
    _o, value, note = _one_market({"country": "Singapore"}, carrier="Nobody At All")
    assert value == []
    assert "No survey peer group" in _rendered(note)


# ── a peer group belongs to the carrier it was chosen for ────────────────────
#
# The survey carrier is BOTH an Output of the panel and a State fed back into it, so the
# form cannot tell "the author overrode the match" from "the form made this match last time
# round" by value alone. Without the trigger, a match made for one carrier outlived it:
# select Zurich, then select AIG, and the survey panel still said Zurich — and still offered
# Zurich's peers.


def _carrier_state(selected, basis="premium_survey", pinned=None, record=None, keep_pin=True):
    from studio.authoring.setup import survey_panel_state

    return survey_panel_state(selected, basis, pinned, record, keep_pin=keep_pin)


def test_changing_the_carrier_rematches_the_survey_carrier(seeded_survey):
    _s, _o, value, note = _carrier_state({"carrier": "AIG", "country": "Singapore"},
                                         pinned="Zurich", keep_pin=False)
    assert value == "AIG"
    assert "Surveyed as AIG" in _rendered(note)


def test_a_pin_still_survives_a_change_that_is_not_the_carrier(seeded_survey):
    """The override is only discarded by the thing that invalidates it."""
    _s, _o, value, _n = _carrier_state({"carrier": "Zurich", "country": "Singapore"},
                                       pinned="Chubb", keep_pin=True)
    assert value == "Chubb"


def test_changing_the_carrier_drops_the_previous_carriers_peers(seeded_survey):
    _o, value, _n = _one_market({"carrier": "AIG", "country": "Singapore"}, carrier="AIG",
                                chosen={"Singapore": ["Chubb"]}, keep_choice=False)
    from studio.data import peer_members

    assert set(value) == set(peer_members("survey", "AIG", country=("Singapore",)))


def test_the_callbacks_ask_dash_what_actually_changed():
    """The rule lives in the callback, so assert it is wired — a `State` cannot carry it."""
    import inspect

    from studio.authoring import setup as S

    src = inspect.getsource(S.register_setup)
    assert 'keep_pin="carrier" not in changed_ids()' in src
    assert "keep_choice=not ({\"carrier\", \"studio-survey-carrier\"} & changed)" in src


# ── several countries, several peer groups ───────────────────────────────────
#
# The Peers table keys a group on (carrier, COUNTRY). A run over several markets therefore
# has several groups, and one flat row of chips claimed a benchmark set that exists in no
# market at all — the author could not see that Singapore and Japan rank against different
# carriers.

_MULTI = {"carrier": "Zurich", "country": ["Singapore", "Japan"]}


def test_the_premium_panel_shows_one_peer_group_per_selected_country(seeded_survey):
    from studio.authoring.setup import peer_panel_state

    _opts, _style, body = peer_panel_state(_MULTI, "existing", None, {})
    rendered = _rendered(body)
    assert "Singapore" in rendered and "Japan" in rendered
    assert "One peer group per market" in rendered


def test_the_premium_groups_are_the_peers_table_s_own_per_country(seeded_survey):
    from studio.authoring.setup import peer_groups
    from studio.data import peer_members

    groups = dict(peer_groups("gpr", "Zurich", _MULTI, None))
    assert set(groups) == {"Singapore", "Japan"}
    for country, members in groups.items():
        assert set(members) <= set(peer_members("gpr", "Zurich", country=country))
    assert groups["Singapore"] != groups["Japan"], "the seed gives these markets two groups"


def test_the_survey_panel_gives_each_selected_country_its_own_peer_group(seeded_survey):
    """One dropdown per market, each pre-filled with THAT market's group.

    The old flat dropdown could pin one set for the whole run, so it had to be left empty
    here — pre-filling the union would have ranked Japan's page against Singapore's peers.
    Per market there is a right answer everywhere, so every market shows one.
    """
    groups, body = _peer_state(_MULTI)
    by_country = {c: list(peers) for c, _o, peers in groups}
    assert set(by_country) == {"Singapore", "Japan"}
    assert all(by_country.values()), "every market must show the group its page will use"
    assert by_country["Singapore"] != by_country["Japan"]
    rendered = _rendered(body)
    assert "Singapore" in rendered and "Japan" in rendered
    assert "each survey page benchmarks against its own" in rendered.lower()


def test_each_survey_market_only_offers_carriers_that_book_surveys_there(seeded_survey):
    """A candidate list drawn from the whole run would offer Japan-only names in
    Singapore — and the page cannot rank a carrier it has no score for."""
    groups, _body = _peer_state(_MULTI)
    for country, options, _peers in groups:
        names = {o["value"] for o in options}
        assert "Zurich" not in names, "a carrier is not its own peer"
        assert names, f"{country} must offer survey peer candidates"


def test_editing_one_survey_market_leaves_the_others_on_their_own_group(seeded_survey):
    """The whole point of the split: overriding Japan must not touch Singapore."""
    groups, _body = _peer_state(_MULTI, chosen={"Japan": ["Chubb"]})
    by_country = {c: list(peers) for c, _o, peers in groups}
    assert by_country["Japan"] == ["Chubb"]
    from studio.authoring.setup import _survey_result, survey_peer_group

    assert by_country["Singapore"] == list(
        survey_peer_group(_survey_result(), "Zurich", "Singapore"))


def test_the_survey_groups_are_the_peers_table_s_own_per_country(seeded_survey):
    from studio.authoring.setup import _survey_result, survey_peer_group

    result = _survey_result()
    singapore = survey_peer_group(result, "Zurich", "Singapore")
    japan = survey_peer_group(result, "Zurich", "Japan")
    assert singapore and japan and singapore != japan
    assert "Zurich" not in singapore + japan


def test_one_country_still_pins_the_group_because_there_is_only_one(seeded_survey):
    _opts, value, _body = _one_market({"carrier": "Zurich", "country": ["Singapore"]})
    assert value, "a single market has one right answer — pinning it changes nothing"


def test_the_ai_assist_control_is_gone_and_generate_pins_it_on():
    """The narrative IS the deck — nobody should have to opt in to the commentary, and
    the checkbox that let them opt out was one more decision on an already long form."""
    import inspect

    from studio.authoring import setup as S

    form = _form()
    assert "studio-ai-toggle" not in form and "AI assist" not in form
    src = inspect.getsource(S.register_setup)
    assert 'State("studio-ai-toggle"' not in src
    assert '"ai": True' in src


# ── peers: existing = names from the Peers table, custom = a scoped dropdown ──


def test_custom_peer_options_are_scoped_and_exclude_the_subject():
    scoped = carriers_in_scope({"carrier": "Zurich", "country": ["Singapore"]}, None)
    names = [o["value"] for o in scoped]
    assert names, "Singapore must offer peer candidates"
    assert "Zurich" not in names           # a carrier is never its own peer
    assert "AXA XL" in names


def test_peer_candidates_ignore_the_year_but_follow_the_scope():
    """Year must not drop a peer for a quiet year; other filters must narrow."""
    base = {"carrier": "Zurich", "country": ["Singapore"]}
    assert carriers_in_scope(base, None) == carriers_in_scope({**base, "year": 2025}, None)

    everywhere = {o["value"] for o in carriers_in_scope({"carrier": "Zurich"}, None)}
    in_country = {o["value"] for o in carriers_in_scope(base, None)}
    assert in_country <= everywhere


def test_existing_peers_come_from_the_peers_table_not_the_carrier_list():
    from studio.data import peer_members

    members = peer_members("gpr", "Zurich", country=["Singapore"])
    assert set(members) == {"AIG", "AXA XL", "Allianz", "Chubb"}  # the seed peer group
    everyone = {o["value"] for o in carriers_in_scope({"carrier": "Zurich"}, None)}
    assert set(members) < everyone  # a strict subset — never the whole list


# ── existing peers are scoped to the selected country ────────────────────────


def _peers_db(tmp_path, rows, *, with_country: bool = True):
    """A tiny GPR + Peers database, shaped like the live one."""
    from sqlalchemy import create_engine, text

    tmp_path.mkdir(parents=True, exist_ok=True)
    engine = create_engine(f"sqlite:///{tmp_path / 'peers.db'}")
    columns = "Carrier_Group TEXT, Country TEXT, Overall_Peer_Group TEXT" if with_country \
        else "Carrier_Group TEXT, Overall_Peer_Group TEXT"
    with engine.begin() as conn:
        conn.execute(text(f"CREATE TABLE Peers ({columns})"))
        for carrier, country, peer in rows:
            if with_country:
                conn.execute(
                    text("INSERT INTO Peers VALUES (:c, :co, :p)"),
                    {"c": carrier, "co": country, "p": peer},
                )
            else:
                conn.execute(
                    text("INSERT INTO Peers (Carrier_Group, Overall_Peer_Group) VALUES (:c, :p)"),
                    {"c": carrier, "p": peer},
                )
    return engine


_PEER_ROWS = [
    ("Zurich", "Singapore", "AIG"),
    ("Zurich", "Singapore", "Chubb"),
    ("Zurich", "Japan", "Tokio Marine"),
    ("Zurich", "Japan", "Sompo"),
]


def test_existing_peers_are_scoped_to_the_selected_country(tmp_path, monkeypatch):
    """The Peers table carries a Country column, so a peer group is per-market —
    Setup must list Japan's peers when Japan is selected, not Singapore's."""
    from studio import data as D

    engine = _peers_db(tmp_path, _PEER_ROWS)
    monkeypatch.setattr(D, "get_engine", lambda: engine)

    assert D.peer_members("gpr", "Zurich", country=["Singapore"]) == ["AIG", "Chubb"]
    assert D.peer_members("gpr", "Zurich", country=["Japan"]) == ["Sompo", "Tokio Marine"]


def test_multiple_countries_union_the_peer_groups(tmp_path, monkeypatch):
    from studio import data as D

    engine = _peers_db(tmp_path, _PEER_ROWS)
    monkeypatch.setattr(D, "get_engine", lambda: engine)

    members = D.peer_members("gpr", "Zurich", country=["Singapore", "Japan"])
    assert members == ["AIG", "Chubb", "Sompo", "Tokio Marine"]


def test_no_country_selected_returns_every_market_s_peers(tmp_path, monkeypatch):
    from studio import data as D

    engine = _peers_db(tmp_path, _PEER_ROWS)
    monkeypatch.setattr(D, "get_engine", lambda: engine)

    assert D.peer_members("gpr", "Zurich") == ["AIG", "Chubb", "Sompo", "Tokio Marine"]
    assert D.peer_members("gpr", "Zurich", country=["All"]) == ["AIG", "Chubb", "Sompo", "Tokio Marine"]


def test_peers_table_without_a_country_column_still_resolves(tmp_path, monkeypatch):
    """The registry declares the column; a database that predates it must fall
    back to the unscoped group rather than silently finding nobody."""
    from studio import data as D

    engine = _peers_db(tmp_path, [("Zurich", None, "AIG"), ("Zurich", None, "Chubb")], with_country=False)
    monkeypatch.setattr(D, "get_engine", lambda: engine)

    assert D.peer_members("gpr", "Zurich", country=["Japan"]) == ["AIG", "Chubb"]


def test_peer_country_column_is_ignored_when_the_table_lacks_it(tmp_path):
    from core.analytics.sql import flow_spec, peer_country_column

    spec = flow_spec("gpr")
    assert spec.peer_columns["country"] == "Country"     # declared in flows.yaml
    with_col = _peers_db(tmp_path / "a", _PEER_ROWS)
    without = _peers_db(tmp_path / "b", [("Zurich", None, "AIG")], with_country=False)
    assert peer_country_column(spec, with_col) == "Country"
    assert peer_country_column(spec, without) is None


def test_peer_set_body_shows_names_only():
    body = _rendered(peer_set_body(["AIG", "Chubb"], "2 peers from the Peers table."))
    assert "AIG" in body and "Chubb" in body
    assert "qs-peer-chip" in body
    assert "Dropdown" not in body  # existing peers are read-only


def test_existing_mode_hides_the_custom_peer_dropdown_on_first_paint():
    from studio.page.authoring.setup import _peers_panel

    wrap = next(c for c in _peers_panel().children if getattr(c, "id", "") == "studio-peer-custom-wrap")
    assert wrap.style == {"display": "none"}


# ── custom peers, market by market ───────────────────────────────────────────
#
# A peer group is a per-country statement. One flat dropdown could only pin ONE set for the
# whole run, so a two-market deck silently benchmarked Japan's page against carriers the
# author chose for Singapore. Each market now gets its own picker, its own candidate list
# and its own minimum.


def _custom_groups(selected, record=None, chosen=None):
    from studio.authoring.setup import custom_peer_groups

    return custom_peer_groups(selected, record, chosen)


def test_custom_peers_get_one_picker_per_selected_market():
    assert [c for c, _o, _v in _custom_groups(_MULTI)] == ["Singapore", "Japan"]


def test_one_country_needs_only_one_picker():
    groups = _custom_groups({"carrier": "Zurich", "country": ["Japan"]})
    assert [c for c, _o, _v in groups] == ["Japan"]
    # …and a run with no country pinned has one market: all of them.
    assert [c for c, _o, _v in _custom_groups({"carrier": "Zurich"})] == [""]


def test_each_market_only_offers_carriers_that_write_there():
    """The candidate list is what makes the choice correct — a carrier absent from Japan
    must not be offered as a Japanese peer."""
    for country, options, _v in _custom_groups(_MULTI):
        names = {o["value"] for o in options}
        expected = {o["value"] for o in
                    carriers_in_scope({**_MULTI, "country": country}, None)}
        assert names == expected
        assert "Zurich" not in names, "a carrier is never its own peer"


def test_a_markets_selection_survives_a_filter_change():
    """The pickers are rebuilt on every filter change, so a rebuild must not be a reset."""
    groups = _custom_groups(_MULTI, chosen={"Singapore": ["AIG", "Chubb"]})
    by_country = {c: list(v) for c, _o, v in groups}
    assert by_country["Singapore"] == ["AIG", "Chubb"]
    assert by_country["Japan"] == []


def test_a_peer_that_left_a_market_is_dropped_from_that_market():
    groups = _custom_groups(_MULTI, chosen={"Singapore": ["AIG", "Not A Carrier"]})
    assert {c: list(v) for c, _o, v in groups}["Singapore"] == ["AIG"]


def test_the_panel_renders_a_picker_and_a_minimum_note_per_market():
    from studio.authoring.setup import peer_panel_state

    picker, style, _msg = peer_panel_state(_MULTI, "custom", None, {})
    rendered = _rendered(picker)
    assert style == {}, "custom mode shows the pickers"
    assert rendered.count("studio-peer-custom") == 2
    assert rendered.count("studio-peer-min") == 2
    assert "Singapore" in rendered and "Japan" in rendered


# ── at least five peers, or no deck ──────────────────────────────────────────
#
# A benchmark of one or two carriers is close enough to naming them, which carrier-facing
# output may not do — and an average over two is not an aggregate. The rule is enforced
# twice: live under each picker, and again at Generate, because a note is only a note.


def test_a_thin_peer_set_is_called_out():
    from studio.page.authoring.setup import MIN_PEERS_MESSAGE, peer_min_note

    assert peer_min_note(["AIG", "Chubb"]) == MIN_PEERS_MESSAGE
    assert peer_min_note([]) == MIN_PEERS_MESSAGE
    assert peer_min_note(None) == MIN_PEERS_MESSAGE


def test_five_peers_clears_the_warning():
    from studio.page.authoring.setup import MIN_CUSTOM_PEERS, peer_min_note

    assert MIN_CUSTOM_PEERS == 5
    assert peer_min_note(["AIG", "Chubb", "Allianz", "AXA XL", "QBE"]) == ""


def test_the_note_sits_in_a_class_the_stylesheet_paints_red():
    from pathlib import Path

    css = Path("assets/studio_authoring_v4.css").read_text(encoding="utf-8")
    block = css.split(".qs-root .qs-peer-min {", 1)[1].split("}", 1)[0]
    assert "--va-danger" in block
    # …and takes no room at all while the rule is satisfied.
    assert ".qs-root .qs-peer-min:empty { display: none; }" in css


def test_the_live_check_answers_only_the_market_that_changed():
    """MATCH, not ALL: rebuilding every picker on each pick would close the open menu."""
    import inspect

    from studio.authoring import setup as S

    src = inspect.getsource(S.register_setup)
    assert 'Output({"type": "studio-peer-min", "country": MATCH}, "children")' in src
    assert 'Input({"type": "studio-peer-custom", "country": MATCH}, "value")' in src


def test_every_market_on_the_form_is_checked_not_only_the_answered_ones():
    """A run over two countries with peers pinned in ONE of them is exactly the case a
    check over the answers alone waves through."""
    from studio.authoring.setup import short_markets

    five = ["AIG", "Chubb", "Allianz", "AXA XL", "QBE"]
    assert short_markets({"Singapore": five}, countries=["Singapore", "Japan"]) == ["Japan"]
    assert short_markets({"Singapore": five, "Japan": five},
                         countries=["Singapore", "Japan"]) == []


def test_no_selection_at_all_is_under_the_minimum():
    from studio.authoring.setup import short_markets

    assert short_markets({}, countries=[""]) == [""]
    assert short_markets({}) == [""]


def test_generate_refuses_a_thin_custom_peer_set():
    """The message the author sees is the one they were asked for, in the warn tone."""
    import inspect

    from studio.authoring import setup as S
    from studio.page.authoring.setup import MIN_PEERS_MESSAGE

    warning = _rendered(S._short_peer_warning(["Japan"]))
    assert MIN_PEERS_MESSAGE in warning and "warn" in warning
    assert "Japan" in warning

    src = inspect.getsource(S.register_setup)
    assert "short_markets(" in src and "_short_peer_warning(short)" in src


def test_the_union_is_what_the_overall_block_benchmarks_against():
    """Per-market pins narrow the country pages; the overall page reports on the whole
    selection, so it benchmarks against the whole pinned field."""
    from studio.authoring.setup import flatten_peers, peers_by_country

    ids = [{"type": "studio-peer-custom", "country": "Singapore"},
           {"type": "studio-peer-custom", "country": "Japan"}]
    by_country = peers_by_country(ids, [["AIG", "Chubb"], ["Sompo", "AIG"]])
    assert by_country == {"Singapore": ["AIG", "Chubb"], "Japan": ["Sompo", "AIG"]}
    assert flatten_peers(by_country) == ["AIG", "Chubb", "Sompo"]   # de-duplicated, in order


def test_an_untouched_picker_contributes_no_market():
    from studio.authoring.setup import peers_by_country

    ids = [{"type": "studio-peer-custom", "country": "Singapore"},
           {"type": "studio-peer-custom", "country": "Japan"}]
    assert peers_by_country(ids, [["AIG"], []]) == {"Singapore": ["AIG"]}


# ── a market's peers reach that market's pages ───────────────────────────────


def test_a_country_sub_deck_benchmarks_against_its_own_pinned_peers():
    """The whole point of the split. Without this the country pages all benchmarked
    against the union — a field that exists in no market at all."""
    from studio.compute import compute_overall
    from studio.template_fill.bindings import scope_to_country

    result = compute_overall(
        filters={"carrier": "Zurich", "country": ["Singapore", "Japan"], "year": 2025},
        peers=["AIG", "Chubb", "Sompo", "Tokio Marine"],
        peers_by_country={"Singapore": ["AIG", "Chubb"], "Japan": ["Sompo", "Tokio Marine"]},
        survey_peers_by_country={"Japan": ["Sompo"]},
    )
    assert scope_to_country(result, "Singapore").peers == ("AIG", "Chubb")
    assert scope_to_country(result, "Japan").peers == ("Sompo", "Tokio Marine")
    assert scope_to_country(result, "Japan").survey_peers == ("Sompo",)
    assert result.peers == ("AIG", "Chubb", "Sompo", "Tokio Marine"), "the union is untouched"


def test_a_market_left_empty_resolves_its_own_group_not_the_union():
    """Once a run pins market by market, a market with no answer must NOT inherit what the
    other markets add up to — that union is a benchmark set no market has."""
    from studio.compute import compute_overall
    from studio.template_fill.bindings import scope_to_country

    result = compute_overall(
        filters={"carrier": "Zurich", "country": ["Singapore", "Japan"]},
        survey_peers=["Chubb", "Sompo"],
        survey_peers_by_country={"Japan": ["Sompo"]},
    )
    assert scope_to_country(result, "Japan").survey_peers == ("Sompo",)
    assert scope_to_country(result, "Singapore").survey_peers is None


def test_a_market_with_no_pin_keeps_the_runs_peers():
    """Nothing changes for a run that pinned no per-market peers — the flat pin (or the
    Peers table, when there is none) still answers."""
    from studio.compute import compute_overall
    from studio.template_fill.bindings import scope_to_country

    flat = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"]},
                           peers=["AIG", "Chubb"])
    assert scope_to_country(flat, "Singapore").peers == ("AIG", "Chubb")

    none = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"]})
    assert scope_to_country(none, "Singapore").peers is None


def test_blank_markets_are_not_pinned_at_all():
    from studio.compute import compute_overall

    result = compute_overall(filters={"carrier": "Zurich"},
                             peers_by_country={"Singapore": [], "": ["AIG"]})
    assert result.peers_by_country is None


# ── the Setup form asks questions, and explains every one of them ────────────


def test_the_deck_shape_questions_come_first():
    """They are answered from the brief, before any filter is touched."""
    first = _rendered(setup_body([], filter_options={}, filter_values={}))
    shape = first.index("studio-audience")
    assert shape < first.index("studio-data-source") < first.index("studio-peer-mode")


def test_the_controls_ask_questions_rather_than_naming_themselves():
    form = _form()
    for question in ("How much of the deck should we build?",
                     "Who is this deck for?",
                     "How should the commentary read?",
                     "Where should the numbers come from?",
                     "Which books should the deck draw on?"):
        assert question in form, question
    for label in ('"SCOPE"', '"AUDIENCE"', '"COMMENTARY STYLE"', '"DATA SOURCE"',
                  '"DATA BASIS"', '"SURVEY CARRIER"', '"SURVEY PEERS"'):
        assert label not in form, label


def test_every_question_carries_an_explanation():
    """A three-word answer to a modelling question is exactly what needs a footnote."""
    form = _form()
    for tip in ("qs-tip-scope", "qs-tip-audience", "qs-tip-style", "qs-tip-source",
                "qs-tip-basis", "qs-tip-sec-peers", "qs-tip-sec-filters",
                "qs-tip-survey-carrier", "qs-tip-survey-peers"):
        assert tip in form, tip
    assert "bi-info-circle" in form              # the ⓘ itself
    assert form.count('"target"') >= 9           # …and a tooltip pointed at each one


def test_the_page_is_named_for_what_it_makes():
    form = _form()
    assert "QBR Creator" in form and "Deck setup" not in form
    assert "Deck sections" not in form
    # …and the paragraph that explained the machinery is gone; the tooltips carry it now.
    assert "no invented numbers" not in form


def test_the_filter_pane_starts_its_second_line_with_year():
    """Auto-fill made the wrap point a function of the window width, so Year — the filter
    that decides the reporting PERIOD — could land anywhere."""
    from pathlib import Path

    from studio.page.layout import GPR_FILTERS

    assert [f["id"] for f in GPR_FILTERS][:4] == ["region", "country", "carrier", "year"]
    css = Path("assets/studio_authoring_v4.css").read_text(encoding="utf-8")
    block = css.split(".qs-root .studio-filter-grid {", 1)[1].split("}", 1)[0]
    assert "repeat(3, minmax(0, 1fr))" in block


# ── deck sections track the assembly scope ───────────────────────────────────


def test_scope_axes_maps_all_to_every_registered_axis():
    assert scope_axes("all") == ("overall", "product", "country")
    assert scope_axes("product") == ("product",)
    assert scope_axes(None) == ("overall", "product", "country")


def test_deck_sections_round_trip_between_all_and_a_single_axis():
    """Regression: switching Scope to Product and back to All used to leave the
    panel showing the overall template alone, so 'All' looked like a dead click."""
    every = _rendered(template_sections_panel("all"))
    product = _rendered(template_sections_panel("product"))

    assert every != product
    assert _rendered(template_sections_panel("all")) == every  # returning restores it

    for axis in ("OVERALL", "PRODUCT", "COUNTRY"):
        assert axis.title() in every
    assert "Overall" not in product and "Country" not in product


def test_deck_sections_page_count_is_the_sum_of_the_axes_it_assembles():
    def pages(scope: str) -> int:
        panel = template_sections_panel(scope)
        summary = panel.children[0].children[0]  # "<b>N</b> base pages"
        return int(summary.children[0].children)

    assert pages("all") == pages("overall") + pages("product") + pages("country")


def test_country_scoped_peer_group_reaches_the_deck_benchmark(tmp_path):
    """End to end over a real seed-shaped database: the peer benchmark the deck
    computes must use the SAME per-country group Setup lists, not a global one."""
    from sqlalchemy import create_engine

    from studio import data as D
    from studio.compute import _resolve_filters, peer_gap
    from studio.seed import build_seed

    seed = build_seed(tmp_path / "seed.db")
    engine = create_engine(f"sqlite:///{seed}")

    def gap(country):
        return peer_gap("gpr", _resolve_filters({"carrier": "Zurich", "country": [country],
                                                 "year": 2025}), engine)

    singapore, japan = gap("Singapore"), gap("Japan")
    assert singapore["n_peers"] == 4 and japan["n_peers"] == 4
    # Different peer sets in the two markets → different benchmarks.
    assert singapore["peer_avg"] != japan["peer_avg"]

    # And the names Setup shows match the group the benchmark resolved.
    original = D.get_engine
    try:
        D.get_engine = lambda: engine
        assert D.peer_members("gpr", "Zurich", country=["Japan"]) == [
            "AIG", "MS&AD", "Sompo", "Tokio Marine",
        ]
        assert D.peer_members("gpr", "Zurich", country=["Singapore"]) == [
            "AIG", "AXA XL", "Allianz", "Chubb",
        ]
    finally:
        D.get_engine = original
    engine.dispose()


# ── end to end: the form's selection still produces a real deck ──────────────


def test_setup_selection_builds_the_assembled_deck(tmp_path):
    """selection → data load → evidence → template binding → merged .pptx."""
    from pptx import Presentation

    from studio.authoring.generate import _assembled_export, _generated_deck

    selection = {
        "report": "qbr",                       # pinned by generate, no longer a control
        "filters": {"carrier": "Zurich", "country": ["Singapore"], "year": 2025},
        "breakdowns": ["Product_Line", "SIC_Major_Class"],
        "cuts": [],
        "peers": None,                          # existing peers → the Peers table
        "audience": "executive",
        "meeting_length": "standard",
        "style": "balanced",
        "ai": False,                            # deterministic path for the test
        "template_scope": "overall",
        "template_path": "template/overall_template.pptx",
        "dataset_id": None,
    }

    deck = _generated_deck(selection)
    assert deck and deck.slides

    path = _assembled_export(selection)
    assert path, "the overall sub-deck must assemble"
    assert len(Presentation(path).slides) > 0


def test_per_market_custom_peers_build_a_two_country_deck():
    """The new selection shape, end to end: two markets, two pinned peer sets, one merged
    deck — and each country block benchmarking against its own market's carriers."""
    from pptx import Presentation

    from studio.authoring.generate import _assembled_export, _generated_deck
    from studio.template_fill.bindings import scope_to_country

    pinned = {"Singapore": ["AIG", "AXA XL", "Allianz", "Chubb", "QBE"],
              "Japan": ["AIG", "MS&AD", "Sompo", "Tokio Marine", "Chubb"]}
    selection = {
        "report": "qbr",
        "filters": {"carrier": "Zurich", "country": ["Singapore", "Japan"], "year": 2025},
        "breakdowns": ["Product_Line", "SIC_Major_Class"],
        "cuts": [],
        "peers": ["AIG", "AXA XL", "Allianz", "Chubb", "QBE", "MS&AD", "Sompo", "Tokio Marine"],
        "peers_by_country": pinned,
        "audience": "executive",
        "meeting_length": "standard",
        "style": "balanced",
        "ai": False,                            # deterministic path for the test
        "template_scope": "country",
        "template_path": "template/country_template.pptx",
        "dataset_id": None,
    }

    assert _generated_deck(selection).slides
    path = _assembled_export(selection)
    assert path and len(Presentation(path).slides) > 0

    # …and the benchmark each country page computes is its own market's, not the union.
    from studio.compute import compute_overall

    result = compute_overall(filters=selection["filters"], peers=selection["peers"],
                             peers_by_country=pinned)
    for country, expected in pinned.items():
        assert scope_to_country(result, country).peers == tuple(expected)


def test_pinned_custom_peers_reach_the_deck():
    """Custom mode writes carriers into the selection; the benchmark must use them."""
    from studio.compute import _resolve_filters, peer_gap
    from studio.data import get_engine

    filters = _resolve_filters({"carrier": "Zurich", "country": ["Singapore"], "year": 2025})
    pinned = peer_gap("gpr", filters, get_engine(), peers=["QBE", "Sompo"])
    assert pinned and pinned["n_peers"] == 2

    from_table = peer_gap("gpr", filters, get_engine())
    assert from_table["n_peers"] == 4
    assert pinned["peer_avg"] != from_table["peer_avg"]


# ── the busy overlay: every filter/selection change is acknowledged ───────────


def test_the_form_carries_the_busy_overlay_and_one_flag_per_callback():
    from studio.page.authoring import setup as P

    form = _form()
    for flag in (P.BUSY_FORM, P.BUSY_PREVIEW, P.BUSY_SECTIONS):
        assert flag in form, f"{flag} must exist for its callback's running= to target it"
    assert "qs-page-loader" in form and "qs-page-spinner" in form


def test_every_setup_callback_raises_a_flag_while_it_works():
    """A change answered by several callbacks must stay covered until the LAST finishes,
    so each one owns its own flag rather than sharing one that races."""
    import inspect

    from studio.authoring import export as E
    from studio.authoring import setup as S

    src = inspect.getsource(S.register_setup) + inspect.getsource(E.register_export)
    assert src.count("running=") >= 3, "the option cascade, scope preview and deck sections"
    for flag in ("BUSY_FORM", "BUSY_PREVIEW", "BUSY_SECTIONS"):
        assert flag in src, f"no callback raises {flag}"


def test_a_busy_flag_is_raised_for_the_call_and_dropped_after():
    from dash import Output

    from studio.authoring.setup import _busy
    from studio.page.authoring import setup as P

    [(target, during, after)] = _busy(P.BUSY_FORM)
    assert target == Output(P.BUSY_FORM, "className")
    assert "is-busy" in during and "is-busy" not in after
    assert after == P.BUSY_FLAG_CLASS      # back to the resting class, not blank


def test_the_overlay_never_swallows_a_click():
    """It is a progress cue, not a modal: a change made while it is up must still land."""
    from pathlib import Path

    css = Path("assets/studio_authoring.css").read_text(encoding="utf-8")
    block = css.split(".qs-page-loader {", 1)[1].split("}", 1)[0]
    assert "pointer-events: none" in block


def test_the_overlay_is_shown_by_display_not_by_a_fade():
    """Reduced motion and background tabs suppress transitions; the overlay must still
    appear the moment work starts."""
    from pathlib import Path

    css = Path("assets/studio_authoring.css").read_text(encoding="utf-8")
    assert ".qs-page-loader.is-on {" in css
    on_block = css.split(".qs-page-loader.is-on {", 1)[1].split("}", 1)[0]
    assert "display: flex" in on_block


# ── the cascade sends only what changed ──────────────────────────────────────


def _diff(options, signature, token="t1"):
    from studio.authoring.setup import changed_options

    return changed_options(options, signature, token)


def test_the_first_cascade_sends_every_list():
    from dash import no_update

    fresh, signature = _diff({"carrier": [{"label": "AIG", "value": "AIG"}], "year": []}, None)
    assert no_update not in fresh.values()
    assert set(signature["digests"]) == {"carrier", "year"} and signature["token"] == "t1"


def test_an_unchanged_list_is_not_sent_again():
    """The saving: a filter change moves two or three of the ten lists, and re-sending the
    rest costs payload and a re-render that both grow with the vocabulary."""
    from dash import no_update

    first = {"carrier": [{"label": "AIG", "value": "AIG"}],
             "country": [{"label": "Japan", "value": "Japan"}]}
    _, signature = _diff(first, None)

    changed = {**first, "country": [{"label": "Japan", "value": "Japan"},
                                    {"label": "Singapore", "value": "Singapore"}]}
    fresh, _ = _diff(changed, signature)
    assert fresh["carrier"] is no_update
    assert fresh["country"] == changed["country"]


def test_a_re_rendered_form_is_sent_everything_again():
    """A re-render rebuilds the dropdowns from the UNFILTERED lists, so a signature from the
    previous rendering describes something that is no longer on screen — skipping a list
    there would leave the user looking at options the selection has ruled out."""
    from dash import no_update

    options = {"carrier": [{"label": "AIG", "value": "AIG"}]}
    _, signature = _diff(options, None, token="render-1")
    fresh, _ = _diff(options, signature, token="render-2")
    assert fresh["carrier"] is not no_update


def test_the_form_carries_a_token_that_changes_with_each_rendering():
    from studio.page.authoring.setup import form_token

    assert form_token().data != form_token().data
