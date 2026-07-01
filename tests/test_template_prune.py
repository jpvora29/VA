"""Pruning surplus per-country pages down to the country count.

Hermetic: builds a deck whose slides enumerate countries the way the product template does
(feedback pages of 2 countries each, plus a product-level page and a 1-4 ranking page).
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from studio.template_fill.analyze import analyze
from studio.template_fill.prune import hidden_country_pages


def _deck_with_country_pages(path: str) -> str:
    prs = Presentation()
    pages = [
        "Portfolio Analysis",                       # 0: no country tokens → never dropped
        "Feedback — Country (1) / Country (2)",     # 1: countries 1-2
        "Feedback — Country (3) / Country (4)",     # 2: countries 3-4
        "Feedback — Country (5) / Country (6)",     # 3: countries 5-6
        "Ranking — Country (1) Country (2) Country (3) Country (4)",  # 4: min idx 1 → never dropped
    ]
    for text in pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = text
    prs.save(path)
    return path


def test_prunes_pages_wholly_beyond_country_count(tmp_path):
    t = analyze(_deck_with_country_pages(str(tmp_path / "product.pptx")))

    assert hidden_country_pages(t, 2) == [2, 3]      # keep countries 1-2 page, drop 3-6
    assert hidden_country_pages(t, 4) == [3]         # keep 1-4, drop 5-6
    assert hidden_country_pages(t, 6) == []          # everything fits
    assert hidden_country_pages(t, 1) == [2, 3]      # page 1 still kept (has country 1)


def test_pages_without_country_tokens_are_never_hidden(tmp_path):
    t = analyze(_deck_with_country_pages(str(tmp_path / "product.pptx")))
    hidden = hidden_country_pages(t, 1)
    assert 0 not in hidden      # portfolio page (no tokens)
    assert 4 not in hidden      # ranking page (min index 1)
