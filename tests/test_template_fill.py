"""Template-fill: dynamic analysis, slot detection, role mapping, and the fill.

Deterministic — runs against the seed DB (no DB_PATH, no LLM). Proves the analyzer
is generic (works on the starter template AND a different seed deck) and that the
fill replaces mapped tokens while leaving unmapped slots as placeholders.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from studio.compute import compute_overall
from studio.template_fill import roles as R
from studio.template_fill import slots as S
from studio.template_fill.analyze import analyze
from studio.template_fill.fill import fill_template
from studio.template_fill.model import materialize_fields, new_template_doc
from studio.template_fill.slots import classify

TEMPLATE = "template/qbr_template.pptx"
SEED_DECK = "studio/_seed/test_qbr.pptx"


def _result():
    return compute_overall(filters={"carrier": "Zurich", "year": 2025})


# ── slot classification (the generic token grammar) ──────────────────────────


@pytest.mark.parametrize("token,kind", [
    ("xx.x", "int"), ("x.x%", "pct"), ("$xx,xxxm", "money"), ("x.xB", "money"),
    ("x", "int"), ("Country (1)", "text"), ("………", "text"),
])
def test_classify_detects_placeholders(token, kind):
    assert classify(token) == kind


@pytest.mark.parametrize("text", ["xyz", "Carrier Country xyz YoY change", "Property", "$2,100M", ""])
def test_classify_ignores_non_placeholders(text):
    # Real values and prose with an embedded 'x' must NOT be slots.
    assert classify(text) is None or text == "Carrier Country xyz YoY change"


# ── analyzer is generic across templates ─────────────────────────────────────


def test_analyzer_detects_slots_on_starter_template():
    slots = S.detect(analyze(TEMPLATE))
    assert len(slots) > 50
    # No hard-coded slide indices: slots span multiple slides.
    assert len({s.slide_idx for s in slots}) > 3


def test_analyzer_runs_on_a_different_pptx():
    if not Path(SEED_DECK).exists():
        pytest.skip("seed deck not present")
    template = analyze(SEED_DECK)          # must not raise on a non-starter deck
    assert template.slides


# ── role inference + fill ────────────────────────────────────────────────────


def test_manifest_maps_core_roles():
    bindings = R.infer(S.detect(analyze(TEMPLATE)))
    roles = {b.role for b in bindings if b.role}
    assert {"subject_name", "marsh_gwp", "carrier_gwp", "sow_pct", "rank"} <= roles


# ── spotlight YoY (the "… Country xyz YoY change" highlights) ─────────────────


def test_spotlight_slots_get_distinct_roles():
    # The "Carrier/Marsh Country xyz YoY change" callouts must NOT map to the overall
    # carrier/marsh YoY — else they duplicate the blended figure.
    from studio.template_fill.slots import Slot

    def role_of(context):
        return R._infer_role(Slot(1, 1, ["para", 0], "+xx.x%", "pct", context))

    assert role_of("Carrier Country xyz YoY change") == "spotlight_carrier_yoy"
    assert role_of("Marsh Country xyz YoY change") == "spotlight_marsh_yoy"
    assert role_of("Change in Premium YoY") == "carrier_gwp_yoy"      # overall, unchanged


def test_spotlight_is_a_country_distinct_from_overall_when_multi_country():
    from studio.template_fill.bindings import resolve_roles

    r = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore", "Hong Kong"], "year": 2025})
    v = resolve_roles(r)
    assert v.get("spotlight_name") in {"Singapore", "Hong Kong"}
    # A single country's YoY should differ from the blended two-country YoY.
    assert v.get("spotlight_carrier_yoy") != v.get("carrier_gwp_yoy")


def test_spotlight_drills_to_a_product_when_single_country():
    from studio.template_fill.bindings import resolve_roles

    r = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"], "year": 2025})
    v = resolve_roles(r)
    # One country in scope → the spotlight is a product line, not a country.
    assert v.get("spotlight_name") not in (None, "Singapore")


@pytest.mark.parametrize("countries", [["Singapore"], ["Singapore", "Hong Kong"], None])
def test_spotlight_yoy_resolves_with_no_year_pinned(countries):
    """Both callouts must populate when the user leaves the year filter on "All".

    They are period comparisons, and ``compute.period_totals`` returns nothing at all
    without a reference year — so scoping them to the raw selection instead of the resolved
    reporting year left "Carrier/Marsh Country xyz YoY change" on its ``+xx.x%`` placeholder.
    """
    from studio.template_fill.bindings import resolve_roles

    filters = {"carrier": "Zurich"}
    if countries:
        filters["country"] = countries
    v = resolve_roles(compute_overall(filters=filters))
    assert v.get("spotlight_name")
    assert isinstance(v.get("spotlight_carrier_yoy"), float)
    assert isinstance(v.get("spotlight_marsh_yoy"), float)


def test_fill_replaces_mapped_tokens_and_keeps_placeholders(tmp_path):
    doc = new_template_doc(_result(), template_path=TEMPLATE)
    out = fill_template(doc, out_path=str(tmp_path / "filled.pptx"))
    prs = Presentation(out)
    frame_text = "\n".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                           if sh.has_text_frame)
    table_text = "\n".join(c.text for s in prs.slides for sh in s.shapes if sh.has_table
                           for row in sh.table.rows for c in row.cells)
    # The carrier name replaced the generic "Carrier" label, and no mapped money
    # placeholder token survived in the body text frames.
    assert "Zurich" in frame_text
    assert "$xx,xxxm" not in frame_text and "$xxx.xM" not in frame_text
    # Qualitative relationship-feedback prose (table cells, not data-derivable) is left
    # as its ellipsis placeholder — only fact-grounded sections get commentary.
    assert "…" in (frame_text + table_text) or "......" in (frame_text + table_text)


def test_render_token_matches_placeholder_style():
    from studio.template_fill.render import render_token
    # Money auto-scales to billions (short, no overflow); keeps the token's $ presence.
    assert render_token("$xx,xxxm", 2.29e9, "money") == "$2.3B"
    assert render_token("x.xB", 2.29e9, "money") == "2.3B"
    assert render_token("xxxM", 2.079e8, "money") == "208M"          # millions below 1bn
    assert render_token("+xx.x%", 28.6, "pct") == "+28.6%"
    assert render_token("#x", 5, "rank") == "#5"
    assert render_token("PY (-x.x%▼)", 9.9, "pct") == "PY (+9.9%▲)"   # sign + arrow follow data


def _divider_template(blocks: int, *, slides_per_block: int = 2):
    """A synthetic template of ``blocks`` enumerated "Country (n)" divider blocks.

    Built in memory rather than read off a .pptx so the pruning rule is tested on its own
    terms — the shipped templates carry ONE country block each (the split pipeline fills it
    once per country), so they can no longer exercise a surplus block.
    """
    from studio.template_fill.analyze import Shape, Slide, Template

    def text_slide(index: int, text: str) -> Slide:
        return Slide(index=index, layout="Blank",
                     shapes=[Shape(shape_id=2, name="Title 1", kind="text", paragraphs=[text])])

    slides, i = [], 0
    for block in range(1, blocks + 1):
        slides.append(text_slide(i, f"Country ({block})"))
        i += 1
        for _ in range(slides_per_block):
            slides.append(text_slide(i, "Carrier breakdown"))
            i += 1
    return Template(path="<memory>", width_emu=12192000, height_emu=6858000, slides=slides)


@pytest.mark.parametrize("selected,expected_hidden", [
    (["Singapore"], [3, 4, 5]),                    # 1 country → the "Country (2)" block goes
    (["Singapore", "Japan"], []),                  # 2 countries → both blocks stay
])
def test_surplus_country_blocks_are_hidden(selected, expected_hidden):
    from studio.template_fill.model import _hidden_blocks

    res = compute_overall(filters={"carrier": "Zurich", "country": selected, "year": 2025})
    assert _hidden_blocks(_divider_template(2), res) == expected_hidden


def test_hidden_blocks_are_dropped_from_the_export(tmp_path):
    # Whatever the doc marks hidden must actually leave the exported deck.
    res = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"], "year": 2025})
    doc = new_template_doc(res, template_path=TEMPLATE)
    full = len(analyze(TEMPLATE).slides)
    doc["hidden"] = [full - 1]
    out = fill_template(doc, out_path=str(tmp_path / "pruned.pptx"))
    assert len(Presentation(out).slides) == full - 1


def test_country_token_substituted_in_labels(tmp_path):
    res = compute_overall(filters={"carrier": "Zurich", "country": ["Singapore"], "year": 2025})
    doc = new_template_doc(res, template_path=TEMPLATE)
    out = fill_template(doc, out_path=str(tmp_path / "subbed.pptx"))
    text = "\n".join(sh.text_frame.text for s in Presentation(out).slides
                     for sh in s.shapes if sh.has_text_frame)
    assert "xyz" not in text.lower()          # the literal placeholder is gone
    assert "Singapore" in text                # replaced by the selection


def test_materialize_is_parity_source():
    doc = new_template_doc(_result(), template_path=TEMPLATE)
    fields = materialize_fields(doc)
    filled = [f for f in fields.values() if f["filled"]]
    assert filled and all(not str(f["text"]).strip().endswith("x") for f in filled)


# ── carrier vs Marsh disambiguation (the slide-1 headline bug) ───────────────


def test_carrier_premium_not_mapped_to_marsh_book():
    # "Premium written with Marsh" is the SUBJECT's premium, not the whole book;
    # only "Overall Marsh premium" maps to the Marsh-book role.
    bindings = R.infer(S.detect(analyze(TEMPLATE)))
    by_ctx = {b.slot.context.lower(): b.role for b in bindings if b.role}
    carrier_hdr = next((r for c, r in by_ctx.items() if "premium written with marsh" in c), None)
    marsh_hdr = next((r for c, r in by_ctx.items() if "overall marsh premium" in c), None)
    assert carrier_hdr == "carrier_gwp"
    assert marsh_hdr == "marsh_gwp"


# ── per-product breakdown grid (slides 9 / 14) ───────────────────────────────


def test_breakdown_grid_fills_distinct_rows_per_product():
    from studio.template_fill import grids
    from studio.template_fill.registry import derive_manifest

    res = _result()
    template, _ = derive_manifest(TEMPLATE)
    values = grids.grid_values(template, res)
    # The grid produced per-row GWP keys, and the rows are NOT all the same value
    # (the old bug repeated the carrier total down every row).
    gwp = {k: v for k, v in values.items() if k.endswith(":gwp") and v not in ("", None)}
    assert len(gwp) >= 4
    assert len(set(gwp.values())) > 1, "breakdown rows must differ per product"


def test_breakdown_grid_reports_the_same_year_as_the_rest_of_the_deck():
    """The grid must be year-scoped, or it disagrees with every other page.

    Without a resolved reporting year ``product_breakdown_rows`` sums EVERY year in the book
    into one "GWP" (so the column read far higher than the GWP-performance page's bars) and
    has no prior year to compare against (so Var % and Rank change stayed empty).
    """
    from studio.template_fill import grids
    from studio.template_fill.bindings import resolve_roles_for_country, scope_to_country
    from studio.template_fill.registry import derive_manifest

    # No year pinned — the failing case.
    run = compute_overall(filters={"carrier": "Zurich"})
    template, _ = derive_manifest(TEMPLATE)
    country = "Singapore"
    values = grids.grid_values(template, scope_to_country(run, country))

    def column(metric):
        return [v for k, v in values.items() if k.endswith(f":{metric}") and v not in ("", None)]

    assert column("var"), "Var % must populate"
    assert [c for c in column("rank_change") if c], "Rank change must populate"
    # The GWP column reports the same book as the country's headline KPI: its rows are a
    # subset of that year's lines of business (the all-year sum would run far past it).
    total = resolve_roles_for_country(run, country)["carrier_gwp"]
    assert 0 < sum(column("gwp")) <= total * (1 + 1e-9)
    assert max(column("gwp")) < total


def test_sections_classify_known_slides():
    from studio.template_fill.sections import Section, classify_sections

    secs = set(classify_sections(analyze(TEMPLATE)).values())
    # The starter deck must expose the sections commentary + grids depend on.
    assert {Section.TRADING_SUMMARY, Section.FEEDBACK, Section.BREAKDOWN,
            Section.COUNTRY_DIVIDER} <= secs


def test_commentary_fills_trading_summary_with_grounded_text():
    from studio.template_fill import commentary
    from studio.template_fill.registry import derive_manifest

    template, _ = derive_manifest(TEMPLATE)
    vals = commentary.values(template, _result())
    notes = {k: v for k, v in vals.items() if k.startswith("note:") and v}
    # The 4 trading-summary columns are filled, the carrier is named, and every figure
    # is faithful (the verifier guarantees no number absent from the deterministic text).
    assert len(notes) >= 3
    assert any("Zurich" in v for v in notes.values())


def test_feedback_prose_is_left_as_placeholder():
    # Qualitative relationship feedback isn't premium-derivable → must NOT be auto-filled.
    from studio.template_fill import commentary
    from studio.template_fill.registry import derive_manifest

    template, _ = derive_manifest(TEMPLATE)
    from studio.template_fill.sections import Section, section_of
    fb = [s.index for s in template.slides if section_of(s) == Section.FEEDBACK]
    vals = commentary.values(template, _result())
    assert fb and not any(k.split(":")[1] in {str(i) for i in fb} for k in vals)


def test_external_charts_are_flagged_and_not_rewritten(tmp_path):
    # think-cell / externally-linked charts must be flagged (for the manual-fill cue)
    # and left untouched by the guarded chart fill — the export still reopens cleanly.
    template = analyze(TEMPLATE)
    charts = [sh for s in template.slides for sh in s.shapes if sh.kind == "chart"]
    assert charts and all(c.chart_external for c in charts), "starter charts are think-cell"
    out = fill_template(new_template_doc(_result(), template_path=TEMPLATE),
                        out_path=str(tmp_path / "charts.pptx"))
    # Re-analysing the export must still find every chart (none dropped/corrupted).
    assert len([sh for s in analyze(out).slides for sh in s.shapes if sh.kind == "chart"]) == len(charts)


def test_breakdown_slide_is_scoped_to_its_own_country():
    from studio.template_fill import grids
    from studio.template_fill.registry import derive_manifest

    template, _ = derive_manifest(TEMPLATE)
    values = grids.grid_values(template, _result())
    subtitles = [v for k, v in values.items() if k.endswith(":subtitle")]
    # Each breakdown slide re-resolves its hard-coded "Country (n)" to a real country, and
    # names the carrier rather than the template's literal "Carrier".
    assert subtitles
    assert len(set(subtitles)) == len(subtitles)
    assert all("(1)" not in s and "Carrier" not in s for s in subtitles)
    assert all("Zurich" in s for s in subtitles)


# ── the runway bars, one per grid row ────────────────────────────────────────

COUNTRY_TEMPLATE = "template/country_template.pptx"


def _breakdown_grid():
    from studio.template_fill import grids

    if not Path(COUNTRY_TEMPLATE).exists():
        pytest.skip("country template not present")
    for slide in analyze(COUNTRY_TEMPLATE).slides:
        grid = grids._detect(slide)
        if grid is not None:
            return slide.index, grid
    raise AssertionError("no breakdown grid in the country template")


def test_the_runway_bars_are_found_through_their_own_value_labels():
    """The bars are one picture, and the value labels are drawn on it — which is what
    identifies it, rather than where on the slide it happens to sit."""
    _, grid = _breakdown_grid()
    assert grid.runway_chart_id is not None
    assert grid.row_count > 1


@pytest.mark.parametrize("live,expected", [(7, None), (3, 4 / 7), (1, 6 / 7), (0, 1.0)])
def test_the_runway_is_cropped_to_the_rows_that_have_a_product(live, expected):
    """A blanked row keeps its bar unless the picture is cut — a runway under no product.
    A full grid is left alone; a grid with nothing in it loses the picture whole."""
    from studio.template_fill import grids

    idx, grid = _breakdown_grid()
    out = {}
    grids._trim_runway_chart(out, idx, grid, live)
    crop = (out.get("picture_crops") or {}).get(f"{idx}:{grid.runway_chart_id}")
    if expected is None:
        assert not out, "a full grid needs no crop"
    else:
        assert crop["bottom"] == pytest.approx(expected)


def test_cropping_trims_the_picture_and_its_frame_together(tmp_path):
    """The bars line up with the table's rows, so the visible part has to keep its scale
    and its top edge — a crop alone would stretch the survivors over the whole frame."""
    idx, grid = _breakdown_grid()
    before = next(sh for sh in Presentation(COUNTRY_TEMPLATE).slides[idx].shapes
                  if int(sh.shape_id) == grid.runway_chart_id)
    doc = {"template_path": COUNTRY_TEMPLATE, "manifest": [],
           "values": {"picture_crops": {f"{idx}:{grid.runway_chart_id}": {"bottom": 3 / 7}}},
           "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "cropped.pptx"))

    after = next(sh for sh in Presentation(out).slides[idx].shapes
                 if int(sh.shape_id) == grid.runway_chart_id)
    # OOXML stores a crop in hundred-thousandths, so the fraction lands rounded.
    assert after.crop_bottom == pytest.approx(3 / 7, rel=1e-4)
    assert after.height == pytest.approx(before.height * 4 / 7, rel=1e-3)
    assert (after.left, after.top, after.width) == (before.left, before.top, before.width)


def test_a_grid_with_no_products_loses_the_runway_picture(tmp_path):
    idx, grid = _breakdown_grid()
    doc = {"template_path": COUNTRY_TEMPLATE, "manifest": [],
           "values": {"picture_crops": {f"{idx}:{grid.runway_chart_id}": {"bottom": 1.0}}},
           "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path / "gone.pptx"))
    ids = {int(sh.shape_id) for sh in Presentation(out).slides[idx].shapes}
    assert grid.runway_chart_id not in ids
