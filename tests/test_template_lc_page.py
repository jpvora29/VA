"""The LC-ranking page: panel detection, responsive layout, and the refilled quadrants.

Two layers:
  * unit — panel ordering, rank bands, the layout plan, the re-anchoring rule and the label
    placement, on in-memory objects, so the logic is pinned independently of any .pptx;
  * integration — the real ``template/overall_template.pptx`` valued from the seed DB and
    then filled, proving a panel populates end-to-end, that the author's painted quadrant
    bands survive onto the new frame, and that the deck STILL OPENS (these charts are
    authored off binary workbooks, which a naive ``replace_data`` corrupts).

Deterministic: seed DB, no LLM.
"""
from __future__ import annotations

from dataclasses import replace
from pathlib import Path

import pytest
from pptx import Presentation

from studio.compute import compute_overall
from studio.template_fill import fill as F
from studio.template_fill import lc_page as L
from studio.template_fill.analyze import Shape, Slide, Template, analyze

OVERALL_TEMPLATE = "template/overall_template.pptx"
IN = 914400                                     # EMU per inch


def _scatter(shape_id: int, x: float, y: float) -> Shape:
    return Shape(shape_id=shape_id, name=f"Chart {shape_id}", kind="chart",
                 chart_type="XY_SCATTER (-4169)",
                 x=int(x * IN), y=int(y * IN), w=int(6.2 * IN), h=int(1.8 * IN))


def _page(*charts: Shape) -> Template:
    title = Shape(shape_id=1, name="Title 1", kind="text",
                  paragraphs=["Marsh Portfolio and LC ranking"])
    return Template(path="synthetic", width_emu=12192000, height_emu=6858000,
                    slides=[Slide(index=0, layout="", shapes=[title, *charts])])


def _grid() -> Template:
    return _page(_scatter(40, x=6.8, y=4.71), _scatter(10, x=0.4, y=2.34),
                 _scatter(30, x=0.4, y=4.71), _scatter(20, x=6.8, y=2.35))


# ── panel detection ──────────────────────────────────────────────────────────


def test_panels_are_ordered_top_row_first_then_left_to_right():
    # Authored out of order, and with the two panels of a row a hair apart vertically —
    # a plain (y, x) sort would read that as four rows of one.
    assert [p.shape_id for p in L.panels(_grid())] == [10, 20, 30, 40]


def test_only_the_ranking_page_is_a_panel():
    page = _page(_scatter(10, x=0.4, y=2.34))
    page.slides[0].shapes[0].paragraphs = ["Carrier vs Marsh growth rates"]
    assert L.panels(page) == []


# ── rank bands ───────────────────────────────────────────────────────────────


@pytest.mark.parametrize("rank,band", [
    (1, L.LEAD), (2, L.LEAD),                   # the position to defend
    (3, L.STRONG), (5, L.STRONG),               # …still inside the deck's top-5 benchmark
    (6, L.CHASING), (8, L.CHASING),             # outside it, within reach
    (9, L.BEHIND), (30, L.BEHIND),              # the ground to make up
])
def test_band_of_follows_the_decks_top_five_benchmark(rank, band):
    assert L.band_of(rank) == band


# ── responsive layout ────────────────────────────────────────────────────────


def _frames():
    return [p.frame for p in L.panels(_grid())]


def test_one_country_takes_the_whole_page():
    [rect] = L.plan_layout(_frames(), 1)
    whole = L._union(_frames())
    assert (rect.x, rect.y, rect.w, rect.h) == (whole.x, whole.y, whole.w, whole.h)
    assert rect.h > _frames()[0].h * 2           # both rows, not a quarter of the page


def test_two_countries_share_the_page_side_by_side_at_full_height():
    frames = _frames()
    left, right = L.plan_layout(frames, 2)
    assert (left.x, left.w) == (frames[0].x, frames[0].w)      # the author's own columns
    assert (right.x, right.w) == (frames[1].x, frames[1].w)
    assert left.y == right.y and left.h == right.h == L._union(frames).h


@pytest.mark.parametrize("live", [3, 4])
def test_three_or_four_countries_keep_the_authored_grid(live):
    frames = _frames()
    assert L.plan_layout(frames, live) == frames[:live]


def test_a_layout_it_cannot_plan_falls_back_to_the_authored_grid():
    frames = _frames()
    assert L.plan_layout(frames, 0) == frames and L.plan_layout(frames, 9) == frames


# ── re-anchoring the panel's own furniture ───────────────────────────────────


def test_furniture_keeps_the_relationship_the_author_gave_it():
    from studio.template_fill.fill import _reanchor

    old = (0, 0, 600, 180)                       # a panel…
    new = (0, 0, 600, 400)                       # …grown to full height
    # A caption hugging the bottom edge stays under the panel.
    assert _reanchor((100, 190, 300, 20), old, new) == (100, 410, 300, 20)
    # A title above it stays above.
    assert _reanchor((10, -40, 300, 20), old, new) == (10, -40, 300, 20)
    # A rotated caption centred beside it stays centred.
    assert _reanchor((-60, 80, 120, 20), old, new) == (-60, 190, 120, 20)
    # A rule spanning the panel spans the new one.
    assert _reanchor((10, -10, 580, 4), old, new) == (10, -10, 580, 4)
    assert _reanchor((10, -10, 580, 4), old, (0, 0, 1200, 400)) == (10, -10, 1180, 4)


def test_the_painted_bands_follow_the_plot_area_not_the_chart_frame():
    # The plot's insets are a FRACTION of the chart, so a panel grown to the whole page
    # moves its plot area further in — bands scaled off the frame would sit adrift of the
    # points they belong behind.
    from studio.template_fill.fill import _move_panel_furniture

    class _Box:
        left, top, width, height = 100, 100, 600, 180

    band = _Box()
    _move_panel_furniture([band], (100, 100, 600, 180), (100, 100, 1200, 400),
                          old_plot=(100, 100, 600, 180), new_plot=(160, 110, 1080, 380))
    assert (band.left, band.top, band.width, band.height) == (160, 110, 1080, 380)


# ── label placement ──────────────────────────────────────────────────────────


def test_labels_stay_inside_the_plot_and_off_each_other():
    from pptx.enum.chart import XL_LABEL_POSITION as POS

    points = [{"name": "Property", "size": 900.0, "rank": 6},   # near the right edge
              {"name": "Casualty", "size": 800.0, "rank": 6},   # …sharing its row
              {"name": "Marine", "size": 300.0, "rank": 6},     # …and a third
              {"name": "Cyber", "size": 200.0, "rank": 1}]
    assert L.band_of(1) == L.LEAD                # sanity: the fixture is a real ranking
    assert _sides(points) == [POS.LEFT, POS.ABOVE, POS.BELOW, POS.RIGHT]


def _sides(points):
    from studio.template_fill.fill import _label_sides

    return _label_sides(points, 1000.0)


# ── integration: the real template, valued from the seed DB ──────────────────


def _result(countries, product=None):
    from studio.template_fill.bindings import scope_to_product

    run = compute_overall(filters={"carrier": "Zurich", "country": countries, "year": 2025})
    run = replace(run, scope_countries=tuple(countries))
    return scope_to_product(run, product) if product else run


@pytest.fixture(scope="module")
def overall_template():
    if not Path(OVERALL_TEMPLATE).exists():
        pytest.skip("overall template not present")
    return analyze(OVERALL_TEMPLATE)


def test_every_panel_gets_a_payload_and_the_ones_in_scope_get_points(overall_template):
    panels = L.values(overall_template, _result(["Singapore", "Japan"]))["lc_ranking"]
    assert len(panels) == len(L.panels(overall_template)) == 4
    filled = [p for p in panels.values() if p["points"]]
    assert len(filled) == 2 and all(p["country"] for p in filled)
    # The panels past the countries in scope carry no points — the fill engine drops them
    # rather than shipping the template's authored example book under an erased title.
    assert [p["points"] for p in panels.values() if not p["country"]] == [[], []]


def test_two_live_panels_are_re_laid_out_to_take_the_page_back(overall_template):
    panels = list(L.values(overall_template, _result(["Singapore", "Japan"]))["lc_ranking"].values())
    authored = L.panels(overall_template)[0].frame
    live = [p["rect"] for p in panels if p["points"]]
    assert len(live) == 2
    for rect in live:
        assert rect["h"] > authored.h * 1.9      # full height, not a quarter of the page


def test_points_run_biggest_pool_first_and_carry_a_rank(overall_template):
    panel = next(p for p in L.values(overall_template, _result(["Singapore"]))["lc_ranking"].values()
                 if p["points"])
    sizes = [p["size"] for p in panel["points"]]
    assert sizes == sorted(sizes, reverse=True) and len(sizes) > 1
    for point in panel["points"]:
        assert point["size"] > 0 and point["rank"] >= 1
        assert point["band"] in {L.LEAD, L.STRONG, L.CHASING, L.BEHIND}


def test_a_product_subdeck_still_ranks_the_whole_line_of_business_mix(overall_template):
    # The page is a PORTFOLIO view: a per-product sub-deck's own pin must not leave every
    # panel with the single point of the product the sub-deck happens to be about.
    scoped = _result(["Singapore"], product="Property")
    panel = next(p for p in L.values(overall_template, scoped)["lc_ranking"].values() if p["points"])
    names = {p["name"] for p in panel["points"]}
    assert len(names) > 1 and "Property" in names


# ── the page reports on the Setup selection ──────────────────────────────────


_PICKED = ("Property", "Financial Lines", "Cyber")


def _points_of(template, result):
    panel = next(p for p in L.values(template, result)["lc_ranking"].values() if p["points"])
    return {p["name"] for p in panel["points"]}


def test_the_page_ranks_only_the_lines_the_run_selected(overall_template):
    """Setup governs the page: pinning three lines of business must not leave the panel
    ranking the carrier's whole book. Regression — the product pin was dropped outright."""
    run = replace(_result(["Singapore"]), scope_products=_PICKED)
    assert _points_of(overall_template, run) == set(_PICKED)


def test_a_product_subdeck_widens_back_to_the_selection_not_the_whole_book(overall_template):
    """A per-product page still needs a portfolio to rank its product against — but that
    portfolio is the run's selection, not every line the carrier writes."""
    run = replace(_result(["Singapore"], product="Property"), scope_products=_PICKED)
    assert _points_of(overall_template, run) == set(_PICKED)


def test_with_no_selection_the_page_still_ranks_the_whole_book(overall_template):
    """Nothing pinned means nothing to narrow to — the page keeps its old behaviour."""
    assert len(_points_of(overall_template, _result(["Singapore"]))) > len(_PICKED)


# ── integration: the written panel ───────────────────────────────────────────


_POINTS = [{"name": "Property", "size": 545e6, "rank": 6, "band": L.CHASING},
           {"name": "Financial Lines", "size": 414e6, "rank": 3, "band": L.STRONG},
           {"name": "Cyber", "size": 300e6, "rank": 1, "band": L.LEAD}]


@pytest.fixture(scope="module")
def filled_page(tmp_path_factory):
    if not Path(OVERALL_TEMPLATE).exists():
        pytest.skip("overall template not present")
    from studio.template_fill.fill import fill_template

    panels = L.panels(analyze(OVERALL_TEMPLATE))
    whole = L.plan_layout([p.frame for p in panels], 1)[0]
    payload = {f"{p.slide_idx}:{p.shape_id}": {"country": None, "points": [],
                                               "rect": p.frame.to_dict()}
               for p in panels}
    payload[f"{panels[0].slide_idx}:{panels[0].shape_id}"] = {
        "country": "Singapore", "points": _POINTS, "rect": whole.to_dict()}
    doc = {"template_path": OVERALL_TEMPLATE, "manifest": [],
           "values": {"lc_ranking": payload}, "overrides": {}, "map_overrides": {}, "added": {}}
    out = fill_template(doc, out_path=str(tmp_path_factory.mktemp("lc") / "lc.pptx"))
    return Presentation(out).slides[panels[0].slide_idx], panels, whole


def _only_chart(slide):
    charts = [sh for sh in slide.shapes if getattr(sh, "has_chart", False)]
    assert len(charts) == 1, "the panels with no country must be removed whole"
    return charts[0]


def _plot_box(shape):
    """Where the chart's plot area sits on the slide, in absolute EMU."""
    from pptx.oxml.ns import qn

    manual = shape.chart._chartSpace.plotArea.find(qn("c:layout")).find(qn("c:manualLayout"))
    fx, fy, fw, fh = (float(manual.find(qn(f"c:{k}")).get("val")) for k in ("x", "y", "w", "h"))
    return (shape.left + fx * shape.width, shape.top + fy * shape.height,
            fw * shape.width, fh * shape.height)


def test_the_panel_stays_the_authored_scatter_over_its_planned_rect(filled_page):
    """The panel's PLOT lands on the planned rect's authored plot area. Its frame is wider
    than that — the money ticks it now states need room outside the plot, and PowerPoint
    drops a manual plot layout it cannot label, so the frame has to hold both."""
    slide, panels, whole = filled_page
    shape = _only_chart(slide)
    assert "SCATTER" in str(shape.chart.chart_type)
    authored = next(sh for sh in Presentation(OVERALL_TEMPLATE).slides[panels[0].slide_idx].shapes
                    if getattr(sh, "has_chart", False)
                    and int(sh.shape_id) == panels[0].shape_id)
    planned = F._plot_frame(authored.chart, (whole.x, whole.y, whole.w, whole.h))
    assert _plot_box(shape) == pytest.approx(planned, abs=2)
    assert shape.left <= planned[0] and shape.left + shape.width >= planned[0] + planned[2]
    assert shape.top + shape.height >= planned[1] + planned[3]


def test_the_plot_lands_exactly_on_the_painted_bands(filled_page):
    """The bands are the panel's backdrop: the plot border drawn over them has to be the
    rectangle they are, or the page shows two frames a quarter-inch apart."""
    slide, _, whole = filled_page
    band = next(sh for sh in slide.shapes if str(sh.shape_type).startswith("GROUP")
                and sh.width > whole.w * 0.8)
    plot = _plot_box(_only_chart(slide))
    assert plot == pytest.approx((band.left, band.top, band.width, band.height),
                                 abs=0.01 * band.width)


def test_each_point_is_a_line_of_business_named_on_the_authored_axes(filled_page):
    chart = _only_chart(filled_page[0]).chart
    series = chart.plots[0].series[0]
    assert [round(v) for v in series.iter_values()] == [p["rank"] for p in _POINTS]
    assert [pt.data_label.text_frame.text for pt in series.points] == [p["name"] for p in _POINTS]
    # The money axis is scaled to this country's largest pool and states its own ticks
    # (the author's hand-drawn broken axis is gone); the rank axis keeps the author's 0–11,
    # which is what puts their painted band boundary on the top-5 line.
    assert chart.category_axis.maximum_scale == pytest.approx(545e6 * 1.04)
    assert chart.value_axis.maximum_scale == 11.0
    assert not chart.has_legend and not chart.has_title


def test_the_plot_keeps_the_authored_ruling(filled_page):
    """The author authored both grids ``noFill`` — the painted priority bands are the
    panel's backdrop, and a grid over them competes with the reading they exist to give.
    The refill must not paint one on."""
    chart = _only_chart(filled_page[0]).chart
    for axis, which in ((chart.category_axis, "money"), (chart.value_axis, "rank")):
        line = axis.major_gridlines.format.line
        assert line.fill.type is None or str(line.fill.type).startswith("BACKGROUND"), \
            f"a grid was painted onto the {which} axis"


def test_the_deck_still_opens_after_the_chart_is_refilled(filled_page):
    # These charts are authored off a BINARY workbook: replace_data writes .xlsx bytes into
    # the .xlsb part unless the chart is detached first, and PowerPoint then refuses the file.
    chart_part = _only_chart(filled_page[0]).chart.part
    targets = [str(rel.target_ref or "").lower() for rel in chart_part.rels.values()]
    assert not [t for t in targets if t.endswith(".xlsb")]
    assert [t for t in targets if t.endswith(".xlsx")], "no workbook backs the written data"


def test_the_painted_bands_survive_onto_the_new_frame(filled_page):
    slide, _, whole = filled_page
    bands = [sh for sh in slide.shapes
             if str(sh.shape_type).startswith("GROUP") and sh.width > whole.w * 0.8
             and sh.height > whole.h * 0.8]
    assert len(bands) == 1, "the author's priority matrix is the page — it must not be stripped"


def test_the_hand_drawn_axis_and_dummy_point_labels_do_not_survive(filled_page):
    slide, _, whole = filled_page
    inside = [sh for sh in slide.shapes
              if not getattr(sh, "has_chart", False) and sh.left is not None
              and whole.x <= sh.left + sh.width // 2 <= whole.x + whole.w
              and whole.y <= sh.top + sh.height // 2 <= whole.y + whole.h]
    texts = [sh.text_frame.text for sh in inside if getattr(sh, "has_text_frame", False)]
    # No "$2,080M" tick faking a broken axis, and no "Product (1)" dummy label.
    assert not [t for t in texts if "Product" in t or "$" in t]


def test_both_axis_captions_survive_and_move_with_the_panel(filled_page):
    slide, panels, whole = filled_page
    captions = {sh.text_frame.text.strip(): sh for sh in slide.shapes
                if getattr(sh, "has_text_frame", False)}
    size = captions["Size of Marsh Portfolio"]
    assert size.top > panels[0].frame.y + panels[0].frame.h    # followed the panel down
    assert size.top < whole.y + whole.h + IN                   # …and stayed under it
    # The rank caption stays: the refilled panel still plots rank up the side.
    assert [t for t in captions if "Rank" in t]
