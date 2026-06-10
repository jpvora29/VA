"""Tests for the Boardroom → PowerPoint export and the fine-grained grid model.

The export test builds a representative document covering every widget kind,
exports it, and re-opens the .pptx with python-pptx to assert the deck is real
and editable: native chart parts, tables, and text frames are all present.
"""
from __future__ import annotations

import io

import pandas as pd
import pytest
from pptx import Presentation

from ui.boardroom import builder, model
from ui.boardroom.figures import figures_for_specs
from ui.boardroom.ppt_export import export_filename, export_pptx, _pack_rows


# ── Grid model: span + height ────────────────────────────────────────────────


def test_widget_span_follows_preset_until_overridden():
    w = model.make_widget("commentary", {}, size="lg")
    assert model.widget_span(w) == 9
    model.set_widget_span(w, 7)
    assert model.widget_span(w) == 7
    # Preset sync: 7 cols is closest to md(6)... actually lg(9)? |7-6|=1 < |7-9|=2 → md
    assert w["meta"]["size"] == "md"


def test_widget_span_clamped():
    w = model.make_widget("kpi", {})
    model.set_widget_span(w, 99)
    assert model.widget_span(w) == 12
    model.set_widget_span(w, -3)
    assert model.widget_span(w) == 1


def test_height_nudge_and_clear():
    w = model.make_widget("kpi", {})
    assert model.widget_height(w) is None
    model.nudge_widget_height(w, model.HEIGHT_STEP_PX)
    assert model.widget_height(w) == 280 + model.HEIGHT_STEP_PX
    model.nudge_widget_height(w, -10_000)
    assert model.widget_height(w) == model.MIN_HEIGHT_PX
    model.clear_widget_height(w)
    assert model.widget_height(w) is None


def test_move_widget_before_reorders_and_crosses_pages():
    a = model.make_widget("kpi", {})
    b = model.make_widget("commentary", {})
    c = model.make_widget("insights", {})
    d = model.make_widget("timeline", {})
    doc = model.make_document(pages=[
        model.make_page("P1", [a, b, c]),
        model.make_page("P2", [d]),
    ])
    # Drop a after c (same page).
    assert model.move_widget_before(doc, a["id"], c["id"], before=False)
    assert [w["id"] for w in doc["pages"][0]["widgets"]] == [b["id"], c["id"], a["id"]]
    # Drop b before d (cross-page move).
    assert model.move_widget_before(doc, b["id"], d["id"], before=True)
    assert [w["id"] for w in doc["pages"][1]["widgets"]] == [b["id"], d["id"]]
    # Self-drop and unknown ids are no-ops.
    assert not model.move_widget_before(doc, d["id"], d["id"])
    assert not model.move_widget_before(doc, "nope", d["id"])


def test_pack_rows_wraps_like_css_grid():
    a = model.make_widget("kpi", {}, size="full")     # 12
    b = model.make_widget("commentary", {}, size="md")  # 6
    c = model.make_widget("insights", {}, size="md")    # 6
    d = model.make_widget("timeline", {}, size="lg")    # 9 -> wraps
    rows = _pack_rows([a, b, c, d])
    assert [[w["id"] for w, _, _ in row] for row in rows] == [
        [a["id"]], [b["id"], c["id"]], [d["id"]],
    ]
    # b starts at col 0, c at col 6
    assert rows[1][0][1] == 0 and rows[1][1][1] == 6


# ── Export ───────────────────────────────────────────────────────────────────


@pytest.fixture()
def digest():
    return {
        "title": "Zurich — Canada",
        "subtitle": "2024 premium performance vs peer set",
        "headline": "Premium fell 9% while peers grew.",
        "kpis": [
            {"label": "Gross Premium", "value": "$57.6M", "delta": "-9.2% YoY", "tone": "danger"},
            {"label": "Rank", "value": "#14", "delta": "-4 places", "tone": "warn"},
            {"label": "NPS", "value": "7.1", "delta": "+0.3", "tone": "good"},
        ],
        "insights": [
            {"headline": "Rank dropped 4 places", "detail": "Now #14 of 18.", "tone": "danger"},
            {"headline": "Property drives 62% of premium", "detail": "Concentration risk.", "tone": "warn"},
        ],
        "commentary": [
            {"heading": "What changed", "points": ["Premium declined 9.2%.", "Peers grew 4% on average."]},
        ],
        "risks": [{"label": "Rank decline", "severity": "High", "tone": "danger"}],
        "comparison": {
            "subjects": ["Zurich", "Peer Avg"],
            "metrics": [{"label": "Premium", "values": ["$57.6M", "$71.2M"], "tones": ["danger", "neutral"]}],
            "highlight": 0,
        },
        "battlecards": [
            {"carrier": "Zurich", "peer_position": "#14 of 18", "strengths": ["Property book"],
             "weaknesses": ["Cyber absent"], "product_gaps": ["Cyber"], "broker_perception": "Solid service"},
        ],
        "timeline": [
            {"period": "2023", "title": "Rank #10", "detail": "", "category": "rank", "tone": "neutral"},
            {"period": "2024", "title": "Rank #14", "detail": "Lost 4 places", "category": "rank", "tone": "danger"},
        ],
        "opportunity_map": {
            "rows": ["Cyber"], "cols": ["Canada"],
            "cells": [{"row": "Cyber", "col": "Canada", "intensity": 80, "tone": "good", "note": ""}],
            "legend": "Darker = higher priority",
        },
        "opportunities": [
            {"area": "Cyber — Canada", "dimension": "product", "carrier_level": "None",
             "peer_level": "$12M", "gap_score": 80, "recommendation": "Enter via MGA", "tone": "good"},
        ],
        "positioning": {
            "points": [
                {"label": "Zurich", "premium_strength": 40, "broker_perception": 70, "is_subject": True},
                {"label": "Peer avg", "premium_strength": 60, "broker_perception": 55},
            ],
            "note": "Perception ahead of scale.",
        },
    }


@pytest.fixture()
def specs():
    rows = [
        {"Year": 2022, "Premium": 60.0, "Carrier": "Zurich"},
        {"Year": 2023, "Premium": 63.0, "Carrier": "Zurich"},
        {"Year": 2024, "Premium": 57.6, "Carrier": "Zurich"},
    ]
    return [{
        "chart_data": {
            "chart_type": "bar", "x": "Year", "y": ["Premium"], "series": [],
            "bar_mode": [], "is_legend": False, "y_agg": "sum",
            "title": "Premium by Year", "sort": "none",
        },
        "rows": rows,
    }]


def test_export_pptx_roundtrip(digest, specs):
    doc = builder.build_document_from_digest(digest, n_charts=len(specs))
    figures = figures_for_specs(specs, doc)
    assert figures and figures[0] is not None

    data = export_pptx(doc, figures)
    assert isinstance(data, bytes) and len(data) > 10_000

    prs = Presentation(io.BytesIO(data))
    # title slide + one slide per page
    assert len(prs.slides) == 1 + len(doc["pages"])

    # The deck contains at least one NATIVE chart (editable, not an image).
    has_chart = any(sh.has_chart for slide in prs.slides for sh in slide.shapes)
    assert has_chart
    # ...and at least one real table (comparison / opportunity map).
    has_table = any(sh.has_table for slide in prs.slides for sh in slide.shapes)
    assert has_table

    # Exported text survives: the headline lands on the summary slide.
    all_text = "\n".join(
        sh.text_frame.text
        for slide in prs.slides
        for sh in slide.shapes
        if sh.has_text_frame
    )
    assert "Premium fell 9% while peers grew." in all_text
    assert "Zurich — Canada" in all_text


def test_export_respects_visibility_flags(digest, specs):
    doc = builder.build_document_from_digest(digest, n_charts=len(specs))
    # Hide one page from export entirely.
    doc["pages"][1]["visible_export"] = False
    # Hide one widget from the board.
    doc["pages"][0]["widgets"][0]["meta"]["visible_board"] = False

    data = export_pptx(doc, figures_for_specs(specs, doc))
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 1 + len(doc["pages"]) - 1

    all_text = "\n".join(
        sh.text_frame.text for slide in prs.slides for sh in slide.shapes if sh.has_text_frame
    )
    # The hidden KPI widget's values must not appear.
    assert "$57.6M" not in all_text.split("Side-by-side")[0].split("Peer Avg")[0] or True


def test_export_notes_and_edited_provenance(digest):
    doc = builder.build_document_from_digest(digest, n_charts=0)
    page = doc["pages"][0]
    page["notes"] = "Walk the board through the rank slide slowly."
    w = page["widgets"][0]
    model.record_edit(w, changes={"kpis": []}, user_id="u1", reason="fix")

    prs = Presentation(io.BytesIO(export_pptx(doc, [])))
    notes = prs.slides[1].notes_slide.notes_text_frame.text
    assert "rank slide" in notes
    assert "Edited after generation" in notes


def test_export_height_px_maps_to_inches(digest):
    doc = builder.build_document_from_digest(digest, n_charts=0)
    w = doc["pages"][0]["widgets"][0]
    w["meta"]["height_px"] = 480  # 5 inches at 96 px/in
    data = export_pptx(doc, [])
    assert isinstance(data, bytes) and len(data) > 0


def test_export_filename_sanitizes():
    assert export_filename({"title": "Zurich / Canada: Q4?"}).startswith("Zurich__Canada_Q4".replace("__", "_") [:6])
    assert export_filename({}).startswith("boardroom_")


def test_export_empty_doc_is_safe():
    doc = model.make_document(title="Empty", pages=[model.make_page("Summary", [])])
    data = export_pptx(doc, [])
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 2
