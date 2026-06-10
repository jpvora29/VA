"""Export an editable Boardroom document to a native, fully-editable .pptx.

Design goals (per the product ask):
  * EXACT layout parity - the deck reuses the same governed 12-column grid the
    on-screen board uses (`model.widget_span` / `model.widget_height`), so a
    widget that spans 9 columns at 320 px tall on screen occupies the same
    relative footprint on the slide. CSS pixels map to inches at 96 px/in.
  * EDITABLE output - everything is real PowerPoint objects: text boxes,
    autoshapes, tables, and NATIVE charts built from the underlying rows (no
    screenshots), so a leader can retitle a chart or edit a KPI in PowerPoint.
  * Honour the document's presentation flags: pages with
    ``visible_export=False`` and widgets hidden from the board or export are
    skipped; page speaker notes land in the slide notes.

Public entry point: :func:`export_pptx`.
"""
from __future__ import annotations

import base64
import io
import math
import re
from datetime import date
from typing import Any, Callable, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from logger import get_logger
from ui.boardroom import catalog, model
from ui.boardroom.themes import theme as bm_theme
from ui.color_pallet import ColorPalette

logger = get_logger(__name__)

# -- Geometry: 16:9 slide, CSS px <-> inches at 96 px/in ----------------------
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.45
HEADER_IN = 0.72
PX_PER_IN = 96.0
GUTTER_IN = 14 / PX_PER_IN  # the on-screen grid gap is 14 px

CONTENT_W = SLIDE_W_IN - 2 * MARGIN_IN
CONTENT_H = SLIDE_H_IN - MARGIN_IN - HEADER_IN - 0.25
COL_W = (CONTENT_W - (model.GRID_COLUMNS - 1) * GUTTER_IN) / model.GRID_COLUMNS

FONT = "Segoe UI"
BULLET = "•  "

# Brand / tone colours (mirrors the UI tone classes + boardroom themes).
NAVY = RGBColor(0x0B, 0x13, 0x20)
GRAY = RGBColor(0x5B, 0x65, 0x77)
LIGHT_BORDER = RGBColor(0xD7, 0xDF, 0xEB)
SOFT_BG = RGBColor(0xF4, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_TONE_HEX = {
    "good": (0x1F, 0x9D, 0x55),
    "warn": (0xB9, 0x81, 0x0A),
    "danger": (0xC5, 0x35, 0x32),
    "neutral": (0x0B, 0x4B, 0xFF),
}


def _tone(value: Optional[str]) -> str:
    v = (value or "neutral").strip().lower()
    return v if v in _TONE_HEX else "neutral"


def _tone_color(value: Optional[str]) -> RGBColor:
    return RGBColor(*_TONE_HEX[_tone(value)])


def _tone_soft(value: Optional[str], alpha: float = 0.12) -> RGBColor:
    """Tone blended towards white - the PPT stand-in for rgba(tone, alpha)."""
    r, g, b = _TONE_HEX[_tone(value)]
    mix = lambda c: int(round(255 + (c - 255) * alpha))  # noqa: E731
    return RGBColor(mix(r), mix(g), mix(b))


def _hex_rgb(value: str, fallback: RGBColor = NAVY) -> RGBColor:
    m = re.fullmatch(r"#?([0-9a-fA-F]{6})", (value or "").strip())
    return RGBColor.from_string(m.group(1)) if m else fallback


# -- Low-level shape helpers ---------------------------------------------------


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return box, tf


def _para(
    tf,
    text: str,
    *,
    size: float = 11,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = NAVY,
    align=PP_ALIGN.LEFT,
    bullet: bool = False,
    space_after: float = 2,
    first: bool = False,
):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = (BULLET + text) if bullet else text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p


def _rounded(slide, x, y, w, h, *, fill: Optional[RGBColor] = WHITE,
             line: Optional[RGBColor] = LIGHT_BORDER, radius: float = 0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    try:  # corner radius (fraction of the smaller side)
        shp.adjustments[0] = radius
    except Exception:  # noqa: BLE001 - cosmetic only
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _rect(slide, x, y, w, h, fill: RGBColor, line: Optional[RGBColor] = None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _dot(slide, x, y, d, fill: RGBColor, line: Optional[RGBColor] = None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    return shp


def _set_cell(cell, text: str, *, size=10, bold=False, color: RGBColor = NAVY,
              fill: Optional[RGBColor] = None, align=PP_ALIGN.LEFT):
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


# -- Height estimation (inches) for auto-height widgets -------------------------


def _text_lines(text: str, w_in: float, size_pt: float = 11) -> int:
    chars_per_line = max(10, int(w_in * 96 / (size_pt * 0.55)))
    lines = 0
    for para in str(text or "").split("\n"):
        lines += max(1, math.ceil(len(para) / chars_per_line))
    return lines


def _estimate_height(widget: Dict[str, Any], w_in: float) -> float:
    """Content-driven height estimate, mirroring how tall each widget renders."""
    kind = widget.get("kind")
    data = widget.get("data") or {}
    content = catalog.content_of(kind)

    if content == "chart":
        return 3.3
    if content == "kpis" or kind == "kpi":
        n = len(data.get("kpis") or []) or 1
        per_row = max(1, int(w_in // 1.85))
        return 0.25 + math.ceil(n / per_row) * 1.05
    if kind == "insights":
        n = len(data.get("insights") or []) or 1
        per_row = 2 if w_in > 5 else 1
        return 0.2 + math.ceil(n / per_row) * 1.05
    if kind == "commentary":
        h = 0.2
        if data.get("headline"):
            h += 0.45
        for s in data.get("sections") or []:
            h += 0.34 + 0.3 * len(s.get("points") or [])
        risks = data.get("risks") or []
        if risks:
            h += 0.34 + 0.3 * len(risks)
        return max(h, 0.8)
    if kind == "comparison":
        comp = data.get("comparison") or {}
        return 0.45 + 0.34 * (len(comp.get("metrics") or []) + 1)
    if kind == "timeline":
        return 0.45 + 0.52 * len(data.get("timeline") or [])
    if kind == "opportunity_map":
        m = data.get("opportunity_map") or {}
        return 0.7 + 0.34 * (len(m.get("rows") or []) + 1)
    if kind == "opportunity_radar":
        return 0.4 + 0.95 * len(data.get("opportunities") or [])
    if kind == "positioning":
        return 3.7
    if kind == "battlecards":
        cards = data.get("battlecards") or []
        per_row = 2 if w_in > 7 else 1
        return 0.2 + math.ceil(max(len(cards), 1) / per_row) * 2.3
    if content == "table":
        return 0.4 + 0.32 * (len(data.get("rows") or []) + 1)
    if content == "kv":
        return 0.3 + 0.3 * len(data.get("rows") or [])
    if content == "text":
        return 0.3 + 0.24 * _text_lines(data.get("text") or "", w_in)
    if content == "callout":
        return 0.4 + 0.24 * _text_lines(data.get("text") or "", w_in, 12)
    if content == "image":
        return 2.6
    if content == "quote":
        return 0.6 + 0.24 * _text_lines(data.get("text") or "", w_in, 12)
    if content == "section_title":
        return 0.5
    if content == "divider":
        return 0.12
    if content == "list":
        return 0.3 + 0.28 * len(data.get("items") or [])
    return 1.2


# -- Native chart construction from a plotly figure ------------------------------


def _trace_color(tr, i: int) -> RGBColor:
    c = getattr(getattr(tr, "marker", None), "color", None)
    if isinstance(c, str) and c.startswith("#"):
        return _hex_rgb(c, _hex_rgb(ColorPalette.color_for(i)))
    line_c = getattr(getattr(tr, "line", None), "color", None)
    if isinstance(line_c, str) and line_c.startswith("#"):
        return _hex_rgb(line_c, _hex_rgb(ColorPalette.color_for(i)))
    return _hex_rgb(ColorPalette.color_for(i))


def _as_float(v) -> Optional[float]:
    try:
        f = float(v)
        return None if math.isnan(f) else f
    except (TypeError, ValueError):
        return None


def _seq(v) -> list:
    """Plotly trace arrays may be tuples or numpy arrays (whose truthiness is
    ambiguous) - normalise to a plain list."""
    if v is None:
        return []
    try:
        return list(v)
    except TypeError:
        return []


def _add_native_chart(slide, x, y, w, h, fig, fallback_title: str = "") -> bool:
    """Build a native, editable PPT chart from a plotly figure's traces.

    Mapping: Bar->clustered/stacked column, Scatter->line (or XY scatter when
    marker-only), Pie->pie/doughnut, Waterfall->column of the movements. A combo
    (bars + secondary-axis line) lands as a clustered column with the line
    measure as an extra series - still one editable chart object in PowerPoint.
    Returns False when the figure has nothing chartable.
    """
    traces = list(getattr(fig, "data", []) or [])
    if not traces:
        return False
    title = ""
    try:
        title = (fig.layout.title.text or "").strip()
    except Exception:  # noqa: BLE001
        pass
    title = title or fallback_title

    pies = [t for t in traces if t.type == "pie"]
    if pies:
        tr = pies[0]
        data = CategoryChartData()
        data.categories = [str(l) for l in _seq(tr.labels)]
        data.add_series(title or "Share", [(_as_float(v) or 0) for v in _seq(tr.values)])
        hole = getattr(tr, "hole", 0) or 0
        ctype = XL_CHART_TYPE.DOUGHNUT if hole else XL_CHART_TYPE.PIE
        gframe = slide.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), data)
        chart = gframe.chart
    else:
        xy_traces = [t for t in traces if t.type in ("bar", "scatter", "waterfall")]
        if not xy_traces:
            return False
        scatter_only = all(t.type == "scatter" for t in xy_traces)
        marker_only = scatter_only and all(
            "lines" not in (getattr(t, "mode", "") or "lines") for t in xy_traces
        )
        if marker_only:
            data = XyChartData()
            for i, tr in enumerate(xy_traces):
                s = data.add_series(str(tr.name or f"Series {i + 1}"))
                for xv, yv in zip(_seq(tr.x), _seq(tr.y)):
                    fx, fy = _as_float(xv), _as_float(yv)
                    if fx is not None and fy is not None:
                        s.add_data_point(fx, fy)
            gframe = slide.shapes.add_chart(
                XL_CHART_TYPE.XY_SCATTER, Inches(x), Inches(y), Inches(w), Inches(h), data
            )
            chart = gframe.chart
        else:
            # Align all traces on the ordered union of category labels.
            cats: List[str] = []
            seen = set()
            for tr in xy_traces:
                for xv in _seq(tr.x):
                    k = str(xv)
                    if k not in seen:
                        seen.add(k)
                        cats.append(k)
            data = CategoryChartData()
            data.categories = cats
            for i, tr in enumerate(xy_traces):
                lut = {str(xv): _as_float(yv) for xv, yv in zip(_seq(tr.x), _seq(tr.y))}
                data.add_series(str(tr.name or f"Series {i + 1}"), [lut.get(c) for c in cats])
            stacked = (getattr(fig.layout, "barmode", None) == "stack")
            has_bars = any(t.type in ("bar", "waterfall") for t in xy_traces)
            if has_bars:
                ctype = XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED
            else:
                ctype = XL_CHART_TYPE.LINE_MARKERS
            gframe = slide.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), data)
            chart = gframe.chart
            # Brand-colour the series to match the on-screen figure.
            for i, (series, tr) in enumerate(zip(chart.series, xy_traces)):
                try:
                    color = _trace_color(tr, i)
                    if has_bars:
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = color
                    else:
                        series.format.line.color.rgb = color
                except Exception:  # noqa: BLE001 - styling is best-effort
                    pass

    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for p in chart.chart_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = NAVY
    multi = len(getattr(chart, "series", [])) > 1 or bool(pies)
    chart.has_legend = multi
    if multi:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        try:
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = FONT
        except Exception:  # noqa: BLE001
            pass
    return True


# -- Widget renderers ------------------------------------------------------------
# Each takes (slide, x, y, w, h, widget, ctx) with the rect in inches.


def _render_kpis(slide, x, y, w, h, widget, ctx):
    kpis = (widget.get("data") or {}).get("kpis") or []
    if not kpis:
        return
    per_row = max(1, min(len(kpis), int(w // 1.85) or 1))
    gap = 0.12
    tile_w = (w - (per_row - 1) * gap) / per_row
    rows = math.ceil(len(kpis) / per_row)
    tile_h = min(1.0, (h - (rows - 1) * gap) / rows)
    for i, card in enumerate(kpis):
        r, c = divmod(i, per_row)
        tx = x + c * (tile_w + gap)
        ty = y + r * (tile_h + gap)
        tone = _tone(card.get("tone"))
        _rounded(slide, tx, ty, tile_w, tile_h, fill=WHITE, line=LIGHT_BORDER)
        _dot(slide, tx + 0.12, ty + tile_h / 2 - 0.11, 0.22, _tone_soft(tone, 0.25))
        _dot(slide, tx + 0.185, ty + tile_h / 2 - 0.045, 0.09, _tone_color(tone))
        _, tf = _textbox(slide, tx + 0.45, ty + 0.08, tile_w - 0.55, tile_h - 0.16)
        _para(tf, card.get("label", ""), size=8.5, color=GRAY, first=True, space_after=0)
        _para(tf, card.get("value", ""), size=15, bold=True, color=NAVY, space_after=0)
        if card.get("delta"):
            _para(tf, card["delta"], size=8.5, bold=True, color=_tone_color(tone), space_after=0)


def _render_commentary(slide, x, y, w, h, widget, ctx):
    data = widget.get("data") or {}
    _, tf = _textbox(slide, x, y, w, h)
    first = True
    if data.get("headline"):
        _para(tf, data["headline"], size=13, bold=True, color=NAVY, first=first, space_after=6)
        first = False
    for s in data.get("sections") or []:
        _para(tf, s.get("heading", ""), size=11, bold=True, color=NAVY, first=first, space_after=2)
        first = False
        for pt in s.get("points") or []:
            _para(tf, pt, size=10, color=GRAY, bullet=True, space_after=2)
    risks = data.get("risks") or []
    if risks:
        _para(tf, "Risks & watch items", size=11, bold=True, color=NAVY, first=first, space_after=2)
        for r in risks:
            _para(
                tf,
                f"{r.get('label', '')} - {r.get('severity', '')}",
                size=10,
                bold=True,
                color=_tone_color(r.get("tone")),
                bullet=True,
                space_after=2,
            )


def _render_insights(slide, x, y, w, h, widget, ctx):
    cards = (widget.get("data") or {}).get("insights") or []
    if not cards:
        return
    per_row = 2 if w > 5 else 1
    gap = 0.12
    card_w = (w - (per_row - 1) * gap) / per_row
    rows = math.ceil(len(cards) / per_row)
    card_h = min(1.05, (h - (rows - 1) * gap) / rows)
    for i, card in enumerate(cards):
        r, c = divmod(i, per_row)
        cx = x + c * (card_w + gap)
        cy = y + r * (card_h + gap)
        tone = _tone(card.get("tone"))
        _rounded(slide, cx, cy, card_w, card_h, fill=_tone_soft(tone, 0.10), line=LIGHT_BORDER)
        _rect(slide, cx, cy, 0.045, card_h, _tone_color(tone))
        _, tf = _textbox(slide, cx + 0.15, cy + 0.08, card_w - 0.28, card_h - 0.16)
        _para(tf, card.get("headline", ""), size=10.5, bold=True, color=NAVY, first=True, space_after=2)
        if card.get("detail"):
            _para(tf, card["detail"], size=9, color=GRAY, space_after=0)


def _render_comparison(slide, x, y, w, h, widget, ctx):
    comp = (widget.get("data") or {}).get("comparison") or {}
    subjects = comp.get("subjects") or []
    metrics = comp.get("metrics") or []
    if not subjects or not metrics:
        return
    highlight = comp.get("highlight", 0)
    rows, cols = len(metrics) + 1, len(subjects) + 1
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    table.first_row = False
    table.horz_banding = True
    _set_cell(table.cell(0, 0), "", fill=SOFT_BG)
    for j, s in enumerate(subjects):
        _set_cell(
            table.cell(0, j + 1), s, size=10, bold=True,
            color=WHITE if j == highlight else NAVY,
            fill=_hex_rgb("#0b4bff") if j == highlight else SOFT_BG,
            align=PP_ALIGN.CENTER,
        )
    for i, m in enumerate(metrics):
        _set_cell(table.cell(i + 1, 0), m.get("label", ""), size=9.5, bold=True)
        values, tones = m.get("values") or [], m.get("tones") or []
        for j in range(len(subjects)):
            val = values[j] if j < len(values) else "-"
            tn = tones[j] if j < len(tones) else None
            _set_cell(
                table.cell(i + 1, j + 1), val, size=9.5,
                bold=j == highlight,
                color=_tone_color(tn) if tn else NAVY,
                fill=_tone_soft("neutral", 0.06) if j == highlight else None,
                align=PP_ALIGN.CENTER,
            )


def _render_timeline(slide, x, y, w, h, widget, ctx):
    events = (widget.get("data") or {}).get("timeline") or []
    if not events:
        return
    row_h = min(0.62, h / len(events))
    rail_x = x + 0.95
    rail_h = row_h * (len(events) - 1) + 0.12 if len(events) > 1 else 0.1
    _rect(slide, rail_x, y + 0.08, 0.018, rail_h, LIGHT_BORDER)
    for i, e in enumerate(events):
        ey = y + i * row_h
        _, tf = _textbox(slide, x, ey, 0.85, row_h)
        _para(tf, e.get("period", ""), size=9, bold=True, color=GRAY, first=True, align=PP_ALIGN.RIGHT)
        _dot(slide, rail_x - 0.05, ey + 0.05, 0.12, _tone_color(e.get("tone")))
        _, tf = _textbox(slide, rail_x + 0.18, ey, w - (rail_x - x) - 0.2, row_h)
        _para(tf, e.get("title", ""), size=10, bold=True, color=NAVY, first=True, space_after=0)
        if e.get("detail"):
            _para(tf, e["detail"], size=8.5, color=GRAY, space_after=0)


def _render_opportunity_map(slide, x, y, w, h, widget, ctx):
    m = (widget.get("data") or {}).get("opportunity_map") or {}
    rows, cols = m.get("rows") or [], m.get("cols") or []
    cells = {(c.get("row"), c.get("col")): c for c in (m.get("cells") or [])}
    if not rows or not cols:
        return
    legend = m.get("legend") or ""
    table_h = h - (0.25 if legend else 0)
    shape = slide.shapes.add_table(
        len(rows) + 1, len(cols) + 1, Inches(x), Inches(y), Inches(w), Inches(table_h)
    )
    table = shape.table
    table.first_row = False
    _set_cell(table.cell(0, 0), "", fill=SOFT_BG)
    for j, c in enumerate(cols):
        _set_cell(table.cell(0, j + 1), c, size=9, bold=True, fill=SOFT_BG, align=PP_ALIGN.CENTER)
    for i, r in enumerate(rows):
        _set_cell(table.cell(i + 1, 0), r, size=9, bold=True)
        for j, c in enumerate(cols):
            cell = cells.get((r, c))
            if cell:
                inten = max(0, min(100, int(cell.get("intensity", 0) or 0)))
                _set_cell(
                    table.cell(i + 1, j + 1), str(inten), size=9, bold=inten >= 60,
                    color=NAVY,
                    fill=_tone_soft(cell.get("tone"), 0.10 + 0.006 * inten),
                    align=PP_ALIGN.CENTER,
                )
            else:
                _set_cell(table.cell(i + 1, j + 1), "", fill=WHITE)
    if legend:
        _, tf = _textbox(slide, x, y + table_h + 0.04, w, 0.2)
        _para(tf, legend, size=8.5, italic=True, color=GRAY, first=True)


def _render_radar(slide, x, y, w, h, widget, ctx):
    ops = (widget.get("data") or {}).get("opportunities") or []
    ops = sorted(ops, key=lambda o: o.get("gap_score", 0), reverse=True)
    if not ops:
        return
    item_h = min(0.95, h / len(ops))
    for i, o in enumerate(ops):
        oy = y + i * item_h
        tone = _tone(o.get("tone"))
        gap = max(0, min(100, int(o.get("gap_score", 0) or 0)))
        _, tf = _textbox(slide, x, oy, w - 0.55, 0.24)
        p = _para(tf, o.get("area", ""), size=10.5, bold=True, color=NAVY, first=True, space_after=0)
        dim = (o.get("dimension") or "").title()
        if dim:
            run = p.add_run()
            run.text = f"   {dim}"
            run.font.name = FONT
            run.font.size = Pt(8.5)
            run.font.color.rgb = GRAY
        _, tf = _textbox(slide, x + w - 0.5, oy, 0.5, 0.24)
        _para(tf, str(gap), size=11, bold=True, color=_tone_color(tone), first=True, align=PP_ALIGN.RIGHT)
        _rounded(slide, x, oy + 0.3, w, 0.1, fill=SOFT_BG, line=None, radius=0.5)
        if gap > 0:
            _rounded(slide, x, oy + 0.3, max(0.12, w * gap / 100), 0.1, fill=_tone_color(tone), line=None, radius=0.5)
        detail = []
        if o.get("carrier_level"):
            detail.append(f"Carrier: {o['carrier_level']}")
        if o.get("peer_level"):
            detail.append(f"Marsh/Peers: {o['peer_level']}")
        line2 = " | ".join(detail)
        rec = o.get("recommendation") or ""
        if line2 or rec:
            _, tf = _textbox(slide, x, oy + 0.46, w, item_h - 0.48)
            if line2:
                _para(tf, line2, size=8.5, color=GRAY, first=True, space_after=1)
            if rec:
                _para(tf, "→ " + rec, size=9, italic=True, color=NAVY, first=not line2, space_after=0)


def _render_positioning(slide, x, y, w, h, widget, ctx):
    mx = (widget.get("data") or {}).get("positioning") or {}
    pts = mx.get("points") or []
    if not pts:
        return
    note = mx.get("note") or ""
    plot_h = h - 0.5 - (0.25 if note else 0)
    plot_w = min(w - 0.3, plot_h * 1.35)
    px0 = x + (w - plot_w) / 2
    py0 = y + 0.06
    _rounded(slide, px0, py0, plot_w, plot_h, fill=SOFT_BG, line=LIGHT_BORDER, radius=0.03)
    _rect(slide, px0 + plot_w / 2, py0 + 0.04, 0.012, plot_h - 0.08, LIGHT_BORDER)
    _rect(slide, px0 + 0.04, py0 + plot_h / 2, plot_w - 0.08, 0.012, LIGHT_BORDER)
    quads = [("Emerging", 0.06, 0.05), ("Strong", plot_w - 1.1, 0.05),
             ("Underperforming", 0.06, plot_h - 0.28), ("Vulnerable", plot_w - 1.1, plot_h - 0.28)]
    for label, qx, qy in quads:
        _, tf = _textbox(slide, px0 + qx, py0 + qy, 1.05, 0.2)
        _para(tf, label, size=8, bold=True, color=GRAY, first=True)
    for p in pts:
        vx = max(0, min(100, int(p.get("premium_strength", 50) or 0)))
        vy = max(0, min(100, int(p.get("broker_perception", 50) or 0)))
        is_subject = bool(p.get("is_subject"))
        d = 0.22 if is_subject else 0.15
        cx = px0 + 0.15 + (plot_w - 0.3) * vx / 100 - d / 2
        cy = py0 + 0.15 + (plot_h - 0.3) * (100 - vy) / 100 - d / 2
        _dot(slide, cx, cy, d, _hex_rgb("#001f52") if is_subject else _tone_color(p.get("tone")),
             line=WHITE if is_subject else None)
        _, tf = _textbox(slide, cx + d + 0.02, cy - 0.02, 1.4, 0.2)
        _para(tf, p.get("label", ""), size=8.5, bold=is_subject, color=NAVY, first=True)
    _, tf = _textbox(slide, px0, py0 + plot_h + 0.04, plot_w, 0.2)
    _para(tf, "Premium strength →   (↑ Broker perception)", size=8, italic=True, color=GRAY,
          first=True, align=PP_ALIGN.CENTER)
    if note:
        _, tf = _textbox(slide, x, y + h - 0.22, w, 0.2)
        _para(tf, note, size=8.5, italic=True, color=GRAY, first=True)


def _render_battlecards(slide, x, y, w, h, widget, ctx):
    cards = (widget.get("data") or {}).get("battlecards") or []
    if not cards:
        return
    per_row = 2 if w > 7 else 1
    gap = 0.15
    card_w = (w - (per_row - 1) * gap) / per_row
    rows = math.ceil(len(cards) / per_row)
    card_h = min(2.35, (h - (rows - 1) * gap) / rows)
    for i, bc in enumerate(cards):
        r, c = divmod(i, per_row)
        cx = x + c * (card_w + gap)
        cy = y + r * (card_h + gap)
        _rounded(slide, cx, cy, card_w, card_h, fill=WHITE, line=LIGHT_BORDER)
        name = bc.get("carrier") or "?"
        _dot(slide, cx + 0.12, cy + 0.12, 0.34, _hex_rgb("#001f52"))
        _, tf = _textbox(slide, cx + 0.13, cy + 0.17, 0.32, 0.24)
        _para(tf, name[:2].upper(), size=10, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)
        _, tf = _textbox(slide, cx + 0.56, cy + 0.1, card_w - 0.7, 0.45)
        _para(tf, name, size=11.5, bold=True, color=NAVY, first=True, space_after=0)
        if bc.get("peer_position"):
            _para(tf, bc["peer_position"], size=8.5, color=GRAY, space_after=0)
        cols = [
            ("Strengths", bc.get("strengths"), "good"),
            ("Weaknesses", bc.get("weaknesses"), "danger"),
            ("Product gaps", bc.get("product_gaps"), "warn"),
        ]
        cols = [(t, items, tn) for t, items, tn in cols if items]
        if cols:
            col_w = (card_w - 0.24 - 0.1 * (len(cols) - 1)) / len(cols)
            footer = 0.3 if bc.get("broker_perception") else 0.08
            for j, (title, items, tn) in enumerate(cols):
                lx = cx + 0.12 + j * (col_w + 0.1)
                _, tf = _textbox(slide, lx, cy + 0.62, col_w, card_h - 0.62 - footer)
                _para(tf, title.upper(), size=7.5, bold=True, color=_tone_color(tn), first=True, space_after=2)
                for it in (items or [])[:4]:
                    _para(tf, it, size=8.5, color=NAVY, bullet=True, space_after=1)
        if bc.get("broker_perception"):
            _, tf = _textbox(slide, cx + 0.12, cy + card_h - 0.3, card_w - 0.24, 0.26)
            _para(tf, "“" + bc["broker_perception"] + "”", size=8.5, italic=True, color=GRAY, first=True)


def _render_chart(slide, x, y, w, h, widget, ctx):
    data = widget.get("data") or {}
    i = data.get("spec_index")
    figures = ctx.get("figures") or []
    fig = figures[i] if isinstance(i, int) and 0 <= i < len(figures) else None
    title = widget.get("title") or ""
    if fig is not None and _add_native_chart(slide, x, y, w, h, fig, title):
        return
    _, tf = _textbox(slide, x, y + h / 2 - 0.15, w, 0.3)
    _para(tf, "Chart unavailable for this view", size=10, italic=True, color=GRAY,
          first=True, align=PP_ALIGN.CENTER)


def _render_generic(slide, x, y, w, h, widget, ctx):
    """Library content types: text / list / table / kv / callout / image / quote /
    section_title / divider - mirrors ``widgets_library.render_content``."""
    content = catalog.content_of(widget.get("kind"))
    data = widget.get("data") or {}

    if content == "text":
        _, tf = _textbox(slide, x, y, w, h)
        first = True
        for parag in str(data.get("text") or "").split("\n"):
            if parag.strip():
                _para(tf, parag, size=10.5, color=NAVY, first=first, space_after=4)
                first = False
        return
    if content == "list":
        _, tf = _textbox(slide, x, y, w, h)
        for i, it in enumerate(data.get("items") or []):
            _para(tf, str(it), size=10.5, color=NAVY, bullet=True, first=i == 0, space_after=2)
        return
    if content == "table":
        cols = data.get("columns") or []
        rows = data.get("rows") or []
        if not cols:
            return
        shape = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        table.first_row = False
        for j, cname in enumerate(cols):
            _set_cell(table.cell(0, j), cname, size=9.5, bold=True, fill=SOFT_BG)
        for i, r in enumerate(rows):
            vals = r if isinstance(r, list) else [r]
            for j in range(len(cols)):
                _set_cell(table.cell(i + 1, j), vals[j] if j < len(vals) else "", size=9.5)
        return
    if content == "kv":
        rows = data.get("rows") or []
        if not rows:
            return
        shape = slide.shapes.add_table(len(rows), 2, Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        table.first_row = False
        for i, r in enumerate(rows):
            _set_cell(table.cell(i, 0), r[0] if r else "", size=9.5, bold=True, color=GRAY)
            _set_cell(table.cell(i, 1), r[1] if len(r) > 1 else "", size=9.5)
        return
    if content == "callout":
        tone = data.get("tone")
        _rounded(slide, x, y, w, h, fill=_tone_soft(tone, 0.10), line=None)
        _rect(slide, x, y, 0.05, h, _tone_color(tone))
        _, tf = _textbox(slide, x + 0.18, y + 0.08, w - 0.3, h - 0.16)
        _para(tf, data.get("text", ""), size=12, bold=True, color=NAVY, first=True)
        return
    if content == "image":
        url = (data.get("url") or "").strip()
        blob = _image_bytes(url)
        if blob is not None:
            cap = data.get("caption") or ""
            pic_h = h - (0.25 if cap else 0)
            try:
                slide.shapes.add_picture(io.BytesIO(blob), Inches(x), Inches(y), height=Inches(pic_h))
            except Exception:  # noqa: BLE001 - bad image data
                blob = None
            if cap and blob is not None:
                _, tf = _textbox(slide, x, y + pic_h + 0.03, w, 0.2)
                _para(tf, cap, size=8.5, italic=True, color=GRAY, first=True)
        if blob is None:
            _, tf = _textbox(slide, x, y + h / 2 - 0.12, w, 0.24)
            _para(tf, "[image unavailable]", size=9.5, italic=True, color=GRAY, first=True, align=PP_ALIGN.CENTER)
        return
    if content == "quote":
        _, tf = _textbox(slide, x + 0.15, y, w - 0.3, h)
        _para(tf, "“" + str(data.get("text", "")) + "”", size=12.5, italic=True, color=NAVY, first=True, space_after=4)
        if data.get("attribution"):
            _para(tf, "- " + data["attribution"], size=9.5, color=GRAY)
        _rect(slide, x, y, 0.045, h, _hex_rgb("#0b4bff"))
        return
    if content == "section_title":
        _rounded(slide, x, y, w, h, fill=_hex_rgb("#001f52"), line=None, radius=0.15)
        _, tf = _textbox(slide, x + 0.2, y, w - 0.4, h)
        box_p = _para(tf, data.get("text", "Section"), size=14, bold=True, color=WHITE, first=True)
        box_p.alignment = PP_ALIGN.LEFT
        return
    if content == "divider":
        _rect(slide, x, y + h / 2, w, 0.015, LIGHT_BORDER)
        return
    # Unknown -> dump as text so nothing is silently lost.
    _, tf = _textbox(slide, x, y, w, h)
    _para(tf, str(data), size=9, color=GRAY, first=True)


def _image_bytes(url: str) -> Optional[bytes]:
    """Decode a data-URI image (the upload path) or fetch http(s) best-effort."""
    if not url:
        return None
    if url.startswith("data:"):
        try:
            return base64.b64decode(url.split(",", 1)[1])
        except Exception:  # noqa: BLE001
            return None
    if url.startswith(("http://", "https://")):
        try:
            import urllib.request

            with urllib.request.urlopen(url, timeout=5) as resp:  # noqa: S310
                return resp.read()
        except Exception:  # noqa: BLE001
            return None
    return None


_RENDERERS: Dict[str, Callable] = {
    "kpi": _render_kpis,
    "commentary": _render_commentary,
    "insights": _render_insights,
    "comparison": _render_comparison,
    "timeline": _render_timeline,
    "opportunity_map": _render_opportunity_map,
    "opportunity_radar": _render_radar,
    "positioning": _render_positioning,
    "battlecards": _render_battlecards,
    "charts": _render_chart,
}


# -- Page layout: pack widgets onto the grid exactly like the CSS grid does -----


def _pack_rows(widgets: List[Dict[str, Any]]) -> List[List[Tuple[Dict[str, Any], int, int]]]:
    """Flow widgets left->right, wrapping at 12 columns (CSS auto-placement).
    Returns rows of (widget, col_start, span)."""
    rows: List[List[Tuple[Dict[str, Any], int, int]]] = []
    current: List[Tuple[Dict[str, Any], int, int]] = []
    cursor = 0
    for w in widgets:
        span = model.widget_span(w)
        if cursor + span > model.GRID_COLUMNS and current:
            rows.append(current)
            current, cursor = [], 0
        current.append((w, cursor, span))
        cursor += span
    if current:
        rows.append(current)
    return rows


def _exportable_widgets(page: Dict[str, Any]) -> List[Dict[str, Any]]:
    out = []
    for w in page.get("widgets", []):
        meta = w.get("meta") or {}
        if meta.get("visible_board", True) and meta.get("visible_export", True):
            out.append(w)
    return out


def _slide_header(slide, doc: Dict[str, Any], page: Dict[str, Any], page_no: int, total: int):
    _, tf = _textbox(slide, MARGIN_IN, 0.22, CONTENT_W - 2.2, 0.4)
    _para(tf, page.get("title") or f"Page {page_no}", size=18, bold=True, color=NAVY, first=True)
    _, tf = _textbox(slide, SLIDE_W_IN - MARGIN_IN - 3.2, 0.27, 3.2, 0.3)
    _para(tf, doc.get("title", ""), size=10, color=GRAY, first=True, align=PP_ALIGN.RIGHT)
    _rect(slide, MARGIN_IN, HEADER_IN - 0.06, CONTENT_W, 0.02, _hex_rgb("#0b4bff"))
    _, tf = _textbox(slide, SLIDE_W_IN - MARGIN_IN - 0.8, SLIDE_H_IN - 0.32, 0.8, 0.25)
    _para(tf, f"{page_no} / {total}", size=9, color=GRAY, first=True, align=PP_ALIGN.RIGHT)


def _render_widget_frame(slide, x, y, w, h, widget, ctx):
    """Widget chrome: themed soft panel + title, then the kind-specific body."""
    meta = widget.get("meta") or {}
    theme_key = meta.get("theme") or "default"
    themed = theme_key != "default"
    pad = 0.0
    if themed:
        th = bm_theme(theme_key)
        _rounded(slide, x, y, w, h, fill=_hex_rgb(th["soft"]), line=None, radius=0.05)
        _rect(slide, x, y, w, 0.035, _hex_rgb(th["accent"]))
        pad = 0.12
    bx, by, bw, bh = x + pad, y + pad, w - 2 * pad, h - 2 * pad
    title = (widget.get("title") or "").strip()
    if title and widget.get("kind") != "charts":  # charts carry their own title
        _, tf = _textbox(slide, bx, by, bw, 0.26)
        _para(tf, title, size=11, bold=True, color=NAVY, first=True)
        by += 0.3
        bh -= 0.3
    renderer = _RENDERERS.get(widget.get("kind"), _render_generic)
    try:
        renderer(slide, bx, by, bw, max(bh, 0.2), widget, ctx)
    except Exception:  # noqa: BLE001 - one bad widget must not kill the export
        logger.exception("PPT export: widget %s (%s) failed to render", widget.get("id"), widget.get("kind"))
        _, tf = _textbox(slide, bx, by, bw, 0.3)
        _para(tf, "[widget could not be exported]", size=9, italic=True, color=GRAY, first=True)


def _add_title_slide(prs: Presentation, doc: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, _hex_rgb("#001538"))
    _rect(slide, MARGIN_IN, 2.6, 2.2, 0.045, _hex_rgb("#009DE0"))
    _, tf = _textbox(slide, MARGIN_IN, 2.85, SLIDE_W_IN - 2 * MARGIN_IN, 1.2)
    _para(tf, doc.get("title") or "Boardroom", size=40, bold=True, color=WHITE, first=True)
    if doc.get("subtitle"):
        _, tf = _textbox(slide, MARGIN_IN, 4.05, SLIDE_W_IN - 2 * MARGIN_IN, 0.6)
        _para(tf, doc["subtitle"], size=16, color=RGBColor(0x9F, 0xE0, 0xFF), first=True)
    _, tf = _textbox(slide, MARGIN_IN, SLIDE_H_IN - 0.75, SLIDE_W_IN - 2 * MARGIN_IN, 0.3)
    _para(tf, f"Virtual Analyst | {date.today():%d %b %Y}", size=10.5,
          color=RGBColor(0x8F, 0xA5, 0xC2), first=True)


def export_pptx(doc: Dict[str, Any], figures: Optional[List[Any]] = None) -> bytes:
    """Render a Boardroom document to .pptx bytes.

    `figures` are the live plotly figures, index-aligned with the message's
    chart specs (same list `render_document` receives), used to build native
    editable charts.
    """
    doc = doc or {}
    ctx = {"figures": figures or []}
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_H_IN * 914400))

    _add_title_slide(prs, doc)

    pages = [p for p in (doc.get("pages") or []) if p.get("visible_export", True)]
    total = len(pages)
    for page_no, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _slide_header(slide, doc, page, page_no, total)

        widgets = _exportable_widgets(page)
        rows = _pack_rows(widgets)

        # Row heights: explicit pixel heights map 1:1 (96 px/in); auto widgets get
        # a content-based estimate. If the page runs taller than the slide, scale
        # every row by the same factor so relative proportions stay exact.
        row_heights: List[float] = []
        for row in rows:
            hs = []
            for w, _c0, span in row:
                fixed = model.widget_height(w)
                w_in = span * COL_W + (span - 1) * GUTTER_IN
                hs.append(fixed / PX_PER_IN if fixed else _estimate_height(w, w_in))
            row_heights.append(max(hs) if hs else 0.5)
        total_h = sum(row_heights) + GUTTER_IN * max(0, len(rows) - 1)
        scale = CONTENT_H / total_h if total_h > CONTENT_H else 1.0

        y = HEADER_IN + 0.12
        for row, rh in zip(rows, row_heights):
            rh_s = max(0.35, rh * scale)
            for w, c0, span in row:
                x = MARGIN_IN + c0 * (COL_W + GUTTER_IN)
                w_in = span * COL_W + (span - 1) * GUTTER_IN
                _render_widget_frame(slide, x, y, w_in, rh_s, w, ctx)
            y += rh_s + GUTTER_IN

        # Speaker notes (+ provenance for transparency on edited evidence).
        notes_lines = []
        if page.get("notes"):
            notes_lines.append(page["notes"])
        edited = [w.get("title") or catalog.meta_of(w.get("kind")).get("label", w.get("kind"))
                  for w in widgets if (w.get("meta") or {}).get("edited")]
        if edited:
            notes_lines.append("Edited after generation: " + ", ".join(edited))
        if notes_lines:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(notes_lines)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def export_filename(doc: Dict[str, Any]) -> str:
    title = (doc or {}).get("title") or "boardroom"
    safe = re.sub(r"[^A-Za-z0-9 _-]+", "", title).strip().replace(" ", "_") or "boardroom"
    return f"{safe}_{date.today():%Y%m%d}.pptx"
